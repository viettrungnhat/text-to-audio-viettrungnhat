#Ứng dụng Text to MP3/M4A đa ngôn ngữ _tích hợp A.I, Discord, Amazon AWS, gTTS.
import time
import warnings
import os
import re
from langdetect import detect
from gtts import gTTS
import tkinter as tk
from tkinter import ttk, messagebox, filedialog, simpledialog
import threading
import datetime
import ctypes
import sys
import json
import base64
import shutil
import requests
from openai import OpenAI
from google import genai
import tempfile
import subprocess
import signal
import random
import socket
from datetime import datetime
from pathlib import Path

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FFMPEG_BUNDLE_DIR = os.path.join(BASE_DIR, "ffmpeg_bin", "bin")
if os.path.isdir(FFMPEG_BUNDLE_DIR):
    os.environ["PATH"] = FFMPEG_BUNDLE_DIR + os.pathsep + os.environ.get("PATH", "")
    os.environ.setdefault("FFMPEG_BINARY", os.path.join(FFMPEG_BUNDLE_DIR, "ffmpeg"))
    os.environ.setdefault("IMAGEIO_FFMPEG_EXE", os.path.join(FFMPEG_BUNDLE_DIR, "ffmpeg"))

TCLTK_ROOT = "/usr/local/Cellar/tcl-tk/9.0.4/lib"
os.environ.setdefault("TK_LIBRARY", os.path.join(TCLTK_ROOT, "tk9.0"))
os.environ.setdefault("TCL_LIBRARY", os.path.join(TCLTK_ROOT, "tcl9.0"))
os.environ.setdefault("TK_SILENCE_DEPRECATION", "1")


def _load_env_file(path):
    if not os.path.exists(path):
        return
    try:
        with open(path, "r", encoding="utf-8") as env_file:
            for raw_line in env_file:
                line = raw_line.strip()
                if not line or line.startswith("#") or "=" not in line:
                    continue
                key, value = line.split("=", 1)
                key = key.strip()
                value = value.strip().strip('"').strip("'")
                if key and key not in os.environ:
                    os.environ[key] = value
    except Exception as exc:
        print(f"⚠ Không đọc được file .env: {exc}")


_load_env_file(os.path.join(BASE_DIR, ".env"))


def _env_or(value, *env_keys):
    for env_key in env_keys:
        env_value = os.environ.get(env_key, "").strip()
        if env_value:
            return env_value
    return value
from cryptography.fernet import Fernet, InvalidToken
from cryptography.hazmat.backends import default_backend
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from config.settings import AWS_REGION, DEFAULT_VOICE
from services.tts_service import (
    export_audio_batch,
    generate_audio_core,
    process_text_lines,
)
from utils.file_reader import read_excel_vocab, read_text_file

#=== chỉ chạy 1 app
def check_already_running(port=65432):
    s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    try:
        s.bind(('127.0.0.1', port))
    except socket.error:
        return

    # Ngăn socket bị đóng: lưu vào biến global
    global singleton_socket
    singleton_socket = s

# Không chặn mở app nếu instance cũ còn giữ cổng; macOS launcher sẽ tự mở bản hiện tại.


#============
warnings.simplefilter("ignore")

#warnings.filterwarnings("ignore", category=UserWarning) #tắt cảnh báo rác
import pygame
current_channel = None
# Disabled: game audio channel is not used in audio-tool version.
# mixer_channel_game = pygame.mixer.Channel(5)
import uuid
import boto3
from pydub import AudioSegment


def _ensure_pygame_mixer():
    global current_channel
    if pygame.mixer.get_init():
        return True
    try:
        if sys.platform == "darwin":
            os.environ.setdefault("SDL_AUDIODRIVER", "coreaudio")
        pygame.mixer.init()
        return True
    except Exception as exc:
        print(f"⚠ Không khởi tạo được pygame mixer: {exc}")
        return False

def _resolve_app_build_time():
    try:
        repo_root = Path(__file__).resolve().parent
        result = subprocess.run(
            ["git", "log", "-1", "--format=%cI", "--", Path(__file__).name],
            cwd=repo_root,
            capture_output=True,
            text=True,
            check=True,
        )
        raw = (result.stdout or "").strip()
        if raw:
            return datetime.fromisoformat(raw).astimezone().strftime("%Y-%m-%d %H:%M:%S")
    except Exception:
        pass
    return datetime.fromtimestamp(os.path.getmtime(__file__)).strftime("%Y-%m-%d %H:%M:%S")


APP_BUILD_TIME = _resolve_app_build_time()
APP_BUILD_TAG = f"Code mới nhất: {APP_BUILD_TIME}"
# Vocab M4A uses 32 kbps unless the user explicitly selects 26 kbps.
DEFAULT_VOCAB_M4A_BITRATE = "32k"


def app_beep(freq=1000, duration=200, widget=None):
    if os.name == "nt":
        import winsound
        winsound.Beep(freq, duration)
        return

    try:
        if widget is not None:
            widget.bell()
            return
        if tk._default_root is not None:
            tk._default_root.bell()
            return
    except Exception:
        pass

    print("\a", end="", flush=True)


def attach_mouse_text_menu(widget):
    menu = tk.Menu(widget, tearoff=0)

    def popup(event):
        try:
            widget.focus_set()
            menu.tk_popup(event.x_root, event.y_root)
        finally:
            menu.grab_release()
        return "break"

    menu.add_command(label="Cắt", command=lambda: widget.event_generate("<<Cut>>"))
    menu.add_command(label="Sao chép", command=lambda: widget.event_generate("<<Copy>>"))
    menu.add_command(label="Dán", command=lambda: widget.event_generate("<<Paste>>"))
    menu.add_separator()
    menu.add_command(label="Chọn tất cả", command=lambda: widget.tag_add("sel", "1.0", "end-1c") if isinstance(widget, tk.Text) else (widget.selection_range(0, tk.END), widget.icursor(tk.END)))

    for sequence in ("<Button-3>", "<Button-2>", "<Control-Button-1>"):
        widget.bind(sequence, popup, add="+")


def open_path_cross_platform(target):
    if not target:
        return

    target = os.fspath(target)
    try:
        if os.name == "nt":
            os.startfile(target)
        elif sys.platform == "darwin":
            subprocess.Popen(["open", target])
        else:
            subprocess.Popen(["xdg-open", target])
    except Exception as e:
        raise Exception(f"Không mở được file: {e}")


dang_doc_game = False

def end_doc_game():
    try:
        if 'current_channel' in globals() and current_channel and current_channel.get_busy():
            current_channel.stop()
    except: pass


def stop_all_doc():
    global current_channel
    try:
        if current_channel and current_channel.get_busy():
            current_channel.stop()
    except:
        pass




#===================

# Xác định thư mục gốc app (chạy EXE hoặc chạy Python)
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(os.path.abspath(sys.executable))
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

ASSETS_DIR = os.path.join(BASE_DIR, "assets")
SOUNDS_DIR = os.path.join(ASSETS_DIR, "sounds")
IMAGES_DIR = os.path.join(ASSETS_DIR, "images")
FONTS_DIR = os.path.join(ASSETS_DIR, "fonts")

# Thư mục dữ liệu app (AppData)
APPDATA_ROOT = os.path.join(BASE_DIR, "AppData")
os.makedirs(APPDATA_ROOT, exist_ok=True)

# Đường dẫn logo mặc định
LOGO_PATH = os.path.join(IMAGES_DIR, "logo.png")

# Vault secret được mã hóa bằng mật khẩu dùng chung
SECRET_VAULT_FILE = os.path.join(BASE_DIR, "secrets.enc")
SECRET_VAULT_PASSWORD_ENV = "TEXTTOMP3_VAULT_PASSWORD"
SECRET_VAULT_ITERATIONS = 390000

# Disabled: Game/Image paths kept as constants only for legacy dead code.
# Game - đường dẫn gốc (bản cài đặt)
EXCEL_GAME_ORIGINAL = os.path.join(BASE_DIR, "Game_doan_chu.xlsx")

# Game - đường dẫn dùng trong AppData (sẽ thao tác thật)
APPDATA_GAME_FILE = os.path.join(APPDATA_ROOT, "Game_doan_chu.xlsx")
APPDATA_IMAGE_FOLDER = os.path.join(APPDATA_ROOT, "images")

# Gán biến dùng chung
EXCEL_GAME_PATH = APPDATA_GAME_FILE
IMAGE_FOLDER = APPDATA_IMAGE_FOLDER
LICH_SU_FILE = os.path.join(APPDATA_ROOT, "lich_su_game.json")

# Disabled: không còn tự động copy Excel game hoặc tạo folder ảnh game.
# if not os.path.exists(APPDATA_GAME_FILE) and os.path.exists(EXCEL_GAME_ORIGINAL):
#     shutil.copy(EXCEL_GAME_ORIGINAL, APPDATA_GAME_FILE)
#     print("✅ Đã copy file Game_doan_chu.xlsx vào AppData")
#
# os.makedirs(APPDATA_IMAGE_FOLDER, exist_ok=True)


def _derive_vault_key(password, salt):
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=SECRET_VAULT_ITERATIONS,
        backend=default_backend(),
    )
    return base64.urlsafe_b64encode(kdf.derive(password.encode("utf-8")))


def _prompt_vault_password():
    env_password = os.environ.get(SECRET_VAULT_PASSWORD_ENV, "").strip()
    if env_password:
        return env_password

    temp_root = None
    try:
        temp_root = tk.Tk()
        temp_root.withdraw()
        dialog = tk.Toplevel(temp_root)
        dialog.title("Nhap mat khau secrets")
        dialog.resizable(False, False)
        dialog.attributes("-topmost", True)
        dialog.grab_set()
        try:
            dialog.lift()
            dialog.focus_force()
        except Exception:
            pass

        result = {"password": None}

        tk.Label(
            dialog,
            text="Nhap mat khau de mo vault bi mat:\nGoi y: Sinh nhật Bảo Khiêm-Linh Dương Call 0916385682",
            justify="left",
            wraplength=360,
        ).pack(padx=16, pady=(16, 8))

        password_var = tk.StringVar()
        entry = tk.Entry(dialog, textvariable=password_var, show="*", width=40)
        entry.pack(padx=16, pady=(0, 12))

        button_frame = tk.Frame(dialog)
        button_frame.pack(pady=(0, 14))

        def submit():
            result["password"] = password_var.get()
            dialog.destroy()

        def cancel():
            result["password"] = None
            dialog.destroy()

        tk.Button(button_frame, text="OK", width=10, command=submit).pack(side="left", padx=8)
        tk.Button(button_frame, text="Huy", width=10, command=cancel).pack(side="left", padx=8)

        dialog.protocol("WM_DELETE_WINDOW", cancel)
        dialog.bind("<Return>", lambda event: submit())
        dialog.bind("<Escape>", lambda event: cancel())
        entry.focus_set()
        dialog.update_idletasks()
        screen_width = dialog.winfo_screenwidth()
        screen_height = dialog.winfo_screenheight()
        width, height = 420, 180
        x = (screen_width // 2) - (width // 2)
        y = (screen_height // 2) - (height // 2)
        dialog.geometry(f"{width}x{height}+{x}+{y}")
        temp_root.wait_window(dialog)

        password = result["password"]
        if password is None:
            raise SystemExit("Nguoi dung da huy nhap mat khau secrets")
        password = password.strip()
        if not password:
            raise SystemExit("Mat khau secrets rong")
        return password
    finally:
        try:
            if temp_root is not None:
                temp_root.destroy()
        except Exception:
            pass


def _load_legacy_bundle():
    legacy_bundle = {}

    legacy_config_path = os.path.join(BASE_DIR, "config_default.json")
    if os.path.exists(legacy_config_path):
        try:
            with open(legacy_config_path, "r", encoding="utf-8") as f:
                legacy_bundle["config_default"] = json.load(f)
        except Exception:
            legacy_bundle["config_default"] = {}

    legacy_client_secret_path = os.path.join(BASE_DIR, "client_secret.json")
    if os.path.exists(legacy_client_secret_path):
        try:
            with open(legacy_client_secret_path, "r", encoding="utf-8") as f:
                legacy_bundle["client_secret_json"] = f.read()
        except Exception:
            legacy_bundle["client_secret_json"] = ""

    return legacy_bundle


def _load_secret_bundle():
    secret_password = os.environ.get(SECRET_VAULT_PASSWORD_ENV, "").strip()
    if not secret_password:
        print("ℹ️ Bỏ qua secrets.enc vì chưa có mật khẩu vault; dùng .env/AppData trước.")
        return _load_legacy_bundle()

    if os.path.exists(SECRET_VAULT_FILE):
        try:
            with open(SECRET_VAULT_FILE, "r", encoding="utf-8") as f:
                vault = json.load(f)

            salt_b64 = vault.get("salt", "")
            token = vault.get("token", "")
            if not salt_b64 or not token:
                raise ValueError("Secret vault thiếu salt hoặc token")

            salt = base64.b64decode(salt_b64)
            key = _derive_vault_key(secret_password, salt)
            plain = Fernet(key).decrypt(token.encode("utf-8")).decode("utf-8")
            bundle = json.loads(plain)
            if not isinstance(bundle, dict):
                raise ValueError("Secret vault không chứa JSON object hợp lệ")
            return bundle
        except (InvalidToken, ValueError, json.JSONDecodeError, OSError) as exc:
            print(f"⚠ Không giải mã được secrets.enc: {exc}")
            raise SystemExit("Không mo duoc secrets.enc. Hay kiem tra lai mat khau.")

    return _load_legacy_bundle()


def _sync_local_secret_files(secret_bundle):
    config_defaults = secret_bundle.get("config_default", {})
    if not isinstance(config_defaults, dict):
        config_defaults = {}

    if config_defaults and not os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                json.dump(config_defaults, f, indent=2, ensure_ascii=False)
        except Exception as exc:
            print(f"⚠ Không tạo được config.json local: {exc}")

    client_secret_payload = secret_bundle.get("client_secret_json", "")
    appdata_client_secret = os.path.join(APPDATA_ROOT, "client_secret.json")
    if client_secret_payload and not os.path.exists(appdata_client_secret):
        try:
            try:
                os.chmod(appdata_client_secret, 0o666)
                ctypes.windll.kernel32.SetFileAttributesW(str(appdata_client_secret), 0)
            except Exception:
                pass
            with open(appdata_client_secret, "w", encoding="utf-8") as f:
                if isinstance(client_secret_payload, dict):
                    json.dump(client_secret_payload, f, indent=2, ensure_ascii=False)
                else:
                    f.write(client_secret_payload)
        except Exception as exc:
            print(f"⚠ Không ghi được client_secret.json local: {exc}")

    return config_defaults




#======================

# === FILE cấu hình ===
CONFIG_FILE = os.path.join(APPDATA_ROOT, "config.json")

# Defer heavy secret loading until after the main window appears.
secret_bundle = {}
config_default = {}

# === Đọc file config chính ===
if os.path.exists(CONFIG_FILE):
    with open(CONFIG_FILE, "r", encoding="utf-8") as f:
        config = json.load(f)
else:
    config = {}


def _normalize_lang_code(lang_value):
    lang_clean = (lang_value or "vi").lower().strip()
    if lang_clean.startswith("vi"):
        return "vi"
    if lang_clean.startswith(("zh", "zh-cn", "zh_tw", "zh-hk")):
        return "zh"
    if lang_clean.startswith("ja"):
        return "ja"
    if lang_clean.startswith("en"):
        return "en"
    return "vi"


def _load_google_tts_profiles_from_config():
    stored = config.get("GOOGLE_TTS_PROFILES", {})
    if not isinstance(stored, dict):
        return
    for lang_code, profile in stored.items():
        if not isinstance(profile, dict):
            continue
        normalized = _normalize_lang_code(lang_code)
        GOOGLE_TTS_PROFILES[normalized] = _google_tts_normalize_profile(profile)


def _save_google_tts_profiles_to_config():
    try:
        config["GOOGLE_TTS_PROFILES"] = {
            lang_code: _google_tts_normalize_profile(profile)
            for lang_code, profile in GOOGLE_TTS_PROFILES.items()
        }
        os.makedirs(APPDATA_ROOT, exist_ok=True)
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
    except Exception as exc:
        print(f"⚠ Không lưu được profile Google TTS vào config: {exc}")


def _load_ui_tts_engine_from_config():
    return (config.get("TTS_ENGINE_SELECTED") or "Google Cloud TTS").strip()


def _save_ui_tts_engine_to_config(engine_name):
    try:
        config["TTS_ENGINE_SELECTED"] = (engine_name or "Google Cloud TTS").strip()
        os.makedirs(APPDATA_ROOT, exist_ok=True)
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
    except Exception as exc:
        print(f"⚠ Không lưu được engine TTS vào config: {exc}")


def _write_app_config():
    os.makedirs(APPDATA_ROOT, exist_ok=True)
    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


def _load_hsk30_recent_selection():
    return (
        str(config.get("HSK30_LAST_EXCEL", "") or ""),
        str(config.get("HSK30_LAST_SHEET", "") or ""),
    )


def _save_hsk30_recent_selection(excel_path, sheet_name):
    try:
        config["HSK30_LAST_EXCEL"] = str(excel_path or "")
        config["HSK30_LAST_SHEET"] = str(sheet_name or "")
        _write_app_config()
    except Exception as exc:
        print(f"⚠ Không lưu được Excel/sheet HSK 3.0 gần nhất: {exc}")


def _load_hsk30_builder_state():
    return {
        "version": str(config.get("HSK30_LAST_VERSION", "HSK 3.0") or "HSK 3.0"),
        "level": str(config.get("HSK30_LAST_LEVEL", "HSK 1") or "HSK 1"),
        "pack_version": str(config.get("HSK30_LAST_PACK_VERSION", "1") or "1"),
        "output_dir": str(config.get("HSK30_LAST_OUTPUT_DIR", os.path.join(BASE_DIR, "output")) or os.path.join(BASE_DIR, "output")),
        "bitrate": str(config.get("HSK30_LAST_BITRATE", "32 kbps") or "32 kbps"),
        "audio_mode": str(config.get("HSK30_LAST_AUDIO_MODE", "Đọc tiếng Trung + Tiếng Việt") or "Đọc tiếng Trung + Tiếng Việt"),
    }


def _save_hsk30_builder_state(*, version=None, level=None, pack_version=None, output_dir=None, bitrate=None, audio_mode=None):
    try:
        changed = False
        if version is not None:
            new_value = str(version or "")
            if config.get("HSK30_LAST_VERSION", "") != new_value:
                config["HSK30_LAST_VERSION"] = new_value
                changed = True
        if level is not None:
            new_value = str(level or "")
            if config.get("HSK30_LAST_LEVEL", "") != new_value:
                config["HSK30_LAST_LEVEL"] = new_value
                changed = True
        if pack_version is not None:
            new_value = str(pack_version or "")
            if config.get("HSK30_LAST_PACK_VERSION", "") != new_value:
                config["HSK30_LAST_PACK_VERSION"] = new_value
                changed = True
        if output_dir is not None:
            new_value = str(output_dir or "")
            if config.get("HSK30_LAST_OUTPUT_DIR", "") != new_value:
                config["HSK30_LAST_OUTPUT_DIR"] = new_value
                changed = True
        if bitrate is not None:
            new_value = str(bitrate or "")
            if config.get("HSK30_LAST_BITRATE", "") != new_value:
                config["HSK30_LAST_BITRATE"] = new_value
                changed = True
        if audio_mode is not None:
            new_value = str(audio_mode or "")
            if config.get("HSK30_LAST_AUDIO_MODE", "") != new_value:
                config["HSK30_LAST_AUDIO_MODE"] = new_value
                changed = True
        if changed:
            _write_app_config()
    except Exception as exc:
        print(f"⚠ Không lưu được state HSK 3.0: {exc}")


def _load_secret_bundle_after_ui():
    global secret_bundle, config_default, config
    try:
        if not os.environ.get(SECRET_VAULT_PASSWORD_ENV, "").strip():
            print("ℹ️ Bỏ qua load secrets.enc lúc khởi động để mở UI nhanh hơn.")
            return
        loaded_secret_bundle = _load_secret_bundle()
        loaded_config_default = _sync_local_secret_files(loaded_secret_bundle)
        secret_bundle = loaded_secret_bundle
        config_default = loaded_config_default

        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                config = json.load(f)
        else:
            config = {}

        # Refresh the most common secrets from the loaded local defaults.
        globals().update({
            "SENDER_EMAIL": _env_or(config.get("SENDER_EMAIL", config_default.get("SENDER_EMAIL", "")), "SENDER_EMAIL"),
            "SENDER_NAME": _env_or(config.get("SENDER_NAME", config_default.get("SENDER_NAME", "")), "SENDER_NAME"),
            "APP_PASSWORD": _env_or(config.get("APP_PASSWORD", config_default.get("APP_PASSWORD", "")), "APP_PASSWORD"),
            "AWS_ACCESS_KEY_ID": _env_or(config.get("AWS_ACCESS_KEY_ID", config_default.get("AWS_ACCESS_KEY_ID", "")), "AWS_ACCESS_KEY_ID"),
            "AWS_SECRET_ACCESS_KEY": _env_or(config.get("AWS_SECRET_ACCESS_KEY", config_default.get("AWS_SECRET_ACCESS_KEY", "")), "AWS_SECRET_ACCESS_KEY"),
            "AWS_REGION": _env_or(config.get("AWS_REGION", config_default.get("AWS_REGION", "ap-southeast-2")), "AWS_REGION"),
            "GEMINI_API_KEY": _env_or(config.get("GEMINI_API_KEY", config_default.get("GEMINI_API_KEY", "")).strip(), "GEMINI_API_KEY"),
            "GOOGLE_TTS_API_KEY": _env_or(config.get("GOOGLE_TTS_API_KEY", "").strip(), "GOOGLE_TTS_API_KEY"),
            "DISCORD_WEBHOOK_URL": _env_or(config.get("DISCORD_WEBHOOK_URL", config_default.get("DISCORD_WEBHOOK_URL", "")).strip(), "DISCORD_WEBHOOK_URL"),
        })
        if "GOOGLE_TTS_PROFILES" in config:
            _load_google_tts_profiles_from_config()
    except Exception as exc:
        print(f"⚠ Không nạp được secrets sau khi mở UI: {exc}")


def _guess_hsk30_sheet_from_state(sheet_names, version, level):
    try:
        from pipelines.vocab_zip_builder import SHEET_SELECTIONS, _normalize_sheet_name
    except Exception:
        return ""
    wanted_version = str(version or "").strip()
    wanted_level = str(level or "").strip().lower()
    for sheet_name in sheet_names or []:
        mapped = SHEET_SELECTIONS.get(_normalize_sheet_name(sheet_name))
        if not mapped:
            continue
        mapped_version, mapped_level = mapped
        if mapped_version == wanted_version and mapped_level == wanted_level:
            return sheet_name
    return ""

# === Các hàm validate ví dụ (anh có thể thay bằng logic riêng nếu muốn) ===
#CHECK ĐÚNG ĐỊNH DẠNG EMAIL, SDT
def is_valid(val, check_func):
    try:
        return check_func(val.strip())
    except:
        return False
#mở file aN TOÀN    
def mo_file_an_toan(path):
    if os.path.exists(path):
        open_path_cross_platform(path)
    else:
        messagebox.showwarning("Không tìm thấy file", f"Không tìm thấy file:\n{path}")
#gỡ ẩn file
def unprotect_file(path):
    try:
        # Gỡ ẩn, gỡ readonly
        os.chmod(path, 0o666)
        ctypes.windll.kernel32.SetFileAttributesW(str(path), 0)
    except:
        pass
# Các hàm kiểm tra cấu hình
is_email = lambda e: "@" in e and "." in e
is_token = lambda t: ":" in t
is_chat_id = lambda c: str(c).lstrip("-").isdigit()
is_webhook = lambda u: "api/webhooks/" in u
is_github_models_token = lambda k: k.startswith(("ghp_", "github_pat_")) and len(k) >= 20
is_gemini = lambda k: k.startswith("AIza")
is_weather = lambda k: len(k) > 15

# === Lấy giá trị từ config (ưu tiên file chính, fallback file mặc định) ===
SENDER_EMAIL = config.get("SENDER_EMAIL", config_default.get("SENDER_EMAIL", ""))
SENDER_NAME = config.get("SENDER_NAME", config_default.get("SENDER_NAME", ""))
APP_PASSWORD = config.get("APP_PASSWORD", config_default.get("APP_PASSWORD", ""))
AWS_ACCESS_KEY_ID = config.get("AWS_ACCESS_KEY_ID", config_default.get("AWS_ACCESS_KEY_ID", ""))
AWS_SECRET_ACCESS_KEY = config.get("AWS_SECRET_ACCESS_KEY", config_default.get("AWS_SECRET_ACCESS_KEY", ""))
AWS_REGION = config.get("AWS_REGION", config_default.get("AWS_REGION", "ap-southeast-2"))
GEMINI_API_KEY = _env_or(config.get("GEMINI_API_KEY", "").strip(), "GEMINI_API_KEY")
GOOGLE_TTS_API_KEY = _env_or(config.get("GOOGLE_TTS_API_KEY", "").strip(), "GOOGLE_TTS_API_KEY")
DISCORD_WEBHOOK_URL = _env_or(config.get("DISCORD_WEBHOOK_URL", "").strip(), "DISCORD_WEBHOOK_URL")
GITHUB_MODELS_TOKEN = _env_or(config.get("GITHUB_MODELS_TOKEN", "").strip(), "GITHUB_MODELS_TOKEN")



#Youtube
APPDATA_CLIENT_SECRET = os.path.join(APPDATA_ROOT, "client_secret.json")
if os.path.exists(APPDATA_CLIENT_SECRET):
    print("✅ AppData đã có sẵn client_secret.json.")
else:
    print("⚠️ client_secret.json chưa được nạp vào AppData.")


# Config AWS Polly=====================================
#==============================
# Email cấu hình
#============================


#====================
def tao_file_mp3(text, lang="vi", voice="Mặc định", toc_do="Bình thường",
                 engine="gTTS", file_out="out.mp3"):
    """
    Tạo file mp3 từ text, hỗ trợ gTTS hoặc Amazon Polly.
    - Làm sạch trước khi đọc.
    - Nếu rỗng: tạo 300ms im lặng để không vỡ pipeline ghép file.
    - Bắt riêng lỗi 'No text to speak' từ gTTS.
    """
    from gtts import gTTS
    from pydub import AudioSegment
    import re

    def _export_silence(path, ms=300):
        AudioSegment.silent(duration=ms).export(path, format="mp3", bitrate="192k")
        print(f"⏭️ Dòng rỗng → tạo {ms}ms im lặng: {path}")

    def _get_polly_voice_id(lang_code, voice_value):
        if lang_code == "ja":
            return "Mizuki" if voice_value in ["Nữ", "Female"] else "Takumi"
        if lang_code == "zh":
            if voice_value in ["Nam", "Male"]:
                print("⚠ Polly tiếng Trung không có giọng Nam, dùng Zhiyu (nữ).")
            return "Zhiyu"
        if lang_code == "en":
            return "Joanna" if voice_value in ["Nữ", "Female"] else "Matthew"
        return None

    def _polly_supports_lang(lang_code):
        return lang_code in ["en", "ja", "zh"]

    def _speak_with_polly(text_value, lang_code, voice_value, speed_value, out_path):
        import boto3

        voice_id = _get_polly_voice_id(lang_code, voice_value)
        if not voice_id:
            raise ValueError("Polly không hỗ trợ ngôn ngữ này")

        polly_client = boto3.Session(
            aws_access_key_id=AWS_ACCESS_KEY_ID,
            aws_secret_access_key=AWS_SECRET_ACCESS_KEY,
            region_name=AWS_REGION
        ).client('polly')

        text_pol = re.sub(r"^\d+\s*\.\s*\.\s*", "", text_value)
        text_pol = re.sub(r"(\s*\.\s*){1,}", ' <break time="400ms"/> ', text_pol).strip()

        if speed_value == "Chậm":
            ssml_text = f"<speak><prosody rate='80%'>{text_pol}</prosody></speak>"
        else:
            ssml_text = f"<speak>{text_pol}</speak>"

        response = polly_client.synthesize_speech(
            VoiceId=voice_id,
            OutputFormat='mp3',
            Text=ssml_text,
            TextType="ssml"
        )

        with open(out_path, 'wb') as f:
            f.write(response['AudioStream'].read())

        print(f"✅ Đã tạo file bằng Amazon Polly: {out_path}")

    def _speak_with_gtts(text_value, lang_code, slow_value, out_path):
        tts = gTTS(text=text_value, lang=lang_code, slow=slow_value)
        tts.save(out_path)
        print(f"✅ Đã tạo file bằng gTTS: {out_path}")

    def _speak_with_google(text_value, lang_code, slow_value, out_path):
        tao_file_google_mp3(
            text_value,
            lang=lang_code,
            gender=voice,
            voice_name="",
            toc_do=("Chậm" if slow_value else "Bình thường"),
            file_out=out_path,
        )

    try:
        #slow = (toc_do == "Chậm")
        # Nếu là tiếng Việt thì luôn tốc độ bình thường
        if (lang or "").lower().startswith("vi"):
            slow = False
        else:
            slow = (toc_do == "Chậm")


        # ✅ Làm sạch văn bản (không để None)
        try:
            text = lam_sach_van_ban(text or "")
        except Exception:
            text = str(text or "").strip()

        # ✅ Chuẩn hoá lang
        lang_clean = _normalize_lang_code(lang)

        # ✅ Nếu sau khi làm sạch mà rỗng → tạo im lặng và thoát
        if not text.strip():
            _export_silence(file_out, 300)
            return

        engine_norm = (engine or "gTTS").strip().lower()

        # ================= Amazon Polly =================
        if engine_norm == "polly":
            try:
                _speak_with_polly(text, lang_clean, voice, toc_do, file_out)
                return
            except Exception as e:
                print(f"⚠ Polly lỗi, thử Google Cloud TTS: {e}")
                thong_bao_loi_api(e, "AWS Polly")
                try:
                    _speak_with_google(text, lang_clean, slow, file_out)
                    return
                except Exception as e2:
                    print(f"⚠ Google Cloud TTS lỗi, thử gTTS: {e2}")
                    _speak_with_gtts(text, lang_clean, slow, file_out)
                    return

        # ================= gTTS =================
        if engine_norm == "gtts":
            # gTTS is explicit: do not contact Google Cloud first.
            _speak_with_gtts(text, lang_clean, slow, file_out)
            return

        # ================= Google Cloud TTS =================
        if engine_norm in {"google cloud tts", "google cloud", "google"}:
            try:
                _speak_with_google(text, lang_clean, slow, file_out)
                return
            except Exception as e:
                print(f"⚠ Google Cloud TTS lỗi hoặc không khả dụng, fallback to gTTS: {e}")
                thong_bao_loi_api(e, "Google Cloud TTS")
            try:
                _speak_with_gtts(text, lang_clean, slow, file_out)
                return
            except Exception as e:
                msg = str(e)
                if "No text to speak" in msg or "No text to send to TTS API" in msg:
                    _export_silence(file_out, 300)
                    return
                print(f"⚠ gTTS lỗi, thử Polly thay thế: {e}")
                if _polly_supports_lang(lang_clean):
                    try:
                        _speak_with_polly(text, lang_clean, voice, toc_do, file_out)
                        return
                    except Exception as polly_error:
                        print(f"❌ Polly cũng lỗi sau khi gTTS lỗi: {polly_error}")
                        raise
                print(f"⚠ Không fallback sang Polly vì ngôn ngữ '{lang_clean}' không được hỗ trợ.")
                raise

        print(f"⚠ Engine '{engine}' không được nhận diện, dùng Google Cloud TTS rồi gTTS.")
        try:
            _speak_with_google(text, lang_clean, slow, file_out)
            return
        except Exception as e:
            print(f"⚠ Google Cloud TTS lỗi, fallback to gTTS: {e}")
        _speak_with_gtts(text, lang_clean, slow, file_out)
        return

    except Exception as e:
        print("❌ Lỗi tạo file mp3:", e)
        thong_bao_loi_api(e, "TTS")
        raise


def _google_tts_lang_code(lang_value):
    lang_clean = (lang_value or "vi").lower().strip()
    if lang_clean.startswith("vi"):
        return "vi-VN"
    if lang_clean.startswith("en"):
        return "en-US"
    if lang_clean.startswith("ja"):
        return "ja-JP"
    if lang_clean.startswith(("zh", "zh-cn", "zh_tw", "zh-hk")):
        return "cmn-CN"
    return "vi-VN"


def _google_tts_gender_enum(label):
    try:
        from google.cloud import texttospeech
    except Exception:
        raise RuntimeError("google-cloud-texttospeech not installed")

    mapping = {
        "Nam": texttospeech.SsmlVoiceGender.MALE,
        "Nữ": texttospeech.SsmlVoiceGender.FEMALE,
        "Trung tính": texttospeech.SsmlVoiceGender.NEUTRAL,
        "Mặc định": texttospeech.SsmlVoiceGender.SSML_VOICE_GENDER_UNSPECIFIED,
        None: texttospeech.SsmlVoiceGender.SSML_VOICE_GENDER_UNSPECIFIED,
        "": texttospeech.SsmlVoiceGender.SSML_VOICE_GENDER_UNSPECIFIED,
    }
    return mapping.get(label, texttospeech.SsmlVoiceGender.SSML_VOICE_GENDER_UNSPECIFIED)


GOOGLE_TTS_SLOT_LABELS = ("Mặc định", "Nam", "Nữ", "Trung tính")


def _google_tts_default_profile():
    return {
        "gender": "Mặc định",
        "voice_name": "",
        "slots": {label: {"gender": label, "voice_name": ""} for label in GOOGLE_TTS_SLOT_LABELS},
    }


def _google_tts_slot_label(label):
    normalized = (label or "Mặc định").strip()
    return normalized if normalized in GOOGLE_TTS_SLOT_LABELS else "Mặc định"


def _google_tts_normalize_profile(profile):
    normalized = _google_tts_default_profile()
    if not isinstance(profile, dict):
        return normalized

    raw_slots = profile.get("slots", {})
    if isinstance(raw_slots, dict):
        for slot_label in GOOGLE_TTS_SLOT_LABELS:
            slot_profile = raw_slots.get(slot_label, {})
            if isinstance(slot_profile, dict):
                normalized["slots"][slot_label] = {
                    "gender": _google_tts_slot_label(slot_profile.get("gender", slot_label)),
                    "voice_name": slot_profile.get("voice_name", "") or "",
                }

    legacy_gender = _google_tts_slot_label(profile.get("gender", "Mặc định"))
    legacy_voice_name = profile.get("voice_name", "") or ""
    if legacy_voice_name:
        normalized["slots"][legacy_gender] = {
            "gender": legacy_gender,
            "voice_name": legacy_voice_name,
        }
        normalized["gender"] = legacy_gender
        normalized["voice_name"] = legacy_voice_name
    else:
        normalized["gender"] = legacy_gender
        normalized["voice_name"] = ""

    if normalized["slots"]["Mặc định"]["voice_name"] == "" and normalized["gender"] == "Mặc định":
        normalized["slots"]["Mặc định"]["gender"] = "Mặc định"

    return normalized


def _google_tts_get_profile(lang_code):
    lang_clean = _normalize_lang_code(lang_code)
    profile = GOOGLE_TTS_PROFILES.setdefault(lang_clean, _google_tts_default_profile())
    normalized = _google_tts_normalize_profile(profile)
    GOOGLE_TTS_PROFILES[lang_clean] = normalized
    return normalized


def _google_tts_get_profile_slot(lang_code, gender="Mặc định"):
    profile = _google_tts_get_profile(lang_code)
    slot_label = _google_tts_slot_label(gender)
    slot = profile.setdefault("slots", {}).setdefault(slot_label, {"gender": slot_label, "voice_name": ""})
    slot["gender"] = _google_tts_slot_label(slot.get("gender", slot_label))
    slot["voice_name"] = slot.get("voice_name", "") or ""
    return slot


def _google_tts_set_profile(lang_code, gender="Mặc định", voice_name=""):
    profile = _google_tts_get_profile(lang_code)
    slot_label = _google_tts_slot_label(gender)
    slot = profile.setdefault("slots", {}).setdefault(slot_label, {"gender": slot_label, "voice_name": ""})
    slot["gender"] = slot_label
    slot["voice_name"] = voice_name or ""
    profile["gender"] = slot_label
    profile["voice_name"] = voice_name or ""
    return profile


def _google_tts_resolve_selection(lang_code, requested_gender="Mặc định"):
    profile = _google_tts_get_profile(lang_code)
    requested_label = _google_tts_slot_label(requested_gender)
    slots = profile.get("slots", {})
    slot = slots.get(requested_label, {})
    voice_name = (slot.get("voice_name") or "").strip()
    resolved_gender = _google_tts_slot_label(slot.get("gender", requested_label))

    if not voice_name and requested_label == "Mặc định":
        voice_name = (profile.get("voice_name") or "").strip()
        resolved_gender = _google_tts_slot_label(profile.get("gender", "Mặc định"))

    return {
        "gender": resolved_gender or requested_label,
        "voice_name": voice_name,
    }


def _google_tts_client(texttospeech_module):
    """Create a Cloud TTS client using the dedicated API key when configured."""
    if GOOGLE_TTS_API_KEY:
        from google.api_core.client_options import ClientOptions
        return texttospeech_module.TextToSpeechClient(
            client_options=ClientOptions(api_key=GOOGLE_TTS_API_KEY)
        )
    return texttospeech_module.TextToSpeechClient()


def _google_tts_list_voices(lang_code, gender=None):
    try:
        from google.cloud import texttospeech
    except Exception:
        raise RuntimeError("google-cloud-texttospeech not installed")

    client = _google_tts_client(texttospeech)
    voices = client.list_voices(language_code=_google_tts_lang_code(lang_code)).voices
    results = []
    for voice in voices:
        voice_gender = getattr(voice.ssml_gender, "name", str(voice.ssml_gender))
        if gender and gender not in ["Mặc định", ""]:
            if gender == "Nam" and voice_gender != "MALE":
                continue
            if gender == "Nữ" and voice_gender != "FEMALE":
                continue
            if gender == "Trung tính" and voice_gender not in ["NEUTRAL", "SSML_VOICE_GENDER_UNSPECIFIED"]:
                continue
        results.append({
            "name": voice.name,
            "gender": voice_gender,
            "display": f"{voice.name} · {voice_gender}",
        })
    return results


def tao_file_google_mp3(text, lang="vi", gender="Mặc định", voice_name="", toc_do="Bình thường", file_out="out.mp3"):
    try:
        from google.cloud import texttospeech
    except Exception:
        raise RuntimeError("google-cloud-texttospeech not installed")

    text = lam_sach_van_ban(text or "")
    if not text.strip():
        from pydub import AudioSegment
        AudioSegment.silent(duration=300).export(file_out, format="mp3", bitrate="192k")
        print(f"⏭️ Dòng rỗng → tạo im lặng: {file_out}")
        return

    requested_gender = _google_tts_slot_label(gender)
    profile_selection = _google_tts_resolve_selection(lang, requested_gender)
    profile_gender = (profile_selection.get("gender") or requested_gender).strip()
    profile_voice_name = (profile_selection.get("voice_name") or "").strip()
    voice_name = (voice_name or profile_voice_name or "").strip()
    gender = requested_gender if requested_gender else "Mặc định"
    if not voice_name:
        gender = profile_gender if profile_gender and profile_gender != "Mặc định" else gender
    lang_code = _google_tts_lang_code(lang)

    client = _google_tts_client(texttospeech)
    synthesis_input = texttospeech.SynthesisInput(text=text)

    voice_kwargs = {"language_code": lang_code}
    if voice_name:
        voice_kwargs["name"] = voice_name
    else:
        gender_enum = _google_tts_gender_enum(gender)
        if gender_enum != texttospeech.SsmlVoiceGender.SSML_VOICE_GENDER_UNSPECIFIED:
            voice_kwargs["ssml_gender"] = gender_enum

    audio_kwargs = {"audio_encoding": texttospeech.AudioEncoding.MP3}
    if toc_do == "Chậm":
        audio_kwargs["speaking_rate"] = 0.8

    response = client.synthesize_speech(
        input=synthesis_input,
        voice=texttospeech.VoiceSelectionParams(**voice_kwargs),
        audio_config=texttospeech.AudioConfig(**audio_kwargs),
    )

    with open(file_out, "wb") as f:
        f.write(response.audio_content)

    print(f"✅ Đã tạo file bằng Google Cloud TTS: {file_out}")

#=========================
    

#=====


def thong_bao_loi_cauhinh(loai, chi_tiet="", hien_popup=True):
    r"""
    Hiển thị lỗi cấu hình + ghi log vào AppData\loi_cauhinh.txt
    """
    try:
        from datetime import datetime
        os.makedirs(APPDATA_ROOT, exist_ok=True)
        log_path = os.path.join(APPDATA_ROOT, "loi_cauhinh.txt")
        unprotect_file(log_path)
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] ❌ {loai} lỗi: {chi_tiet}\n")

        print(f"🛑 Lỗi {loai}: {chi_tiet}")
        if hien_popup:
            messagebox.showerror("Lỗi cấu hình", f"{loai} bị lỗi:\n{chi_tiet}")
    except Exception as e:
        print("‼ Không thể ghi log lỗi:", e)

#==================

# === GitHub Models ===
# GitHub Models exposes an OpenAI-compatible chat-completions endpoint.  This is
# intentionally separate from the old OpenAI key so an `sk-...` key is never
# sent to GitHub by accident.
GITHUB_MODELS_BASE_URL = "https://models.github.ai/inference"
GITHUB_MODELS_CATALOG_URL = "https://models.github.ai/catalog/models"
GITHUB_MODELS_API_VERSION = "2022-11-28"
GITHUB_MODELS_MODEL = "openai/gpt-4.1"


def _create_github_models_client(token):
    """Return a client only after a token has been configured.

    Recent OpenAI SDK versions reject an empty api_key during construction, so
    the app must remain usable before the user opens Settings for the first
    time.
    """
    if not token:
        return None
    try:
        return OpenAI(base_url=GITHUB_MODELS_BASE_URL, api_key=token)
    except ImportError as exc:
        print(f"⚠ Không khởi tạo được GitHub Models client: {exc}")
        return None


if not is_valid(GITHUB_MODELS_TOKEN, is_github_models_token):
    thong_bao_loi_cauhinh(
        "GitHub Models token",
        "Sai định dạng hoặc rỗng. Hãy nhập GitHub PAT (ghp_... hoặc github_pat_...).",
        hien_popup=False,
    )
    GITHUB_MODELS_TOKEN = ""

if not is_valid(GEMINI_API_KEY, is_gemini):
    thong_bao_loi_cauhinh("Gemini API Key", "Sai định dạng hoặc rỗng → dùng mặc định.", hien_popup=False)
    GEMINI_API_KEY = config_default.get("GEMINI_API_KEY", "")

# Google Cloud Text-to-Speech dùng key riêng với key Gemini.
if not is_valid(DISCORD_WEBHOOK_URL, is_webhook):
    thong_bao_loi_cauhinh("Discord Webhook", "Không chứa /api/webhooks → dùng mặc định.", hien_popup=False)
    DISCORD_WEBHOOK_URL = config_default.get("DISCORD_WEBHOOK_URL", "")

#===========================BIẾN TOÀN CỤC ============================================================BIẾN TOÀN CỤC===============
client = _create_github_models_client(GITHUB_MODELS_TOKEN)
#Biến điều khiển đọc âm thanh
doc_thread = None
dang_doc = False
file_am_thanh = ""
noi_dung_cuoi = ""
dang_tam_dung = False
dung_doc_ngay = False
tam_dung = False  # biến toàn cục
tam_dung_doc = False
chup_anh_luu = None

#Kênh âm thanh bằng pygame
current_seek_pos = 0  # vị trí phát hiện tại (giây)
seek_start_time = 0   # thời điểm bắt đầu phát, tính bằng time.time()
GOOGLE_TTS_LANG_MAP = {
    "vi": "vi-VN",
    "en": "en-US",
    "ja": "ja-JP",
    "zh": "cmn-CN",
}
GOOGLE_TTS_GENDER_LABELS = ["Mặc định", "Nam", "Nữ", "Trung tính"]
GOOGLE_TTS_PROFILES = {
    "vi": {"gender": "Mặc định", "voice_name": ""},
    "en": {"gender": "Mặc định", "voice_name": ""},
    "ja": {"gender": "Mặc định", "voice_name": ""},
    "zh": {"gender": "Mặc định", "voice_name": ""},
}

_load_google_tts_profiles_from_config()
#==================================================
###====CÁC FILE KHÁC HỆ THỐNG====
#Sound , font các biến âm thanh=============================================================================

def pick_existing_asset(*names):
    for name in names:
        candidate = os.path.join(SOUNDS_DIR, name)
        if os.path.exists(candidate):
            return candidate
    for name in names:
        candidate = os.path.join(BASE_DIR, name)
        if os.path.exists(candidate):
            return candidate
    # Fallback: check assets_legacy for legacy sound files
    for name in names:
        candidate = os.path.join(BASE_DIR, "assets_legacy", name)
        if os.path.exists(candidate):
            return candidate
    # If not found, print warning and return None for graceful fallback
    print(f"⚠️  Asset not found: {', '.join(names)}")
    return None


def pick_existing_executable(*names):
    for name in names:
        candidate = os.path.join(BASE_DIR, name)
        if os.path.isfile(candidate) and os.access(candidate, os.X_OK):
            return candidate
    return None

CLICK_SOUND = pick_existing_asset("click.ogg", "click.wav")
WARNING_SOUND = pick_existing_asset("warning.ogg", "warning.wav")
SUCCESS_SOUND = pick_existing_asset("Success.ogg", "Success.wav", "Success1.wav")
TRUE_SOUND = pick_existing_asset("True.ogg", "True.wav", "True1.wav")
FALSE_SOUND = pick_existing_asset("False.ogg", "False.wav", "False1.wav")
WIN_SOUND = pick_existing_asset("Win.ogg", "Win.wav", "Win1.wav")
#background_music
BACKGROUND_MUSIC = pick_existing_asset("background_music.ogg", "background_music.mp3", "background_music.MP3")

#============font chữ đa ngôn ngữ

FONT_LATIN = os.path.join(FONTS_DIR, "DejaVuSans.ttf")
FONT_CJK = os.path.join(FONTS_DIR, "SamsungGothicKorean.ttf")

#=======================================

# Tài nguyên tĩnh
##LOGO_PATH = os.path.join(IMAGES_DIR, "logo.jpg")
##if not os.path.exists(LOGO_PATH):
##    LOGO_PATH = os.path.join(IMAGES_DIR, "logo.png")
ICON_ICO_PATH = os.path.join(IMAGES_DIR, "logo.ico")

#==================

# ==== Thiết lập ffmpeg cho pydub xuất mp3 ====


FOLDER = os.path.dirname(os.path.abspath(sys.argv[0]))

FFMPEG_DOWNLOAD_URL = "https://www.gyan.dev/ffmpeg/builds/ffmpeg-release-essentials.zip"

def resolve_ffmpeg_paths():
    bundle_folder = os.path.join(FOLDER, "ffmpeg_bin", "bin")
    bundle_ffmpeg = pick_existing_executable(os.path.join("ffmpeg_bin", "bin", "ffmpeg"), os.path.join("ffmpeg_bin", "bin", "ffmpeg.exe"))
    bundle_ffprobe = pick_existing_executable(os.path.join("ffmpeg_bin", "bin", "ffprobe"), os.path.join("ffmpeg_bin", "bin", "ffprobe.exe"))
    if bundle_ffmpeg:
        return bundle_ffmpeg, bundle_ffprobe or bundle_ffmpeg, bundle_folder

    system_ffmpeg = shutil.which("ffmpeg")
    system_ffprobe = shutil.which("ffprobe")
    if system_ffmpeg:
        return system_ffmpeg, system_ffprobe or system_ffmpeg, None

    cache_root = os.path.join(APPDATA_ROOT, "ffmpeg_bin")
    cache_folder = os.path.join(cache_root, "bin")
    cache_ffmpeg = pick_existing_executable("ffmpeg", "ffmpeg.exe")
    cache_ffprobe = pick_existing_executable("ffprobe", "ffprobe.exe")
    if cache_ffmpeg:
        return cache_ffmpeg, cache_ffprobe or cache_ffmpeg, cache_folder

    def find_ffmpeg_folder(search_root):
        for current_root, _, files in os.walk(search_root):
            if "ffmpeg" in files or "ffmpeg.exe" in files:
                ffmpeg_name = "ffmpeg" if "ffmpeg" in files else "ffmpeg.exe"
                ffprobe_name = "ffprobe" if "ffprobe" in files else "ffprobe.exe"
                ffmpeg_path = os.path.join(current_root, ffmpeg_name)
                ffprobe_path = os.path.join(current_root, ffprobe_name)
                return ffmpeg_path, ffprobe_path if os.path.isfile(ffprobe_path) else ffmpeg_path, current_root
        return None, None, None

    try:
        import zipfile
        from urllib.request import urlopen

        os.makedirs(cache_root, exist_ok=True)
        zip_path = os.path.join(cache_root, "ffmpeg_release.zip")
        print("ℹ️ Không thấy ffmpeg trong PATH, đang tải bản portable về AppData...")

        with urlopen(FFMPEG_DOWNLOAD_URL, timeout=120) as response, open(zip_path, "wb") as output_file:
            shutil.copyfileobj(response, output_file)

        with zipfile.ZipFile(zip_path, "r") as archive:
            archive.extractall(cache_root)

        try:
            os.remove(zip_path)
        except:
            pass

        downloaded = find_ffmpeg_folder(cache_root)
        if downloaded[0]:
            return downloaded
    except Exception as e:
        print(f"⚠️ Tải ffmpeg tự động thất bại: {e}")

    return None, None, None


FFMPEG_PATH, ffprobe_path, FFMPEG_FOLDER = resolve_ffmpeg_paths()
M4A_VOICE_BITRATE = "48k"
M4A_VOICE_FALLBACK_BITRATE = "40k"
M4A_VOICE_SAMPLE_RATE = "22050"

if FFMPEG_PATH:
    AudioSegment.converter = FFMPEG_PATH
    AudioSegment.ffmpeg = FFMPEG_PATH
    AudioSegment.ffprobe = ffprobe_path
#=======Gửi discor chung
def gui_discord_thong_bao(msg=""):
    try:
        if not DISCORD_WEBHOOK_URL:
            raise ValueError("Webhook URL trống")

        import requests
        requests.post(DISCORD_WEBHOOK_URL, json={"content": msg})
    except Exception as e:
        print("Lỗi gửi Discord:", e)
        thong_bao_loi_api(e, "Discord")
#==========Set logo icon cho toàn thông báo app
def set_popup_icon(win):
    try:
        win.iconbitmap(ICON_ICO_PATH)
    except Exception:
        pass

#=======
#===Nhạc nền:

def play_background_music():
    """Phát nhạc nền lặp vô hạn"""
    try:
        pygame.mixer.music.load(BACKGROUND_MUSIC)
        pygame.mixer.music.play(-1)  # -1: loop vô hạn
        print("🎵 Đang phát nhạc nền...")
    except Exception as e:
        print(f"⚠️ Lỗi phát nhạc nền: {e}")

def stop_background_music():
    """Dừng nhạc nền"""
    try:
        pygame.mixer.music.stop()
        print("🛑 Đã dừng nhạc nền.")
    except Exception as e:
        print(f"⚠️ Lỗi dừng nhạc nền: {e}")


#=======Mật khẩu , phím mini, gui mail khi quên
def ask_password_with_keyboard(callback):
    pw_window = tk.Toplevel(root)
    set_popup_icon(pw_window)
    pw_window.title("Nhập mật khẩu")
    pw_window.geometry("400x440")
    pw_window.resizable(False, False)
    pw_window.grab_set()
    pw_window.attributes('-topmost', True)

    tk.Label(pw_window, text="Nhập mật khẩu:", font=("Arial", 20)).pack(pady=10)
    entry_pw = tk.Entry(pw_window, show="*", font=("Arial", 16))
    entry_pw.pack()

    def insert_char(ch): entry_pw.insert(tk.END, ch)
    def delete_last(): entry_pw.delete(len(entry_pw.get()) - 1, tk.END)
    def clear_entry(): entry_pw.delete(0, tk.END)

    def submit():
        password = entry_pw.get()
        pw_window.destroy()
        try:
            with open(os.path.join(APPDATA_ROOT, "pass.txt"), "r") as f:
                mat_khau_dung = f.read().strip()
        except:
            mat_khau_dung = "1234"

        if password == mat_khau_dung:
            callback(True)
        else:
            messagebox.showerror("Sai mật khẩu", "❌ Mật khẩu không đúng.")
            callback(False)

    tk.Button(pw_window, text="OK", width=10, command=submit).pack(pady=5)

    # === Bàn phím mini ===
    frm_kb = tk.Frame(pw_window)
    frm_kb.pack()
    btns = [['1', '2', '3'], ['4', '5', '6'], ['7', '8', '9'], ['←', '0', 'C']]
    for r, row in enumerate(btns):
        for c, char in enumerate(row):
            if char == '←':
                cmd = delete_last
            elif char == 'C':
                cmd = clear_entry
            else:
                cmd = lambda ch=char: insert_char(ch)
            tk.Button(frm_kb, text=char, width=5, height=2, command=cmd).grid(row=r, column=c, padx=2, pady=2)

    # === Quên mật khẩu ===
    def quen_mat_khau():
     
        def send_password_request():
            phone = entry_phone.get().strip()
            email_user = entry_email.get().strip()
            diachi = entry_diachi.get().strip()

            if not phone.isdigit() or len(phone) < 9:
                messagebox.showwarning("Lỗi", "Số điện thoại không hợp lệ.", parent=popup)
                return
            if "@" not in email_user or "." not in email_user:
                messagebox.showwarning("Lỗi", "Email không hợp lệ.", parent=popup)
                return
            if not diachi:
                messagebox.showwarning("Lỗi", "Vui lòng nhập địa chỉ liên hệ.", parent=popup)
                return

            try:
                with open(os.path.join(APPDATA_ROOT, "pass.txt"), "r") as f:
                    current_pw = f.read().strip()
            except:
                current_pw = "1234"

            import smtplib
            from email.message import EmailMessage

            msg1 = EmailMessage()
            msg1["Subject"] = "🔐 Mật khẩu Máy Học Tập"
            msg1["From"] = f"{SENDER_NAME} <{SENDER_EMAIL}>"
            msg1["To"] = email_user
            msg1.set_content(
                f"Chào bạn,\n\nMật khẩu hiện tại của bạn là:\n🔑 {current_pw}\n\n"
                f"Vui lòng bảo mật thông tin này.\n\n-- Máy Học Tập --"
            )

            msg2 = EmailMessage()
            msg2["Subject"] = "📩 Yêu cầu quên mật khẩu"
            msg2["From"] = f"{SENDER_NAME} <{SENDER_EMAIL}>"
            msg2["To"] = SENDER_EMAIL

            noi_dung = (
                f"📝 YÊU CẦU QUÊN MẬT KHẨU:\n\n"
                f"📞 Số điện thoại: {phone}\n"
                f"📧 Email người dùng: {email_user}\n"
                f"🏠 Địa chỉ liên hệ: {diachi}\n"
                f"🔐 Mật khẩu hiện tại: {current_pw}"
            )

            msg2.set_content(noi_dung + "\n\n" + get_system_info())

            with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
                smtp.login(SENDER_EMAIL, APP_PASSWORD)
                smtp.send_message(msg1)
                smtp.send_message(msg2)

            messagebox.showinfo("Thành công", f"📧 Đã gửi mật khẩu tới {email_user}", parent=popup)
            popup.destroy()

        # ========== Cửa sổ to đẹp giữa màn hình ==========
        popup = tk.Toplevel(pw_window)
        set_popup_icon(popup)
        popup.title("🔐 Quên mật khẩu?")
        popup.geometry("550x360")
        popup.resizable(False, False)
        popup.transient(pw_window)
        popup.grab_set()
        popup.attributes("-topmost", True)

        popup.update_idletasks()
        w, h = 550, 360
        x = root.winfo_x() + (root.winfo_width() - w) // 2
        y = root.winfo_y() + (root.winfo_height() - h) // 2
        popup.geometry(f"{w}x{h}+{x}+{y}")

        tk.Label(popup, text="Vui lòng điền thông tin khôi phục:", font=("Arial", 13, "bold")).pack(pady=12)

        def add_input(label_text):
            frame = tk.Frame(popup)
            frame.pack(pady=6)
            tk.Label(frame, text=label_text + ":", font=("Arial", 11), width=18, anchor="e").pack(side="left", padx=5)
            entry = tk.Entry(frame, font=("Arial", 13), width=35)
            entry.pack(side="left", ipady=6)
            return entry

        entry_phone = add_input("📞 Số điện thoại")
        entry_email = add_input("📧 Email của bạn")
        entry_diachi = add_input("🏠 Địa chỉ liên hệ")

        tk.Button(
            popup, text="📤 Gửi yêu cầu", font=("Arial", 12, "bold"),
            bg="#4caf50", fg="white", width=18, height=1,
            command=send_password_request
        ).pack(pady=20)



    tk.Button(pw_window, text="❓ Quên mật khẩu", fg="blue", command=quen_mat_khau).pack(pady=5)

def doi_mat_khau():
    def thuc_hien(ok):
        if not ok:
            return

        win = tk.Toplevel(root)
        set_popup_icon(win)
        win.title("🔐 Đổi mật khẩu ứng dụng")
        win.geometry("650x600")
        win.grab_set()
        win.resizable(False, False)

        tk.Label(win, text="🔒 Nhập mật khẩu cũ:", font=("Arial", 11)).pack(pady=(10, 2))
        entry_old = tk.Entry(win, show="*", font=("Arial", 13), width=35)
        entry_old.pack()

        tk.Label(win, text="🔐 Mật khẩu mới:", font=("Arial", 11)).pack(pady=(8, 2))
        entry1 = tk.Entry(win, show="*", font=("Arial", 13), width=35)
        entry1.pack()

        tk.Label(win, text="🔁 Nhập lại mật khẩu:", font=("Arial", 11)).pack(pady=(8, 2))
        entry2 = tk.Entry(win, show="*", font=("Arial", 13), width=35)
        entry2.pack()

        tk.Label(win, text="📧 Email của bạn:", font=("Arial", 11)).pack(pady=(8, 2))
        entry_email = tk.Entry(win, font=("Arial", 13), width=35)
        entry_email.pack()

        tk.Label(win, text="📞 Số điện thoại:", font=("Arial", 11)).pack(pady=(8, 2))
        entry_sdt = tk.Entry(win, font=("Arial", 13), width=35)
        entry_sdt.pack()

        def insert(ch):
            for e in [entry_old, entry1, entry2]:
                if e.focus_get() == e:
                    e.insert(tk.END, ch)

        def backspace():
            for e in [entry_old, entry1, entry2]:
                if e.focus_get() == e:
                    e.delete(len(e.get()) - 1, tk.END)

        def clear_all():
            for e in [entry_old, entry1, entry2]:
                if e.focus_get() == e:
                    e.delete(0, tk.END)

        frm_kb = tk.Frame(win)
        frm_kb.pack(pady=8)
        keys = [['1', '2', '3'], ['4', '5', '6'], ['7', '8', '9'], ['←', '0', 'C']]
        for r, row in enumerate(keys):
            for c, ch in enumerate(row):
                cmd = backspace if ch == '←' else (clear_all if ch == 'C' else lambda x=ch: insert(x))
                tk.Button(frm_kb, text=ch, width=5, height=2, command=cmd).grid(row=r, column=c, padx=2, pady=2)

        def xac_nhan():
            pass_path = os.path.join(APPDATA_ROOT, "pass.txt")

            # Nếu chưa có file pass.txt → tự tạo với mật khẩu mặc định "1234"
            if not os.path.exists(pass_path):
                try:
                    unprotect_file(pass_path)
                    with open(pass_path, "w", encoding="utf-8") as f:
                        f.write("1234")
                except Exception as e:
                    messagebox.showerror("Lỗi", f"Không thể tạo file pass.txt: {e}", parent=win)
                    return

            try:
                with open(pass_path, "r", encoding="utf-8") as f:
                    mk_cu = f.read().strip()
            except:
                mk_cu = "1234"

            old = entry_old.get().strip()
            new1 = entry1.get().strip()
            new2 = entry2.get().strip()
            email = entry_email.get().strip()
            sdt = entry_sdt.get().strip()

            if not all([old, new1, new2, email, sdt]):
                messagebox.showwarning("Thiếu", "Vui lòng điền đầy đủ thông tin.", parent=win)
                return
            if old != mk_cu:
                messagebox.showerror("Sai mật khẩu", "❌ Mật khẩu cũ không đúng.", parent=win)
                return
            if new1 != new2:
                messagebox.showerror("Không khớp", "❌ Hai mật khẩu mới không khớp.", parent=win)
                return
            if "@" not in email or "." not in email:
                messagebox.showerror("Email không hợp lệ", "Vui lòng nhập đúng định dạng email.", parent=win)
                return

            # ✅ Ghi mật khẩu mới vào pass.txt
            try:
                with open(pass_path, "w", encoding="utf-8") as f:
                    f.write(new1)
            except Exception as e:
                messagebox.showerror("Lỗi ghi file", f"Không thể lưu mật khẩu mới: {e}", parent=win)
                return

            try:
                import smtplib
                from email.message import EmailMessage

                msg_user = EmailMessage()
                msg_user['Subject'] = "🔐 Mật khẩu mới của bạn"
                msg_user['From'] = SENDER_EMAIL
                msg_user['To'] = email
                msg_user.set_content(f"""
Chào bạn,

Bạn vừa đổi mật khẩu thành công ứng dụng text to mp3 đa ngôn ngữ.

🔐 Mật khẩu mới: {new1}
📞 Số điện thoại: {sdt}

Cảm ơn bạn đã sử dụng ứng dụng.
""")

                with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
                    smtp.login(SENDER_EMAIL, APP_PASSWORD)
                    smtp.send_message(msg_user)

                messagebox.showinfo("Thành công", f"📧 Đã gửi mật khẩu mới tới {email}", parent=win)
                win.destroy()

            except Exception as e:
                thong_bao_loi_api(e, "Email")
                messagebox.showerror("Lỗi gửi mail", f"Không thể gửi email:\n{e}", parent=win)

        # ✅ Thêm nút Đổi mật khẩu đầy đủ
        tk.Button(win, text="✅ Đổi mật khẩu", bg="green", fg="white", font=("Arial", 12, "bold"), command=xac_nhan).pack(pady=12)

    ask_password_with_keyboard(thuc_hien)

#=======================NHÓM  - HỎI A.I , ĐỌC ĐỀ, XUẤT MP3====
#hàm hỏi A.I qua GitHub Models
def goi_github_models_cau_hoi(prompt):
    try:
        if client is None:
            raise RuntimeError(
                "GitHub Models API key chưa được cấu hình. "
                "Mở Cài đặt → Sửa GitHub Models Token để nhập GitHub PAT."
            )
        response = client.chat.completions.create(
            model=GITHUB_MODELS_MODEL,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        thong_bao_loi_api(e, "GitHub Models")
        return f"Lỗi GitHub Models: {e}"
##có chọn ngôn ngữ
def goi_gemini_cau_hoi(prompt):
    from langdetect import detect
    import re

    LANGUAGE_KEYWORDS = {
        "tiếng Trung": "zh",
        "tiếng Nhật": "ja",
        "tiếng Anh": "en",
        "tiếng Việt": "vi"
    }

    def get_expected_lang(text):
        for key, lang_code in LANGUAGE_KEYWORDS.items():
            if key.lower() in text.lower():
                return lang_code
        return "vi"

    def contains_pinyin_or_english(text):
        return bool(re.search(r"[a-zA-Z]{2,}", text))

    try:
        client = genai.Client(api_key=GEMINI_API_KEY)

        expected_lang = get_expected_lang(prompt)
        original_lang = detect(prompt)
        response = client.models.generate_content(
            model="gemini-1.5-flash",
            contents=prompt,
        )
        text = (response.text or "").strip()

        # Nếu có pinyin/English hoặc sai ngôn ngữ → hỏi lại
        if detect(text) != expected_lang or contains_pinyin_or_english(text):
            prompt_fix = f"""
Chỉ trả lời bằng đúng ngôn ngữ '{expected_lang}', không thêm phiên âm, không thêm chú giải hay dịch. Trả lời đơn giản, chỉ liệt kê.
Nội dung yêu cầu là:
{prompt}
"""
            text = (client.models.generate_content(
                model="gemini-1.5-flash",
                contents=prompt_fix,
            ).text or "").strip()

        # Nếu prompt gốc là tiếng Việt → chèn dòng tiêu đề tiếng Việt lên đầu
        if original_lang == "vi":
            text = f"{prompt}:\n\n{text}"

        return text

    except Exception as e:
        thong_bao_loi_api(e, "Gemini")
        return f"Lỗi Gemini: {e}"



def gui_hoi_github_models():
    cau_hoi = entry_cau_hoi.get().strip()
    if not cau_hoi:
        messagebox.showwarning("Chưa nhập câu hỏi", "Hãy nhập nội dung để hỏi.")
        return

    txt_de.delete("1.0", tk.END)
    txt_de.insert(tk.END, "⏳ Đang gửi tới GitHub Models ...\n")

    def call():
        try:
            ket_qua = goi_github_models_cau_hoi(cau_hoi)
            if any(x in ket_qua.lower() for x in ["github models", "quota", "error", "api key", "invalid", "401", "403"]):
                raise Exception("GitHub Models có thể đang bận")

            txt_de.delete("1.0", tk.END)
            txt_de.insert(tk.END, ket_qua)
        except:
            txt_de.delete("1.0", tk.END)
            txt_de.insert(tk.END, "⚠ GitHub Models có thể đang bận, chuyển sang Gemini ...\n⏳ Đang hỏi Gemini ...\n")
            ket_qua = goi_gemini_cau_hoi(cau_hoi)
            txt_de.delete("1.0", tk.END)
            txt_de.insert(tk.END, ket_qua)

    threading.Thread(target=call, daemon=True).start()
def gui_hoi_gemini(on_done=None):
    cau_hoi = entry_cau_hoi.get().strip()
    if not cau_hoi:
        messagebox.showwarning("Chưa nhập câu hỏi", "Hãy nhập nội dung để hỏi cho các nhà thông thái!")
        return
    txt_de.delete("1.0", tk.END)
    txt_de.insert(tk.END, "⏳ Đang hỏi mẹ ...")
    def call():
        ket_qua = goi_gemini_cau_hoi(cau_hoi)
        txt_de.delete("1.0", tk.END)
        txt_de.insert(tk.END, ket_qua)
        if on_done:
            on_done()
    threading.Thread(target=call, daemon=True).start()
#=========================
def nhap_giong_noi_advanced(entry_cau_hoi, on_text_got=None, root=None):
    
    import os, sys, threading

    wave_running = {"on": True}
    record_flag = {"is_recording": True}

    try:
        global dung_doc_ngay
        dung_doc_ngay = True
        if "pygame" in sys.modules:
            import pygame
            if pygame.mixer.get_init():
                pygame.mixer.music.stop()
    except Exception:
        pass

    def show_countdown_and_record():
        popup = tk.Toplevel(root)
        set_popup_icon(popup)
        popup.title("Ghi âm")

        # 👉 Căn giữa màn hình và cao lên
        screen_w = popup.winfo_screenwidth()
        screen_h = popup.winfo_screenheight()
        win_w, win_h = 370, 200
        x = (screen_w - win_w) // 2
        y = (screen_h - win_h) // 3
        popup.geometry(f"{win_w}x{win_h}+{x}+{y}")

        popup.grab_set()
        popup.resizable(False, False)

        tk.Label(popup, text="🎤 Hãy nói vào micro...", font=("Arial", 13, "bold")).pack(pady=6)
        lbl_count = tk.Label(popup, text="", font=("Arial", 32, "bold"), fg="red")
        lbl_count.pack(pady=6)

        canvas_wave = tk.Canvas(popup, width=300, height=60, bg="white")
        canvas_wave.pack(pady=6)
        progress_am = ttk.Progressbar(popup, orient="horizontal", length=280, mode="determinate", maximum=100)
        progress_am.pack(pady=(0, 5))

        root.update_idletasks()

        def countdown(n=3):
            if n > 0:
                lbl_count.config(text=str(n))
                popup.after(700, countdown, n - 1)
            else:
                lbl_count.config(text="Bắt đầu!")
                app_beep(1200, 200, popup)

                # ✅ Giữ lại popup một chút để sóng âm hiện rõ rồi đóng 6 giây
                def stop_wave_and_close():
                    wave_running["on"] = False
                    if popup.winfo_exists():
                        popup.destroy()
                popup.after(6000, stop_wave_and_close)

                threading.Thread(target=run_record, daemon=True).start()
                threading.Thread(target=vong_lap_suong_am, args=(canvas_wave, progress_am), daemon=True).start()


        countdown(3)

    #  vong_lap_suong_am(canvas):
    def vong_lap_suong_am(canvas, progressbar):
        import sounddevice as sd
        import numpy as np
        import colorsys

        w, h = 300, 60
        middle = h // 2
        scale = h / 2
        latest_data = {"data": None}

        def get_color(volume):
            volume = min(1.0, max(0.0, volume))
            hue = 0.33 * (1 - volume)
            r, g, b = colorsys.hsv_to_rgb(hue, 1.0, 1.0)
            return '#%02x%02x%02x' % (int(r * 255), int(g * 255), int(b * 255))

        def audio_callback(indata, frames, time_, status):
            latest_data["data"] = indata[:, 0].copy()

        def draw_loop():
           
            if not wave_running["on"] or not canvas.winfo_exists():
                return

            data = latest_data.get("data")
            if data is not None:
                try:
                    canvas.delete("all")
                    step = max(1, len(data) // w)
                    volume_raw = np.mean(np.abs(data))  # độ nhạy ổn định hơn norm
                    color = get_color(volume_raw)

                    for x in range(w):
                        y = int(data[x * step] * scale)
                        canvas.create_line(x, middle, x, middle - y, fill=color, width=1)

                    # 💡 Làm mượt dao động bằng buffer trung bình
                    if not hasattr(progressbar, "_vol_buffer"):
                        progressbar._vol_buffer = []

                    buffer = progressbar._vol_buffer
                    buffer.append(volume_raw)
                    if len(buffer) > 5:
                        buffer.pop(0)

                    volume_smoothed = sum(buffer) / len(buffer)
                    volume_scaled = min(100, max(0, int(volume_smoothed * 1000)))
                    progressbar["value"] = volume_scaled

                except Exception as e:
                    print("⛔ Lỗi vẽ sóng âm:", e)
                    wave_running["on"] = False
                    return

            canvas.after(50, draw_loop)


        try:
            stream = sd.InputStream(callback=audio_callback, channels=1, samplerate=44100, blocksize=512)
            stream.start()
            canvas.after(50, draw_loop)
            time.sleep(10)
            stream.stop()
        except Exception as e:
            print("⛔ Lỗi mở stream sóng âm:", e)




    def run_record():
        import sounddevice as sd
        import speech_recognition as sr
        import numpy as np
        fs = 44100
        recognizer = sr.Recognizer()
        try:
            recording = sd.rec(int(5 * fs), samplerate=fs, channels=1, dtype='int16') #thời gian ghi âm chỉnh int(10 * fs)
            sd.wait()

            audio_np = np.squeeze(recording)
            audio_bytes = audio_np.tobytes()
            audio = sr.AudioData(audio_bytes, sample_rate=fs, sample_width=2)

            detected_lang = "vi"
            text = ""

            for lang_code in ["vi-VN", "en-US", "ja-JP", "zh-CN"]:
                try:
                    text = recognizer.recognize_google(audio, language=lang_code)
                    detected_lang = lang_code.split("-")[0]
                    break
                except sr.UnknownValueError:
                    continue

            if not text.strip():
                raise sr.UnknownValueError

            # 🔒 Chống lỗi gTTS không hỗ trợ ngôn ngữ lạ
            if detected_lang not in ["vi", "en", "ja", "zh"]:
                detected_lang = "vi"

            entry_cau_hoi.delete(0, tk.END)
            entry_cau_hoi.insert(tk.END, text)
            entry_cau_hoi.focus_set()
            print(f"Nhận: [{detected_lang}] {text}")

            def send_and_read():
                gui_hoi_gemini()
                # Đợi tới khi có nội dung thực sự (hoặc timeout tối đa 20 giây)
                bat_dau = time.time()
                while True:
                    content = txt_de.get("1.0", tk.END).strip()
                    if content and "⏳" not in content:
                        break
                    if time.time() - bat_dau > 20:  # Sau 20 giây vẫn chưa có thì bỏ
                        print("⚠ Quá thời gian chờ nội dung GitHub Models/Gemini.")
                        return
                    time.sleep(0.2)

                                    
                doc_noi_dung_de()

            threading.Thread(target=send_and_read, daemon=True).start()
            if on_text_got:
                on_text_got(text, detected_lang)

        except sr.UnknownValueError:
            def canh_bao_va_doc_lai():
                def show_popup():
                    try:
                        if root and root.winfo_exists():
                            popup = tk.Toplevel(root)
                            set_popup_icon(popup)
                            popup.title("Lỗi")
                            popup.geometry("340x120+{}+{}".format(root.winfo_x() + 200, root.winfo_y() + 150))
                            popup.configure(bg="white")
                            popup.resizable(False, False)
                            popup.attributes("-topmost", True)
                            tk.Label(popup, text="⚠ Không nhận diện được giọng nói.", font=("Arial", 12), bg="white", fg="red").pack(pady=15)
                            popup.after(5000, popup.destroy)
                    except Exception as e:
                        print("⚠ Không thể tạo popup:", e)

                try:
                    if root and root.winfo_exists():
                        root.after(0, show_popup)
                except Exception as e:
                    print("⚠ Không thể gọi lại GUI sau khi nhận diện lỗi:", e)

                def phat_loi():
                    try:
                        
                        import pygame
                        if not os.path.exists(APPDATA_ROOT):
                            os.makedirs(APPDATA_ROOT, exist_ok=True)
                        path = os.path.join(APPDATA_ROOT, f"mic_prompt_{int(time.time())}.mp3")
                        # Prefer Google Cloud TTS, fallback to gTTS
                        try:
                            _speak_with_google("Hãy nói vào mic để hỏi tôi!", 'vi', False, path)
                        except Exception:
                            tts = gTTS("Hãy nói vào mic để hỏi tôi!", lang="vi")
                            tts.save(path)
                        pygame.mixer.init()
                        pygame.mixer.music.load(path)
                        pygame.mixer.music.play()
                    except Exception as e:
                        print("Không phát được nhắc mic:", e)

                threading.Thread(target=phat_loi, daemon=True).start()

            try:
                if root and root.winfo_exists():
                    root.after(0, canh_bao_va_doc_lai)
            except Exception as e:
                print("⚠ Không thể gọi lại GUI sau khi nhận diện lỗi:", e)

        except sr.RequestError:
            try:
                entry_cau_hoi.after(0, lambda: messagebox.showerror("Lỗi", "Không thể kết nối dịch vụ nhận dạng."))
            except:
                pass
        except Exception as e:
            try:
                entry_cau_hoi.after(0, lambda: messagebox.showwarning("Lỗi khác", str(e)))
            except:
                pass
        finally:
            record_flag["is_recording"] = False

    show_countdown_and_record()

##===============
# ===== DEAD CODE: VIDEO / YOUTUBE UPLOAD =====
# DEAD CODE - remove later. Feature disabled in audio-tool version.

def popup_google_login(file_path):
    print("Disabled: YouTube/video upload feature removed in audio-tool version.")
    try:
        messagebox.showinfo("Đã tắt", "Tính năng upload video YouTube đã được tắt.")
    except Exception:
        pass
    return

    import tkinter as tk
    from tkinter import messagebox
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.oauth2.credentials import Credentials
    import os

    APPDATA_ROOT = os.path.join(BASE_DIR, "AppData")
    TOKEN_JSON = os.path.join(APPDATA_ROOT, "token.json")
    CLIENT_SECRET_FILE = os.path.join(APPDATA_ROOT, "client_secret.json")

    SCOPES = ["https://www.googleapis.com/auth/youtube.upload"]

    # ==== Kiểm tra token trước ====
    if os.path.exists(TOKEN_JSON):
        creds = Credentials.from_authorized_user_file(TOKEN_JSON, SCOPES)
        if creds and creds.valid:
            popup_youtube_upload_v2(file_path)
            return

    # ==== Nếu chưa có token hoặc token không hợp lệ, hiện popup login ====
    popup = tk.Toplevel()
    popup.title("Đăng nhập Google")
    popup.geometry("400x180")
    popup.grab_set()
    set_popup_icon(popup)

    tk.Label(popup, text="Đăng nhập tài khoản Google\nđể chuẩn bị upload YouTube", font=("Arial", 11)).pack(pady=15)

    def start_login():
        try:
            flow = InstalledAppFlow.from_client_secrets_file(CLIENT_SECRET_FILE, SCOPES)
            creds = flow.run_local_server(port=0)

            unprotect_file(TOKEN_JSON)
            with open(TOKEN_JSON, "w", encoding="utf-8") as f:
                f.write(creds.to_json())

            messagebox.showinfo("Thành công", "✅ Đăng nhập thành công và đã lưu!")
            popup.destroy()
            popup_youtube_upload_v2(file_path)

        except Exception as e:
            messagebox.showerror("Lỗi", f"Đăng nhập lỗi:\n{e}")

    tk.Button(popup, text="🔑 Đăng nhập Google", bg="green", fg="white", command=start_login).pack(pady=10)
    tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=5)



def popup_youtube_upload_v2(file_path):
    print("Disabled: YouTube/video upload feature removed in audio-tool version.")
    try:
        messagebox.showinfo("Đã tắt", "Tính năng upload video YouTube đã được tắt.")
    except Exception:
        pass
    return

    import tkinter as tk
    from tkinter import ttk, messagebox
    import threading
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload
    import os
    import webbrowser
    import pyperclip
    import shutil
    import tempfile

    APPDATA_ROOT = os.path.join(BASE_DIR, "AppData")
    TOKEN_JSON = os.path.join(APPDATA_ROOT, "token.json")

    popup = tk.Toplevel()
    popup.title("Upload video YouTube")
    popup.geometry("460x520")
    popup.grab_set()
    set_popup_icon(popup)

    tk.Label(popup, text="Tiêu đề:", font=("Arial", 10, "bold")).pack(pady=3)
    entry_title = tk.Entry(popup, width=55)

    video_name = os.path.basename(file_path)
    title_default = f"{video_name} from App Text to MP3 \nVCJ School International.(84)986183806 VCJ Co.Ltd"
    entry_title.insert(0, title_default)
    entry_title.pack(pady=2)

    tk.Label(popup, text="Mô tả:", font=("Arial", 10, "bold")).pack(pady=3)
    text_desc = tk.Text(popup, width=55, height=5)
    text_desc.insert("1.0", " Đây là video {video_name} xuất tự động từ phần mềm Text to MP3 đa ngôn ngữ. \n Thuộc dự án Máy học tập Thông minh Smart Learning, \n VCJ School International. (84)986183806 \n VCj Co., Ltd")
    text_desc.pack(pady=2)

    tk.Label(popup, text="Chế độ:", font=("Arial", 10, "bold")).pack(pady=3)
    combo_privacy = ttk.Combobox(popup, values=["public", "unlisted", "private"])
    combo_privacy.current(1) # mặc định là không côNG KHAI, unlisted (1), công khai thì là (0) public, riêng tư là (2) pribvate
    combo_privacy.pack(pady=2)

    tk.Label(popup, text="Đường dẫn video:", font=("Arial", 10, "bold")).pack(pady=3)
    entry_file = tk.Entry(popup, width=55)
    entry_file.insert(0, file_path)
    entry_file.pack(pady=2)

    progress_var = tk.StringVar()
    progress_var.set("Chưa bắt đầu")
    tk.Label(popup, textvariable=progress_var, fg="green").pack(pady=5)

    progress_bar = ttk.Progressbar(popup, orient="horizontal", length=350, mode="determinate")
    progress_bar.pack(pady=5)

    def start_upload_thread():
        def run_upload():
            try:
                creds = Credentials.from_authorized_user_file(TOKEN_JSON, ["https://www.googleapis.com/auth/youtube.upload"])
                youtube = build("youtube", "v3", credentials=creds)

                original_path = entry_file.get()
                file_name = os.path.basename(original_path)

                # Nếu tên file dài, copy sang tên tạm
                if len(file_name) > 90:
                    short_name = "temp_upload_video.mp4"
                    temp_dir = tempfile.gettempdir()
                    temp_path = os.path.join(temp_dir, short_name)
                    shutil.copy(original_path, temp_path)
                    entry_file.delete(0, tk.END)
                    entry_file.insert(0, temp_path)
                else:
                    temp_path = original_path

                # Giới hạn title
                title = entry_title.get()
                if len(title) > 100:
                    title = title[:100]

                media = MediaFileUpload(temp_path, chunksize=-1, resumable=True, mimetype="video/*")

                request = youtube.videos().insert(
                    part="snippet,status",
                    body={
                        "snippet": {
                            "title": title,
                            "description": text_desc.get("1.0", "end-1c"),
                            "tags": ["text-to-mp3", "multi-language"],
                            "categoryId": "27"
                        },
                        "status": {
                            "privacyStatus": combo_privacy.get()
                        }
                    },
                    media_body=media
                )

                response = None
                while response is None:
                    status, response = request.next_chunk()
                    if status:
                        percent = int(status.progress() * 100)
                        progress_var.set(f"Đang tải: {percent}%")
                        progress_bar["value"] = percent
                        popup.update_idletasks()

                progress_var.set("✅ Đã upload xong!")
                progress_bar["value"] = 100

                video_id = response['id']
                video_url = f"https://www.youtube.com/watch?v={video_id}"

                def open_and_copy():
                    webbrowser.open(video_url)
                    pyperclip.copy(video_url)
                    messagebox.showinfo("Đã copy", "🔗 Link video đã được copy vào clipboard!")

                tk.Button(popup, text="🔗 Mở & Copy Link", bg="green", fg="white", command=open_and_copy).pack(pady=5)
                pygame.mixer.init()
                pygame.mixer.music.load(SUCCESS_SOUND)
                pygame.mixer.music.play()
                gui_discord_thong_bao(f"🎙️ [TextToMp3] ✅ Video đã upload Youtube thành công: {video_url}")
                messagebox.showinfo("Hoàn tất", f"✅ Video đã upload thành công!\nID: {video_id}")
                
            except Exception as ex:
                progress_var.set("❌ Lỗi")
                messagebox.showerror("Lỗi", f"Lỗi upload: {ex}")
                gui_discord_thong_bao(f"🎙️ [TextToMp3] Lỗi upload: {ex}")

        threading.Thread(target=run_upload, daemon=True).start()
        progress_var.set("🔄 Đang khởi tạo...")

    def logout_google():
        try:
            if os.path.exists(TOKEN_JSON):
                unprotect_file(TOKEN_JSON)
                os.remove(TOKEN_JSON)
                messagebox.showinfo("Đăng xuất", "✅ Đã xoá token. Lần sau sẽ yêu cầu đăng nhập lại.")
                popup.destroy()
            else:
                messagebox.showinfo("Đăng xuất", "Không có token để xoá (chưa từng đăng nhập).")
        except Exception as e:
            messagebox.showerror("Lỗi", f"Lỗi khi xoá token: {e}")

    tk.Button(popup, text="🚀 Upload", bg="blue", fg="white", command=start_upload_thread).pack(pady=6)
    tk.Button(popup, text="Đăng xuất Google", bg="red", fg="white", command=logout_google).pack(pady=3)
    tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=3)

# ===== END DEAD CODE: VIDEO / YOUTUBE UPLOAD =====
# ==========

#==============XỬ LÝ VĂN BẢN TRƯỚC KHI ĐỌC=====
def lam_sach_van_ban(text):
    import re
    if not text:
        return ""

    # Quy ước ký tự bullet / đánh dấu cần loại
    BULLETS = r'•◦‣▪▫■□◆◇♦⬤●○★☆✓✔✗✘☑✅❌➤➔▶►→←·※■□▪▫❖✦✧'
    BULLETS_CLASS = f"[{BULLETS}]"

    # 1) Chuẩn hoá phép toán thường gặp
    text = re.sub(r'(?<=\d)\s*-\s*(?=\d)', ' trừ ', text)
    text = re.sub(r'(?<=\d)\s*\*\s*(?=\d)', ' nhân ', text)
    text = text.replace("×", " nhân ").replace("÷", " chia ").replace("/", " chia ")
    text = text.replace("+", " cộng ").replace("=", " bằng ")

    # 2) Bỏ các ký tự trang trí/markdown
    text = re.sub(r'[*_`]+', '', text).replace('"', '').replace("“", '').replace("”", '')
    text = re.sub(r'_+', ' ', text)

    # 3) Loại bullet đầu dòng, dấu gạch đầu dòng thừa
    text = re.sub(rf'^\s*(?:{BULLETS_CLASS}|[-–—])+(\s+|$)', '', text)

    # 4) Loại mọi bullet còn sót lại ở giữa câu
    text = re.sub(rf'{BULLETS_CLASS}', ' ', text)

    # 5) Dấu câu tiếng Trung -> chuẩn
    text = text.replace("。", ".").replace("，", ",").replace("、", ",").replace("？", "?")

    # 6) Rút gọn khoảng trắng
    text = re.sub(r'\s+', ' ', text).strip()

    # 7) Nếu chỉ còn ký tự không phải chữ/số thì coi như rỗng (bỏ qua khi TTS)
    #   Bao phủ Latin có dấu, CJK, Hiragana/Katakana, Hangul, và số
    if not re.search(r'[A-Za-zÀ-ỿ\u0100-\u024F\u1E00-\u1EFF\u4E00-\u9FFF\u3040-\u30FF\uAC00-\uD7A3\d]', text):
        return ""

    return text


def tach_de_thanh_danh_sach_da_ngon_ngu():
    
    from langdetect import detect
    noi_dung = txt_de.get("1.0", tk.END).strip()
    danh_sach = []
    for dong in noi_dung.split('\n'):
        dong = dong.strip()
        if not dong:
            continue
        try:
            lang = detect(dong)
        except:
            lang = "vi"
        danh_sach.append((dong, lang))
    return danh_sach


def doan_ngon_ngu_theo_ky_tu(text):
    text = text.strip()
    count_zh = sum(0x4E00 <= ord(c) <= 0x9FFF for c in text)
    count_ja = sum(0x3040 <= ord(c) <= 0x30FF for c in text)
    count_ko = sum(0xAC00 <= ord(c) <= 0xD7AF for c in text)
    count_vi = sum(c in 'ăâêôơưđáàảãạấầẩẫậắằẳẵặéèẻẽẹếềểễệíìỉĩịóòỏõọốồổỗộớờởỡợúùủũụứừửữựỳýỵỷỹ' for c in text.lower())

    if count_ja > 0:
        return "ja"
    elif count_zh > 0:
        return "zh"
    elif count_ko > 0:
        return "ko"
    elif count_vi > 0:
        return "vi"
    else:
        try:
            from langdetect import detect
            return detect(text)
        except:
            return "vi"


def _get_runtime_tts_selection():
    """Read the active popup selection, falling back to the saved app config."""
    engine = _load_ui_tts_engine_from_config()
    voice = config.get("VOCAB_TTS_VOICE", "Nam") or "Nam"
    try:
        engine = combo_engine.get().strip() or engine
    except Exception:
        pass
    try:
        voice = combo_giong_popup.get().strip() or voice
    except Exception:
        pass
    return engine, voice


def doc_noi_dung_de():
    import threading, tempfile, os, time
    import tkinter as tk
    from tkinter import messagebox
    import pygame
    from langdetect import detect

    global dung_doc_ngay, dang_doc, channel_doc
    dung_doc_ngay = False
    engine_selected, voice_selected = _get_runtime_tts_selection()

    def run():
        global noi_dung_cuoi, dang_doc, channel_doc
        if dang_doc:
            print("⛔ Đang đọc, không thể đọc mới.")
            return

        noi_dung = txt_de.get("1.0", tk.END).strip()
        if not noi_dung:
            messagebox.showinfo("Trống", "Không có nội dung để đọc.", parent=root)
            return

        noi_dung_cuoi = noi_dung
        cac_dong = noi_dung.split('\n')

        try:
            pygame.mixer.init()
        except:
            pass

        try:
            channel_doc = pygame.mixer.Channel(1)
        except Exception as e:
            print("❌ Lỗi tạo kênh âm thanh:", e)
            return

        dang_doc = True

        gtts_supported = ["vi", "en", "ja", "zh", "ko", "fr", "de", "es", "it", "pt"]

        chon_combo = combo_ngon_ngu.get().strip()
        lang_from_combo = {
            "Tiếng Việt": "vi",
            "Tiếng Anh": "en",
            "Tiếng Nhật": "ja",
            "Tiếng Trung": "zh"
        }.get(chon_combo, None)

        dialogue_count = 0
        for dong in cac_dong:
            if dung_doc_ngay:
                print("⛔ Dừng đọc ngay.")
                break

            dong = dong.strip()
            if not dong:
                continue

            try:
                dong_sach = lam_sach_van_ban(dong)

                # Ưu tiên ngôn ngữ do người dùng chọn
                if lang_from_combo:
                    lang = lang_from_combo
                else:
                    try:
                        lang = detect(dong_sach)
                    except:
                        lang = "vi"

                lang = lang.lower()
                if lang in ["zh-cn", "zh-tw", "zh-hk"]:
                    lang = "zh"

                if lang not in gtts_supported:
                    print(f"⚠ Ngôn ngữ '{lang}' không hỗ trợ, dùng tiếng Việt.")
                    lang = "vi"

                dialogue_pairs = {
                    "Hội thoại 1 câu nam - 1 câu nữ": ("Nam", "Nữ"),
                    "Hội thoại 1 câu nữ - 1 câu nam": ("Nữ", "Nam"),
                }
                if voice_selected in dialogue_pairs:
                    first_voice, second_voice = dialogue_pairs[voice_selected]
                    voice = first_voice if dialogue_count % 2 == 0 else second_voice
                    dialogue_count += 1
                else:
                    voice = voice_selected

                print(f"📢 Đọc ({lang}) [{voice}] Engine: {engine_selected}: {dong_sach}")
                with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tf:
                    file_path = tf.name

                toc_do = combo_toc_do.get()
                tao_file_mp3(dong_sach, lang=lang, voice=voice, toc_do=toc_do, engine=engine_selected, file_out=file_path)

                sound = pygame.mixer.Sound(file_path)
                channel_doc.play(sound)
                time.sleep(0.2)  # chờ để âm thanh bắt đầu phát

                while channel_doc.get_busy():
                    if dung_doc_ngay:
                        print("⛔ Dừng đọc giữa dòng.")
                        channel_doc.stop()
                        break
                    time.sleep(0.1)

                os.remove(file_path)

                if dung_doc_ngay:
                    break

            except Exception as e:
                print(f"❌ Lỗi đọc dòng: {e}")
                continue

        dang_doc = False
        print("✅ Kết thúc đọc nội dung.")

    threading.Thread(target=run).start()

#=================

def toggle_tam_dung():
    import pygame
    global dung_doc_ngay
    dung_doc_ngay = True
    if not pygame.mixer.get_init():
        return
    if pygame.mixer.music.get_busy():
        pygame.mixer.music.stop()
        dung_doc_ngay = True
        btn_tam_dung.config(text="▶️ Dừng Đọc")
    else:
        btn_tam_dung.config(text="⏸ Dừng Đọc")


#Share Zalo


def open_zalo_and_folder(file_path):
    try:
        # Mở thư mục chứa file
        folder_path = os.path.dirname(file_path)
        open_path_cross_platform(folder_path)

        # Mở Zalo Desktop
        zalo_path = r"C:\Users\{}\AppData\Local\Programs\Zalo\Zalo.exe".format(os.getlogin())
        if os.path.exists(zalo_path):
            open_path_cross_platform(zalo_path)
        else:
            tk.messagebox.showwarning("Zalo", "⚠ Không tìm thấy Zalo Desktop. Vui lòng kiểm tra đường dẫn hoặc mở thủ công.")
    except Exception as e:
        tk.messagebox.showerror("Lỗi", f"Không mở được Zalo hoặc thư mục:\n{e}")

#=================

def convert_seconds_to_timestamp(seconds):
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    ms = int((seconds - int(seconds)) * 1000)
    return f"{h:02}:{m:02}:{s:02},{ms:03}"


#===============
def build_danh_sach_doc_from_text():
    lines = txt_de.get("1.0", tk.END).split("\n")
    danh_sach = []
    for line in lines:
        line = line.strip()
        if line:
            danh_sach.append((line, "vi"))
    return danh_sach
# ===== DEAD CODE: GAME DATA / HISTORY =====
# DEAD CODE - remove later. Game Đoán Chữ disabled in audio-tool version.


def tai_lich_su():
    print("Disabled: game history removed in audio-tool version.")
    return []

    try:
        if os.path.exists(LICH_SU_FILE):
            with open(LICH_SU_FILE, "r", encoding="utf-8") as f:
                du_lieu = json.load(f)
                for ng in du_lieu:
                    if "thoigian" not in ng:
                        ng["thoigian"] = datetime.now().strftime("%H:%M %d/%m/%Y")
                    if "ngonngu" not in ng:
                        ng["ngonngu"] = "Ja"  # mã rút gọn mặc định
                return du_lieu
        else:
            return []
    except json.JSONDecodeError:
        print("⚠️ File lịch sử bị lỗi JSON. Đang thử sửa...")
        try:
            with open(LICH_SU_FILE, "r", encoding="utf-8") as f:
                content = f.read()
            if not content.strip().endswith("]"):
                content += "]"
            with open(LICH_SU_FILE, "w", encoding="utf-8") as f:
                f.write(content)
            with open(LICH_SU_FILE, "r", encoding="utf-8") as f:
                du_lieu = json.load(f)
            for ng in du_lieu:
                if "thoigian" not in ng:
                    ng["thoigian"] = datetime.now().strftime("%H:%M %d/%m/%Y")
                if "ngonngu" not in ng:
                    ng["ngonngu"] = "Ja"
            print("✅ Đã sửa lỗi JSON tự động.")
            return du_lieu
        except Exception as e:
            print("❌ Không thể sửa file JSON:", e)
            return []
    except Exception as e:
        print("❌ Lỗi đọc lịch sử game:", e)
        return []



def doc_du_lieu_game(ngonngu):
    print("Disabled: Game Đoán Chữ data loading removed in audio-tool version.")
    return None

    if ngonngu == "Dnn":
        df_ja = pd.read_excel(EXCEL_GAME_PATH, sheet_name="Ja")
        df_cn = pd.read_excel(EXCEL_GAME_PATH, sheet_name="Cn")
        df_en = pd.read_excel(EXCEL_GAME_PATH, sheet_name="En")
        df = pd.concat([df_ja, df_cn, df_en], ignore_index=True)
    else:
        df = pd.read_excel(EXCEL_GAME_PATH, sheet_name=ngonngu)

    df = df.dropna(subset=["Câu hỏi", "Câu rút gọn", "Đáp án", "Nghĩa TV"])
    return df


def play_sound(sound_path):
    pygame.mixer.init()
    pygame.mixer.music.load(sound_path)
    pygame.mixer.music.play()

def luu_diem(ten, level, diem, ngonngu=""):
    print("Disabled: game score saving removed in audio-tool version.")
    return

    from datetime import datetime
    try:
        # 🔁 Ánh xạ ngôn ngữ đầy đủ → mã rút gọn
        ma_hoa = {
            "Tiếng Nhật": "Ja",
            "Tiếng Trung": "Cn",
            "Tiếng Anh": "En",
            "Đa ngôn ngữ": "Dnn"
        }
        ngonngu_ma = ma_hoa.get(ngonngu, ngonngu)  # nếu đã là mã rút gọn thì giữ nguyên

        lich_su = tai_lich_su()
        lich_su.append({
            "ten": ten,
            "level": level,
            "diem": diem,
            "ngonngu": ngonngu_ma,
            "thoigian": datetime.now().strftime("%H:%M %d/%m/%Y")
        })
        if len(lich_su) > 100:
            lich_su = lich_su[-100:]  # Giữ lại 100 người gần nhất

        unprotect_file(LICH_SU_FILE)
        with open(LICH_SU_FILE, "w", encoding="utf-8") as f:
            json.dump(lich_su, f, ensure_ascii=False, indent=2)

    except Exception as e:
        print("❌ Lỗi lưu điểm:", e)

        
#====
# ===== END DEAD CODE: GAME DATA / HISTORY =====


#===========
#Chọn ngôn ngữ thủ công=========================================
#==========================

background_path = LOGO_PATH
count = 0
popup_lang_open = False  # Chỉ mở 1 cửa sổ


#===============================================CỬA SỔ CHỌN NGÔN NGỮ====================================
def mo_popup_chon_lang(mo_tu_ben_ngoai=False):
#def mo_popup_chon_lang():
    from tkinter import ttk, filedialog
    from langdetect import detect
    import tempfile, os, threading, time
    import pygame
    global popup_lang_open, popup, background_path
    global combo_engine, combo_giong_popup, combo_toc_do_popup

    if popup_lang_open:
        tk.messagebox.showwarning("Đang mở", "Cửa sổ chọn ngôn ngữ đã được mở rồi!")
        return

    popup_lang_open = True
    popup = tk.Toplevel(root)
    set_popup_icon(popup)
    popup.title("Chọn ngôn ngữ từng dòng")
    popup.geometry("1420x860+40+20")
    popup.minsize(1240, 760)
    popup.grab_set()
    popup.transient(root)

    body_frame = tk.Frame(popup, bg="#eef3ee")
    body_frame.pack(fill="both", expand=True, padx=10, pady=10)
    body_frame.columnconfigure(0, weight=1)
    body_frame.columnconfigure(1, weight=0)
    body_frame.rowconfigure(0, weight=1)

    left_frame = tk.Frame(body_frame, bg="#ffffff")
    left_frame.grid(row=0, column=0, sticky="nsew", padx=(0, 8))

    right_frame = tk.Frame(body_frame, bg="#f8fff8", width=400)
    right_frame.grid(row=0, column=1, sticky="ns")
    right_frame.pack_propagate(False)

    danh_sach = []
    

    
    for dong in txt_de.get("1.0", tk.END).split('\n'):
        dong = dong.strip()
        if not dong:
            continue
        try:
            lang = detect(dong)
            if lang == "ko" and any('\u4e00' <= c <= '\u9fff' for c in dong):
                lang = "zh-cn"
        except:
            lang = "vi"
        if lang == "zh":
            lang = "zh-cn"
        danh_sach.append((dong, lang))

        
    selected_lang_list = [lang for _, lang in danh_sach]
    
    canvas = tk.Canvas(left_frame, bg="#ffffff", highlightthickness=0)
    frame = tk.Frame(canvas)
    vsb = tk.Scrollbar(left_frame, orient="vertical", command=canvas.yview)
    canvas.configure(yscrollcommand=vsb.set)

    canvas.pack(side="left", fill="both", expand=True)
    vsb.pack(side="right", fill="y")
    list_window = canvas.create_window((0, 0), window=frame, anchor="nw")

    def on_list_canvas_configure(event):
        canvas.itemconfigure(list_window, width=event.width)

    canvas.bind("<Configure>", on_list_canvas_configure)

    lang_options = ['vi', 'en', 'ja', 'zh', 'zh-cn']
    combos = []
    allowed_langs = ["vi", "zh", "ja", "en", "zh-cn"]

    # Frame tùy chọn bên phải (scroll được để không bị che nút)
    option_canvas = tk.Canvas(right_frame, bg="#f8fff8", highlightthickness=0)
    option_vsb = tk.Scrollbar(right_frame, orient="vertical", command=option_canvas.yview)
    option_canvas.configure(yscrollcommand=option_vsb.set)
    option_canvas.pack(side="left", fill="both", expand=True)
    option_vsb.pack(side="right", fill="y")

    option_frame = tk.Frame(option_canvas, bg="#f8fff8")
    option_window = option_canvas.create_window((0, 0), window=option_frame, anchor="nw")

    def on_option_frame_configure(event):
        option_canvas.configure(scrollregion=option_canvas.bbox("all"))

    option_frame.bind("<Configure>", on_option_frame_configure)

    def on_option_canvas_configure(event):
        option_canvas.itemconfigure(option_window, width=event.width)

    option_canvas.bind("<Configure>", on_option_canvas_configure)

    tk.Label(option_frame, text="Cấu hình xuất", font=("Arial", 12, "bold"), bg="#f8fff8").pack(anchor="w", padx=10, pady=(4, 6))
    
    # ==== Chọn ngôn ngữ muốn giữ lại ====
    tk.Label(option_frame, text="Ngôn ngữ sử dụng:", font=("Arial", 10, "bold"), bg="#f8fff8").pack(anchor="w", padx=10, pady=(6, 2))

    stored_vocab_languages = config.get("VOCAB_TTS_LANGUAGES", ["vi", "en", "ja", "zh"])
    if not isinstance(stored_vocab_languages, list):
        stored_vocab_languages = ["vi", "en", "ja", "zh"]
    ngon_ngu_flags = {
        "vi": tk.BooleanVar(value="vi" in stored_vocab_languages),
        "en": tk.BooleanVar(value="en" in stored_vocab_languages),
        "ja": tk.BooleanVar(value="ja" in stored_vocab_languages),
        "zh": tk.BooleanVar(value="zh" in stored_vocab_languages),
    }
    for code, name in [("vi", "Tiếng Việt"), ("en", "Tiếng Anh"), ("ja", "Tiếng Nhật"), ("zh", "Tiếng Trung")]:
        cb = tk.Checkbutton(option_frame, text=name, variable=ngon_ngu_flags[code], bg="#f8fff8")
        cb.pack(anchor="w", padx=10)


    tk.Label(option_frame, text="Engine (Google Cloud TTS / gTTS / Polly):", bg="#f8fff8").pack(anchor="w", padx=10, pady=(8, 2))
    combo_engine = ttk.Combobox(option_frame, values=["Google Cloud TTS", "gTTS", "Polly"], state="readonly")
    combo_engine.set(_load_ui_tts_engine_from_config())
    combo_engine.pack(fill="x", padx=10, pady=2)
    engine_status_var = tk.StringVar(value=f"Đang lưu: {combo_engine.get()}")
    tk.Label(
        option_frame,
        textvariable=engine_status_var,
        bg="#f8fff8",
        fg="#2f5f2f",
        wraplength=360,
        justify="left",
    ).pack(anchor="w", padx=10, pady=(0, 2))
    

    tk.Label(option_frame, text="Tốc độ:", bg="#f8fff8").pack(anchor="w", padx=10, pady=(10, 2))
    combo_toc_do_popup = ttk.Combobox(option_frame, values=["Chậm", "Bình thường"], state="readonly")
    combo_toc_do_popup.set(config.get("VOCAB_TTS_SPEED", "Chậm"))
    combo_toc_do_popup.pack(fill="x", padx=10, pady=2)

    tk.Label(option_frame, text="Giọng đọc:", bg="#f8fff8").pack(anchor="w", padx=10, pady=(10, 2))
    def _dialogue_voice_pair(voice_value):
        if voice_value == "Hội thoại 1 câu nữ - 1 câu nam":
            return ("Nữ", "Nam")
        return ("Nam", "Nữ")

    combo_giong_popup = ttk.Combobox(
        option_frame,
        values=["Nam", "Nữ", "Hội thoại 1 câu nam - 1 câu nữ", "Hội thoại 1 câu nữ - 1 câu nam"],
        state="readonly",
    )
    combo_giong_popup.set(config.get("VOCAB_TTS_VOICE", "Hội thoại 1 câu nam - 1 câu nữ"))
    combo_giong_popup.pack(fill="x", padx=10, pady=2)

    vocab_tts_confirmed_var = tk.BooleanVar(
        value=str(config.get("VOCAB_TTS_CONFIG_CONFIRMED", "false")).strip().lower() in {"1", "true", "yes"}
    )
    tk.Checkbutton(
        option_frame,
        text="Xác nhận dùng cấu hình này cho M4A vocab HSK 2.0 và HSK 3.0",
        variable=vocab_tts_confirmed_var,
        bg="#f8fff8",
        wraplength=360,
        justify="left",
    ).pack(anchor="w", padx=10, pady=(4, 4))

    google_voice_status_var = tk.StringVar(value="Google Cloud TTS: mặc định theo ngôn ngữ")
    tk.Label(option_frame, textvariable=google_voice_status_var, fg="#2f5f2f", bg="#f8fff8", wraplength=360,
             justify="left").pack(anchor="w", padx=10, pady=(4, 2))

    def update_google_voice_status():
        try:
            summary = []
            for code, label in [("vi", "VI"), ("en", "EN"), ("ja", "JA"), ("zh", "ZH")]:
                profile = _google_tts_get_profile(code)
                active_gender = profile.get("gender", "Mặc định")
                slot_bits = []
                for slot_label in GOOGLE_TTS_SLOT_LABELS:
                    slot_profile = profile.get("slots", {}).get(slot_label, {})
                    slot_voice = (slot_profile.get("voice_name") or "").strip()
                    slot_bits.append(f"{slot_label}={slot_voice or 'mặc định'}")
                summary.append(f"{label}[{active_gender}]: " + "; ".join(slot_bits))
            google_voice_status_var.set(" | ".join(summary))
        except Exception as exc:
            google_voice_status_var.set(f"Google Cloud TTS: không đọc được trạng thái ({exc})")

    def open_google_voice_popup():
        popup_google = tk.Toplevel(option_frame)
        set_popup_icon(popup_google)
        popup_google.title("Giọng Google Cloud TTS")
        popup_google.geometry("940x650")
        popup_google.resizable(True, True)
        popup_google.transient(root)
        popup_google.grab_set()
        popup_google.attributes("-topmost", True)

        popup_google.update_idletasks()
        x = root.winfo_x() + 80
        y = root.winfo_y() + 80
        popup_google.geometry(f"940x650+{x}+{y}")

        main = tk.Frame(popup_google, bg="white")
        main.pack(fill="both", expand=True, padx=16, pady=16)

        tk.Label(main, text="Chọn giọng Google Cloud TTS", font=("Arial", 14, "bold"), bg="white").pack(anchor="w")
        tk.Label(
            main,
            text="Mỗi ngôn ngữ có thể lưu riêng giọng Nam/Nữ/Trung tính. Nếu chưa chọn thì dùng mặc định của Google.",
            bg="white",
            fg="#555",
            wraplength=880,
            justify="left",
        ).pack(anchor="w", pady=(4, 10))

        controls = tk.Frame(main, bg="white")
        controls.pack(fill="x", pady=(10, 6))

        tk.Label(controls, text="Ngôn ngữ", bg="white").grid(row=0, column=0, sticky="w")
        google_lang_var = tk.StringVar(value="Tiếng Việt")
        google_lang_combo = ttk.Combobox(controls, textvariable=google_lang_var,
                                         values=["Tiếng Việt", "Tiếng Anh", "Tiếng Nhật", "Tiếng Trung"], state="readonly", width=28)
        google_lang_combo.grid(row=1, column=0, sticky="we", padx=(0, 10), pady=(2, 8))

        tk.Label(controls, text="Giới tính", bg="white").grid(row=0, column=1, sticky="w")
        google_gender_var = tk.StringVar(value="Mặc định")
        google_gender_combo = ttk.Combobox(controls, textvariable=google_gender_var,
                                           values=GOOGLE_TTS_GENDER_LABELS, state="readonly", width=18)
        google_gender_combo.grid(row=1, column=1, sticky="we", padx=(0, 10), pady=(2, 8))

        tk.Label(controls, text="Giọng cụ thể", bg="white").grid(row=0, column=2, sticky="w")
        google_voice_name_var = tk.StringVar(value="(Mặc định)")
        google_voice_name_combo = ttk.Combobox(controls, textvariable=google_voice_name_var,
                                               values=["(Mặc định)"], state="readonly", width=42)
        google_voice_name_combo.grid(row=1, column=2, sticky="we", pady=(2, 8))

        controls.grid_columnconfigure(0, weight=1)
        controls.grid_columnconfigure(1, weight=1)
        controls.grid_columnconfigure(2, weight=2)

        tk.Label(main, text="Văn bản nghe thử", bg="white").pack(anchor="w")
        google_tts_sample_texts = {
            "vi": "Xin chào, tôi là Phương Anh - trợ lý Google TTS.",
            "en": "Hello, I am Phuong Anh - your Google TTS assistant.",
            "ja": "こんにちは、私はフォン・アイン、Google TTS アシスタントです。",
            "zh": "你好，我是芳英，Google TTS 助理。",
        }

        sample_var = tk.StringVar(value=google_tts_sample_texts["vi"])
        sample_entry = tk.Entry(main, textvariable=sample_var)
        sample_entry.pack(fill="x", pady=(2, 8))

        voice_map = {}
        info_var = tk.StringVar(value="Bấm 'Tải giọng' để lấy danh sách voice từ Google Cloud.")
        info_label = tk.Label(main, textvariable=info_var, bg="white", fg="#555", wraplength=610, justify="left")
        info_label.pack(anchor="w", pady=(0, 8))

        def _lang_code_from_label(label):
            return {
                "Tiếng Việt": "vi",
                "Tiếng Anh": "en",
                "Tiếng Nhật": "ja",
                "Tiếng Trung": "zh",
            }.get(label, "vi")

        def load_google_voices():
            nonlocal voice_map
            code = _lang_code_from_label(google_lang_var.get())
            gender = google_gender_var.get()
            try:
                voices = _google_tts_list_voices(code, gender if gender != "Mặc định" else None)
                voice_map = {"(Mặc định)": ""}
                values = ["(Mặc định)"]
                for item in voices:
                    display = item["display"]
                    voice_map[display] = item["name"]
                    values.append(display)
                google_voice_name_combo["values"] = values
                if google_voice_name_var.get() not in values:
                    google_voice_name_var.set("(Mặc định)")
                info_var.set(f"Đã tải {len(voices)} giọng cho {code}. Nếu để '(Mặc định)' thì Google tự chọn giọng chuẩn của ngôn ngữ.")
            except Exception as exc:
                google_voice_name_combo["values"] = ["(Mặc định)"]
                google_voice_name_var.set("(Mặc định)")
                info_var.set(f"Không tải được danh sách giọng: {exc}")
                thong_bao_loi_api(exc, "Google Cloud TTS")

        def sync_google_voice_selection():
            code = _lang_code_from_label(google_lang_var.get())
            profile = _google_tts_get_profile(code)
            slot_label = _google_tts_slot_label(google_gender_var.get())
            stored_voice = ""
            if slot_label == "Mặc định":
                stored_voice = (profile.get("slots", {}).get("Mặc định", {}).get("voice_name") or profile.get("voice_name", "") or "").strip()
            else:
                stored_voice = (profile.get("slots", {}).get(slot_label, {}).get("voice_name") or "").strip()
            load_google_voices()
            if stored_voice:
                for display, actual_name in voice_map.items():
                    if actual_name == stored_voice:
                        google_voice_name_var.set(display)
                        return
            google_voice_name_var.set("(Mặc định)")

        def save_google_profile():
            code = _lang_code_from_label(google_lang_var.get())
            voice_name = voice_map.get(google_voice_name_var.get(), "")
            _google_tts_set_profile(code, gender=google_gender_var.get(), voice_name=voice_name)
            _save_google_tts_profiles_to_config()
            update_google_voice_status()
            messagebox.showinfo("Đã lưu", f"Đã lưu giọng Google cho {code}.", parent=popup_google)
            popup_google.destroy()

        def test_google_voice():
            code = _lang_code_from_label(google_lang_var.get())
            voice_name = voice_map.get(google_voice_name_var.get(), "")
            sample_text = sample_var.get().strip()
            if not sample_text:
                sample_text = google_tts_sample_texts.get(code, google_tts_sample_texts["vi"])
            selected_gender = google_gender_var.get()

            def _worker():
                try:
                    import pygame
                    test_path = os.path.join(APPDATA_ROOT, f"google_tts_test_{uuid.uuid4().hex}.mp3")
                    tao_file_google_mp3(
                        sample_text,
                        lang=code,
                        gender=selected_gender,
                        voice_name=voice_name,
                        toc_do="Bình thường",
                        file_out=test_path,
                    )
                    pygame.mixer.init()
                    pygame.mixer.music.load(test_path)
                    pygame.mixer.music.play()
                    print(f"✅ Đang nghe thử Google Cloud TTS: {test_path}")
                except Exception as exc:
                    print("❌ Không nghe thử được Google Cloud TTS:", exc)
                    thong_bao_loi_api(exc, "Google Cloud TTS")
                    error_message = str(exc)
                    try:
                        popup_google.after(
                            0,
                            lambda msg=error_message: messagebox.showerror(
                                "Lỗi nghe thử",
                                msg,
                                parent=popup_google,
                            ),
                        )
                    except Exception:
                        pass

            threading.Thread(target=_worker, daemon=True).start()

        def on_language_change(event=None):
            code = _lang_code_from_label(google_lang_var.get())
            profile = _google_tts_get_profile(code)
            google_gender_var.set(_google_tts_slot_label(profile.get("gender", "Mặc định")))
            sample_var.set(google_tts_sample_texts.get(code, google_tts_sample_texts["vi"]))
            sync_google_voice_selection()

        def on_gender_change(event=None):
            sync_google_voice_selection()

        def sync_from_selected_lines():
            try:
                current_codes = []
                for _, combo in combos[:5]:
                    val = combo.get().strip()
                    if val:
                        current_codes.append(val)
                if not current_codes:
                    return
                first_code = current_codes[0]
                mapped = {
                    "vi": "Tiếng Việt",
                    "en": "Tiếng Anh",
                    "ja": "Tiếng Nhật",
                    "zh": "Tiếng Trung",
                    "zh-cn": "Tiếng Trung",
                }.get(first_code, "Tiếng Việt")
                google_lang_var.set(mapped)
                on_language_change()
            except Exception:
                pass

        google_lang_combo.bind("<<ComboboxSelected>>", on_language_change)
        google_gender_combo.bind("<<ComboboxSelected>>", on_gender_change)

        btn_row = tk.Frame(main, bg="white")
        btn_row.pack(fill="x", pady=(6, 0))

        tk.Button(btn_row, text="Tải giọng", command=sync_google_voice_selection, width=12).pack(side="left", padx=(0, 8))
        tk.Button(btn_row, text="Nghe thử", command=test_google_voice, width=12).pack(side="left", padx=(0, 8))
        tk.Button(btn_row, text="Lấy theo dòng hiện tại", command=sync_from_selected_lines, width=18).pack(side="left", padx=(0, 8))
        tk.Button(btn_row, text="Lưu", command=save_google_profile, width=12).pack(side="left", padx=(0, 8))
        tk.Button(btn_row, text="Đóng", command=popup_google.destroy, width=12).pack(side="right")

        load_google_voices()
        on_language_change()

    tk.Button(option_frame, text="🎙 Giọng Google Cloud…", command=open_google_voice_popup).pack(fill="x", padx=10, pady=(4, 2))
    update_google_voice_status()

    def _canonical_vocab_bitrate(value):
        return "26k" if str(value or "").strip().lower() in {"26k", "26 kbps"} else DEFAULT_VOCAB_M4A_BITRATE

    def _vocab_bitrate_display(value):
        return "26 kbps" if _canonical_vocab_bitrate(value) == "26k" else "32 kbps"

    def _canonical_vocab_audio_mode(value):
        return "zh_only" if str(value or "").strip().lower() in {
            "zh_only", "zh-only", "chinese_only", "chỉ đọc tiếng trung"
        } else "zh_vi"

    def _vocab_audio_mode_display(value):
        return "Chỉ đọc tiếng Trung" if _canonical_vocab_audio_mode(value) == "zh_only" else "Đọc tiếng Trung + Tiếng Việt"

    def collect_vocab_tts_config(parent, bitrate_override=None, audio_mode_override=None):
        """Take one confirmed, serialisable snapshot for both vocab workflows."""
        selected_languages = [code for code in ("vi", "en", "ja", "zh") if ngon_ngu_flags[code].get()]
        if not vocab_tts_confirmed_var.get():
            messagebox.showwarning(
                "Chưa xác nhận cấu hình",
                "Hãy tích 'Xác nhận dùng cấu hình này cho M4A vocab HSK 2.0 và HSK 3.0'.",
                parent=parent,
            )
            return None
        if not {"vi", "zh"}.issubset(set(selected_languages)):
            messagebox.showwarning(
                "Thiếu ngôn ngữ vocab",
                "Vocab HSK cần chọn cả Tiếng Việt và Tiếng Trung.",
                parent=parent,
            )
            return None
        # The workflow-local selector is authoritative. Blank/invalid means 32 kbps.
        bitrate = _canonical_vocab_bitrate(bitrate_override)
        audio_mode = _canonical_vocab_audio_mode(audio_mode_override)
        snapshot = {
            "engine": combo_engine.get().strip(),
            "speed": combo_toc_do_popup.get().strip(),
            "voice": combo_giong_popup.get().strip(),
            "bitrate": bitrate,
            "audio_mode": audio_mode,
            "languages": selected_languages,
            "confirmed": True,
        }
        try:
            config["VOCAB_TTS_SPEED"] = snapshot["speed"]
            config["VOCAB_TTS_VOICE"] = snapshot["voice"]
            config["VOCAB_M4A_BITRATE"] = snapshot["bitrate"]
            config["VOCAB_AUDIO_MODE"] = snapshot["audio_mode"]
            config["VOCAB_TTS_LANGUAGES"] = snapshot["languages"]
            config["VOCAB_TTS_CONFIG_CONFIRMED"] = True
            _write_app_config()
        except Exception as exc:
            print(f"⚠ Không lưu được cấu hình vocab TTS: {exc}")
        return snapshot

    def apply_vocab_tts_env(environment, snapshot):
        environment.update(
            {
                "TTS_ENGINE": snapshot["engine"],
                "TTS_SPEED": snapshot["speed"],
                "TTS_VOICE": snapshot["voice"],
                "M4A_BITRATE": snapshot["bitrate"],
                "TTS_AUDIO_MODE": snapshot["audio_mode"],
                "TTS_LANGUAGES": ",".join(snapshot["languages"]),
                "TTS_CONFIG_CONFIRMED": "true" if snapshot["confirmed"] else "false",
            }
        )

    


    
    #==== phân trang khi văn bản dài:
    current_page = 0
    ITEMS_PER_PAGE = 710 # 710 dòng max 1 trang, giới hạn có thể hiện thị

    #Tạo hàm tính danh sách theo trang
    def get_page_data(danh_sach, page, items_per_page=710):
        start = page * items_per_page
        end = start + items_per_page
        return danh_sach[start:end]

    #Thêm nút điều hướng
    def next_page():
        nonlocal current_page
        if (current_page + 1) * ITEMS_PER_PAGE < len(danh_sach):
            current_page += 1
            rebuild_page()

    def prev_page():
        nonlocal current_page
        if current_page > 0:
            current_page -= 1
            rebuild_page()

    #Viết hàm rebuild
    def rebuild_page():
        for widget in frame.winfo_children():
            widget.destroy()

        combos.clear()

        page_data = get_page_data(danh_sach, current_page, ITEMS_PER_PAGE)
        for i, (dong, lang) in enumerate(page_data):
            global_index = i + current_page * ITEMS_PER_PAGE

            label = tk.Label(frame, text=f"{i+1 + current_page * ITEMS_PER_PAGE}. {dong}", anchor='w', justify='left', wraplength=840)

            if lang not in allowed_langs:
                label.configure(bg="yellow")
                loi_path = os.path.join(APPDATA_ROOT, "loi_ngon_ngu.txt")
                unprotect_file(loi_path)
                with open(loi_path, "a", encoding="utf-8") as f:
                    f.write(f"Dòng {i+1 + current_page * ITEMS_PER_PAGE}: {dong}\n")

            label.pack(fill="x", padx=5)

            combo = ttk.Combobox(frame, values=lang_options)
            combo.set(selected_lang_list[global_index])  # ✅ Lấy ngôn ngữ đang chọn từ list

            combo.pack(fill="x", padx=5, pady=2)
            combos.append((dong, combo))

            def on_combo_change(event, idx=global_index, lbl=label):
                selected_lang_list[idx] = event.widget.get()
                lbl.configure(bg="lightgreen")

            combo.bind("<<ComboboxSelected>>", on_combo_change)
        
        tk.Label(frame, text="", height=2).pack()
        frame.update_idletasks()
        canvas.config(scrollregion=canvas.bbox("all"))



    #===
    rebuild_page() # gọi lại trang

    #==============
    #===Báo chưa sửa lỗi ngôn ngữ

    def kiem_tra_ngon_ngu_hop_le(selected_lang_list, allowed_langs):
        """
        Kiểm tra danh sách ngôn ngữ được chọn, nếu có dòng chưa đúng, hiện popup báo lỗi.
        Trả về True nếu hợp lệ, False nếu có lỗi.
        """
        # Tìm những dòng không hợp lệ
        danh_sach_loi = [i + 1 for i, lang in enumerate(selected_lang_list) if lang not in allowed_langs]

        if danh_sach_loi:
            msg = "⚠ Có dòng chưa được sửa đúng ngôn ngữ!\n\nCác dòng lỗi: " + ", ".join(map(str, danh_sach_loi))
            tk.messagebox.showerror("Lỗi", msg)
            return False

        return True

    #======
    def chon_bg(): #chọn ảnh nền video
        print("Disabled: chọn background/video removed in audio-tool version.")
        try:
            messagebox.showinfo("Đã tắt", "Tính năng chọn background video đã được tắt.", parent=popup)
        except Exception:
            pass
        return

        global background_path
        from tkinter import filedialog, messagebox
        f = filedialog.askopenfilename(filetypes=[("Image files", "*.png;*.jpg;*.jpeg")])
        if f:
            background_path = f
            messagebox.showinfo("OK", "Đã chọn background:\n" + f)

    def get_engine_giong_tocdo():
        try:
            return (
                combo_engine.get(),
                combo_giong_popup.get(),
                combo_toc_do_popup.get()
            )
        except Exception as e:
            print("⚠️ Không lấy được engine/giong/tốc độ:", e)
            return ("gTTS", "Nam", "Bình thường")




    def on_engine_change(event=None):
        engine = combo_engine.get()
        _save_ui_tts_engine_to_config(engine)
        engine_status_var.set(f"Đang lưu: {engine}")

        # ✅ Cập nhật toàn bộ list
        for idx, lang in enumerate(selected_lang_list):
            if lang == "zh-cn" and engine == "Polly":
                selected_lang_list[idx] = "zh"
                print(f"⚠ Đã tự động chuyển zh-cn ➜ zh (Polly) để đọc chuẩn hơn.")

        # ✅ Cập nhật lại combobox đang hiện
        for i, (dong, combo) in enumerate(combos):
            global_index = i + current_page * ITEMS_PER_PAGE
            combo.set(selected_lang_list[global_index])
    combo_engine.bind("<<ComboboxSelected>>", on_engine_change)

    def doc_popup():
        danh_sach_doc = process_text_lines(danh_sach, selected_lang_list)
        #danh_sach_doc = [(dong, combo.get().strip()) for dong, combo in combos]
        toc_do = combo_toc_do_popup.get()
        giong = combo_giong_popup.get()
        engine = combo_engine.get()

        def run_doc():
            global dung_doc_ngay, dang_doc
            dung_doc_ngay = False
            dang_doc = True
            count = 0

            for dong, lang in danh_sach_doc:
                if dung_doc_ngay: break
                try:
                    dong_sach = lam_sach_van_ban(dong)
                    if giong in {"Hội thoại 1 câu nam - 1 câu nữ", "Hội thoại 1 câu nữ - 1 câu nam"}:
                        first_voice, second_voice = _dialogue_voice_pair(giong)
                        voice = first_voice if count % 2 == 0 else second_voice
                        count += 1
                    else:
                        voice = giong

                    print(f"📢 Đọc ({lang}) [{voice}] Engine: {engine}")
                    temp_mp3 = os.path.join(tempfile.gettempdir(), f"temp_{uuid.uuid4().hex}.mp3")
                    tao_file_mp3(dong_sach, lang=lang, voice=voice, toc_do=toc_do, engine=engine, file_out=temp_mp3)

                    pygame.mixer.init()
                    sound = pygame.mixer.Sound(temp_mp3)
                    channel = pygame.mixer.find_channel()
                    if channel:
                        channel.play(sound)
                        while channel.get_busy():
                            if dung_doc_ngay:
                                channel.stop()
                                break
                            time.sleep(0.1)
                    os.remove(temp_mp3)
                except Exception as e:
                    print(f"❌ Lỗi đọc: {e}")
                    continue

            dang_doc = False
            print("✅ Kết thúc đọc nội dung.")

        threading.Thread(target=run_doc).start()

    def dung_doc():
        global dung_doc_ngay
        dung_doc_ngay = True

    def doc_lai_popup():
        dung_doc()
        doc_popup()

    def export_m4a_with_fallback(full_audio, output_path, progress_callback=None):
        import tempfile, os, subprocess, uuid

        temp_wav = os.path.join(tempfile.gettempdir(), f"temp_{uuid.uuid4().hex}.wav")
        try:
            # Update UI before exporting
            if progress_callback:
                progress_callback("Đang chuẩn bị audio...", 50)
            
            voice_audio = full_audio.set_channels(1).set_frame_rate(int(M4A_VOICE_SAMPLE_RATE))
            voice_audio.export(temp_wav, format="wav")

            if progress_callback:
                progress_callback("Đang chuyển đổi sang M4A...", 75)

            # Mono AAC bitrate thấp đủ rõ cho giọng nói và giảm mạnh dung lượng M4A.
            primary_cmd = [
                FFMPEG_PATH, "-y",
                "-i", temp_wav,
                "-c:a", "aac",
                "-b:a", M4A_VOICE_BITRATE,
                "-ac", "1",
                "-ar", M4A_VOICE_SAMPLE_RATE,
                "-movflags", "+faststart",
                output_path,
            ]
            result = subprocess.run(primary_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            if result.returncode == 0:
                if progress_callback:
                    progress_callback("Hoàn tất!", 100)
                return

            fallback_cmd = [
                FFMPEG_PATH, "-y",
                "-i", temp_wav,
                "-c:a", "aac",
                "-b:a", M4A_VOICE_FALLBACK_BITRATE,
                "-ac", "1",
                "-ar", M4A_VOICE_SAMPLE_RATE,
                "-movflags", "+faststart",
                output_path,
            ]
            result2 = subprocess.run(fallback_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            if result2.returncode != 0:
                raise Exception(result.stderr or result2.stderr or "FFmpeg không xuất được M4A.")
            if progress_callback:
                progress_callback("Hoàn tất!", 100)
        finally:
            try:
                if os.path.exists(temp_wav):
                    os.remove(temp_wav)
            except:
                pass

    def xuat_popup_audio(output_kind="mp3"):
        if not kiem_tra_ngon_ngu_hop_le(selected_lang_list, allowed_langs):
            return

        export_profiles = {
            "mp3": {
                "extension": ".mp3",
                "filetypes": [("MP3 files", "*.mp3")],
                "save_title": "Lưu file MP3",
                "progress_title": "Đang xuất MP3",
                "done_title": "Hoàn tất xuất MP3",
                "done_label": "🎉 Đã xuất file MP3:\n",
                "discord_message": "🎧 File MP3 mới được xuất từ Máy Học Tập",
                "discord_mime": "audio/mpeg",
                "export_kwargs": {"format": "mp3", "bitrate": "192k"},
            },
            "m4a": {
                "extension": ".m4a",
                "filetypes": [("M4A files", "*.m4a")],
                "save_title": "Lưu file M4A",
                "progress_title": "Đang xuất M4A",
                "done_title": "Hoàn tất xuất M4A",
                "done_label": "🎉 Đã xuất file M4A:\n",
                "discord_message": "🎧 File M4A mới được xuất từ Máy Học Tập",
                "discord_mime": "audio/mp4",
                "export_kwargs": None,
            },
        }

        spec = export_profiles.get(output_kind, export_profiles["mp3"])

        danh_sach_doc = process_text_lines(danh_sach, selected_lang_list)
        toc_do = combo_toc_do_popup.get()
        giong = combo_giong_popup.get()
        engine = combo_engine.get()

        file_path = filedialog.asksaveasfilename(
            defaultextension=spec["extension"],
            filetypes=spec["filetypes"],
            title=spec["save_title"],
        )
        if not file_path:
            return

        popup_progress = tk.Toplevel(popup)
        set_popup_icon(popup_progress)
        popup_progress.title(spec["progress_title"])
        popup_progress.geometry("420x120")
        popup_progress.grab_set()
        popup_progress.transient(popup)

        label_status = tk.Label(popup_progress, text="Bắt đầu...", font=("Arial", 11))
        label_status.pack(pady=5)
        progress = ttk.Progressbar(popup_progress, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=5)

        # 👉 Bắt đầu phát nhạc nền
        play_background_music()

        def thread_xuat():
            tong_dong = len(danh_sach_doc)

            try:
                def update_audio_progress(done, total):
                    progress["value"] = int(done / total * 100) if total else 100
                    label_status.config(text=f"Đang xuất dòng {done}/{tong_dong}")
                    popup_progress.update_idletasks()

                def update_export_progress(msg, pct):
                    """Update progress during export (M4A conversion)."""
                    progress["value"] = pct
                    label_status.config(text=msg)
                    popup_progress.update_idletasks()

                full_audio, _ = generate_audio_core(
                    danh_sach_doc,
                    giong=giong,
                    toc_do=toc_do,
                    engine=engine,
                    progress_callback=update_audio_progress,
                    clean_text_func=lam_sach_van_ban,
                    tts_func=tao_file_mp3,
                )
                export_audio_batch(full_audio, file_path, output_kind, export_m4a_with_fallback, export_progress_callback=update_export_progress)
                stop_background_music()
                popup_progress.destroy()

                def open_file():
                    try:
                        open_path_cross_platform(file_path)
                    except Exception as e:
                        tk.messagebox.showerror("Lỗi", f"Không mở được file:\n{e}")

                def open_folder():
                    try:
                        open_path_cross_platform(os.path.dirname(file_path))
                    except Exception as e:
                        tk.messagebox.showerror("Lỗi", f"Không mở được thư mục:\n{e}")

                def gui_discord():
                    try:
                        import requests
                        webhook_url = DISCORD_WEBHOOK_URL
                        if not webhook_url or "api/webhooks/" not in webhook_url:
                            tk.messagebox.showerror("Lỗi", "Webhook Discord không hợp lệ.")
                            return

                        with open(file_path, "rb") as f:
                            files = {"file": (os.path.basename(file_path), f, spec["discord_mime"])}
                            data = {"content": spec["discord_message"]}
                            response = requests.post(webhook_url, data=data, files=files)

                        if response.status_code in [200, 204]:
                            tk.messagebox.showinfo("OK", "✅ Đã gửi file lên Discord thành công!")
                        else:
                            tk.messagebox.showerror("Lỗi", f"Không gửi được Discord.\nHTTP {response.status_code}\n{response.text}")
                    except Exception as e:
                        thong_bao_loi_api(e, "Discord")
                        tk.messagebox.showerror("Lỗi", f"Không gửi Discord:\n{e}")

                popup_done = tk.Toplevel(popup)
                set_popup_icon(popup_done)

                pygame.mixer.init()
                pygame.mixer.music.load(SUCCESS_SOUND)
                pygame.mixer.music.play()
                gui_discord_thong_bao(f"🟢 Đã xuất xong file {output_kind.upper()}: {file_path}")

                popup_done.title(spec["done_title"])
                popup_done.geometry("480x180")
                popup_done.grab_set()
                popup_done.transient(popup)

                tk.Label(popup_done, text=spec["done_label"] + file_path, font=("Arial", 11, "bold"), fg="green").pack(pady=13)
                frm = tk.Frame(popup_done); frm.pack(pady=3)
                tk.Button(frm, text="Mở file", width=10, command=lambda: [popup_done.destroy(), open_file()]).pack(side="left", padx=6)
                tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup_done.destroy(), open_folder()]).pack(side="left", padx=6)
                tk.Button(frm, text="Gửi Discord", width=12, command=lambda: [popup_done.destroy(), gui_discord()]).pack(side="left", padx=6)
                tk.Button(frm, text="Gửi Zalo", width=12, command=lambda: [popup_done.destroy(), open_zalo_and_folder(file_path)]).pack(side="left", padx=6)
                tk.Button(popup_done, text="Đóng", command=popup_done.destroy).pack(pady=8)

            except Exception as e:
                stop_background_music()
                popup_progress.destroy()
                pygame.mixer.init()
                pygame.mixer.music.load(WARNING_SOUND)
                pygame.mixer.music.play()
                tk.messagebox.showerror("Lỗi", f"Xuất {output_kind.upper()} bị lỗi:\n{e}")
                gui_discord_thong_bao(f"🟢 Xuất {output_kind.upper()} bị lỗi: {file_path}")

        threading.Thread(target=thread_xuat, daemon=True).start()

    def xuat_popup_mp3():
        xuat_popup_audio("mp3")

    def xuat_popup_m4a():
        xuat_popup_audio("m4a")

    def xuat_popup_m4a_multifiles():
        if not kiem_tra_ngon_ngu_hop_le(selected_lang_list, allowed_langs):
            return

        use_default = messagebox.askyesno(
            "M4A MultiFiles",
            "Dùng mặc định 2 dòng cho mỗi file M4A?\nChọn No để nhập số dòng theo ý bạn.",
            parent=popup,
        )
        if use_default:
            so_dong_moi_file = 2
        else:
            so_dong_moi_file = simpledialog.askinteger(
                "M4A MultiFiles",
                "Nhập số dòng cho mỗi file M4A:",
                parent=popup,
                minvalue=1,
                initialvalue=2,
            )
            if not so_dong_moi_file:
                return

        parent_dir = filedialog.askdirectory(title="Chọn thư mục lưu M4A MultiFiles")
        if not parent_dir:
            return

        danh_sach_doc = [(dong, selected_lang_list[idx]) for idx, (dong, _) in enumerate(danh_sach)]
        if not danh_sach_doc:
            messagebox.showwarning("Trống", "Không có nội dung để xuất.", parent=popup)
            return

        toc_do = combo_toc_do_popup.get()
        giong = combo_giong_popup.get()
        engine = combo_engine.get()
        count = 0

        from datetime import datetime
        folder_name = f"M4A_MutiFiles_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        output_dir = os.path.join(parent_dir, folder_name)
        os.makedirs(output_dir, exist_ok=True)

        popup_progress = tk.Toplevel(popup)
        set_popup_icon(popup_progress)
        popup_progress.title("Đang xuất M4A MultiFiles")
        popup_progress.geometry("420x130")
        popup_progress.grab_set()
        popup_progress.transient(popup)

        label_status = tk.Label(popup_progress, text="Bắt đầu...", font=("Arial", 11))
        label_status.pack(pady=5)
        progress = ttk.Progressbar(popup_progress, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=5)

        play_background_music()

        def thread_multi():
            nonlocal count

            tong_dong = len(danh_sach_doc)
            tong_file = (tong_dong + so_dong_moi_file - 1) // so_dong_moi_file
            dong_da_xu_ly = 0
            files_da_tao = []

            try:
                for file_idx in range(tong_file):
                    bat_dau = file_idx * so_dong_moi_file
                    ket_thuc = min(bat_dau + so_dong_moi_file, tong_dong)
                    group = danh_sach_doc[bat_dau:ket_thuc]

                    start_processed = dong_da_xu_ly

                    def update_multi_progress(done, total):
                        current_done = start_processed + done
                        progress["value"] = int(current_done / tong_dong * 100) if tong_dong else 100
                        label_status.config(text=f"Đang xuất file {file_idx + 1}/{tong_file} - dòng {current_done}/{tong_dong}")
                        popup_progress.update_idletasks()

                    full_audio, count = generate_audio_core(
                        group,
                        giong=giong,
                        toc_do=toc_do,
                        engine=engine,
                        start_count=count,
                        initial_silence_ms=250,
                        progress_callback=update_multi_progress,
                        clean_text_func=lam_sach_van_ban,
                        tts_func=tao_file_mp3,
                    )
                    dong_da_xu_ly += len(group)

                    file_name = f"{folder_name}_{file_idx + 1:03d}.m4a"
                    file_path = os.path.join(output_dir, file_name)
                    
                    def update_export_progress(msg, pct):
                        progress["value"] = int(pct)
                        label_status.config(text=f"Đang xuất file {file_idx + 1}/{tong_file} - {msg}")
                        popup_progress.update_idletasks()
                    
                    export_m4a_with_fallback(full_audio, file_path, progress_callback=update_export_progress)
                    files_da_tao.append(file_path)

                stop_background_music()
                popup_progress.destroy()

                def open_folder():
                    try:
                        open_path_cross_platform(output_dir)
                    except Exception as e:
                        tk.messagebox.showerror("Lỗi", f"Không mở được thư mục:\n{e}")

                pygame.mixer.init()
                pygame.mixer.music.load(SUCCESS_SOUND)
                pygame.mixer.music.play()
                gui_discord_thong_bao(f"🟢 Đã xuất xong M4A MultiFiles: {output_dir}")

                popup_done = tk.Toplevel(popup)
                set_popup_icon(popup_done)
                popup_done.title("Hoàn tất xuất M4A MultiFiles")
                popup_done.geometry("500x190")
                popup_done.grab_set()
                popup_done.transient(popup)

                tk.Label(
                    popup_done,
                    text=f"🎉 Đã xuất {len(files_da_tao)} file M4A vào:\n" + output_dir,
                    font=("Arial", 11, "bold"),
                    fg="green",
                ).pack(pady=13)
                frm = tk.Frame(popup_done); frm.pack(pady=3)
                tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup_done.destroy(), open_folder()]).pack(side="left", padx=6)
                tk.Button(frm, text="Đóng", width=10, command=popup_done.destroy).pack(side="left", padx=6)

            except Exception as e:
                stop_background_music()
                popup_progress.destroy()
                pygame.mixer.init()
                pygame.mixer.music.load(WARNING_SOUND)
                pygame.mixer.music.play()
                tk.messagebox.showerror("Lỗi", f"Xuất M4A MultiFiles bị lỗi:\n{e}")

        threading.Thread(target=thread_multi, daemon=True).start()
    def run_xuat_srt_thread():
        print("Disabled: MP3+SRT subtitle export removed in audio-tool version.")
        try:
            messagebox.showinfo("Đã tắt", "Tính năng tạo phụ đề + MP3 đã được tắt.", parent=popup)
        except Exception:
            pass
        return

        if not kiem_tra_ngon_ngu_hop_le(selected_lang_list, allowed_langs):
            return

        danh_sach_doc = [(dong, selected_lang_list[idx]) for idx, (dong, _) in enumerate(danh_sach)]
        engine = combo_engine.get()
        giong = combo_giong_popup.get()
        toc_do = combo_toc_do_popup.get()
        count = 0

        file_path = filedialog.asksaveasfilename(defaultextension=".mp3", filetypes=[("MP3 files", "*.mp3")], title="Lưu file MP3 + SRT")
        if not file_path:
            return

        popup_progress = tk.Toplevel(popup)
        set_popup_icon(popup_progress)
        popup_progress.title("Đang xuất MP3 + SRT")
        popup_progress.geometry("420x120")
        popup_progress.grab_set()
        popup_progress.transient(popup)

        label_status = tk.Label(popup_progress, text="Bắt đầu...", font=("Arial", 11))
        label_status.pack(pady=5)
        progress = ttk.Progressbar(popup_progress, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=5)

        # 👉 Bắt đầu phát nhạc nền
        play_background_music()

        def thread_export():
            nonlocal count
            from pydub import AudioSegment
            import tempfile, os, uuid

            full_audio = AudioSegment.silent(duration=500)
            temp_dir = tempfile.gettempdir()
            current_time = 0.0
            srt_entries = []
            tong_dong = len(danh_sach_doc)

            try:
                for i, (dong, lang) in enumerate(danh_sach_doc):
                    dong_sach = lam_sach_van_ban(dong)

                    # Giọng
                    if giong in {"Hội thoại 1 câu nam - 1 câu nữ", "Hội thoại 1 câu nữ - 1 câu nam"}:
                        first_voice, second_voice = _dialogue_voice_pair(giong)
                        voice = first_voice if count % 2 == 0 else second_voice
                        count += 1
                    else:
                        voice = giong

                    temp_mp3 = os.path.join(temp_dir, f"temp_{uuid.uuid4().hex}.mp3")

                    try:
                        # Nếu rỗng → chèn im lặng 300ms, KHÔNG tạo SRT cho dòng này
                        if not dong_sach.strip():
                            segment = AudioSegment.silent(duration=300)
                            full_audio += segment + AudioSegment.silent(duration=300)
                            current_time += 0.3 + 0.3
                        else:
                            # Tạo file TTS
                            tao_file_mp3(dong_sach, lang=lang, voice=voice, toc_do=toc_do, engine=engine, file_out=temp_mp3)

                            # Nạp đoạn, lỗi → im lặng 300ms
                            try:
                                segment = AudioSegment.from_mp3(temp_mp3)
                            except Exception as e2:
                                print(f"⚠ Không nạp được mp3 tạm: {e2} → chèn im lặng 300ms")
                                segment = AudioSegment.silent(duration=300)

                            # Ghép âm
                            full_audio += segment + AudioSegment.silent(duration=300)

                            # Thời lượng & SRT (chỉ tạo khi có text hợp lệ)
                            duration_sec = getattr(segment, "duration_seconds", 0.3)
                            start_str = convert_seconds_to_timestamp(current_time)
                            end_str = convert_seconds_to_timestamp(current_time + duration_sec)
                            srt_entries.append(f"{len(srt_entries)+1}\n{start_str} --> {end_str}\n{dong_sach}\n\n")
                            current_time += duration_sec + 0.3
                    except Exception as e:
                        msg = str(e)
                        if "No text to speak" in msg or "No text to send to TTS API" in msg:
                            # Im lặng 300ms, không SRT
                            segment = AudioSegment.silent(duration=300)
                            full_audio += segment + AudioSegment.silent(duration=300)
                            current_time += 0.3 + 0.3
                        else:
                            raise
                    finally:
                        try:
                            if os.path.exists(temp_mp3):
                                os.remove(temp_mp3)
                        except:
                            pass

                    progress["value"] = int((i + 1) / tong_dong * 100)
                    label_status.config(text=f"Đang xuất {i+1}/{tong_dong}")
                    popup_progress.update_idletasks()

                # Xuất MP3 + SRT
                full_audio.export(file_path, format="mp3", bitrate="192k")
                file_srt = os.path.splitext(file_path)[0] + ".srt"
                with open(file_srt, "w", encoding="utf-8") as f:
                    f.writelines(srt_entries)

                stop_background_music()
                popup_progress.destroy()

                # 🎉 Popup xong
                def open_file():
                    try:
                        open_path_cross_platform(file_path)
                    except Exception as e:
                        tk.messagebox.showerror("Lỗi", f"Không mở được file:\n{e}")

                def open_folder():
                    try:
                        open_path_cross_platform(os.path.dirname(file_path))
                    except Exception as e:
                        tk.messagebox.showerror("Lỗi", f"Không mở được thư mục:\n{e}")

                def gui_discord():
                    try:
                        import requests
                        webhook_url = DISCORD_WEBHOOK_URL
                        if not webhook_url or "api/webhooks/" not in webhook_url:
                            tk.messagebox.showerror("Lỗi", "Webhook Discord không hợp lệ.")
                            return

                        with open(file_path, "rb") as f:
                            files = {"file": (os.path.basename(file_path), f, "audio/mpeg")}
                            data = {"content": "🎧 File MP3 mới được xuất từ Máy Học Tập"}
                            response = requests.post(webhook_url, data=data, files=files)

                        if response.status_code in [200, 204]:
                            tk.messagebox.showinfo("OK", "✅ Đã gửi file lên Discord thành công!")
                        else:
                            tk.messagebox.showerror("Lỗi", f"Không gửi được Discord.\nHTTP {response.status_code}\n{response.text}")
                    except Exception as e:
                        tk.messagebox.showerror("Lỗi", f"Không gửi Discord:\n{e}")

                popup_done = tk.Toplevel(popup)
                set_popup_icon(popup_done)
                pygame.mixer.init()
                pygame.mixer.music.load(SUCCESS_SOUND)
                pygame.mixer.music.play()
                gui_discord_thong_bao(f"🎙️ [TextToMp3] Đã xuất xong mp3 và phụ đề : {file_path}")

                popup_done.title("Hoàn tất xuất MP3 + SRT")
                popup_done.geometry("480x180")
                popup_done.grab_set()
                popup_done.transient(popup)

                tk.Label(popup_done, text="🎉 Đã xuất file MP3 + SRT:\n" + file_path, font=("Arial", 11, "bold"), fg="green").pack(pady=13)
                frm = tk.Frame(popup_done); frm.pack(pady=3)
                tk.Button(frm, text="Mở file", width=10, command=lambda: [popup_done.destroy(), open_file()]).pack(side="left", padx=6)
                tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup_done.destroy(), open_folder()]).pack(side="left", padx=6)
                tk.Button(frm, text="Gửi Discord", width=12, command=lambda: [popup_done.destroy(), gui_discord()]).pack(side="left", padx=6)
                tk.Button(frm, text="Gửi Zalo", width=12, command=lambda: [popup_done.destroy(), open_zalo_and_folder(file_path)]).pack(side="left", padx=6)
                tk.Button(popup_done, text="Đóng", command=popup_done.destroy).pack(pady=8)

            except Exception as e:
                stop_background_music()
                popup_progress.destroy()
                pygame.mixer.init()
                pygame.mixer.music.load(WARNING_SOUND)
                pygame.mixer.music.play()
                gui_discord_thong_bao(f"🎙️ [TextToMp3] Lỗi mp3 và phụ đề!!! : {file_path}")
                tk.messagebox.showerror("Lỗi", f"Xuất MP3 + SRT lỗi:\n{e}")

        threading.Thread(target=thread_export, daemon=True).start()


    def run_xuat_video_thread():
        print("Disabled: popup video export removed in audio-tool version.")
        try:
            messagebox.showinfo("Đã tắt", "Tính năng tạo video đã được tắt.", parent=popup)
        except Exception:
            pass
        return

        if not kiem_tra_ngon_ngu_hop_le(selected_lang_list, allowed_langs):
            return

        danh_sach_doc = [(dong, selected_lang_list[idx]) for idx, (dong, _) in enumerate(danh_sach)]
        engine = combo_engine.get()
        giong = combo_giong_popup.get()
        toc_do = combo_toc_do_popup.get()
        count = 0

        global background_path
        if not background_path or not os.path.isfile(background_path):
            background_path = LOGO_PATH

        file_path = filedialog.asksaveasfilename(defaultextension=".mp4", filetypes=[("MP4 files", "*.mp4")], title="Lưu file video")
        if not file_path:
            return

        popup_progress = tk.Toplevel(popup)
        set_popup_icon(popup_progress)
        popup_progress.title("Đang xuất Video")
        popup_progress.geometry("420x150")
        popup_progress.grab_set()
        popup_progress.transient(popup)

        label_status = tk.Label(popup_progress, text="Bắt đầu...", font=("Arial", 11))
        label_status.pack(pady=5)
        progress = ttk.Progressbar(popup_progress, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=5)

        stop_video_export = {"stop": False}
        def cancel_export_video():
            stop_video_export["stop"] = True
            label_status.config(text="⚠️ Đã yêu cầu huỷ, đang dừng...")

        btn_cancel = tk.Button(popup_progress, text="❌ Huỷ xuất video", fg="red", command=cancel_export_video)
        btn_cancel.pack(pady=5)

        def tao_anh_text(background_path, text, output_path):
            import textwrap
            from PIL import Image, ImageDraw, ImageFont

            img = Image.open(background_path).convert("RGB")
            draw = ImageDraw.Draw(img)
            font_size = 42
            try:
                font = ImageFont.truetype(FONT_CJK, font_size)
            except Exception as e:
                print(f"⚠ Không tải được font: {e}, dùng default")
                font = ImageFont.load_default()

            text = text or ""  # phụ đề rỗng vẫn render fine
            wrapped_text = "\n".join(textwrap.wrap(text, width=40)) if text else ""
            w, h = img.size
            if wrapped_text:
                lines = wrapped_text.split("\n")
                line_height = draw.textbbox((0, 0), "A", font=font)[3] + 10
                total_text_height = line_height * len(lines)
                y = h - total_text_height - 50
                for line in lines:
                    line_width = draw.textbbox((0, 0), line, font=font)[2]
                    x = (w - line_width) // 2
                    draw.text((x, y), line, font=font, fill="white", stroke_width=2, stroke_fill="black")
                    y += line_height

            # Fix kích thước chẵn cho x264
            w_new, h_new = img.size
            if w_new % 2 != 0 or h_new % 2 != 0:
                img = img.resize((w_new - w_new % 2, h_new - h_new % 2))

            img.save(output_path)

        def thread_video():
            import tempfile, uuid, subprocess, shutil
            from pydub import AudioSegment

            full_video_list = []
            temp_dir = tempfile.gettempdir()
            tong_dong = len(danh_sach_doc)
            nonlocal count

            # Chốt engine và startupinfo
            engine_now = combo_engine.get().lower()
            startupinfo = subprocess.STARTUPINFO()
            startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW

            # Helper: chuẩn hoá lang
            def normalize_lang(l):
                l = (l or "vi").lower()
                if l.startswith("vi"): return "vi"
                if l.startswith(("zh", "zh-cn", "zh_tw", "zh-hk")): return "zh"
                if l.startswith("ja"): return "ja"
                if l.startswith("en"): return "en"
                return "vi"

            # Helper: tạo mp3 an toàn từ text (có chèn im lặng nếu cần)
                def safe_tts_to_mp3(text_in, lang_in, voice_in, engine_in, out_mp3):
                    # Use central tao_file_mp3 which prefers Google Cloud for Vietnamese
                    lang_norm = normalize_lang(lang_in)
                    text_clean = lam_sach_van_ban(text_in or "")
                    if not text_clean.strip():
                        AudioSegment.silent(duration=300).export(out_mp3, format="mp3", bitrate="192k")
                        return "silence"

                    try:
                        # forward to tao_file_mp3; engine_in preserved so fallback rules apply
                        tao_file_mp3(text_clean, lang=lang_norm, voice=voice_in, toc_do=toc_do, engine=engine_in if engine_in else combo_engine.get(), file_out=out_mp3)
                        return "ok"
                    except Exception as e:
                        msg = str(e)
                        if "No text to speak" in msg or "No text to send to TTS API" in msg:
                            AudioSegment.silent(duration=300).export(out_mp3, format="mp3", bitrate="192k")
                            return "silence"
                        raise

            try:
                for i, (dong, lang) in enumerate(danh_sach_doc):
                    if stop_video_export["stop"]:
                        print("⚠️ Đã huỷ xuất video.")
                        break

                    dong_sach = lam_sach_van_ban(dong)

                    # Chọn giọng
                    if giong in {"Hội thoại 1 câu nam - 1 câu nữ", "Hội thoại 1 câu nữ - 1 câu nam"}:
                        first_voice, second_voice = _dialogue_voice_pair(giong)
                        voice_video = first_voice if count % 2 == 0 else second_voice
                        count += 1
                    else:
                        voice_video = giong

                    temp_mp3 = os.path.join(temp_dir, f"temp_{uuid.uuid4().hex}.mp3")
                    temp_img = os.path.join(temp_dir, f"temp_{uuid.uuid4().hex}.png")
                    temp_mp4 = os.path.join(temp_dir, f"temp_{uuid.uuid4().hex}.mp4")

                    try:
                        # 1) Âm thanh: an toàn với rỗng
                        status = safe_tts_to_mp3(dong_sach, lang, voice_video, engine_now, temp_mp3)

                        # +1.5s im lặng cuối đoạn cho dễ ghép
                        try:
                            segment = AudioSegment.from_mp3(temp_mp3)
                        except Exception as e2:
                            print(f"⚠ Không nạp được mp3 tạm: {e2} → chèn im lặng 300ms")
                            segment = AudioSegment.silent(duration=300)
                        segment += AudioSegment.silent(duration=1500)
                        segment.export(temp_mp3, format="mp3", bitrate="192k")

                        # 2) Ảnh phụ đề (nếu text rỗng → render ảnh không chữ)
                        tao_anh_text(background_path, dong_sach if status != "silence" else "", temp_img)

                        # 3) Tạo video đoạn
                        cmd = [
                            FFMPEG_PATH, "-y",
                            "-loop", "1",
                            "-i", temp_img.replace(os.sep, "/"),
                            "-i", temp_mp3.replace(os.sep, "/"),
                            "-c:v", "libx264", "-tune", "stillimage", "-c:a", "aac", "-b:a", "192k",
                            "-pix_fmt", "yuv420p",
                            "-shortest", temp_mp4
                        ]
                        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, startupinfo=startupinfo)
                        if result.returncode != 0:
                            print("⚠️ FFmpeg stderr:\n", result.stderr)
                            raise Exception(f"❌ Lỗi ffmpeg: {result.stderr}")

                        full_video_list.append(temp_mp4)
                    finally:
                        # Dọn rác tạm
                        for p in [temp_mp3, temp_img]:
                            try:
                                if os.path.exists(p): os.remove(p)
                            except: pass

                    progress["value"] = int((i + 1) / tong_dong * 100)
                    label_status.config(text=f"Đã ghép {i+1}/{tong_dong}")
                    popup_progress.update_idletasks()

                if stop_video_export["stop"]:
                    popup_progress.destroy()
                    tk.messagebox.showinfo("Đã huỷ", "Đã huỷ xuất video thành công.")
                    for vid in full_video_list:
                        try: os.remove(vid)
                        except: pass
                    return

                # Nối các đoạn
                list_file = os.path.join(temp_dir, f"list_{uuid.uuid4().hex}.txt")
                with open(list_file, "w", encoding="utf-8") as f:
                    for vid in full_video_list:
                        f.write(f"file '{vid.replace(os.sep, '/')}'\n")

                temp_output = tempfile.mktemp(suffix=".mp4")
                cmd_concat = [
                    FFMPEG_PATH, "-y",
                    "-f", "concat", "-safe", "0",
                    "-i", list_file,
                    "-c:v", "libx264",
                    "-c:a", "aac",
                    "-b:a", "192k",
                    "-pix_fmt", "yuv420p",
                    temp_output
                ]
                result = subprocess.run(cmd_concat, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, startupinfo=startupinfo)
                if result.returncode != 0:
                    print("⚠️ FFmpeg stderr:\n", result.stderr)
                    raise Exception(f"❌ Lỗi nối video: {result.stderr}")

                shutil.move(temp_output, file_path)
                try: os.remove(list_file)
                except: pass
                for vid in full_video_list:
                    try: os.remove(vid)
                    except: pass

                popup_progress.destroy()

                try:
                    pygame.mixer.init()
                    pygame.mixer.music.load(SUCCESS_SOUND)
                    pygame.mixer.music.play()
                except: pass

                popup_done = tk.Toplevel()
                pygame.mixer.init()
                pygame.mixer.music.load(SUCCESS_SOUND)
                pygame.mixer.music.play()
                gui_discord_thong_bao(f"🎙️ [TextToMp3] Đã xuất xong video: {file_path}")

                popup_done.title("✅ Hoàn tất xuất video")
                popup_done.geometry("480x180")
                popup_done.grab_set()
                set_popup_icon(popup_done)

                tk.Label(popup_done, text="✅ Đã xuất video:\n" + file_path, font=("Arial", 11, "bold"), fg="green").pack(pady=13)
                frm = tk.Frame(popup_done); frm.pack(pady=3)
                tk.Button(frm, text="Mở file", width=10, command=lambda: mo_file_an_toan(file_path)).pack(side="left", padx=6)
                tk.Button(frm, text="Mở thư mục", width=12, command=lambda: open_path_cross_platform(os.path.dirname(file_path))).pack(side="left", padx=6)
                tk.Button(frm, text="Đăng YouTube", width=12, command=lambda: popup_google_login(file_path)).pack(side="left", padx=6)
                tk.Button(popup_done, text="Đóng", command=popup_done.destroy).pack(pady=8)

            except Exception as e:
                popup_progress.destroy()
                try:
                    if os.path.exists(WARNING_SOUND):
                        pygame.mixer.music.load(WARNING_SOUND)
                        pygame.mixer.music.play()
                except: pass
                gui_discord_thong_bao(f"🎙️ [TextToMp3] Xuất video lỗi")
                tk.messagebox.showerror("Lỗi", f"Xuất video lỗi:\n{e}")

        threading.Thread(target=thread_video, daemon=True).start()

    def lay_sheet_names_tu_excel(excel_path):
        from openpyxl import load_workbook
        wb = load_workbook(excel_path, read_only=True, data_only=True)
        try:
            return list(wb.sheetnames)
        finally:
            wb.close()

    def tim_node_executable():
        system_node = shutil.which("node")
        if system_node:
            return system_node

        candidates = [
            os.path.join(BASE_DIR, ".tools", "node", "bin", "node"),
            os.path.join(BASE_DIR, ".tools", "node-v20.20.2-darwin-x64", "bin", "node"),
        ]
        for cand in candidates:
            if os.path.isfile(cand):
                return cand
        return None

    def import_excel_va_deploy_supabase():
        env_path = os.path.join(BASE_DIR, ".env")
        profile_state_path = os.path.join(APPDATA_ROOT, "supabase_import_profiles.json")

        profile_keys = [
            "SUPABASE_URL",
            "SUPABASE_SERVICE_ROLE_KEY",
            "SUPABASE_BUCKET",
            "SUPABASE_TABLE",
            "SUPABASE_STORAGE_FOLDER",
            "INSERT_BATCH_SIZE",
            "SKIP_EXISTING_UPLOAD",
            "CREATE_BUCKET_IF_MISSING",
            "ALLOW_INSERT_WHEN_UPLOAD_FAILED",
        ]

        def doc_env_dict():
            data = {}
            lines = []
            if os.path.exists(env_path):
                with open(env_path, "r", encoding="utf-8") as f:
                    lines = f.read().splitlines()

            for line in lines:
                stripped = line.strip()
                if not stripped or stripped.startswith("#") or "=" not in stripped:
                    continue
                key, value = stripped.split("=", 1)
                data[key.strip()] = value.strip()
            return data, lines

        def ghi_env_values(updates):
            _, lines = doc_env_dict()
            key_pos = {}
            for idx, line in enumerate(lines):
                stripped = line.strip()
                if not stripped or stripped.startswith("#") or "=" not in stripped:
                    continue
                key = stripped.split("=", 1)[0].strip()
                if key not in key_pos:
                    key_pos[key] = idx

            for key, value in updates.items():
                new_line = f"{key}={value}"
                if key in key_pos:
                    lines[key_pos[key]] = new_line
                else:
                    lines.append(new_line)

            with open(env_path, "w", encoding="utf-8") as f:
                f.write("\n".join(lines).rstrip() + "\n")

        def tao_profile_mac_dinh(env_values):
            return {
                "SUPABASE_URL": env_values.get("SUPABASE_URL", ""),
                "SUPABASE_SERVICE_ROLE_KEY": env_values.get("SUPABASE_SERVICE_ROLE_KEY", ""),
                "SUPABASE_BUCKET": env_values.get("SUPABASE_BUCKET", "audio"),
                "SUPABASE_TABLE": env_values.get("SUPABASE_TABLE", "vocab"),
                "SUPABASE_STORAGE_FOLDER": env_values.get("SUPABASE_STORAGE_FOLDER", ""),
                "INSERT_BATCH_SIZE": env_values.get("INSERT_BATCH_SIZE", "100"),
                "SKIP_EXISTING_UPLOAD": env_values.get("SKIP_EXISTING_UPLOAD", "true"),
                "CREATE_BUCKET_IF_MISSING": env_values.get("CREATE_BUCKET_IF_MISSING", "true"),
                "ALLOW_INSERT_WHEN_UPLOAD_FAILED": env_values.get("ALLOW_INSERT_WHEN_UPLOAD_FAILED", "false"),
            }

        def doc_profile_state(default_profile):
            if not os.path.exists(profile_state_path):
                return {
                    "active_profile": "dev",
                    "profiles": {
                        "dev": dict(default_profile),
                        "prod": dict(default_profile),
                    },
                    "last_excel": "",
                    "last_sheet": "",
                }

            try:
                with open(profile_state_path, "r", encoding="utf-8") as f:
                    raw = json.load(f)
            except Exception:
                raw = {}

            profiles = raw.get("profiles") if isinstance(raw, dict) else {}
            if not isinstance(profiles, dict):
                profiles = {}

            for pname in ["dev", "prod"]:
                if pname not in profiles or not isinstance(profiles[pname], dict):
                    profiles[pname] = dict(default_profile)

            for pname, cfg in list(profiles.items()):
                if not isinstance(cfg, dict):
                    profiles[pname] = dict(default_profile)
                    cfg = profiles[pname]
                for k in profile_keys:
                    if k not in cfg:
                        cfg[k] = default_profile.get(k, "")

            active_profile = raw.get("active_profile", "dev") if isinstance(raw, dict) else "dev"
            if active_profile not in profiles:
                active_profile = "dev"

            return {
                "active_profile": active_profile,
                "profiles": profiles,
                "last_excel": raw.get("last_excel", "") if isinstance(raw, dict) else "",
                "last_sheet": raw.get("last_sheet", "") if isinstance(raw, dict) else "",
            }

        def ghi_profile_state(state):
            try:
                with open(profile_state_path, "w", encoding="utf-8") as f:
                    json.dump(state, f, ensure_ascii=False, indent=2)
            except Exception as e:
                messagebox.showerror("Lỗi", f"Không lưu được profile state:\n{e}", parent=cfg_win)

        env_data, _ = doc_env_dict()
        default_profile = tao_profile_mac_dinh(env_data)
        profile_state = doc_profile_state(default_profile)

        cfg_win = tk.Toplevel(popup)
        set_popup_icon(cfg_win)
        cfg_win.title("Import Excel + Deploy Supabase")
        cfg_win.geometry("760x600")
        cfg_win.transient(popup)
        cfg_win.grab_set()

        main_frame = tk.Frame(cfg_win)
        main_frame.pack(fill="both", expand=True, padx=12, pady=10)

        excel_var = tk.StringVar()
        sheet_var = tk.StringVar()
        legacy_bitrate_var = tk.StringVar(value=_vocab_bitrate_display(config.get("VOCAB_M4A_BITRATE", DEFAULT_VOCAB_M4A_BITRATE)))
        legacy_audio_mode_var = tk.StringVar(value=_vocab_audio_mode_display(config.get("VOCAB_AUDIO_MODE", "zh_vi")))
        profile_var = tk.StringVar(value=profile_state["active_profile"])
        show_key_var = tk.BooleanVar(value=False)

        active_profile_cfg = profile_state["profiles"].get(profile_var.get(), default_profile)

        supabase_url_var = tk.StringVar(value=active_profile_cfg.get("SUPABASE_URL", ""))
        supabase_key_var = tk.StringVar(value=active_profile_cfg.get("SUPABASE_SERVICE_ROLE_KEY", ""))
        supabase_bucket_var = tk.StringVar(value=active_profile_cfg.get("SUPABASE_BUCKET", "audio"))
        supabase_table_var = tk.StringVar(value=active_profile_cfg.get("SUPABASE_TABLE", "vocab"))
        supabase_storage_folder_var = tk.StringVar(value=active_profile_cfg.get("SUPABASE_STORAGE_FOLDER", ""))
        insert_batch_var = tk.StringVar(value=active_profile_cfg.get("INSERT_BATCH_SIZE", "100"))

        skip_existing_var = tk.BooleanVar(value=str(active_profile_cfg.get("SKIP_EXISTING_UPLOAD", "true")).lower() != "false")
        create_bucket_var = tk.BooleanVar(value=str(active_profile_cfg.get("CREATE_BUCKET_IF_MISSING", "true")).lower() != "false")
        allow_insert_failed_var = tk.BooleanVar(
            value=str(active_profile_cfg.get("ALLOW_INSERT_WHEN_UPLOAD_FAILED", "false")).lower() == "true"
        )

        profile_row = tk.Frame(main_frame)
        profile_row.pack(fill="x", pady=(0, 8))
        tk.Label(profile_row, text="Profile project:", width=16, anchor="w", font=("Arial", 10, "bold")).pack(side="left")
        profile_combo = ttk.Combobox(
            profile_row,
            textvariable=profile_var,
            state="readonly",
            values=sorted(profile_state["profiles"].keys()),
            width=16,
        )
        profile_combo.pack(side="left")
        tk.Label(
            profile_row,
            text="(dev: chỉ pipeline+validate | prod: deploy đầy đủ)",
            fg="#666",
        ).pack(side="left", padx=(8, 0))

        tk.Label(main_frame, text="1) Chọn file Excel", font=("Arial", 11, "bold")).pack(anchor="w")
        file_row = tk.Frame(main_frame)
        file_row.pack(fill="x", pady=(4, 8))
        tk.Entry(file_row, textvariable=excel_var).pack(side="left", fill="x", expand=True)

        sheet_combo = ttk.Combobox(main_frame, textvariable=sheet_var, state="readonly")

        def derive_level_from_sheet(sheet_name):
            sheet_name = (sheet_name or "").strip().lower()
            match = re.search(r"hsk\d+", sheet_name)
            if match:
                return match.group(0)
            if "_" in sheet_name:
                return sheet_name.split("_", 1)[0]
            return sheet_name

        def sync_auto_fields_from_sheet(sheet_name):
            if advanced_mode_var.get():
                return
            level_name = derive_level_from_sheet(sheet_name)
            if level_name:
                supabase_storage_folder_var.set(level_name)

        def cap_nhat_sheet_list(excel_path, preferred_sheet=None):
            try:
                sheet_names = lay_sheet_names_tu_excel(excel_path)
            except Exception as e:
                messagebox.showerror("Lỗi", f"Không đọc được sheet từ Excel:\n{e}", parent=cfg_win)
                return

            if not sheet_names:
                messagebox.showwarning("Không có sheet", "File Excel không có sheet nào.", parent=cfg_win)
                return

            sheet_combo["values"] = sheet_names
            if preferred_sheet and preferred_sheet in sheet_names:
                sheet_var.set(preferred_sheet)
            else:
                sheet_var.set(sheet_names[0])
            sync_auto_fields_from_sheet(sheet_var.get())

        def browse_excel_file():
            selected_path = filedialog.askopenfilename(
                parent=cfg_win,
                title="Chọn file Excel",
                filetypes=[("Excel files", "*.xlsx;*.xlsm;*.xls")],
            )
            if not selected_path:
                return
            excel_var.set(selected_path)
            cap_nhat_sheet_list(selected_path)

        def on_sheet_change(event=None):
            sync_auto_fields_from_sheet(sheet_var.get())

        sheet_combo.bind("<<ComboboxSelected>>", on_sheet_change)

        tk.Button(file_row, text="Chọn Excel", width=12, command=browse_excel_file).pack(side="left", padx=(8, 0))

        tk.Label(main_frame, text="Sheet:").pack(anchor="w")
        sheet_combo.pack(fill="x", pady=(2, 10))

        quality_row = tk.Frame(main_frame)
        quality_row.pack(fill="x", pady=(0, 8))
        tk.Label(quality_row, text="Chất lượng M4A:", width=18, anchor="w").pack(side="left")
        ttk.Combobox(
            quality_row,
            textvariable=legacy_bitrate_var,
            values=("26 kbps", "32 kbps"),
            state="readonly",
            width=14,
        ).pack(side="left")
        tk.Label(quality_row, text="AAC-LC · mono · 22.05 kHz", fg="#666").pack(side="left", padx=(8, 0))

        audio_mode_row = tk.Frame(main_frame)
        audio_mode_row.pack(fill="x", pady=(0, 8))
        tk.Label(audio_mode_row, text="Nội dung audio:", width=18, anchor="w").pack(side="left")
        ttk.Combobox(
            audio_mode_row,
            textvariable=legacy_audio_mode_var,
            values=("Chỉ đọc tiếng Trung", "Đọc tiếng Trung + Tiếng Việt"),
            state="readonly",
            width=30,
        ).pack(side="left")

        tk.Label(main_frame, text=f"2) Cấu hình Supabase (lưu tại {env_path})", font=("Arial", 11, "bold")).pack(
            anchor="w"
        )

        form = tk.Frame(main_frame)
        form.pack(fill="x", pady=(6, 6))

        def add_field(label, var, show=None):
            row = tk.Frame(form)
            row.pack(fill="x", pady=2)
            tk.Label(row, text=label, width=24, anchor="w").pack(side="left")
            entry = tk.Entry(row, textvariable=var, show=show)
            entry.pack(side="left", fill="x", expand=True)
            return entry

        add_field("SUPABASE_URL", supabase_url_var)

        key_row = tk.Frame(form)
        key_row.pack(fill="x", pady=2)
        tk.Label(key_row, text="SERVICE_ROLE_KEY", width=24, anchor="w").pack(side="left")
        key_entry = tk.Entry(key_row, textvariable=supabase_key_var, show="*")
        key_entry.pack(side="left", fill="x", expand=True)

        def toggle_key_visibility():
            if show_key_var.get():
                key_entry.config(show="")
            else:
                key_entry.config(show="*")

        tk.Checkbutton(
            key_row,
            text="Hiện",
            variable=show_key_var,
            command=toggle_key_visibility,
        ).pack(side="left", padx=(8, 0))

        bucket_entry = add_field("SUPABASE_BUCKET", supabase_bucket_var)
        table_entry = add_field("SUPABASE_TABLE", supabase_table_var)
        storage_entry = add_field("STORAGE_FOLDER", supabase_storage_folder_var)
        batch_entry = add_field("INSERT_BATCH_SIZE", insert_batch_var)

        opt_row = tk.Frame(main_frame)
        opt_row.pack(fill="x", pady=(4, 4))
        cb_skip = tk.Checkbutton(opt_row, text="SKIP_EXISTING_UPLOAD", variable=skip_existing_var)
        cb_skip.pack(side="left", padx=(0, 12))
        cb_create_bucket = tk.Checkbutton(opt_row, text="CREATE_BUCKET_IF_MISSING", variable=create_bucket_var)
        cb_create_bucket.pack(side="left", padx=(0, 12))
        cb_allow_insert_failed = tk.Checkbutton(opt_row, text="ALLOW_INSERT_WHEN_UPLOAD_FAILED", variable=allow_insert_failed_var)
        cb_allow_insert_failed.pack(side="left")

        advanced_mode_var = tk.BooleanVar(value=False)

        def apply_advanced_mode():
            state = "normal" if advanced_mode_var.get() else "disabled"
            for widget in [bucket_entry, table_entry, storage_entry, batch_entry]:
                widget.config(state=state)
            for widget in [cb_skip, cb_create_bucket, cb_allow_insert_failed]:
                widget.config(state=state)

        tk.Checkbutton(
            main_frame,
            text="Bật chỉnh sửa nâng cao (mặc định để tự động an toàn)",
            variable=advanced_mode_var,
            command=apply_advanced_mode,
            fg="#444",
        ).pack(anchor="w", pady=(0, 2))
        apply_advanced_mode()

        hint_lbl = tk.Label(
            main_frame,
            text="Mặc định app dùng cấu hình tự động an toàn. STORAGE_FOLDER là folder trên Supabase, không phải folder output local.",
            fg="#666",
            anchor="w",
            justify="left",
        )
        hint_lbl.pack(fill="x", pady=(2, 8))

        def collect_profile_config_from_ui():
            return {
                "SUPABASE_URL": supabase_url_var.get().strip(),
                "SUPABASE_SERVICE_ROLE_KEY": supabase_key_var.get().strip(),
                "SUPABASE_BUCKET": supabase_bucket_var.get().strip(),
                "SUPABASE_TABLE": supabase_table_var.get().strip(),
                "SUPABASE_STORAGE_FOLDER": supabase_storage_folder_var.get().strip(),
                "INSERT_BATCH_SIZE": insert_batch_var.get().strip(),
                "SKIP_EXISTING_UPLOAD": "true" if skip_existing_var.get() else "false",
                "CREATE_BUCKET_IF_MISSING": "true" if create_bucket_var.get() else "false",
                "ALLOW_INSERT_WHEN_UPLOAD_FAILED": "true" if allow_insert_failed_var.get() else "false",
            }

        def apply_profile_to_ui(profile_cfg):
            supabase_url_var.set(profile_cfg.get("SUPABASE_URL", ""))
            supabase_key_var.set(profile_cfg.get("SUPABASE_SERVICE_ROLE_KEY", ""))
            supabase_bucket_var.set(profile_cfg.get("SUPABASE_BUCKET", "audio"))
            supabase_table_var.set(profile_cfg.get("SUPABASE_TABLE", "vocab"))
            supabase_storage_folder_var.set(profile_cfg.get("SUPABASE_STORAGE_FOLDER", ""))
            insert_batch_var.set(str(profile_cfg.get("INSERT_BATCH_SIZE", "100")))
            skip_existing_var.set(str(profile_cfg.get("SKIP_EXISTING_UPLOAD", "true")).lower() != "false")
            create_bucket_var.set(str(profile_cfg.get("CREATE_BUCKET_IF_MISSING", "true")).lower() != "false")
            allow_insert_failed_var.set(
                str(profile_cfg.get("ALLOW_INSERT_WHEN_UPLOAD_FAILED", "false")).lower() == "true"
            )
            sync_auto_fields_from_sheet(sheet_var.get())

        def luu_profile_hien_tai(show_ok=True):
            profile_name = profile_var.get().strip()
            if not profile_name:
                messagebox.showwarning("Thiếu profile", "Vui lòng chọn profile.", parent=cfg_win)
                return

            profile_state["profiles"][profile_name] = collect_profile_config_from_ui()
            profile_state["active_profile"] = profile_name
            ghi_profile_state(profile_state)
            if show_ok:
                messagebox.showinfo("Đã lưu", f"Đã lưu profile {profile_name}", parent=cfg_win)

        def on_profile_change(event=None):
            pname = profile_var.get().strip()
            cfg = profile_state["profiles"].get(pname)
            if not cfg:
                cfg = dict(default_profile)
                profile_state["profiles"][pname] = cfg
            apply_profile_to_ui(cfg)

        profile_combo.bind("<<ComboboxSelected>>", on_profile_change)
        tk.Button(profile_row, text="💾 Lưu profile", width=12, command=luu_profile_hien_tai).pack(side="left", padx=8)

        def validate_config_only():
            cfg = collect_profile_config_from_ui()
            insert_batch = cfg["INSERT_BATCH_SIZE"]
            if not insert_batch.isdigit() or int(insert_batch) <= 0:
                messagebox.showwarning("Sai cấu hình", "INSERT_BATCH_SIZE phải là số nguyên dương.", parent=cfg_win)
                return None

            required = {
                "SUPABASE_URL": cfg["SUPABASE_URL"],
                "SUPABASE_SERVICE_ROLE_KEY": cfg["SUPABASE_SERVICE_ROLE_KEY"],
                "SUPABASE_BUCKET": cfg["SUPABASE_BUCKET"],
                "SUPABASE_TABLE": cfg["SUPABASE_TABLE"],
            }
            missing = [k for k, v in required.items() if not v]
            if missing:
                messagebox.showwarning("Thiếu cấu hình", "Thiếu giá trị: " + ", ".join(missing), parent=cfg_win)
                return None
            return cfg

        def validate_and_collect():
            cfg = validate_config_only()
            if not cfg:
                return None

            excel_path = excel_var.get().strip()
            sheet_name = sheet_var.get().strip()

            if not excel_path or not os.path.isfile(excel_path):
                messagebox.showwarning("Thiếu file", "Vui lòng chọn file Excel hợp lệ.", parent=cfg_win)
                return None
            if not sheet_name:
                messagebox.showwarning("Thiếu sheet", "Vui lòng chọn sheet cần import.", parent=cfg_win)
                return None

            return {
                "excel_path": excel_path,
                "sheet_name": sheet_name,
                "PROFILE_MODE": profile_var.get().strip().lower(),
                **cfg,
            }

        def save_supabase_config(show_ok=True):
            collected = validate_and_collect()
            if not collected:
                return None

            updates = {
                "SUPABASE_URL": collected["SUPABASE_URL"],
                "SUPABASE_SERVICE_ROLE_KEY": collected["SUPABASE_SERVICE_ROLE_KEY"],
                "SUPABASE_BUCKET": collected["SUPABASE_BUCKET"],
                "SUPABASE_TABLE": collected["SUPABASE_TABLE"],
                "SUPABASE_STORAGE_FOLDER": collected["SUPABASE_STORAGE_FOLDER"],
                "INSERT_BATCH_SIZE": collected["INSERT_BATCH_SIZE"],
                "SKIP_EXISTING_UPLOAD": collected["SKIP_EXISTING_UPLOAD"],
                "CREATE_BUCKET_IF_MISSING": collected["CREATE_BUCKET_IF_MISSING"],
                "ALLOW_INSERT_WHEN_UPLOAD_FAILED": collected["ALLOW_INSERT_WHEN_UPLOAD_FAILED"],
            }

            try:
                ghi_env_values(updates)
                profile_state["active_profile"] = profile_var.get().strip() or "dev"
                profile_state["profiles"][profile_state["active_profile"]] = collect_profile_config_from_ui()
                profile_state["last_excel"] = excel_var.get().strip()
                profile_state["last_sheet"] = sheet_var.get().strip()
                ghi_profile_state(profile_state)
            except Exception as e:
                messagebox.showerror("Lỗi", f"Không lưu được file .env:\n{e}", parent=cfg_win)
                return None

            if show_ok:
                messagebox.showinfo("Đã lưu", "Đã lưu cấu hình Supabase vào .env", parent=cfg_win)
            return collected

        def test_ket_noi_supabase():
            cfg = validate_config_only()
            if not cfg:
                return

            import ssl
            from urllib import request as urlrequest
            from urllib import error as urlerror

            ssl_context = None
            try:
                import certifi

                ssl_context = ssl.create_default_context(cafile=certifi.where())
            except Exception:
                ssl_context = ssl.create_default_context()

            base_url = cfg["SUPABASE_URL"].rstrip("/")
            table_name = cfg["SUPABASE_TABLE"]
            test_url = f"{base_url}/rest/v1/{table_name}?select=*&limit=1"
            headers = {
                "apikey": cfg["SUPABASE_SERVICE_ROLE_KEY"],
                "Authorization": f"Bearer {cfg['SUPABASE_SERVICE_ROLE_KEY']}",
                "Accept": "application/json",
            }

            try:
                req = urlrequest.Request(test_url, headers=headers, method="GET")
                with urlrequest.urlopen(req, timeout=12, context=ssl_context) as resp:
                    code = resp.getcode()
                    if code in (200, 206):
                        messagebox.showinfo("Kết nối OK", "✅ Kết nối Supabase thành công.", parent=cfg_win)
                    else:
                        messagebox.showwarning("Kết nối", f"Nhận mã phản hồi {code}", parent=cfg_win)
            except urlerror.HTTPError as e:
                try:
                    body = e.read().decode("utf-8", errors="ignore")
                except Exception:
                    body = ""
                messagebox.showerror(
                    "Test kết nối thất bại",
                    f"HTTP {e.code}: {e.reason}\n{body[:400]}",
                    parent=cfg_win,
                )
            except Exception as e:
                msg = str(e)
                if "CERTIFICATE_VERIFY_FAILED" in msg or "certificate verify failed" in msg.lower():
                    msg += "\n\nGợi ý: Lỗi SSL cục bộ trên máy (CA store), không phải do sai key Supabase."
                messagebox.showerror("Test kết nối thất bại", msg, parent=cfg_win)

        def hoi_che_do_deploy(parent):
            selected = {"mode": None}

            chooser = tk.Toplevel(parent)
            set_popup_icon(chooser)
            chooser.title("Chọn chế độ deploy")
            chooser.geometry("520x210")
            chooser.transient(parent)
            chooser.grab_set()

            tk.Label(
                chooser,
                text="Bản này đã giống bản cũ. Bạn muốn làm gì với những phần đã có?",
                font=("Arial", 11, "bold"),
                wraplength=470,
                justify="center",
            ).pack(pady=(18, 8), padx=16)
            tk.Label(
                chooser,
                text="Ghi đè = tạo lại file local và upload lại storage.\nChỉ thêm phần chưa có = bỏ qua file/record đã tồn tại.",
                fg="#555",
                justify="center",
            ).pack(pady=(0, 14), padx=16)

            btn_row = tk.Frame(chooser)
            btn_row.pack(pady=4)

            def choose(mode):
                selected["mode"] = mode
                chooser.destroy()

            tk.Button(btn_row, text="Chỉ thêm phần chưa có", width=22, command=lambda: choose("add_missing")).pack(
                side="left", padx=6
            )
            tk.Button(btn_row, text="Ghi đè bản đã có", width=18, command=lambda: choose("overwrite")).pack(
                side="left", padx=6
            )
            tk.Button(chooser, text="Huỷ", width=10, command=lambda: choose(None)).pack(pady=(12, 0))
            chooser.protocol("WM_DELETE_WINDOW", lambda: choose(None))
            chooser.wait_window()
            return selected["mode"]

        def mo_cua_so_log_va_chay(collected, deploy_mode):
            sheet_name = collected["sheet_name"]
            excel_path = collected["excel_path"]
            level_match = re.search(r"hsk\d+", sheet_name.lower())
            level_name = level_match.group(0) if level_match else sheet_name.split("_")[0].lower()
            deck_name = sheet_name

            output_dir = os.path.join(BASE_DIR, "output", sheet_name)
            audio_dir = os.path.join(output_dir, "audio")
            json_path = os.path.join(output_dir, "output_vocab.json")
            storage_folder = collected["SUPABASE_STORAGE_FOLDER"] or level_name

            log_win = tk.Toplevel(cfg_win)
            set_popup_icon(log_win)
            log_win.title(f"Import Excel + Deploy: {sheet_name}")
            log_win.geometry("940x640")
            log_win.transient(cfg_win)
            log_win.grab_set()

            status_frame = tk.Frame(log_win)
            status_frame.pack(fill="x", padx=10, pady=(10, 6))

            def make_status_row(label_text):
                row = tk.Frame(status_frame)
                row.pack(fill="x", pady=1)
                tk.Label(row, text=label_text, width=28, anchor="w", font=("Arial", 10, "bold")).pack(side="left")
                v = tk.StringVar(value="Chưa chạy")
                lbl = tk.Label(row, textvariable=v, anchor="w", fg="#666")
                lbl.pack(side="left", fill="x", expand=True)
                return v, lbl

            st_m4a, lb_m4a = make_status_row("M4A:")
            st_json, lb_json = make_status_row("JSON:")
            st_check, lb_check = make_status_row("Check:")
            st_deploy, lb_deploy = make_status_row("Supabase Storage + DB:")

            progress_row = tk.Frame(log_win)
            progress_row.pack(fill="x", padx=10, pady=(0, 8))
            progress_label_var = tk.StringVar(value="Tiến trình: 0%")
            progress_count_var = tk.StringVar(value="Item: 0/0")
            tk.Label(progress_row, textvariable=progress_label_var, width=18, anchor="w", font=("Arial", 10, "bold")).pack(side="left")
            progress_bar = ttk.Progressbar(progress_row, orient="horizontal", length=650, mode="determinate", maximum=100)
            progress_bar.pack(side="left", fill="x", expand=True, padx=(8, 0))
            tk.Label(progress_row, textvariable=progress_count_var, width=14, anchor="e", font=("Arial", 10, "bold")).pack(side="right")

            log_text = tk.Text(log_win, wrap="word", height=26)
            log_text.pack(fill="both", expand=True, padx=10, pady=(0, 10))

            log_text.tag_configure("info", foreground="#0b4f8a")
            log_text.tag_configure("success", foreground="#0b7a28")
            log_text.tag_configure("warning", foreground="#b26a00")
            log_text.tag_configure("error", foreground="#b00020")
            log_text.tag_configure("header", foreground="#5b2c83", font=("Menlo", 12, "bold"))
            log_text.tag_configure("plain", foreground="#000000")
            log_text.tag_configure("phase_pipeline", foreground="#0b4f8a", background="#eaf4ff", font=("Menlo", 11, "bold"))
            log_text.tag_configure("phase_validate", foreground="#7a4f00", background="#fff4db", font=("Menlo", 11, "bold"))
            log_text.tag_configure("phase_deploy", foreground="#0b7a28", background="#e8f9ec", font=("Menlo", 11, "bold"))
            log_text.tag_configure("phase_done", foreground="#5b2c83", background="#f0e8ff", font=("Menlo", 11, "bold"))

            scroll = tk.Scrollbar(log_text, orient="vertical", command=log_text.yview)
            log_text.configure(yscrollcommand=scroll.set)
            scroll.pack(side="right", fill="y")

            control_row = tk.Frame(log_win)
            control_row.pack(fill="x", padx=10, pady=(0, 8))

            def append_log(msg):
                def _apply():
                    line = msg.rstrip("\n")
                    tag = "plain"
                    if line.startswith("==="):
                        tag = "header"
                    elif line.startswith("--- RUN: vocab_pipeline.py ---"):
                        tag = "phase_pipeline"
                    elif line.startswith("--- RUN: import_hsk1_to_supabase.js ---"):
                        tag = "phase_deploy"
                    elif line.startswith("[Pipeline] Running validator") or line.startswith("[Pipeline] Validator"):
                        tag = "phase_validate"
                    elif line.startswith("[Pipeline] Finished successfully") or line.startswith("✅ Hoàn tất"):
                        tag = "phase_done"
                    elif line.startswith("---"):
                        tag = "header"
                    elif line.startswith("❌") or "FAIL" in line or "lỗi" in line.lower() or "error" in line.lower():
                        tag = "error"
                    elif line.startswith("⚠") or line.startswith("[WARN]"):
                        tag = "warning"
                    elif line.startswith("✅") or "PASS" in line or "đã tạo đủ" in line.lower() or "hoàn tất" in line.lower():
                        tag = "success"
                    elif line.startswith("ℹ️") or line.startswith("[Pipeline]"):
                        tag = "info"

                    log_text.insert("end", line + "\n", tag)
                    log_text.see("end")

                log_text.after(0, _apply)

            def update_status(var, label_widget, text, color):
                def _apply():
                    var.set(text)
                    label_widget.config(fg=color)
                log_win.after(0, _apply)

            progress_state = {"total": 0, "processed": 0}

            def update_progress(percent, label=None):
                percent = max(0, min(100, int(percent)))

                def _apply():
                    progress_bar["value"] = percent
                    progress_label_var.set(label or f"Tiến trình: {percent}%")
                    total = progress_state.get("total", 0)
                    processed = progress_state.get("processed", 0)
                    if total:
                        progress_count_var.set(f"Item: {processed}/{total}")
                    else:
                        progress_count_var.set("Item: 0/0")

                log_win.after(0, _apply)

            def update_pipeline_progress_from_line(line):
                total_match = re.search(r"\[Pipeline\] Valid rows:\s*(\d+)", line)
                if total_match:
                    progress_state["total"] = int(total_match.group(1))
                    update_progress(2, f"Tiến trình: 2% - {progress_state['total']} dòng hợp lệ")
                    return

                gen_match = re.search(r"\[Pipeline\] Progress:\s*processed\s*(\d+)\/(\d+)\s*items", line)
                if gen_match:
                    progress_state["processed"] = int(gen_match.group(1))
                    progress_state["total"] = max(progress_state["total"], int(gen_match.group(2)))
                    pct = int(5 + (progress_state["processed"] / max(1, progress_state["total"])) * 80)
                    update_progress(pct, f"Tiến trình: {pct}% - {progress_state['processed']}/{progress_state['total']} items")
                    return

                if "Writing JSON" in line:
                    update_progress(88, "Tiến trình: 88% - đang ghi JSON")
                elif "Running validator" in line:
                    progress_count_var.set(f"Item: {progress_state.get('processed', 0)}/{progress_state.get('total', 0)}")
                    update_progress(92, "Tiến trình: 92% - đang validate")
                elif "Validator passed" in line:
                    update_progress(95, "Tiến trình: 95% - validate PASS")
                elif "Finished successfully" in line:
                    update_progress(100, "Tiến trình: 100% - hoàn tất")

            def update_phase_progress(percent, label):
                update_progress(percent, label)

            cancel_state = {"cancelled": False}
            active_process = {"proc": None}

            def request_cancel():
                cancel_state["cancelled"] = True
                proc = active_process["proc"]
                if proc is not None:
                    try:
                        proc.terminate()
                    except Exception:
                        try:
                            proc.kill()
                        except Exception:
                            pass
                update_status(st_deploy, lb_deploy, "Đã huỷ", "red")
                append_log("\n⛔ Người dùng đã bấm Huỷ. Đang dừng tiến trình...")

            tk.Button(control_row, text="❌ Huỷ", width=12, fg="white", bg="#cc4444", command=request_cancel).pack(
                side="left"
            )
            tk.Label(control_row, text="Nút Huỷ sẽ dừng pipeline/deploy đang chạy.", fg="#666").pack(side="left", padx=8)

            def run_cmd_and_stream(cmd, env=None):
                output_lines = []
                try:
                    process = subprocess.Popen(
                        cmd,
                        cwd=BASE_DIR,
                        env=env,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.STDOUT,
                        text=True,
                        bufsize=1,
                        universal_newlines=True,
                    )
                    active_process["proc"] = process
                except Exception as e:
                    append_log(f"❌ Không chạy được lệnh: {e}")
                    return 1, ""

                if process.stdout is not None:
                    for line in process.stdout:
                        if cancel_state["cancelled"]:
                            break
                        clean = line.rstrip("\n")
                        output_lines.append(clean)
                        append_log(clean)
                        update_pipeline_progress_from_line(clean)

                code = process.wait()
                active_process["proc"] = None
                if cancel_state["cancelled"]:
                    return -1, "\n".join(output_lines)
                return code, "\n".join(output_lines)

            def open_output_folder():
                try:
                    open_path_cross_platform(output_dir)
                except Exception as e:
                    messagebox.showerror("Lỗi", f"Không mở được folder kết quả:\n{e}", parent=log_win)

            def run_deploy_phase(profile_mode, deploy_mode):
                node_exe = tim_node_executable()
                if not node_exe:
                    update_status(st_deploy, lb_deploy, "FAIL (không tìm thấy Node.js)", "red")
                    append_log("\n❌ Không tìm thấy Node.js để chạy script deploy Supabase.")
                    return

                env = os.environ.copy()
                env.update(
                    {
                        "SUPABASE_URL": collected["SUPABASE_URL"],
                        "SUPABASE_SERVICE_ROLE_KEY": collected["SUPABASE_SERVICE_ROLE_KEY"],
                        "SUPABASE_BUCKET": collected["SUPABASE_BUCKET"],
                        "SUPABASE_TABLE": collected["SUPABASE_TABLE"],
                        "INSERT_BATCH_SIZE": collected["INSERT_BATCH_SIZE"],
                        "OVERWRITE_EXISTING_UPLOAD": "true" if deploy_mode == "overwrite" else "false",
                        "SKIP_EXISTING_UPLOAD": collected["SKIP_EXISTING_UPLOAD"],
                        "CREATE_BUCKET_IF_MISSING": collected["CREATE_BUCKET_IF_MISSING"],
                        "ALLOW_INSERT_WHEN_UPLOAD_FAILED": collected["ALLOW_INSERT_WHEN_UPLOAD_FAILED"],
                        "JSON_PATH": json_path,
                        "AUDIO_DIR": audio_dir,
                        "LEVEL": level_name,
                        "DECK": deck_name,
                        "SUPABASE_STORAGE_FOLDER": storage_folder,
                    }
                )

                import_script = os.path.join(BASE_DIR, "scripts", "import_hsk1_to_supabase.js")
                import_cmd = [node_exe, import_script]

                mode_label = "ghi đè" if deploy_mode == "overwrite" else "chỉ thêm phần chưa có"
                update_status(st_deploy, lb_deploy, f"Đang deploy ({mode_label})...", "#cc8800")
                append_log("\n--- RUN: import_hsk1_to_supabase.js ---")
                deploy_code, deploy_output = run_cmd_and_stream(import_cmd, env=env)

                if deploy_code == 0:
                    if cancel_state["cancelled"]:
                        update_status(st_deploy, lb_deploy, "Đã huỷ", "red")
                        append_log("\n⛔ Deploy đã bị huỷ trước khi hoàn tất.")
                        return

                    no_new_records = deploy_mode != "overwrite" and "STATUS: NO_NEW_RECORDS" in deploy_output
                    if no_new_records:
                        update_status(st_deploy, lb_deploy, "Bản này đã deploy đủ rồi", "#b26a00")
                        append_log("\nℹ️ Bản này đã deploy đủ rồi, không có gì mới để đẩy lên.")
                    elif deploy_mode == "overwrite":
                        update_status(st_deploy, lb_deploy, "Đã ghi đè xong", "green")
                        append_log("\n✅ Đã ghi đè xong bản đã có.")
                    else:
                        update_status(st_deploy, lb_deploy, "đã deploy", "green")
                        append_log("\n✅ Hoàn tất toàn bộ: M4A + JSON + PASS + Supabase deploy")

                    done_popup = tk.Toplevel(log_win)
                    set_popup_icon(done_popup)
                    done_popup.title("Hoàn tất import")
                    done_popup.geometry("460x190")
                    done_popup.transient(log_win)
                    done_popup.grab_set()

                    if no_new_records:
                        message = "Bản này đã deploy đủ rồi, không có gì mới để đẩy lên."
                        message_color = "#b26a00"
                    else:
                        message = "✅ Đã deploy xong Supabase." if profile_mode == "prod" else "✅ Đã deploy xong sau khi kiểm tra."
                        message_color = "green"

                    tk.Label(done_popup, text=message, font=("Arial", 11, "bold"), fg=message_color).pack(pady=(16, 10))
                    if profile_mode == "dev":
                        if deploy_mode == "overwrite":
                            hint_text = "Bạn đang ở chế độ ghi đè, có thể mở folder để kiểm tra lại kết quả."
                        else:
                            hint_text = "Bạn hãy mở folder để kiểm tra trước khi deploy." if not no_new_records else "Bạn có thể mở folder để đối chiếu lại kết quả."
                        tk.Label(done_popup, text=hint_text, fg="#444").pack(pady=(0, 8))

                    btn_frame = tk.Frame(done_popup)
                    btn_frame.pack(pady=6)

                    if profile_mode == "dev":
                        tk.Button(btn_frame, text="Mở folder", width=12, command=open_output_folder).pack(side="left", padx=6)
                        tk.Button(
                            btn_frame,
                            text="Đã kiểm tra xong- Deploy ngay",
                            width=22,
                            command=lambda: start_manual_deploy(done_popup),
                        ).pack(side="left", padx=6)
                        tk.Button(done_popup, text="Huỷ", width=10, command=done_popup.destroy).pack(pady=(8, 0))
                    else:
                        tk.Button(btn_frame, text="Mở folder", width=12, command=open_output_folder).pack(side="left", padx=6)
                        tk.Button(done_popup, text="Đóng", width=10, command=done_popup.destroy).pack(pady=(8, 0))
                else:
                    if cancel_state["cancelled"]:
                        update_status(st_deploy, lb_deploy, "Đã huỷ", "red")
                        append_log("\n⛔ Đã huỷ bởi người dùng.")
                        return
                    update_status(st_deploy, lb_deploy, "deploy lỗi", "red")
                    append_log("\n❌ Deploy Supabase bị lỗi.")

            def worker():
                profile_mode = (collected.get("PROFILE_MODE") or "dev").lower()
                append_log("=== BẮT ĐẦU IMPORT EXCEL + DEPLOY SUPABASE ===")
                append_log(f"Excel: {excel_path}")
                append_log(f"Sheet: {sheet_name}")
                append_log(f"Level: {level_name}")
                append_log(f"Deck: {deck_name}")
                append_log(f"Profile mode: {profile_mode}")
                append_log(f"Storage folder: {storage_folder}")
                append_log(f"Config file: {env_path}")
                vocab_tts = collected["VOCAB_TTS_CONFIG"]
                append_log(
                    "Vocab TTS: "
                    f"engine={vocab_tts['engine']} | speed={vocab_tts['speed']} | "
                    f"voice={vocab_tts['voice']} | M4A=AAC-LC mono 22050Hz {vocab_tts['bitrate']} | "
                    f"languages={','.join(vocab_tts['languages'])}"
                )

                update_status(st_m4a, lb_m4a, "Đang chạy...", "#cc8800")
                update_status(st_json, lb_json, "Đang chạy...", "#cc8800")
                update_status(st_check, lb_check, "Đang chạy...", "#cc8800")
                update_status(st_deploy, lb_deploy, "Đang chờ pipeline xong...", "#666")
                update_progress(1, "Tiến trình: 1% - khởi động")

                pipeline_cmd = [
                    sys.executable,
                    "-u",
                    os.path.join(BASE_DIR, "pipelines", "vocab_pipeline.py"),
                    excel_path,
                    "--sheet",
                    sheet_name,
                ]

                pipeline_env = os.environ.copy()
                apply_vocab_tts_env(pipeline_env, vocab_tts)
                pipeline_env["GOOGLE_TTS_PROFILES_JSON"] = json.dumps(GOOGLE_TTS_PROFILES, ensure_ascii=False)
                pipeline_env["GOOGLE_TTS_API_KEY"] = GOOGLE_TTS_API_KEY
                pipeline_env["AWS_ACCESS_KEY_ID"] = AWS_ACCESS_KEY_ID
                pipeline_env["AWS_SECRET_ACCESS_KEY"] = AWS_SECRET_ACCESS_KEY
                pipeline_env["AWS_REGION"] = AWS_REGION
                pipeline_env["AWS_DEFAULT_REGION"] = AWS_REGION

                append_log("\n--- RUN: vocab_pipeline.py ---")
                pipe_code, pipe_output = run_cmd_and_stream(pipeline_cmd, env=pipeline_env)
                thong_bao_loi_api(pipe_output, "pipeline TTS/AWS")
                if cancel_state["cancelled"]:
                    update_status(st_m4a, lb_m4a, "Đã huỷ", "red")
                    update_status(st_json, lb_json, "Đã huỷ", "red")
                    update_status(st_check, lb_check, "Đã huỷ", "red")
                    update_status(st_deploy, lb_deploy, "Đã huỷ", "red")
                    append_log("\n⛔ Quy trình bị huỷ giữa chừng.")
                    return
                if pipe_code != 0:
                    update_status(st_m4a, lb_m4a, "Lỗi tạo M4A", "red")
                    update_status(st_json, lb_json, "Lỗi tạo JSON", "red")
                    update_status(st_check, lb_check, "FAIL", "red")
                    update_status(st_deploy, lb_deploy, "Chưa deploy do pipeline lỗi", "red")
                    update_progress(100, "Tiến trình: lỗi")
                    append_log("\n❌ Pipeline thất bại. Dừng quy trình.")
                    return

                m4a_count = 0
                if os.path.isdir(audio_dir):
                    try:
                        m4a_count = len([f for f in os.listdir(audio_dir) if f.lower().endswith(".m4a")])
                    except Exception:
                        m4a_count = 0

                if m4a_count > 0:
                    update_status(st_m4a, lb_m4a, f"đã tạo đủ ({m4a_count} file)", "green")
                else:
                    update_status(st_m4a, lb_m4a, "không tìm thấy file m4a", "red")

                if os.path.isfile(json_path):
                    update_status(st_json, lb_json, "đã tạo đủ", "green")
                else:
                    update_status(st_json, lb_json, "không tìm thấy output_vocab.json", "red")

                if "STATUS: PASS" in pipe_output:
                    update_status(st_check, lb_check, "PASS", "green")
                    update_progress(90, "Tiến trình: 90% - validate PASS")
                else:
                    update_status(st_check, lb_check, "FAIL (validator không PASS)", "red")
                    update_status(st_deploy, lb_deploy, "Chưa deploy do validate FAIL", "red")
                    update_progress(100, "Tiến trình: lỗi validate")
                    append_log("\n❌ Validation không PASS. Dừng trước bước deploy.")
                    return

                if profile_mode == "dev":
                    update_status(st_deploy, lb_deploy, "Chờ bạn kiểm tra folder", "#2a6")
                    update_progress(95, "Tiến trình: 95% - chờ bạn kiểm tra")
                    append_log("\nℹ️ Dev mode: chỉ chạy pipeline + validate, không deploy Supabase.")
                    append_log("\n✅ Hoàn tất phần kiểm tra dev.")

                    done_popup = tk.Toplevel(log_win)
                    set_popup_icon(done_popup)
                    done_popup.title("Đã kiểm tra xong")
                    done_popup.geometry("480x190")
                    done_popup.transient(log_win)
                    done_popup.grab_set()

                    tk.Label(done_popup, text="✅ Dev đã chạy xong pipeline + validate.", font=("Arial", 11, "bold"), fg="green").pack(pady=(16, 10))
                    tk.Label(done_popup, text="Hãy mở folder để kiểm tra kết quả trước khi deploy.", fg="#444").pack(pady=(0, 8))

                    btn_frame = tk.Frame(done_popup)
                    btn_frame.pack(pady=6)
                    tk.Button(btn_frame, text="Mở folder", width=12, command=open_output_folder).pack(side="left", padx=6)
                    tk.Button(
                        btn_frame,
                        text="Đã kiểm tra xong- Deploy ngay",
                        width=22,
                            command=lambda: [done_popup.destroy(), run_deploy_phase("prod", deploy_mode)],
                    ).pack(side="left", padx=6)
                    tk.Button(done_popup, text="Huỷ", width=10, command=done_popup.destroy).pack(pady=(8, 0))
                    return

                run_deploy_phase(profile_mode, deploy_mode)

            def start_manual_deploy(done_popup=None):
                if done_popup is not None and done_popup.winfo_exists():
                    done_popup.destroy()
                update_status(st_deploy, lb_deploy, "Đang khởi động deploy...", "#cc8800")
                append_log("\n🚀 Bấm Deploy ngay: bắt đầu deploy Supabase ở background.")
                threading.Thread(target=lambda: run_deploy_phase("prod", deploy_mode), daemon=True).start()

            threading.Thread(target=worker, daemon=True).start()

        def on_run():
            vocab_tts = collect_vocab_tts_config(cfg_win, legacy_bitrate_var.get(), legacy_audio_mode_var.get())
            if not vocab_tts:
                return
            deploy_mode = hoi_che_do_deploy(cfg_win)
            if not deploy_mode:
                return
            collected = save_supabase_config(show_ok=False)
            if not collected:
                return
            collected["VOCAB_TTS_CONFIG"] = vocab_tts
            mo_cua_so_log_va_chay(collected, deploy_mode)

        # Nạp lại lựa chọn Excel/sheet gần nhất
        last_excel = profile_state.get("last_excel", "")
        last_sheet = profile_state.get("last_sheet", "")
        if last_excel and os.path.isfile(last_excel):
            excel_var.set(last_excel)
            cap_nhat_sheet_list(last_excel, preferred_sheet=last_sheet)

        btn_row = tk.Frame(main_frame)
        btn_row.pack(fill="x", pady=(8, 0))
        tk.Button(btn_row, text="💾 Lưu cấu hình", width=14, bg="#e6ffe6", command=save_supabase_config).pack(
            side="left"
        )
        tk.Button(btn_row, text="🔌 Test kết nối", width=14, bg="#fff2cc", command=test_ket_noi_supabase).pack(
            side="left", padx=(8, 0)
        )
        tk.Button(btn_row, text="🚀 Chạy Import + Deploy", width=20, bg="#cce6ff", command=on_run).pack(
            side="left", padx=8
        )
        tk.Button(btn_row, text="Đóng", width=10, command=cfg_win.destroy).pack(side="right")


    def hsk30_vocab_zip_builder():
        """Local-only generic HSK 2.0/3.0 vocab pack builder."""
        from pipelines.vocab_zip_deploy import (
            CATALOG_PUBLISH_CONFIRMATION,
            DeployValidationError,
            SupabaseStorageRestClient,
            STAGING_BUCKET,
            build_plan,
            collect_deploy_receipts,
            catalog_matches_receipt,
            input_fingerprint,
            prepare_catalog_publish,
            stage_confirmation_phrase,
            stage_packs_with_client,
        )
        from pipelines.vocab_catalog_publish import (
            DEFAULT_KEY_ID,
            DEFAULT_KEY_PATH,
            initialize_signing_key,
            initialize_pointer_with_client,
            publish_signed_catalog_with_client,
            read_verified_pointer_status,
            signing_status,
        )
        from pipelines.hsk30_help import load_help_text

        builder_win = tk.Toplevel(popup)
        set_popup_icon(builder_win)
        builder_win.title("HSK 2.0 / 3.0 Vocab ZIP Builder")
        # Give the builder enough vertical room for the phase status, receipt,
        # pointer controls and the three workflow buttons on macOS.
        builder_win.geometry("1460x1100")
        builder_win.minsize(1180, 860)
        builder_win.minsize(1280, 860)
        builder_win.resizable(True, True)
        builder_win.transient(popup)
        builder_win.grab_set()

        last_hsk30_excel, last_hsk30_sheet = _load_hsk30_recent_selection()
        last_hsk30_builder_state = _load_hsk30_builder_state()
        excel_var = tk.StringVar(value=last_hsk30_excel if os.path.isfile(last_hsk30_excel) else "")
        sheet_var = tk.StringVar()
        version_label_to_code = {"HSK 2.0": "2.0", "HSK 3.0": "3.0"}
        version_display_var = tk.StringVar(value=last_hsk30_builder_state["version"] if last_hsk30_builder_state["version"] in version_label_to_code else "HSK 3.0")
        level_label_to_code = {
            "HSK 1": "hsk1", "HSK 2": "hsk2", "HSK 3": "hsk3", "HSK 4": "hsk4",
            "HSK 5": "hsk5", "HSK 6": "hsk6", "HSK 7–9": "hsk7_9",
        }
        level_display_var = tk.StringVar(value=last_hsk30_builder_state["level"] if last_hsk30_builder_state["level"] in level_label_to_code else "HSK 1")
        pack_version_var = tk.StringVar(value=last_hsk30_builder_state["pack_version"] if str(last_hsk30_builder_state["pack_version"]).isdigit() and int(last_hsk30_builder_state["pack_version"]) >= 1 else "1")
        pack_version_label_var = tk.StringVar(value=f"Pack version: v{pack_version_var.get()}")
        output_var = tk.StringVar(value=last_hsk30_builder_state["output_dir"] or os.path.join(BASE_DIR, "output"))
        builder_bitrate_var = tk.StringVar(value=last_hsk30_builder_state["bitrate"] or _vocab_bitrate_display(config.get("VOCAB_M4A_BITRATE", DEFAULT_VOCAB_M4A_BITRATE)))
        builder_audio_mode_var = tk.StringVar(value=last_hsk30_builder_state["audio_mode"] or _vocab_audio_mode_display(config.get("VOCAB_AUDIO_MODE", "zh_vi")))
        force_audio_var = tk.BooleanVar(value=False)
        status_var = tk.StringVar(value="Sẵn sàng build local. Phase 2 chỉ chạy sau local PASS và xác nhận.")
        summary_var = tk.StringVar(value="Chưa đọc Excel")
        workflow_state_var = tk.StringVar(value="LOCAL NOT BUILT | PACKS NOT VERIFIED | CATALOG NOT PUBLISHED")
        compatibility_var = tk.StringVar(value="Compatibility hash: chưa verify")
        signing_key_state_var = tk.StringVar(value="SIGNING KEY NOT INITIALIZED")
        signing_public_key_var = tk.StringVar(value="PUBLIC KEY B64: —")
        pointer_state_var = tk.StringVar(value="POINTER STATUS NOT REFRESHED")
        pointer_revision_var = tk.StringVar(value="CURRENT POINTER REVISION: — | CURRENT CATALOG REVISION: —")
        publish_gate_var = tk.StringVar(value="Publish disabled: pointer status chưa được refresh")
        receipt_state_var = tk.StringVar(value="Receipt: chưa có | Remote verification: chưa có")
        artifact_state = {"result": None, "fingerprint": None}
        pointer_status_cache = {"status": "NOT_REFRESHED"}
        build_runtime = {
            "process": None,
            "preview_process": None,
            "paused": False,
            "cancel_requested": False,
        }

        def open_hsk30_help(anchor_phrase=None):
            help_text, help_path, loaded = load_help_text(Path(__file__).resolve().parent)
            help_win = tk.Toplevel(builder_win)
            set_popup_icon(help_win)
            help_win.title("Hướng dẫn HSK 2.0 / 3.0 Vocab ZIP Builder")
            help_win.geometry("900x720")
            help_win.minsize(620, 420)
            help_win.transient(builder_win)
            help_win.grab_set()
            header = tk.Frame(help_win)
            header.pack(fill="x", padx=12, pady=(10, 4))
            tk.Label(header, text="Tài liệu nguồn:", font=("Arial", 10, "bold")).pack(side="left")
            tk.Label(header, text=str(help_path), fg="#555").pack(side="left", padx=(6, 0))
            viewer_frame = tk.Frame(help_win)
            viewer_frame.pack(fill="both", expand=True, padx=12, pady=6)
            viewer_scroll = tk.Scrollbar(viewer_frame, orient="vertical")
            viewer_scroll.pack(side="right", fill="y")
            viewer = tk.Text(
                viewer_frame,
                wrap="word",
                font=("Arial", 11),
                undo=False,
                yscrollcommand=viewer_scroll.set,
            )
            viewer.pack(side="left", fill="both", expand=True)
            viewer_scroll.config(command=viewer.yview)
            viewer.insert("1.0", help_text)
            viewer.configure(state="disabled")

            def select_all(_event=None):
                viewer.configure(state="normal")
                viewer.tag_add("sel", "1.0", "end")
                viewer.configure(state="disabled")
                return "break"

            def copy_selection(_event=None):
                try:
                    selected = viewer.get("sel.first", "sel.last")
                except tk.TclError:
                    return "break"
                help_win.clipboard_clear()
                help_win.clipboard_append(selected)
                return "break"

            def on_mousewheel(event):
                delta = getattr(event, "delta", 0)
                if delta:
                    step = -1 if delta > 0 else 1
                    viewer.yview_scroll(step, "units")
                    return "break"
                return None

            def on_linux_scroll(event):
                if getattr(event, "num", None) == 4:
                    viewer.yview_scroll(-1, "units")
                    return "break"
                if getattr(event, "num", None) == 5:
                    viewer.yview_scroll(1, "units")
                    return "break"
                return None

            viewer.bind("<Command-a>", select_all)
            viewer.bind("<Command-c>", copy_selection)
            viewer.bind("<Control-a>", select_all)
            viewer.bind("<Control-c>", copy_selection)
            viewer.bind("<MouseWheel>", on_mousewheel)
            viewer.bind("<Button-4>", on_linux_scroll)
            viewer.bind("<Button-5>", on_linux_scroll)
            help_win.bind("<MouseWheel>", on_mousewheel)
            help_win.bind("<Button-4>", on_linux_scroll)
            help_win.bind("<Button-5>", on_linux_scroll)

            if anchor_phrase:
                marker = viewer.search(anchor_phrase, "1.0", stopindex="end")
                if marker:
                    viewer.see(marker)
                    viewer.tag_remove("sel", "1.0", "end")
                    viewer.tag_add("sel", marker, f"{marker}+{len(anchor_phrase)}c")
                    viewer.focus_set()

            def jump_to_top(_event=None):
                viewer.yview_moveto(0.0)
                viewer.focus_set()
                return "break"

            def jump_to_bottom(_event=None):
                viewer.yview_moveto(1.0)
                viewer.focus_set()
                return "break"

            nav_row = tk.Frame(help_win)
            nav_row.pack(fill="x", padx=12, pady=(0, 4))
            tk.Button(nav_row, text="Lên đầu", width=12, command=jump_to_top).pack(side="left")
            tk.Button(nav_row, text="Xuống cuối", width=12, command=jump_to_bottom).pack(side="left", padx=(8, 0))
            if not loaded:
                append_log(f"⚠ Không tìm thấy tài liệu Help: {help_path}; đang dùng fallback local.")
            tk.Button(help_win, text="Đóng", width=12, command=help_win.destroy).pack(anchor="e", padx=12, pady=(0, 10))

        frame = tk.Frame(builder_win)
        frame.pack(fill="both", expand=True, padx=14, pady=12)
        title_row = tk.Frame(frame)
        title_row.pack(fill="x")
        tk.Label(title_row, text="HSK 2.0 / 3.0 Vocab ZIP Builder", font=("Arial", 15, "bold")).pack(side="left")
        tk.Button(title_row, text="Help", width=10, command=open_hsk30_help).pack(side="right")
        tk.Label(
            frame,
            text="Build local → stage ZIP packs → publish catalog revision riêng. Dữ liệu HSK 2.0/3.0 dùng namespace độc lập.",
            fg="#555",
        ).pack(anchor="w", pady=(2, 12))

        def form_row(label, variable, browse=None):
            row = tk.Frame(frame)
            row.pack(fill="x", pady=3)
            tk.Label(row, text=label, width=18, anchor="w").pack(side="left")
            tk.Entry(row, textvariable=variable).pack(side="left", fill="x", expand=True)
            if browse:
                tk.Button(row, text=browse[0], width=12, command=browse[1]).pack(side="left", padx=(8, 0))

        def read_sheets(path):
            from openpyxl import load_workbook
            workbook = load_workbook(path, read_only=True, data_only=True)
            try:
                return list(workbook.sheetnames)
            finally:
                workbook.close()

        def choose_excel():
            selected = filedialog.askopenfilename(
                parent=builder_win,
                title="Chọn Excel vocab HSK 2.0 / 3.0",
                filetypes=[("Excel files", "*.xlsx;*.xlsm;*.xls")],
            )
            if not selected:
                return
            try:
                sheets = read_sheets(selected)
            except Exception as exc:
                messagebox.showerror("Không đọc được Excel", str(exc), parent=builder_win)
                return
            excel_var.set(selected)
            sheet_combo["values"] = sheets
            if sheets:
                sheet_var.set(sheets[0])
                sync_version_level_from_sheet()
            _save_hsk30_recent_selection(excel_var.get(), sheet_var.get())
            clear_build_state()

        def remember_hsk30_selection(*_):
            excel_path = excel_var.get().strip()
            if excel_path and os.path.isfile(excel_path):
                _save_hsk30_recent_selection(excel_path, sheet_var.get().strip())

        def choose_output():
            selected = filedialog.askdirectory(parent=builder_win, title="Chọn thư mục output local")
            if selected:
                output_var.set(selected)
                clear_build_state()

        form_row("Excel path:", excel_var, ("Chọn Excel", choose_excel))
        sheet_row = tk.Frame(frame)
        sheet_row.pack(fill="x", pady=3)
        tk.Label(sheet_row, text="Sheet:", width=18, anchor="w").pack(side="left")
        sheet_combo = ttk.Combobox(sheet_row, textvariable=sheet_var, state="readonly")
        sheet_combo.pack(side="left", fill="x", expand=True)
        sheet_combo.bind("<<ComboboxSelected>>", remember_hsk30_selection)

        if excel_var.get():
            try:
                recent_sheets = read_sheets(excel_var.get())
                sheet_combo["values"] = recent_sheets
                preferred_sheet = last_hsk30_sheet if last_hsk30_sheet in recent_sheets else ""
                if not preferred_sheet:
                    preferred_sheet = _guess_hsk30_sheet_from_state(
                        recent_sheets,
                        last_hsk30_builder_state["version"],
                        last_hsk30_builder_state["level"],
                    )
                sheet_var.set(preferred_sheet if preferred_sheet in recent_sheets else (recent_sheets[0] if recent_sheets else ""))
            except Exception:
                excel_var.set("")
                sheet_var.set("")

        def sync_version_level_from_sheet(*_):
            from pipelines.vocab_zip_builder import SHEET_SELECTIONS, _normalize_sheet_name
            selected = _normalize_sheet_name(sheet_var.get())
            mapped = SHEET_SELECTIONS.get(selected)
            if not mapped:
                return
            mapped_version, mapped_level = mapped
            version_display_var.set("HSK 2.0" if mapped_version == "2.0" else "HSK 3.0")
            label = next((name for name, code in level_label_to_code.items() if code == mapped_level), "HSK 1")
            level_display_var.set(label)
            clear_build_state()

        version_row = tk.Frame(frame)
        version_row.pack(fill="x", pady=3)
        tk.Label(version_row, text="Vocab version:", width=18, anchor="w").pack(side="left")
        ttk.Combobox(version_row, textvariable=version_display_var, values=list(version_label_to_code), state="readonly", width=20).pack(side="left")
        tk.Label(version_row, text="Sheet sẽ tự đồng bộ version/level; build vẫn hard-fail nếu tuple không khớp.", fg="#666").pack(side="left", padx=8)

        level_row = tk.Frame(frame)
        level_row.pack(fill="x", pady=3)
        tk.Label(level_row, text="Level:", width=18, anchor="w").pack(side="left")
        ttk.Combobox(level_row, textvariable=level_display_var, values=list(level_label_to_code), state="readonly", width=20).pack(side="left")
        tk.Label(level_row, textvariable=pack_version_label_var, fg="#245a24").pack(side="left", padx=8)
        tk.Spinbox(level_row, from_=1, to=999, width=5, textvariable=pack_version_var).pack(side="left")
        tk.Label(level_row, text="HSK 7–9 build local dùng canonical code hsk7_9; chưa publish pilot.", fg="#666").pack(side="left", padx=8)
        form_row("Output directory:", output_var, ("Chọn thư mục", choose_output))
        sheet_combo.bind("<<ComboboxSelected>>", sync_version_level_from_sheet, add="+")

        quality_row = tk.Frame(frame)
        quality_row.pack(fill="x", pady=3)
        tk.Label(quality_row, text="Chất lượng M4A:", width=18, anchor="w").pack(side="left")
        ttk.Combobox(
            quality_row,
            textvariable=builder_bitrate_var,
            values=("26 kbps", "32 kbps"),
            state="readonly",
            width=14,
        ).pack(side="left")
        tk.Label(quality_row, text="AAC-LC · mono · 22.05 kHz", fg="#666").pack(side="left", padx=(8, 0))

        audio_mode_row = tk.Frame(frame)
        audio_mode_row.pack(fill="x", pady=3)
        tk.Label(audio_mode_row, text="Nội dung audio:", width=18, anchor="w").pack(side="left")
        ttk.Combobox(
            audio_mode_row,
            textvariable=builder_audio_mode_var,
            values=("Chỉ đọc tiếng Trung", "Đọc tiếng Trung + Tiếng Việt"),
            state="readonly",
            width=30,
        ).pack(side="left")
        tk.Checkbutton(audio_mode_row, text="Force regenerate audio for selected level", variable=force_audio_var).pack(side="left", padx=12)

        tts_summary_var = tk.StringVar()

        def refresh_tts_summary(*_):
            tts_summary_var.set(
                f"Engine: {combo_engine.get()} | Giọng: {combo_giong_popup.get()} | "
                f"Tốc độ: {combo_toc_do_popup.get()} | M4A: {builder_bitrate_var.get() or '32 kbps'} | "
                f"Audio: {builder_audio_mode_var.get()}"
            )

        builder_bitrate_var.trace_add("write", refresh_tts_summary)
        builder_audio_mode_var.trace_add("write", refresh_tts_summary)
        tts_row = tk.Frame(frame)
        tts_row.pack(fill="x", pady=(8, 5))
        tk.Label(tts_row, text="TTS dùng chung:", width=18, anchor="w", font=("Arial", 10, "bold")).pack(side="left")
        tk.Label(
            tts_row,
            textvariable=tts_summary_var,
            fg="#245a24",
        ).pack(side="left")
        refresh_tts_summary()
        tk.Label(
            frame,
            text="Builder tái dùng core TTS/M4A và pinyin filename của vocab_pipeline; profile Google/AWS hiện tại được truyền vào subprocess.",
            fg="#666", wraplength=760, justify="left",
        ).pack(anchor="w", pady=(0, 8))
        tk.Label(frame, textvariable=compatibility_var, fg="#245a24", anchor="w").pack(fill="x", pady=(0, 5))
        signing_box = tk.LabelFrame(frame, text="Signed catalog pointer")
        signing_box.pack(fill="x", pady=(0, 6))
        tk.Label(signing_box, textvariable=signing_key_state_var, anchor="w", fg="#245a24").pack(fill="x", padx=8, pady=(4, 0))
        tk.Label(signing_box, textvariable=signing_public_key_var, anchor="w", justify="left", wraplength=760).pack(fill="x", padx=8)
        tk.Label(signing_box, textvariable=pointer_state_var, anchor="w", fg="#245a24").pack(fill="x", padx=8)
        tk.Label(signing_box, textvariable=pointer_revision_var, anchor="w").pack(fill="x", padx=8, pady=(0, 4))
        tk.Label(signing_box, textvariable=publish_gate_var, anchor="w", fg="#9b1c1c", wraplength=1100).pack(fill="x", padx=8, pady=(0, 4))

        phase_box = tk.LabelFrame(frame, text="Trạng thái nấc 1 — Build + Validate Local")
        phase_box.pack(fill="x", pady=(4, 8))
        tk.Label(phase_box, textvariable=summary_var, anchor="w", justify="left", wraplength=750).pack(fill="x", padx=8, pady=(5, 2))
        tk.Label(phase_box, textvariable=status_var, anchor="w", fg="#1d4f91", justify="left", wraplength=750).pack(fill="x", padx=8, pady=(0, 6))
        tk.Label(phase_box, textvariable=receipt_state_var, anchor="w", fg="#555", justify="left", wraplength=1100).pack(fill="x", padx=8, pady=(0, 6))
        tk.Label(phase_box, textvariable=workflow_state_var, anchor="w", fg="#245a24", justify="left", wraplength=1100).pack(fill="x", padx=8, pady=(0, 6))

        log_text = tk.Text(frame, height=18, wrap="word")
        log_text.pack(fill="both", expand=True)

        controls = tk.Frame(frame)
        controls.pack(fill="x", pady=(10, 0))
        pack_controls = tk.Frame(controls)
        pack_controls.pack(fill="x", pady=(0, 4))
        pointer_controls = tk.Frame(controls)
        pointer_controls.pack(fill="x", pady=(0, 4))
        publish_controls = tk.Frame(controls)
        publish_controls.pack(fill="x")

        def append_log(message):
            log_text.insert("end", message.rstrip() + "\n")
            log_text.see("end")

        def signal_build_process(process, sig):
            """Signal the builder and its TTS/FFmpeg children as one process group."""
            if process is None or process.poll() is not None:
                return
            try:
                if os.name == "nt":
                    process.send_signal(sig)
                else:
                    os.killpg(process.pid, sig)
            except (ProcessLookupError, OSError):
                # The worker may have exited between poll() and signalling.
                pass

        def toggle_pause_build():
            process = build_runtime.get("process")
            if process is None or process.poll() is not None:
                return
            if os.name == "nt":
                messagebox.showinfo(
                    "Tạm dừng chưa hỗ trợ",
                    "Tạm dừng builder hiện chỉ hỗ trợ macOS/Linux. Bạn vẫn có thể huỷ build.",
                    parent=builder_win,
                )
                return
            if build_runtime["paused"]:
                signal_build_process(process, signal.SIGCONT)
                build_runtime["paused"] = False
                pause_btn.config(text="⏸ Tạm dừng")
                status_var.set("Đã tiếp tục build local…")
                append_log("▶ Tiếp tục build local")
            else:
                signal_build_process(process, signal.SIGSTOP)
                build_runtime["paused"] = True
                pause_btn.config(text="▶ Tiếp tục")
                status_var.set("Đã tạm dừng build local.")
                append_log("⏸ Tạm dừng build local")

        def cancel_build():
            process = build_runtime.get("process")
            if process is None or process.poll() is not None:
                return
            if not messagebox.askyesno(
                "Huỷ build?",
                "Dừng build vocab hiện tại? Pack chưa hoàn tất sẽ không được coi là PASS.",
                parent=builder_win,
            ):
                return
            build_runtime["cancel_requested"] = True
            if build_runtime["paused"]:
                # SIGTERM is pending for a stopped process; continue first so
                # the process group can exit promptly.
                signal_build_process(process, signal.SIGCONT)
                build_runtime["paused"] = False
            signal_build_process(process, signal.SIGTERM)
            cancel_btn.config(state="disabled")
            pause_btn.config(state="disabled")
            status_var.set("Đang huỷ build local…")
            append_log("✖ Yêu cầu huỷ build local")

        def selected_config():
            from pipelines.vocab_zip_builder import resolve_sheet_selection
            excel_path = excel_var.get().strip()
            sheet = sheet_var.get().strip()
            output = output_var.get().strip()
            if not excel_path or not os.path.isfile(excel_path):
                raise ValueError("Chọn file Excel hợp lệ.")
            if not sheet:
                raise ValueError("Chọn sheet.")
            if not output:
                raise ValueError("Chọn output directory.")
            try:
                if int(pack_version_var.get()) < 1:
                    raise ValueError
            except ValueError:
                raise ValueError("Pack version phải là số nguyên dương.")
            version = version_label_to_code[version_display_var.get()]
            level = level_label_to_code[level_display_var.get()]
            try:
                version, level = resolve_sheet_selection(sheet, version, level)
            except Exception as exc:
                raise ValueError(str(exc)) from exc
            return excel_path, sheet, version, level, output

        def refresh_pack_version_label(*_):
            try:
                pack_version_label_var.set(f"Pack version: v{int(pack_version_var.get())}")
            except ValueError:
                pack_version_label_var.set("Pack version: không hợp lệ")

        def persist_hsk30_builder_state(*_):
            _save_hsk30_builder_state(
                version=version_display_var.get().strip(),
                level=level_display_var.get().strip(),
                pack_version=pack_version_var.get().strip(),
                output_dir=output_var.get().strip(),
                bitrate=builder_bitrate_var.get().strip(),
                audio_mode=builder_audio_mode_var.get().strip(),
            )

        stage_btn = None
        publish_btn = None

        def clear_build_state(*_):
            artifact_state["result"] = None
            artifact_state["fingerprint"] = None
            receipt_state_var.set("Receipt: chưa có | Remote verification: chưa có")
            workflow_state_var.set("LOCAL NOT BUILT | PACKS NOT VERIFIED | CATALOG NOT PUBLISHED")
            summary_var.set("Chưa đọc Excel")
            status_var.set("Sẵn sàng build local. Phase 2 chỉ chạy sau local PASS và xác nhận.")
            if stage_btn is not None:
                stage_btn.config(state="disabled")
            if publish_btn is not None:
                publish_btn.config(state="disabled")

        version_display_var.trace_add("write", clear_build_state)
        level_display_var.trace_add("write", clear_build_state)
        version_display_var.trace_add("write", persist_hsk30_builder_state)
        level_display_var.trace_add("write", persist_hsk30_builder_state)
        pack_version_var.trace_add("write", lambda *_: (refresh_pack_version_label(), persist_hsk30_builder_state()))
        output_var.trace_add("write", persist_hsk30_builder_state)
        builder_bitrate_var.trace_add("write", persist_hsk30_builder_state)
        builder_audio_mode_var.trace_add("write", persist_hsk30_builder_state)

        if sheet_var.get():
            sync_version_level_from_sheet()

        refresh_pack_version_label()

        def refresh_signing_status():
            status = signing_status(DEFAULT_KEY_PATH, DEFAULT_KEY_ID)
            signing_key_state_var.set(str(status.get("status", "SIGNING KEY NOT INITIALIZED")))
            public_key = status.get("publicKeyB64")
            signing_public_key_var.set(f"PUBLIC KEY B64: {public_key or '—'}")
            return status

        def refresh_pointer_status_pending():
            """GET-verify current.json and its catalog; never performs a write."""
            if refresh_pointer_btn.winfo_exists():
                refresh_pointer_btn.config(state="disabled")
            pointer_state_var.set("POINTER STATUS CHECKING…")
            publish_gate_var.set("Publish disabled: đang verify pointer production…")

            def worker():
                try:
                    profile_name, profile = load_active_supabase_profile()
                    project_url = str(profile.get("SUPABASE_URL", "") or "").strip()
                    service_key = str(profile.get("SUPABASE_SERVICE_ROLE_KEY", "") or "").strip()
                    # HSK 3.0 deploys and the signed pointer always use the
                    # staging bucket; keep the legacy profile bucket untouched.
                    bucket = STAGING_BUCKET
                    if not project_url or not service_key:
                        raise DeployValidationError("Thiếu Supabase URL hoặc service-role key trong profile hiện tại.")
                    client = SupabaseStorageRestClient(project_url, service_key, network_enabled=True)
                    result = read_verified_pointer_status(
                        client,
                        bucket=bucket,
                        private_key_path=DEFAULT_KEY_PATH,
                        key_id=DEFAULT_KEY_ID,
                    )
                    result["profileName"] = profile_name
                except Exception as exc:
                    message = str(exc)

                    def show_error():
                        pointer_status_cache["status"] = "UNKNOWN"
                        pointer_state_var.set("POINTER STATUS UNKNOWN — chưa thể verify")
                        pointer_revision_var.set("CURRENT POINTER REVISION: — | CURRENT CATALOG REVISION: —")
                        publish_gate_var.set("Publish disabled: không verify được pointer; thử Refresh Pointer Status lại")
                        append_log(f"⚠ Refresh pointer thất bại (không coi là chưa initialize): {message}")
                        refresh_pointer_btn.config(state="normal")
                        initialize_pointer_btn.config(state="disabled")
                        refresh_catalog_gate()

                    builder_win.after(0, show_error)
                    return

                def show_result():
                    status = str(result.get("status", ""))
                    pointer_status_cache.clear()
                    pointer_status_cache.update(result)
                    if status == "POINTER ACTIVE":
                        pointer_state_var.set("POINTER ACTIVE")
                        pointer_revision_var.set(
                            f"CURRENT POINTER REVISION: {result['pointerRevision']} | "
                            f"CURRENT CATALOG REVISION: {result['catalogRevision']} | "
                            f"CATALOG ENTRIES: {result['entryCount']}"
                        )
                        initialize_pointer_btn.config(state="disabled")
                        append_log(
                            f"Pointer ACTIVE: pointerRevision={result['pointerRevision']} "
                            f"catalogRevision={result['catalogRevision']} entries={result['entryCount']}"
                        )
                    else:
                        pointer_state_var.set("POINTER NOT INITIALIZED")
                        pointer_revision_var.set("CURRENT POINTER REVISION: — | CURRENT CATALOG REVISION: —")
                        initialize_pointer_btn.config(state="normal")
                        append_log("Pointer chưa tồn tại: current.json ABSENT")
                    refresh_pointer_btn.config(state="normal")
                    refresh_catalog_gate()

                builder_win.after(0, show_result)

            threading.Thread(target=worker, daemon=True).start()

        def initialize_signing_key_pending():
            confirm = tk.Toplevel(builder_win)
            set_popup_icon(confirm)
            confirm.title("Initialize Production Signing Key")
            confirm.geometry("620x280")
            confirm.transient(builder_win)
            confirm.grab_set()
            tk.Label(confirm, text=(
                "Tạo Ed25519 raw seed 32 byte ngoài repo tại:\n"
                f"{DEFAULT_KEY_PATH}\n\n"
                "Mất seed sẽ không thể ký pointer mới cho app đang tin key này. "
                "Hãy sao lưu riêng an toàn; tool không log/clipboard/upload seed."
            ), anchor="w", justify="left", wraplength=580).pack(fill="both", expand=True, padx=14, pady=14)
            tk.Label(confirm, text="Nhập chính xác: INITIALIZE VOCAB SIGNING KEY", fg="#9b1c1c", font=("Arial", 10, "bold")).pack(anchor="w", padx=14)
            phrase_var = tk.StringVar()
            tk.Entry(confirm, textvariable=phrase_var, width=44).pack(anchor="w", padx=14, pady=6)

            def initialize():
                try:
                    result = initialize_signing_key(DEFAULT_KEY_PATH, confirmation=phrase_var.get())
                except Exception as exc:
                    messagebox.showerror("Không khởi tạo signing key", str(exc), parent=confirm)
                    return
                confirm.destroy()
                refresh_signing_status()
                append_log(f"{result['status']} | keyId={result['keyId']} | publicKeyB64={result['publicKeyB64']}")
                refresh_catalog_gate()

            buttons = tk.Frame(confirm)
            buttons.pack(fill="x", padx=14, pady=(0, 14))
            tk.Button(buttons, text="Huỷ", width=12, command=confirm.destroy).pack(side="right", padx=(8, 0))
            tk.Button(buttons, text="Khởi tạo key", width=18, command=initialize).pack(side="right")

        def initialize_pointer_pending():
            if pointer_status_cache.get("status") == "POINTER ACTIVE":
                messagebox.showinfo("Pointer đã tồn tại", "POINTER ALREADY INITIALIZED — không tạo lại revision 1.", parent=builder_win)
                return
            if pointer_status_cache.get("status") != "POINTER NOT INITIALIZED":
                messagebox.showwarning(
                    "Chưa xác minh pointer",
                    "Hãy bấm Refresh Pointer Status và chỉ initialize khi current.json thật sự ABSENT.",
                    parent=builder_win,
                )
                return
            status = refresh_signing_status()
            if status.get("status") != "SIGNING KEY READY":
                messagebox.showwarning("Chưa có signing key", "Hãy Initialize Production Signing Key trước.", parent=builder_win)
                return
            confirm = tk.Toplevel(builder_win)
            set_popup_icon(confirm)
            confirm.title("Initialize Signed Pointer")
            confirm.geometry("680x330")
            confirm.transient(builder_win)
            confirm.grab_set()
            tk.Label(confirm, text=(
                "Khởi tạo pointer production cho catalog combined v1 hiện hành.\n"
                "Thao tác này tạo pointer archive rồi cập nhật current.json cuối cùng; đây là remote write.\n\n"
                "Catalog: catalogs/vocab/combined/v1/vocab_pack_catalog_20_30_v1.json\n"
                "Bytes: 7280\nSHA-256: 593d2f8846a6b0ccfca9589512d56b7579f43144814a21ae60509692e9d24413"
            ), anchor="w", justify="left", wraplength=640).pack(fill="both", expand=True, padx=14, pady=14)
            tk.Label(confirm, text="Nhập chính xác: INITIALIZE VOCAB POINTER", fg="#9b1c1c", font=("Arial", 10, "bold")).pack(anchor="w", padx=14)
            phrase_var = tk.StringVar()
            tk.Entry(confirm, textvariable=phrase_var, width=42).pack(anchor="w", padx=14, pady=6)

            def initialize_pointer():
                try:
                    _, profile = load_active_supabase_profile()
                    client = SupabaseStorageRestClient(str(profile.get("SUPABASE_URL", "") or "").strip(), str(profile.get("SUPABASE_SERVICE_ROLE_KEY", "") or "").strip(), network_enabled=True)
                    result = initialize_pointer_with_client(
                        client, catalog_revision=1,
                        catalog_object="catalogs/vocab/combined/v1/vocab_pack_catalog_20_30_v1.json",
                        expected_bytes=7280,
                        expected_sha256="593d2f8846a6b0ccfca9589512d56b7579f43144814a21ae60509692e9d24413",
                        confirmation=phrase_var.get(), private_key_path=DEFAULT_KEY_PATH,
                    )
                except Exception as exc:
                    messagebox.showerror("Không khởi tạo pointer", str(exc), parent=confirm)
                    return
                confirm.destroy()
                pointer_state_var.set("POINTER ACTIVE")
                pointer_revision_var.set("CURRENT POINTER REVISION: 1 | CURRENT CATALOG REVISION: 1")
                append_log(f"Pointer ACTIVE: {result['archivePath']}")

            buttons = tk.Frame(confirm)
            buttons.pack(fill="x", padx=14, pady=(0, 14))
            tk.Button(buttons, text="Huỷ", width=12, command=confirm.destroy).pack(side="right", padx=(8, 0))
            tk.Button(buttons, text="Khởi tạo pointer", width=20, command=initialize_pointer).pack(side="right")

        def fingerprint(config):
            excel_path, sheet, version, level, output = config
            pack_version = 1
            if isinstance(artifact_state.get("result"), dict):
                try:
                    pack_version = int(artifact_state["result"].get("packVersion", 1) or 1)
                except Exception:
                    pack_version = 1
            return input_fingerprint(
                excel_path,
                sheet,
                level,
                output,
                version=version,
                bitrate=_canonical_vocab_bitrate(builder_bitrate_var.get()),
                pack_version=pack_version,
            )

        def load_active_supabase_profile():
            """Read the existing profile/key without changing or displaying the key."""
            env_values = {}
            env_path = os.path.join(BASE_DIR, ".env")
            if os.path.isfile(env_path):
                try:
                    for line in Path(env_path).read_text(encoding="utf-8").splitlines():
                        stripped = line.strip()
                        if stripped and not stripped.startswith("#") and "=" in stripped:
                            key, value = stripped.split("=", 1)
                            env_values[key.strip()] = value.strip()
                except Exception:
                    env_values = {}
            default_profile = {
                "SUPABASE_URL": env_values.get("SUPABASE_URL", ""),
                "SUPABASE_SERVICE_ROLE_KEY": env_values.get("SUPABASE_SERVICE_ROLE_KEY", ""),
                "SUPABASE_BUCKET": env_values.get("SUPABASE_BUCKET", ""),
            }
            profile_path = os.path.join(APPDATA_ROOT, "supabase_import_profiles.json")
            active_name = "dev"
            profile = dict(default_profile)
            try:
                raw = json.loads(Path(profile_path).read_text(encoding="utf-8"))
                active_name = str(raw.get("active_profile", "dev"))
                profiles = raw.get("profiles", {})
                if isinstance(profiles, dict) and isinstance(profiles.get(active_name), dict):
                    profile.update(profiles[active_name])
            except Exception:
                pass
            return active_name, profile

        def current_deploy_plan():
            result = artifact_state.get("result")
            current_config = selected_config()
            profile_name, profile = load_active_supabase_profile()
            return build_plan(
                result,
                current_config,
                artifact_state.get("fingerprint"),
                profile,
                profile_name=profile_name,
            )

        def restore_existing_build_report():
            """Restore a PASS build after the builder window/app was restarted."""
            if isinstance(artifact_state.get("result"), dict):
                return True
            try:
                config = selected_config()
                pack_version = int(pack_version_var.get())
                report_path = os.path.join(
                    config[4], "vocab", config[2], config[3], "builds",
                    f"v{pack_version}", "build_report.json",
                )
                with open(report_path, "r", encoding="utf-8") as report_file:
                    result = json.load(report_file)
                if result.get("status") != "PASS":
                    return False
                if str(result.get("version")) != config[2] or str(result.get("level")) != config[3]:
                    return False
                if int(result.get("packVersion", 0) or 0) != pack_version:
                    return False
                artifact_state["result"] = result
                artifact_state["fingerprint"] = fingerprint(config)
                segments = ("base", "plus1", "plus2") if "plus1" in result else ("base", "plus")
                segment_summary = " | ".join(f"{segment.upper()}={result.get(segment, {}).get('manifest', {}).get('vocabCount', '—')}" for segment in segments)
                summary_var.set(
                    f"Version: HSK {result['version']} | Level: {result['level'].upper()} | Pack version: v{result['packVersion']} | "
                    f"rows={result.get('totalRows', '—')} | audio đã có={result.get('audioReused', '—')} | "
                    f"audio cần tạo/thiếu ban đầu={result.get('audioGenerated', '—')} | {segment_summary}"
                )
                status_var.set("Đã khôi phục local PASS từ build_report.json. Có thể Upload + Verify Packs.")
                workflow_state_var.set("LOCAL PASS | PACKS NOT VERIFIED | CATALOG NOT PUBLISHED")
                return True
            except Exception:
                return False

        def refresh_deploy_gate():
            restore_existing_build_report()
            try:
                plan = current_deploy_plan()
            except (DeployValidationError, ValueError, OSError):
                stage_btn.config(state="disabled")
                compatibility_var.set("Compatibility hash: chưa PASS")
                return False
            compatibility_var.set(f"Compatibility hash: {plan.compatibility_hash}")
            stage_btn.config(state="normal")
            return True

        def refresh_catalog_gate():
            signing = refresh_signing_status()
            if signing.get("status") != "SIGNING KEY READY":
                publish_gate_var.set("Publish disabled: SIGNING KEY READY chưa sẵn sàng")
                publish_btn.config(state="disabled")
                return False
            if pointer_status_cache.get("status") != "POINTER ACTIVE":
                publish_gate_var.set("Publish disabled: pointer status chưa được refresh/ACTIVE")
                publish_btn.config(state="disabled")
                return False
            try:
                selected_version = version_label_to_code[version_display_var.get()]
                selected_level = level_label_to_code[level_display_var.get()]
                # Do not let an unrelated legacy receipt block the selected
                # HSK 3.0 pack.  The publish gate is scoped to the current
                # version/level; legacy HSK 2.0 receipts are handled by their
                # own workflow.
                receipts = collect_deploy_receipts(
                    output_var.get().strip(),
                    levels={selected_level},
                    versions={selected_version},
                )
                # A level can have multiple immutable pack versions.  The
                # catalog gate must only consider the receipt for the exact
                # build currently selected in the window, never a stale
                # receipt from another pack version.
                selected_pack_version = int(pack_version_var.get())
                selected = next((item for item in receipts
                                 if item.get("version") == selected_version
                                 and item.get("level") == selected_level
                                 and int(item.get("packVersion", -1)) == selected_pack_version), None)
            except (DeployValidationError, ValueError, OSError) as exc:
                publish_gate_var.set(f"Publish disabled: receipt không hợp lệ ({exc})")
                publish_btn.config(state="disabled")
                return False
            if selected is None:
                publish_gate_var.set(
                    f"Publish disabled: chưa có receipt remote-verified cho "
                    f"{selected_level.upper()} {selected_version} pack v{selected_pack_version}"
                )
                publish_btn.config(state="disabled")
                return False
            receipt = selected.get("receipt", {})
            # A previous publish may have completed remotely but crashed while
            # writing the local receipt.  The pointer refresh already fetched
            # and verified the active catalog; reconcile that local bookkeeping
            # only when every selected descriptor matches exactly.
            active_catalog = pointer_status_cache.get("catalog")
            if receipt.get("catalogPublished") is not True and isinstance(active_catalog, dict):
                if catalog_matches_receipt(receipt, active_catalog):
                    pointer_revision = pointer_status_cache.get("pointerRevision")
                    catalog_revision = pointer_status_cache.get("catalogRevision")
                    try:
                        from pipelines.vocab_zip_deploy import mark_receipts_catalog_published
                        mark_receipts_catalog_published(
                            [selected],
                            catalog_revision=int(catalog_revision),
                            pointer_revision=int(pointer_revision),
                        )
                        receipt = dict(receipt)
                        receipt["catalogPublished"] = True
                        receipt["catalogRevision"] = int(catalog_revision)
                        receipt["pointerRevision"] = int(pointer_revision)
                        selected["receipt"] = receipt
                        selected["reconciled"] = True
                    except (DeployValidationError, ValueError, OSError):
                        # Never enable/disable publish based on a failed local
                        # bookkeeping write; the verified remote state remains
                        # authoritative and the operator can retry refresh.
                        pass
            if receipt.get("catalogPublished") is True:
                revision = pointer_status_cache.get("catalogRevision", receipt.get("catalogRevision", "—"))
                publish_gate_var.set(
                    f"ALREADY PUBLISHED: {selected_level.upper()} {selected_version} v{receipt.get('packVersion', '—')} "
                    f"| catalogRevision={revision} | không cần Publish lại"
                )
                workflow_state_var.set("LOCAL PASS | REMOTE PACKS VERIFIED | CATALOG PUBLISHED")
                publish_btn.config(state="disabled")
                return False
            publish_gate_var.set(f"Publish READY: {selected_level.upper()} {selected_version} v{receipt.get('packVersion', '—')} receipt đã remote verify; nhập PUBLISH VOCAB CATALOG")
            publish_btn.config(state="normal")
            return True

        def play_first_audio():
            result = artifact_state.get("result")
            if not result:
                messagebox.showinfo("Chưa có audio", "Hãy build local trước.", parent=builder_win)
                return
            audio_dir = result.get("sourceAudioRoot") or os.path.join(
                os.path.dirname(os.path.dirname(os.path.dirname(result["base"]["unpacked"]))),
                "source_audio",
            )
            files = sorted(Path(audio_dir).glob("*.m4a"))
            if not files:
                messagebox.showwarning("Không có audio", "Không tìm thấy M4A local để phát thử.", parent=builder_win)
                return
            try:
                # pygame's bundled SDL_mixer may try to load M4A through
                # libmodplug on macOS.  That optional dylib is not present in
                # many installations, while the system afplay utility plays
                # AAC/M4A directly.  Prefer it for this local preview only;
                # the export/build pipeline is unchanged.
                preview_process = build_runtime.get("preview_process")
                if preview_process is not None and preview_process.poll() is None:
                    preview_process.terminate()
                afplay_path = shutil.which("afplay")
                if not afplay_path and sys.platform == "darwin" and os.path.isfile("/usr/bin/afplay"):
                    afplay_path = "/usr/bin/afplay"
                if sys.platform == "darwin" and afplay_path:
                    build_runtime["preview_process"] = subprocess.Popen(
                        [afplay_path, str(files[0])],
                        stdout=subprocess.DEVNULL,
                        stderr=subprocess.DEVNULL,
                    )
                    append_log(f"▶ Phát thử local (afplay): {files[0].name}")
                else:
                    pygame.mixer.music.load(str(files[0]))
                    pygame.mixer.music.play()
                    append_log(f"▶ Phát thử local (pygame): {files[0].name}")
            except Exception as exc:
                messagebox.showerror("Không phát được audio", str(exc), parent=builder_win)

        def run_local_build():
            vocab_tts = collect_vocab_tts_config(builder_win, builder_bitrate_var.get(), builder_audio_mode_var.get())
            if not vocab_tts:
                return
            try:
                config = selected_config()
            except ValueError as exc:
                messagebox.showwarning("Thiếu cấu hình", str(exc), parent=builder_win)
                return
            clear_build_state()
            _save_hsk30_recent_selection(config[0], config[1])
            build_btn.config(state="disabled")
            pause_btn.config(state="normal", text="⏸ Tạm dừng")
            cancel_btn.config(state="normal")
            build_runtime["paused"] = False
            build_runtime["cancel_requested"] = False
            status_var.set("Đang chạy build local…")
            append_log(f"=== HSK {config[2]} {config[3].upper()} BUILD + VALIDATE LOCAL ===")

            def worker():
                command = [
                    sys.executable, "-u", os.path.join(BASE_DIR, "pipelines", "vocab_zip_builder.py"),
                    config[0], "--sheet", config[1], "--version", config[2], "--level", config[3], "--output", config[4],
                    "--pack-version", pack_version_var.get(),
                    "--engine", vocab_tts["engine"], "--speed", vocab_tts["speed"],
                    "--voice", vocab_tts["voice"], "--bitrate", vocab_tts["bitrate"],
                    "--audio-mode", vocab_tts["audio_mode"],
                    "--languages", ",".join(vocab_tts["languages"]),
                    "--config-confirmed", "true",
                ]
                if force_audio_var.get():
                    command.append("--force-regenerate-audio")
                cache_profile = {
                    "selection": vocab_tts,
                    "google_profiles": GOOGLE_TTS_PROFILES,
                }
                command.extend(["--tts-profile", json.dumps(cache_profile, ensure_ascii=False, sort_keys=True)])
                env = os.environ.copy()
                env.update({
                    "GOOGLE_TTS_PROFILES_JSON": json.dumps(GOOGLE_TTS_PROFILES, ensure_ascii=False),
                    "GOOGLE_TTS_API_KEY": GOOGLE_TTS_API_KEY,
                    "AWS_ACCESS_KEY_ID": AWS_ACCESS_KEY_ID,
                    "AWS_SECRET_ACCESS_KEY": AWS_SECRET_ACCESS_KEY,
                    "AWS_REGION": AWS_REGION,
                    "AWS_DEFAULT_REGION": AWS_REGION,
                })
                apply_vocab_tts_env(env, vocab_tts)
                process_kwargs = {
                    "cwd": BASE_DIR,
                    "env": env,
                    "stdout": subprocess.PIPE,
                    "stderr": subprocess.STDOUT,
                    "text": True,
                }
                if os.name == "nt":
                    process_kwargs["creationflags"] = getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0)
                else:
                    process_kwargs["start_new_session"] = True
                process = subprocess.Popen(command, **process_kwargs)
                build_runtime["process"] = process
                lines = []
                assert process.stdout is not None
                for line in process.stdout:
                    lines.append(line.rstrip())
                    builder_win.after(0, append_log, line)
                code = process.wait()
                cancelled = build_runtime.get("cancel_requested", False)
                build_runtime["process"] = None

                def done():
                    build_btn.config(state="normal")
                    pause_btn.config(state="disabled", text="⏸ Tạm dừng")
                    cancel_btn.config(state="disabled")
                    build_runtime["paused"] = False
                    build_runtime["cancel_requested"] = False
                    if cancelled:
                        status_var.set("Đã huỷ build local — không có upload/publish.")
                        summary_var.set("Build đã huỷ. Hãy chạy lại từ đầu nếu cần.")
                        append_log("✖ Build đã được huỷ")
                        return
                    if code != 0:
                        status_var.set("FAIL — không có upload/publish nào được chạy.")
                        summary_var.set("Build local thất bại. Xem log và build_report.json.")
                        return
                    pack_version = int(pack_version_var.get())
                    report_path = os.path.join(config[4], "vocab", config[2], config[3], "builds", f"v{pack_version}", "build_report.json")
                    try:
                        with open(report_path, "r", encoding="utf-8") as report_file:
                            result = json.load(report_file)
                        if result.get("status") != "PASS":
                            raise ValueError(result.get("error", "Build report không PASS"))
                    except Exception as exc:
                        status_var.set("FAIL — report không hợp lệ.")
                        summary_var.set(str(exc))
                        return
                    artifact_state["result"] = result
                    artifact_state["fingerprint"] = fingerprint(config)
                    base = result["base"]
                    split_hsk79 = result.get("version") == "3.0" and result.get("level") == "hsk7_9" and "plus1" in result
                    segments = ("base", "plus1", "plus2") if split_hsk79 else ("base", "plus")
                    segment_summary = " | ".join(f"{segment.upper()}={result[segment]['manifest']['vocabCount']} ({result[segment]['bytes']} bytes, {result[segment]['sha256']})" for segment in segments)
                    summary_var.set(
                        f"Version: HSK {result['version']} | Level: {result['level'].upper()} | Pack version: v{result['packVersion']} | "
                        f"rows={result['totalRows']} | audio đã có={result['audioReused']} | "
                        f"audio cần tạo/thiếu ban đầu={result['audioGenerated']} | {segment_summary}"
                    )
                    status_var.set("Nấc 1 PASS. Nấc 2 Upload + Verify Packs đã mở; Nấc 3 Publish Combined Catalog vẫn pending.")
                    workflow_state_var.set("LOCAL PASS | PACKS NOT VERIFIED | CATALOG NOT PUBLISHED")
                    receipt_state_var.set(f"Receipt: {config[4]}/vocab/{config[2]}/{config[3]}/deploy_receipt.json | Remote verification: chưa upload | Version: HSK {config[2]}")
                    for segment in segments:
                        pack = result[segment]
                        append_log(f"{segment.upper()}: {pack['zip']} ({pack['bytes']} bytes, {pack['sha256']})")
                    refresh_deploy_gate()
                builder_win.after(0, done)

            threading.Thread(target=worker, daemon=True).start()

        def stage_packs_pending():
            try:
                plan = current_deploy_plan()
                _, profile = load_active_supabase_profile()
            except Exception as exc:
                stage_btn.config(state="disabled")
                messagebox.showwarning("Stage chưa đủ điều kiện", str(exc), parent=builder_win)
                return
            phrase = stage_confirmation_phrase(plan.level, plan.version)
            confirm = tk.Toplevel(builder_win)
            set_popup_icon(confirm)
            confirm.title(f"Xác nhận stage {plan.level.upper()} {plan.version}")
            confirm.geometry("760x520")
            confirm.transient(builder_win)
            confirm.grab_set()
            split_hsk79 = plan.version == "3.0" and plan.level == "hsk7_9" and "plus1" in plan.segment_packs
            stage_segments = ("base", "plus1", "plus2") if split_hsk79 else ("base", "plus")
            pack_details = "\n\n".join(f"{segment.upper()}: {plan.segment_packs.get(segment, {}).get('objectPath', plan.base_object_path if segment == 'base' else plan.plus_object_path)}\n{plan.segment_packs.get(segment, {}).get('bytes', plan.base_bytes if segment == 'base' else plan.plus_bytes)} bytes | {plan.segment_packs.get(segment, {}).get('sha256', plan.base_sha256 if segment == 'base' else plan.plus_sha256)}" for segment in stage_segments)
            details = (
                f"Profile: {plan.profile_name}\nProject URL: {plan.project_url}\nBucket: {plan.bucket}\n"
                f"Version: {plan.version} | Level: {plan.level} | packVersion: v{plan.pack_version}\n\n"
                f"{pack_details}\n\n"
                f"Thao tác này chỉ upload + GET-verify {len(stage_segments)} ZIP bằng create-only. Catalog không được tạo/publish."
            )
            tk.Label(confirm, text=details, anchor="w", justify="left", wraplength=720).pack(fill="both", expand=True, padx=14, pady=(14, 8))
            tk.Label(confirm, text="Nhập chính xác (có thể bôi đen và copy):", fg="#9b1c1c", font=("Arial", 10, "bold")).pack(anchor="w", padx=14)
            phrase_entry = tk.Entry(confirm, width=42, fg="#9b1c1c", readonlybackground="#f8eeee")
            phrase_entry.insert(0, phrase)
            phrase_entry.configure(state="readonly")
            phrase_entry.pack(anchor="w", padx=14, pady=(4, 10))
            phrase_var = tk.StringVar()
            tk.Entry(confirm, textvariable=phrase_var, width=42).pack(anchor="w", padx=14, pady=(0, 10))

            def confirm_stage():
                if phrase_var.get() != phrase:
                    messagebox.showwarning("Xác nhận sai", "Chuỗi xác nhận không khớp; chưa có remote request nào được gọi.", parent=confirm)
                    return
                confirmation = phrase_var.get()
                confirm.destroy()
                stage_btn.config(state="disabled")
                build_btn.config(state="disabled")
                status_var.set(f"Đang stage {plan.level.upper()} {plan.version} ZIP packs…")
                append_log(f"=== STAGE {plan.level.upper()} {plan.version} ===")

                def worker():
                    try:
                        client = SupabaseStorageRestClient(plan.project_url, str(profile.get("SUPABASE_SERVICE_ROLE_KEY", "") or "").strip(), network_enabled=True)
                        result = stage_packs_with_client(client, plan, confirmation=confirmation, output_directory=output_var.get().strip(), progress=lambda message: builder_win.after(0, append_log, message))
                    except Exception as exc:
                        message = str(exc)
                        builder_win.after(0, lambda: status_var.set(f"STAGE FAIL/PARTIAL — {message}"))
                        builder_win.after(0, append_log, f"✖ Stage thất bại: {message}")
                    else:
                        builder_win.after(0, lambda: status_var.set("REMOTE PACKS VERIFIED | CATALOG NOT PUBLISHED"))
                        builder_win.after(0, lambda: workflow_state_var.set("LOCAL PASS | REMOTE PACKS VERIFIED | CATALOG NOT PUBLISHED (không upload lại nếu SHA đã khớp)"))
                        builder_win.after(0, append_log, f"Receipt: {result['receiptPath']}")
                        builder_win.after(0, lambda: receipt_state_var.set(f"Receipt: {result['receiptPath']} | Remote verification: BASE+PLUS PASS"))
                        builder_win.after(0, lambda: summary_var.set(
                            f"Version: HSK {plan.version} | Level: {plan.level.upper()} | Pack version: v{plan.pack_version} | "
                            f"Receipt: {result['receiptPath']} | Remote verification: BASE+PLUS PASS"
                        ))
                        builder_win.after(0, refresh_catalog_gate)
                        # Pointer refresh and receipt write complete on
                        # separate UI callbacks. Re-evaluate once more after
                        # the callbacks have settled so the Publish button
                        # cannot remain visually stale after a successful
                        # stage.
                        builder_win.after(300, refresh_catalog_gate)
                    finally:
                        builder_win.after(0, lambda: build_btn.config(state="normal"))

                threading.Thread(target=worker, daemon=True).start()

            buttons = tk.Frame(confirm)
            buttons.pack(fill="x", padx=14, pady=(0, 14))
            tk.Button(buttons, text="Huỷ", width=12, command=confirm.destroy).pack(side="right", padx=(8, 0))
            tk.Button(buttons, text="Xác nhận stage packs", width=20, command=confirm_stage).pack(side="right")

        def publish_catalog_pending():
            # Re-check the local/remote readiness gate at click time.  This
            # closes the race where a stale enabled button survives a pointer
            # refresh and would otherwise open a new-revision confirmation.
            if not refresh_catalog_gate():
                if str(publish_gate_var.get()).startswith("ALREADY PUBLISHED"):
                    messagebox.showinfo("Catalog đã publish", publish_gate_var.get(), parent=builder_win)
                return
            try:
                selected_version = version_label_to_code[version_display_var.get()]
                selected_level = level_label_to_code[level_display_var.get()]
                publish_plan = prepare_catalog_publish(output_var.get().strip(), levels={selected_level}, versions={selected_version})
                _, profile = load_active_supabase_profile()
            except Exception as exc:
                publish_btn.config(state="disabled")
                messagebox.showwarning("Publish catalog chưa đủ điều kiện", str(exc), parent=builder_win)
                return
            confirm = tk.Toplevel(builder_win)
            set_popup_icon(confirm)
            confirm.title(f"Xác nhận publish catalog v{publish_plan.target_revision}")
            confirm.geometry("760x480")
            confirm.transient(builder_win)
            confirm.grab_set()
            details = (
                f"Catalog source: v{publish_plan.source.revision} | {publish_plan.source.object_path}\n"
                f"Source SHA: {publish_plan.source.sha256}\n"
                f"Catalog target: {publish_plan.target_object_path}\n"
                f"Entries source: {publish_plan.source.entry_count}\n"
                f"Descriptors replace/add từ HSK {selected_version} {selected_level.upper()}: {len(publish_plan.additions)}\n\n"
                "CẢNH BÁO: remote write tạo catalog revision mới bằng create-only. Catalog cũ không bị overwrite."
            )
            tk.Label(confirm, text=details, anchor="w", justify="left", wraplength=720).pack(fill="both", expand=True, padx=14, pady=(14, 8))
            tk.Label(confirm, text="Nhập chính xác (có thể bôi đen và copy):", fg="#9b1c1c", font=("Arial", 10, "bold")).pack(anchor="w", padx=14)
            phrase_entry = tk.Entry(confirm, width=42, fg="#9b1c1c", readonlybackground="#f8eeee")
            phrase_entry.insert(0, CATALOG_PUBLISH_CONFIRMATION)
            phrase_entry.configure(state="readonly")
            phrase_entry.pack(anchor="w", padx=14, pady=(4, 10))
            phrase_var = tk.StringVar()
            tk.Entry(confirm, textvariable=phrase_var, width=42).pack(anchor="w", padx=14, pady=(0, 10))

            def confirm_publish():
                if phrase_var.get() != CATALOG_PUBLISH_CONFIRMATION:
                    messagebox.showwarning("Xác nhận sai", "Chuỗi xác nhận không khớp; chưa có remote request nào được gọi.", parent=confirm)
                    return
                confirmation = phrase_var.get()
                confirm.destroy()
                publish_btn.config(state="disabled")
                status_var.set("Đang publish combined catalog revision mới…")
                append_log("=== PUBLISH COMBINED VOCAB CATALOG ===")

                def worker():
                    try:
                        client = SupabaseStorageRestClient(str(profile.get("SUPABASE_URL", "") or "").strip(), str(profile.get("SUPABASE_SERVICE_ROLE_KEY", "") or "").strip(), network_enabled=True)
                        result = publish_signed_catalog_with_client(
                            client,
                            output_directory=output_var.get().strip(),
                            confirmation=confirmation,
                            private_key_path=DEFAULT_KEY_PATH,
                            levels={selected_level},
                            versions={selected_version},
                            progress=lambda message: builder_win.after(0, append_log, message),
                        )
                    except Exception as exc:
                        message = str(exc)
                        builder_win.after(0, lambda: status_var.set(f"CATALOG FAIL — {message}"))
                        builder_win.after(0, append_log, f"✖ Publish catalog thất bại: {message}")
                    else:
                        result_status = str(result.get("status", "PUBLISHED"))
                        already = result_status == "ALREADY PUBLISHED"
                        builder_win.after(0, lambda: status_var.set("CATALOG ALREADY PUBLISHED" if already else "CATALOG PUBLISHED"))
                        builder_win.after(0, lambda: workflow_state_var.set("LOCAL PASS | REMOTE PACKS VERIFIED | CATALOG PUBLISHED"))
                        builder_win.after(0, lambda: pointer_state_var.set("POINTER ACTIVE"))
                        builder_win.after(0, lambda: pointer_revision_var.set(
                            f"CURRENT POINTER REVISION: {result.get('pointerRevision', '—')} | "
                            f"CURRENT CATALOG REVISION: {result.get('catalogRevision', '—')} | "
                            f"CATALOG ENTRIES: {result.get('entryCount', '—')}"
                        ))
                        builder_win.after(0, lambda: publish_gate_var.set(
                            "ALREADY PUBLISHED | GET VERIFY: PASS | SIGNATURE VERIFY: PASS" if already
                            else "CATALOG PUBLISHED | GET VERIFY: PASS | SIGNATURE VERIFY: PASS"
                        ))
                        builder_win.after(0, lambda: publish_btn.config(state="disabled"))
                        builder_win.after(0, append_log, f"Catalog URL: {result.get('catalogUrl', result.get('publicUrl', '—'))}")
                        builder_win.after(0, append_log, f"Catalog: {result.get('catalogBytes', result.get('bytes', '—'))} bytes | SHA {result.get('catalogSha256', result.get('sha256', '—'))} | entries={result.get('entryCount', '—')}")

                threading.Thread(target=worker, daemon=True).start()

            buttons = tk.Frame(confirm)
            buttons.pack(fill="x", padx=14, pady=(0, 14))
            tk.Button(buttons, text="Huỷ", width=12, command=confirm.destroy).pack(side="right", padx=(8, 0))
            tk.Button(buttons, text="Xác nhận publish catalog", width=23, command=confirm_publish).pack(side="right")

        build_btn = tk.Button(pack_controls, text="1. Build + Validate Local", width=25, bg="#cce6ff", command=run_local_build)
        build_btn.pack(side="left")
        pause_btn = tk.Button(pack_controls, text="⏸ Tạm dừng", width=11, state="disabled", command=toggle_pause_build)
        pause_btn.pack(side="left", padx=(6, 0))
        cancel_btn = tk.Button(pack_controls, text="✖ Huỷ build", width=11, state="disabled", command=cancel_build)
        cancel_btn.pack(side="left", padx=(6, 0))
        tk.Button(pack_controls, text="▶ Phát thử audio local", width=18, command=play_first_audio).pack(side="left", padx=8)
        stage_btn = tk.Button(pack_controls, text="2. Upload + Verify Packs", width=24, state="disabled", command=stage_packs_pending)
        stage_btn.pack(side="left")
        tk.Button(pointer_controls, text="Initialize Production Signing Key", width=25, command=initialize_signing_key_pending).pack(side="left")
        initialize_pointer_btn = tk.Button(pointer_controls, text="Initialize VOCAB POINTER", width=22, command=initialize_pointer_pending)
        initialize_pointer_btn.pack(side="left", padx=(6, 0))
        refresh_pointer_btn = tk.Button(pointer_controls, text="Refresh Pointer Status", width=22, command=refresh_pointer_status_pending)
        refresh_pointer_btn.pack(side="left", padx=(6, 0))
        publish_btn = tk.Button(publish_controls, text="3. Publish Catalog + Signed Pointer", width=30, state="disabled", command=publish_catalog_pending)
        publish_btn.pack(side="left")
        tk.Button(publish_controls, text="? Publish Help", width=14, command=lambda: open_hsk30_help("Cách kích hoạt cấp độ mới trong app")).pack(side="left", padx=(6, 0))

        pack_version_var.trace_add("write", refresh_pack_version_label)
        for var in (excel_var, sheet_var, level_display_var, output_var, builder_bitrate_var, builder_audio_mode_var, pack_version_var):
            var.trace_add("write", clear_build_state)
        refresh_pack_version_label()
        refresh_deploy_gate()
        refresh_signing_status()
        refresh_catalog_gate()
        refresh_pointer_status_pending()


            
    #===Ép ngôn ngữ đã chọn
    def ep_toan_bo_dong_ve_lang():
        lang_da_chon = [k for k, v in ngon_ngu_flags.items() if v.get()]
        if not lang_da_chon:
            tk.messagebox.showwarning("Thiếu lựa chọn", "Bạn cần chọn ít nhất 1 ngôn ngữ.")
            return

        def doan_lang_theo_kytu_va_checkbox(text):
            text = text.strip()

            count_han = sum(0x4E00 <= ord(c) <= 0x9FFF for c in text)
            count_kana = sum(0x3040 <= ord(c) <= 0x30FF for c in text)
            count_vi = sum(c in 'ăâêôơưđáàảãạấầẩẫậắằẳẵặéèẻẽẹếềểễệíìỉĩịóòỏõọốồổỗộớờởỡợúùủũụứừửữựỳýỵỷỹ' for c in text.lower())
            count_ascii = sum('a' <= c.lower() <= 'z' for c in text)

            # Ưu tiên theo kiểu ký tự nếu có trong checkbox
            if count_kana > 0 and "ja" in lang_da_chon:
                return "ja"
            if count_han > 0:
                if "ja" in lang_da_chon:
                    return "ja"
                elif "zh" in lang_da_chon:
                    return "zh"
            if count_vi > 0 and "vi" in lang_da_chon:
                return "vi"
            if count_ascii > 0 and "en" in lang_da_chon:
                return "en"
            # Mặc định nếu không xác định được
            return lang_da_chon[0]

        # Thực hiện ép
        for idx, item in enumerate(danh_sach):
            dong = item[0]
            lang_ep = doan_lang_theo_kytu_va_checkbox(dong)
            selected_lang_list[idx] = lang_ep

        rebuild_page()
        pygame.mixer.init()
        pygame.mixer.music.load(SUCCESS_SOUND)
        pygame.mixer.music.play()
        tk.messagebox.showinfo("Đã ép", "✅ Đã gán ngôn ngữ cho từng dòng theo ký tự & lựa chọn.")

    #===================
    # ====== Toàn cục để dừng đọc âm thanh cũ
# ===== DEAD CODE: GAME UI / GAME AUDIO / GAME IMAGE =====
# DEAD CODE - remove later. Các hàm game bên dưới đã bị ngắt entry point.
    def doc_noi_dung_tung_dong_popup(text, engine, giong, toc_do, callback=None, lang=None):
        # DEAD CODE - remove later. Chỉ còn phục vụ Game Đoán Chữ đã disable.
        import tempfile, os, uuid, pygame, time
        import threading

        def run_doc():
            global current_channel
            try:
                dong_doc = lam_sach_van_ban(text)

                # 🔁 Chỉ tự đoán ngôn ngữ nếu lang=None
                lang_doc = lang
                if not lang:
                    if any('\u3040' <= c <= '\u30ff' or '\u4e00' <= c <= '\u9faf' for c in dong_doc):
                        lang_doc = "ja"
                    elif any('\u4e00' <= c <= '\u9fff' for c in dong_doc):
                        lang_doc = "zh"
                    else:
                        lang_doc = "vi"

                temp_mp3 = os.path.join(tempfile.gettempdir(), f"game_{uuid.uuid4().hex}.mp3")
                tao_file_mp3(dong_doc, lang=lang_doc, voice=giong, toc_do=toc_do, engine=engine, file_out=temp_mp3)

                if not pygame.mixer.get_init():
                    pygame.mixer.init()

                if current_channel and current_channel.get_busy():
                    current_channel.stop()

                sound = pygame.mixer.Sound(temp_mp3)
                current_channel = pygame.mixer.find_channel()
                if current_channel:
                    current_channel.play(sound)
                    while current_channel.get_busy():
                        time.sleep(0.1)

                if callback:
                    callback()

                try:
                    os.remove(temp_mp3)
                except:
                    pass
            except Exception as e:
                print("❌ Lỗi đọc nội dung:", e)

        threading.Thread(target=run_doc, daemon=True).start()



#==============Game
    def doc_dap_an_chinh(thong_bao_text, dap_an_dung, ngonngu, them_nghia_la=False):

    #def doc_dap_an_chinh(thong_bao_text, dap_an_dung, ngonngu):
        global dang_doc_game
        try:
            end_doc_game()

            engine = combo_engine.get()
            giong = combo_giong_popup.get()
            toc_do = combo_toc_do_popup.get()

            import re
            match = re.match(r"^(.*?)[。.．]?\s*(\(|（)(.+?)(\)|）)?$", dap_an_dung.strip())
            if match:
                phan1 = match.group(1).strip()  # phần tiếng gốc
                phan2 = match.group(3).strip()  # nghĩa tiếng Việt
            else:
                phan1 = dap_an_dung.strip()
                phan2 = ""

            # ✅ Xác định mã ngôn ngữ để đọc phần tiếng gốc
            if ngonngu == "Ja":
                lang_phan1 = "ja"
            elif ngonngu == "Cn":
                lang_phan1 = "zh"
            elif ngonngu == "En":
                lang_phan1 = "en"
            elif ngonngu == "Dnn":
                try:
                    from langdetect import detect
                    lang_phan1 = detect(phan1)
                except:
                    lang_phan1 = "vi"
            else:
                lang_phan1 = "vi"


            def ket_thuc_doc():
                global dang_doc_game
                dang_doc_game = False

            def doc_nghia_tv():
                if phan2:
                    text = f"nghĩa là {phan2}" if them_nghia_la else phan2
                    doc_noi_dung_tung_dong_popup(text, engine, giong, toc_do, lang="vi", callback=ket_thuc_doc)
                else:
                    ket_thuc_doc()



            def doc_phan1():
                doc_noi_dung_tung_dong_popup(phan1, engine, giong, toc_do, lang=lang_phan1, callback=doc_nghia_tv)


            def doc_thong_bao():
                doc_noi_dung_tung_dong_popup(thong_bao_text, engine, giong, toc_do, lang="vi", callback=doc_phan1)

            threading.Thread(target=doc_thong_bao, daemon=True).start()

        except Exception as e:
            print("❌ Lỗi đọc đáp án:", e)
            dang_doc_game = False



    def thong_bao_dung(game, dap_an_dung, cau_hoi_full, tiep_cau_tiep, ch=None, ngonngu=None):
        import tkinter as tk
        from PIL import Image, ImageTk
        end_doc_game()  # 🛑 Dừng âm đọc cũ ngay

        popup = tk.Toplevel(game)
        popup.title("🎓 Chính xác!")
        popup.configure(bg="white")
        set_popup_icon(popup)

        popup.update_idletasks()
        w, h = 740, 520
        x = (popup.winfo_screenwidth() // 2) - (w // 2)
        y = (popup.winfo_screenheight() // 2) - (h // 2)
        popup.geometry(f"{w}x{h}+{x}+{y}")
        popup.grab_set()

        tk.Label(popup, text="✅ Đúng rồi!", fg="green", bg="white",
                 font=("Arial", 26, "bold")).pack(pady=10)

        tk.Label(popup, text="🎯 Đáp án đúng là:", fg="green", bg="white",
                 font=("Arial", 16, "bold")).pack()

        tk.Label(popup, text=dap_an_dung, fg="green", bg="white",
                 font=("Arial", 18, "bold")).pack()

        # 🖼️ Hiển thị ảnh nếu có
        if isinstance(ch, dict) and ch and "images" in ch:

            img_path = ""
            if ngonngu == "Dnn":
                for subfolder in ["Ja", "Cn", "En"]:
                    path = os.path.join(IMAGE_FOLDER, subfolder, ch["images"])
                    if os.path.exists(path):
                        img_path = path
                        break
            else:
                thu_muc_anh = os.path.join(IMAGE_FOLDER, ngonngu)
                path = os.path.join(thu_muc_anh, ch["images"])
                if os.path.exists(path):
                    img_path = path

            if img_path:
                try:
                    img = Image.open(img_path).resize((180, 180))
                    tk_img = ImageTk.PhotoImage(img)
                    lbl_img = tk.Label(popup, image=tk_img, bg="white")
                    lbl_img.image = tk_img
                    lbl_img.pack(pady=5)
                except Exception as e:
                    print("⚠ Không thể hiển thị ảnh trong popup:", e)


        tk.Label(popup, text=cau_hoi_full, fg="black", bg="white",
                 font=("Arial", 14), justify="left", wraplength=w - 40).pack(pady=10)

        tk.Button(popup, text="Tiếp tục", font=("Arial", 14),
                  command=lambda: [popup.destroy(), end_doc_game(), tiep_cau_tiep()]).pack(pady=15)

        global dang_doc_game
        dang_doc_game = True
        doc_dap_an_chinh("Chính xác! Đáp án là:", dap_an_dung, ngonngu)


    

    def thong_bao_sai(game, tieu_de, dap_an_dung, cau_hoi_full, on_close, ch=None, ngonngu=None):
        import tkinter as tk
        from PIL import Image, ImageTk
        import os

        end_doc_game()

        popup = tk.Toplevel(game)
        popup.title("❌ Sai rồi!")
        popup.configure(bg="white")
        set_popup_icon(popup)

        popup.update_idletasks()
        w, h = 740, 520
        x = (popup.winfo_screenwidth() // 2) - (w // 2)
        y = (popup.winfo_screenheight() // 2) - (h // 2)
        popup.geometry(f"{w}x{h}+{x}+{y}")
        popup.grab_set()
        popup.protocol("WM_DELETE_WINDOW", lambda: None)

        # Tiêu đề sai
        tk.Label(popup, text=tieu_de, fg="red", bg="white",
                 font=("Arial", 24, "bold")).pack(pady=10)

        # Đáp án đúng
        tk.Label(popup, text="🟢 Đáp án đúng phải là:", font=("Arial", 16, "bold"),
                 fg="green", bg="white").pack()
        tk.Label(popup, text=dap_an_dung, font=("Arial", 18, "bold"),
                 fg="green", bg="white").pack()

        # Ảnh nếu có
        if isinstance(ch, dict) and ch and "images" in ch:
            img_path = ""
            if ngonngu == "Dnn":
                for subfolder in ["Ja", "Cn", "En"]:
                    path = os.path.join(IMAGE_FOLDER, subfolder, ch["images"])
                    if os.path.exists(path):
                        img_path = path
                        break
            else:
                thu_muc_anh = os.path.join(IMAGE_FOLDER, ngonngu)
                path = os.path.join(thu_muc_anh, ch["images"])
                if os.path.exists(path):
                    img_path = path

            if img_path:
                try:
                    img = Image.open(img_path).resize((180, 180))
                    tk_img = ImageTk.PhotoImage(img)
                    lbl_img = tk.Label(popup, image=tk_img, bg="white")
                    lbl_img.image = tk_img
                    lbl_img.pack(pady=5)
                except Exception as e:
                    print("⚠ Không thể hiển thị ảnh trong popup:", e)


        # Câu hỏi đầy đủ
        tk.Label(popup, text=cau_hoi_full, font=("Arial", 14), bg="white",
                 justify="left", wraplength=w - 40).pack(pady=10)

        # Nút tiếp tục
        def dong_popup():
            popup.destroy()
            end_doc_game()
            on_close()

        tk.Button(popup, text="▶ Tiếp tục", font=("Arial", 14),
                  command=dong_popup).pack(pady=15)

        # 🔊 Đọc nội dung giống như popup đúng
        global dang_doc_game
        dang_doc_game = True
        #doc_dap_an_chinh(tieu_de + " Đáp án đúng phải là:", dap_an_dung, ngonngu)
        doc_dap_an_chinh(f"{tieu_de} Đáp án đúng phải là:", dap_an_dung, ngonngu, them_nghia_la=True)





#==============

    def choi_game_ngay_trong_popup(df, ten, level, ngonngu):
        import tempfile, threading, time
        from pydub import AudioSegment
        import pygame
        timer_thread = None
        
        game = tk.Toplevel(popup)
        set_popup_icon(game)
        game.title("Game Đoán Chữ")
        game.state('zoomed')     # ✅ Full màn hình an toàn
        game.grab_set()
        # Khởi tạo điểm
        diem = 0

        # Hiển thị thông tin người chơi ở góc trái
        label_thong_tin = tk.Label(game,

            text=f"{ten} \n Lv{level} – {diem}đ",
            font=("Arial", 15, "bold"), fg="green")
        label_thong_tin.place(x=5, y=5)

        label_cau = tk.Label(game, text="", font=("Arial", 28, "bold"), fg="red", wraplength=900, justify="center") # Tăng cỡ chữ câu hỏi, đáp án, và đổi màu

        label_cau.pack(pady=20)

        frame_btn = tk.Frame(game)
        frame_btn.pack()

        #timer_label = tk.Label(game, text="", font=("Arial", 14), fg="blue") #Đổi màu đếm ngược và cỡ chữ
        timer_label = tk.Label(game, text="", font=("Arial", 20, "bold"), fg="red", bg="white", bd=2, relief="solid")
        timer_label.pack()


        result_label = tk.Label(game, text="", font=("Arial", 14), fg="blue")
        result_label.pack(pady=10)

        cau_hoi = df.sample(level * 5).to_dict(orient="records")
        idx = 0
        diem = 0
        stop_timer = [False]
        #========
        #Gửi discor điểm game
        def lay_top_diem(ngonngu):
            try:
                lich_su_daydu = tai_lich_su()
                if not lich_su_daydu:
                    return "Không rõ", 0, 0

                # Lọc theo mã ngôn ngữ
                top = max(
                    [r for r in lich_su_daydu if r.get("ngonngu") == ngonngu],
                    key=lambda r: (r.get("level", 0), r.get("diem", 0)),
                    default=None
                )

                if top:
                    return top["ten"], top["level"], top["diem"]
            except Exception as e:
                print("❌ Lỗi khi lấy top điểm:", e)

            return "Không rõ", 0, 0

        #======
        def gui_discord_gui_diem(ten, level, diem, ngonngu):
            import requests
            from datetime import datetime
            import tempfile, cv2, os, threading
            from tkinter import messagebox
            from pygrabber.dshow_graph import FilterGraph

            webhook_url = DISCORD_WEBHOOK_URL
            if not webhook_url or "api/webhooks/" not in webhook_url:
                messagebox.showerror("Lỗi", "Webhook Discord không hợp lệ.")
                return

            # 👉 Hiện popup đang gửi
            sending_popup = tk.Toplevel()
            set_popup_icon(sending_popup)
            sending_popup.title("Đang gửi Discord...")
            sending_popup.geometry("360x100")
            sending_popup.attributes("-topmost", True)
            tk.Label(sending_popup, text="🔄 Đang gửi điểm và ảnh lên Discord...", font=("Arial", 12)).pack(pady=20)
            sending_popup.update()

            try:
                # 🔠 Tên ngôn ngữ
                ngon_ngu_text = {
                    "Ja": "Tiếng Nhật",
                    "Cn": "Tiếng Trung",
                    "En": "Tiếng Anh",
                    "Dnn": "Đa ngôn ngữ"
                }.get(ngonngu, "Không xác định")

                # 🏆 Lấy top điểm
                top_name, top_level, top_diem = lay_top_diem(ngonngu)
                thoigian = datetime.now().strftime("Lúc %Hh%M ngày %d/%m/%Y")

                # 📸 Lấy danh sách tên camera
                graph = FilterGraph()
                danh_sach_ten_cam = graph.get_input_devices()  # List[str]

                # 📷 Chụp ảnh từ nhiều camera
                image_files = []
                image_cam_names = {}  # Map từ ảnh sang tên cam
                threads = []

                def chup_anh_tu_cam(cam_index):
                    try:
                        cap = cv2.VideoCapture(cam_index, cv2.CAP_DSHOW)
                        if cap.isOpened():
                            ret, frame = cap.read()
                            if ret:
                                img_path = os.path.join(tempfile.gettempdir(), f"webcam_{cam_index}.jpg")
                                cv2.imwrite(img_path, frame)
                                cam_name = danh_sach_ten_cam[cam_index] if cam_index < len(danh_sach_ten_cam) else f"Camera {cam_index}"
                                print(f"✅ Cam {cam_index}: {cam_name} – đã chụp {img_path}")
                                image_files.append(img_path)
                                image_cam_names[img_path] = cam_name
                            else:
                                print(f"⚠️ Cam {cam_index} mở được nhưng không lấy được frame.")
                        else:
                            print(f"❌ Cam {cam_index} không mở được.")
                        cap.release()
                    except Exception as e:
                        print(f"❌ Cam {cam_index} lỗi: {e}")

                # 👁 Danh sách camera theo thứ tự ưu tiên (C270 là cam 1)
                indices = [1, 0, 2, 3, 4]

                for i in indices:
                    t = threading.Thread(target=chup_anh_tu_cam, args=(i,))
                    t.start()
                    threads.append(t)

                for t in threads:
                    t.join(timeout=5)

                # 📑 Tạo nội dung gửi Discord
                noi_dung = (
                    f"📢 **{thoigian}**\n\n"
                    f"🎉 __**{ten.upper()}**__ đã hoàn thành **Level {level} {ngon_ngu_text}** với số điểm 👉 __**{diem * 10} điểm**__ 🎯 trong Game Đoán Chữ!\n\n"
                    f"🏆 Kỷ lục {ngon_ngu_text} hiện tại: {top_name} – Level {top_level} – {top_diem} điểm!\n"
                )

                if image_files:
                    noi_dung += "\n📷 Ảnh từ các camera:\n"
                    for img in image_files:
                        ten_cam = image_cam_names.get(img, "Webcam")
                        noi_dung += f"- {ten_cam}\n"
                else:
                    noi_dung += "\n⚠️ Không có camera khả dụng để chụp ảnh."

                noi_dung += "\n\n© Cty TNHH Du Lịch & Thương Mại Quốc Tế Việt Trung Nhật\n📞 0986183806"

                # 📤 Chuẩn bị file gửi
                files = {}
                for idx, img in enumerate(image_files):
                    try:
                        f = open(img, "rb")
                        files[f"file{idx+1}"] = (os.path.basename(img), f, "image/jpeg")
                        print(f"📤 Chuẩn bị gửi ảnh: {img}")
                    except Exception as e:
                        print(f"⚠️ Không thể mở ảnh {img} để gửi: {e}")

                # 🔗 Gửi
                response = requests.post(webhook_url, data={"content": noi_dung}, files=files if files else None)

                # 🧹 Đóng và xoá ảnh tạm
                for f in files.values():
                    try:
                        f[1].close()
                        os.remove(os.path.join(tempfile.gettempdir(), f[0]))
                        print(f"🧹 Đã xoá ảnh tạm: {f[0]}")
                    except Exception as e:
                        print(f"⚠️ Không xoá được ảnh {f[0]}: {e}")

                sending_popup.destroy()
                if response.status_code in [200, 204]:
                    messagebox.showinfo("✅ Thành công", "Đã gửi điểm và ảnh lên Discord!")
                else:
                    messagebox.showerror("Lỗi", f"❌ Không gửi được lên Discord.\nHTTP {response.status_code}\n{response.text}")

            except Exception as e:
                sending_popup.destroy()
                messagebox.showerror("Lỗi", f"❌ Lỗi gửi Discord:\n{e}")

        #===================



        def hien_cau_hoi():
            nonlocal idx, diem, timer_label, timer_thread, stop_timer
            global dang_doc_game
            nonlocal frame_btn, ngonngu  # ✅ thêm nonlocal ngonngu

            if dang_doc_game:
                game.after(500, hien_cau_hoi)
                return

            for widget in frame_btn.winfo_children():
                widget.destroy()
            if timer_label:
                timer_label.destroy()

            if idx >= len(cau_hoi):
                #if diem == len(cau_hoi):  # ✅ Đúng hết
                if diem >= int(len(cau_hoi) * 0.7):  # ✅ Qua level nếu đúng ≥ 70%

                    result_label.config(text=f"Hoàn thành level {level}!", fg="green")
                    play_sound(WIN_SOUND)
                    luu_diem(ten, level, diem * 10, ngonngu=ngonngu)

                    popup = tk.Toplevel(game)
                    popup.title("🎉 Hoàn thành level!")
                    popup.configure(bg="white")
                    w, h = 740, 420
                    x = (popup.winfo_screenwidth() // 2) - (w // 2)
                    y = (popup.winfo_screenheight() // 2) - (h // 2)
                    popup.geometry(f"{w}x{h}+{x}+{y}")
                    popup.transient(game)
                    popup.grab_set()
                    set_popup_icon(popup)
                    popup.protocol("WM_DELETE_WINDOW", lambda: None)

                    tk.Label(popup, text=f"🎉 Bạn đã hoàn thành Level {level}!", font=("Arial", 16, "bold"),
                             fg="green", bg="white").pack(pady=10)


                    
                    # === Thông báo chúc mừng 2 ngôn ngữ ===
                    ngon_ngu_text = {"Ja": "tiếng Nhật", "Cn": "tiếng Trung", "En": "tiếng Anh", "Dnn": "đa ngôn ngữ"}.get(ngonngu, "không xác định")
                    thong_bao_tv = f"🎉 Chúc mừng {ten} đã qua level {level} ({ngon_ngu_text}) với số điểm {diem * 10}!"

                    thong_bao_ngonngu = ""

                    if ngonngu == "Ja":
                        thong_bao_ngonngu = f"🎉 {ten}さん、おめでとうございます！レベル{level}（日本語）を{diem*10}点で合格しました！"
                    elif ngonngu == "Cn":
                        thong_bao_ngonngu = f"🎉 恭喜 {ten} 通过了第 {level} 级（中文）考试，得分为 {diem * 10} 分！"
                    elif ngonngu == "En":
                        thong_bao_ngonngu = f"🎉 Congratulations {ten} for passing level {level} (English) with score {diem * 10}!"
                    elif ngonngu == "Dnn":
                        thong_bao_tv = f"🎉 Chúc mừng {ten} đã qua level {level} phần đa ngôn ngữ với số điểm {diem * 10}!"

                    tk.Label(popup, text=thong_bao_tv, font=("Arial", 14, "bold"),
                             fg="green", bg="white", wraplength=460, justify="center").pack(pady=10)

                    if thong_bao_ngonngu:
                        tk.Label(popup, text=thong_bao_ngonngu, font=("Arial", 13),
                                 fg="black", bg="white", wraplength=460, justify="center").pack(pady=5)

                    # === Đọc 2 dòng thông báo ===
                    try:
                        engine = combo_engine.get()
                        giong = combo_giong_popup.get()
                        toc_do = combo_toc_do_popup.get()

                        def doc_tv():
                            level_text = f"level {level}"
                            if ngonngu == "Ja":
                                level_text += " tiếng Nhật"
                            elif ngonngu == "Cn":
                                level_text += " tiếng Trung"
                            elif ngonngu == "En":
                                level_text += " tiếng Anh"
                            elif ngonngu == "Dnn":
                                level_text += " phần đa ngôn ngữ"

                            text_full = f"{thong_bao_tv} ({level_text})"
                            doc_noi_dung_tung_dong_popup(text_full, engine, giong, toc_do, lang="vi")


                        if thong_bao_ngonngu:
                            lang_map = {"Ja": "ja", "Cn": "zh", "En": "en"}
                            lang_code = lang_map.get(ngonngu, "vi")
                            doc_noi_dung_tung_dong_popup(thong_bao_ngonngu, engine, giong, toc_do, lang=lang_code, callback=doc_tv)
                        else:
                            doc_tv()
                    except Exception as e:
                        print("❌ Lỗi đọc thông báo chúc mừng:", e)




                    tk.Label(popup, text="Bạn muốn tiếp tục chơi hay lưu kết quả lại?", font=("Arial", 13),
                             bg="white").pack(pady=5)

                    btn_frame = tk.Frame(popup, bg="white")
                    btn_frame.pack(pady=15)

                    def tiep():
                        popup.destroy()
                        game.destroy()
                        choi_game_ngay_trong_popup(df, ten, level + 1, ngonngu=ngonngu)

                    def luu_ket_qua():
                        popup.destroy()
                        game.destroy()

                    def thoat():
                        popup.destroy()
                        game.destroy()

                    tk.Button(btn_frame,
                              text="📤 Gửi Discord",
                              font=("Arial", 13),
                              width=18,
                              bg="#7289DA",
                              fg="white",
                              command=lambda: threading.Thread(
                                  target=gui_discord_gui_diem,
                                  args=(ten, level, diem, ngonngu),
                                  daemon=True).start()
                              ).grid(row=3, column=0, padx=10, pady=5)


                    tk.Button(btn_frame, text="▶ Tiếp tục chơi", font=("Arial", 13), width=18,
                              bg="lightgreen", command=tiep).grid(row=0, column=0, padx=10, pady=5)
                    tk.Button(btn_frame, text="💾 Lưu kết quả", font=("Arial", 13), width=18,
                              bg="lightblue", command=luu_ket_qua).grid(row=1, column=0, padx=10, pady=5)
                    tk.Button(btn_frame, text="❌ Thoát", font=("Arial", 13), width=18,
                              bg="tomato", command=thoat).grid(row=2, column=0, padx=10, pady=5)

                else:
                    # ❌ Sai ≥ 1 câu
                    play_sound(FALSE_SOUND)
                    popup = tk.Toplevel(game)
                    popup.title("Game Over")
                    popup.geometry("620x420")
                    popup.transient(game)
                    popup.grab_set()
                    set_popup_icon(popup)
                    popup.protocol("WM_DELETE_WINDOW", lambda: None)

                    tk.Label(popup, text="Game over! Bạn chưa trả lời đúng tất cả!", font=("Arial", 18), fg="red").pack(pady=10)
                    tk.Label(popup, text=f"Điểm của bạn: {diem}/{len(cau_hoi)}", font=("Arial", 14)).pack(pady=5)
                    tk.Label(popup, text="Bạn có muốn chơi lại level này không?", font=("Arial", 14)).pack(pady=5)

                    def choi_lai():
                        popup.destroy()
                        game.destroy()
                        choi_game_ngay_trong_popup(df, ten, level, ngonngu=ngonngu)

                    def thoat():
                        popup.destroy()
                        game.destroy()

                    tk.Button(popup, text="🔁 Chơi lại", font=("Arial", 13), width=18,
                              bg="orange", command=choi_lai).pack(pady=5)
                    tk.Button(popup, text="❌ Thoát", font=("Arial", 13), width=18,
                              bg="tomato", command=thoat).pack(pady=5)

                return


            
            ch = cau_hoi[idx]
            cau_so = idx + 1
            label_cau.config(text=f"Câu hỏi số {cau_so}/{len(cau_hoi)} – Level {level}\n\n{ch['Câu hỏi']}\n{ch['Câu rút gọn']}")

            dap_an_dung = ch["Đáp án"] + f" ({ch['Nghĩa TV']})"
            cac_dap_an = df.sample(3).apply(lambda r: r["Đáp án"] + f" ({r['Nghĩa TV']})", axis=1).tolist()
            cac_dap_an.append(dap_an_dung)
            random.shuffle(cac_dap_an)

            global current_channel
            if current_channel and current_channel.get_busy():
                current_channel.stop()

            try:
                cau_chinh = ch["Câu hỏi"]
                cau_rut_gon = ch["Câu rút gọn"]
                engine = combo_engine.get()
                toc_do = combo_toc_do_popup.get()
                giong_mac_dinh = combo_giong_popup.get()

                def doc_cau_hoi_dung_polly():
                    import pygame, tempfile, uuid
                    from langdetect import detect

                    if ngonngu == "Dnn":
                        try: lang_code = detect(cau_chinh)
                        except: lang_code = "vi"
                    else:
                        lang_code = {"Ja": "ja", "Cn": "zh", "En": "en"}.get(ngonngu, "vi")

                    path_nam = os.path.join(tempfile.gettempdir(), f"game_{uuid.uuid4().hex}_nam.mp3")
                    path_nu = os.path.join(tempfile.gettempdir(), f"game_{uuid.uuid4().hex}_nu.mp3")

                    tao_file_mp3(cau_chinh, lang=lang_code, voice="Nam", toc_do=toc_do, engine="Polly", file_out=path_nam)
                    pygame.mixer.init()
                    pygame.mixer.music.load(path_nam)
                    pygame.mixer.music.play()
                    while pygame.mixer.music.get_busy(): time.sleep(0.1)
                    time.sleep(0.5)
                    tao_file_mp3(cau_rut_gon, lang=lang_code, voice="Nữ", toc_do=toc_do, engine="Polly", file_out=path_nu)
                    pygame.mixer.music.load(path_nu)
                    pygame.mixer.music.play()

                if engine == "Polly":
                    threading.Thread(target=doc_cau_hoi_dung_polly, daemon=True).start()
                else:
                    if ngonngu == "Dnn":
                        from langdetect import detect
                        try: lang_cau_hoi = detect(cau_chinh)
                        except: lang_cau_hoi = "vi"
                    else:
                        lang_cau_hoi = {"Ja": "ja", "Cn": "zh", "En": "en"}.get(ngonngu, "vi")
                    doc_noi_dung_tung_dong_popup(cau_chinh + "\n" + cau_rut_gon, engine, giong_mac_dinh, toc_do, lang=lang_cau_hoi)

            except Exception as e:
                print("❌ Lỗi đọc câu hỏi:", e)





            stop_timer = [False]

            def dem_nguoc(ngon_ngu_chon):  # truyền vào thay vì nonlocal
                nonlocal ch
                tg = 25  # thời gian đếm ngược
                while tg > 0 and not stop_timer[0]:
                    try:
                        timer_label.config(text=f"⏳ {tg}s", fg="blue")
                    except:
                        return
                    tg -= 1
                    time.sleep(1)

                if not stop_timer[0]:
                    if current_channel and current_channel.get_busy():
                        current_channel.stop()
                    app_beep(1000, 200)
                    app_beep(1000, 200)

                    try:
                        if ch:
                            thong_bao_sai(game, "⏰ Hết giờ!", dap_an_dung,
                                          f"{ch.get('Câu hỏi', '')}\n{ch.get('Câu rút gọn', '')}",
                                          tiep_cau_tiep, ch, ngonngu=ngon_ngu_chon)
                        else:
                            print("⚠️ Không có câu hỏi hiện tại để thông báo sai.")
                    except Exception as e:
                        print("❌ Lỗi khi thông báo hết giờ:", e)




            def tiep_cau_tiep():
                nonlocal idx
                idx += 1
                hien_cau_hoi()

            timer_label = tk.Label(game, text="", font=("Arial", 68, "bold"), fg="red", bg="white", bd=2, relief="solid")
            timer_label.place(relx=1.0, rely=0.0, anchor="ne")  # Góc phải trên
           
            timer_thread = threading.Thread(target=lambda: dem_nguoc(ngonngu))

            timer_thread.daemon = True
            timer_thread.start()

            def xu_ly_chon(ans):
                nonlocal idx, diem
                stop_timer[0] = True

                if current_channel and current_channel.get_busy():
                    current_channel.stop()

                if ans == dap_an_dung:
                    play_sound(TRUE_SOUND)
                    diem += 1
                    label_thong_tin.config(text=f"{ten} – Lv{level} – {diem * 10}đ")
                    thong_bao_dung(game, dap_an_dung, f"{ch['Câu hỏi']}\n{ch['Câu rút gọn']}", tiep_cau_tiep, ch, ngonngu)
                else:
                    play_sound(FALSE_SOUND)
                    thong_bao_sai(game, "Sai rồi!", dap_an_dung, f"{ch['Câu hỏi']}\n{ch['Câu rút gọn']}", tiep_cau_tiep, ch, ngonngu)






            #=========
            # Vùng hiển thị ảnh và nút đáp án
            # 🧹 Xoá vùng đáp án cũ (nếu có)
            for widget in game.winfo_children():
                if widget not in [label_thong_tin, label_cau, timer_label, result_label]:
                    widget.destroy()

            # 🧱 Frame chứa các đáp án
            frame_btn = tk.Frame(game)
            frame_btn.pack(pady=10)

            for i, da in enumerate(cac_dap_an):
                frame_cot = tk.Frame(frame_btn)
                frame_cot.grid(row=0, column=i, padx=10, pady=10)

                ten_dap_an_goc = da.split(" (")[0].strip()
                dong_anh = df[df["Đáp án"].str.strip() == ten_dap_an_goc]

                duong_dan_anh = ""
                if not dong_anh.empty:
                    ten_file_anh = dong_anh.iloc[0].get("images", "")
                    duong_dan_anh = ""
                    if ten_file_anh:
                        # Nếu là đa ngôn ngữ => dò trong 3 thư mục
                        if ngonngu == "Dnn":
                            for subfolder in ["Ja", "Cn", "En"]:
                                path = os.path.join(IMAGE_FOLDER, subfolder, ten_file_anh)
                                if os.path.exists(path):
                                    duong_dan_anh = path
                                    break
                        else:
                            thu_muc_anh = os.path.join(IMAGE_FOLDER, ngonngu)
                            path = os.path.join(thu_muc_anh, ten_file_anh)
                            if os.path.exists(path):
                                duong_dan_anh = path



                print(f"🔍 Đang tìm ảnh cho đáp án: {ten_dap_an_goc}")
                print("📂 Ảnh tìm được:", duong_dan_anh)

                if duong_dan_anh and os.path.exists(duong_dan_anh):
                    try:
                        img = Image.open(duong_dan_anh).resize((180, 180))
                        tk_img = ImageTk.PhotoImage(img)
                        label_img = tk.Label(frame_cot, image=tk_img)
                        label_img.image = tk_img  # giữ tham chiếu
                        label_img.pack()
                    except Exception as e:
                        print("❌ Không thể hiển thị ảnh đáp án:", e)

                btn = tk.Button(frame_cot, text=da, width=30, font=("Arial", 13), bg="#f0f8ff", fg="darkblue",
                                command=lambda a=da: xu_ly_chon(a))
                btn.pack(pady=4)

                
        #===
        hien_cau_hoi()
        
        # Khi người dùng đóng cửa sổ, dừng đọc
        def on_close():
            try:
                if current_channel and current_channel.get_busy():
                    current_channel.stop()
            except:
                pass
            game.destroy()

        game.protocol("WM_DELETE_WINDOW", on_close)
        

            
#===Cửa sổ bắt đầu chọn chơi game đa ngôn ngữ
        
    def bat_dau_game_popup():
        # DEAD CODE - remove later. Entry point game đã bị ẩn khỏi UI.
        print("Disabled: Game Đoán Chữ removed in audio-tool version.")
        try:
            messagebox.showinfo("Đã tắt", "Game Đoán Chữ đã được tắt.", parent=popup)
        except Exception:
            pass
        return

        from datetime import datetime

        popup_game_start = tk.Toplevel(popup)
        popup_game_start.title("Bắt đầu Game Đoán Chữ")
        popup_game_start.attributes("-fullscreen", True)
        set_popup_icon(popup_game_start)
        popup_game_start.grab_set()

        tk.Button(popup_game_start, text="🗕", command=lambda: popup_game_start.iconify()).place(relx=0.95, rely=0.01)
        tk.Button(popup_game_start, text="❌", fg="red", command=popup_game_start.destroy).place(relx=0.98, rely=0.01)

        tk.Label(popup_game_start, text="Nhập tên của bạn:").pack(pady=5)
        entry_ten = tk.Entry(popup_game_start, font=("Arial", 14))
        entry_ten.insert(0, "Phương Anh")
        entry_ten.pack(pady=5)

        tk.Label(popup_game_start, text="Chọn ngôn ngữ:").pack(pady=5)
        combo_ngon_ngu = ttk.Combobox(popup_game_start, font=("Arial", 14), state="readonly")
        combo_ngon_ngu['values'] = ["Tiếng Nhật", "Tiếng Trung", "Tiếng Anh", "Đa ngôn ngữ"]
        combo_ngon_ngu.set("Tiếng Nhật")
        combo_ngon_ngu.pack(pady=5)

        frame_level = tk.Frame(popup_game_start)
        frame_level.pack(pady=5)
        tk.Label(frame_level, text="Chọn level (1 = 5 câu, 2 = 10 câu,...): 10 điểm mỗi câu").pack()
        entry_level = tk.Entry(frame_level, font=("Arial", 14))
        entry_level.insert(0, "1")
        entry_level.pack()

        def hien_level_combobox(level_max):
            for widget in frame_level.winfo_children():
                widget.destroy()
            tk.Label(frame_level, text=f"Chọn level (tối đa Lv{level_max}):").pack()
            combo = ttk.Combobox(frame_level, font=("Arial", 14), state="readonly")
            combo['values'] = list(range(1, level_max + 1))
            combo.set(level_max)
            combo.pack()
            return combo

        try:
            lich_su_daydu = tai_lich_su()
            lich_su = lich_su_daydu[-100:][::-1]

            def loc_top(lst, ma_ngon_ngu):
                return max([r for r in lst if r.get("ngonngu") == ma_ngon_ngu], key=lambda r: (r["level"], r["diem"]), default=None)

            top_ja = loc_top(lich_su_daydu, "Ja")
            top_cn = loc_top(lich_su_daydu, "Cn")
            top_en = loc_top(lich_su_daydu, "En")
            top_dnn = loc_top(lich_su_daydu, "Dnn")
            top_recent = max(lich_su, key=lambda r: (r["level"], r["diem"]), default=None)

            def to_text_top(top, label):
                if not top: return None
                ma = top.get("ngonngu", "?")
                tg = top.get("thoigian", "")
                return f"🏆 Top 1 {label}: {top['ten']} ({ma}: {top['diem']}đ Lv{top['level']} {tg})"

            frame_top4 = tk.Frame(popup_game_start)
            frame_top4.pack(pady=5)

            top_ngon_ngu = [
                (top_ja, "Tiếng Nhật"),
                (top_cn, "Tiếng Trung"),
                (top_en, "Tiếng Anh"),
                (top_dnn, "Đa ngôn ngữ")
            ]

            for i, (top, label) in enumerate(top_ngon_ngu):
                text = to_text_top(top, label)
                if text:
                    row, col = i // 2, i % 2
                    tk.Label(frame_top4, text=text, font=("Arial", 11, "bold"), fg="red", anchor="w", justify="left").grid(row=row, column=col, sticky="w", padx=10, pady=2)

            if top_recent:
                ma = top_recent.get("ngonngu", "?")
                tg = top_recent.get("thoigian", "")
                tk.Label(popup_game_start,
                         text=f"📈 Top 1 gần nhất: {top_recent['ten']} ({ma}: {top_recent['diem']}đ Lv{top_recent['level']} {tg})",
                         font=("Arial", 10), fg="blue").pack()

            btn_bat_dau = tk.Button(popup_game_start, text="Bắt đầu chơi", font=("Arial", 12),
                                    bg="lightgreen", command=lambda: start_game(entry_ten.get(), entry_level.get(), combo_ngon_ngu.get()))
            btn_bat_dau.pack(pady=10)

            khung_luoi = tk.Frame(popup_game_start)
            khung_luoi.pack(pady=10)

            nut_rong_pixel = 200
            man_hinh_rong = popup_game_start.winfo_screenwidth()
            so_cot = max(1, man_hinh_rong // nut_rong_pixel)

            for i, ng in enumerate(lich_su):
                tg = ng.get("thoigian", "")
                ma = ng.get("ngonngu", "?")
                text = f"{ng['ten']} ({ma}: {ng['diem']}đ Lv{ng['level']} {tg})"
                fg = "blue" if ng == top_recent else "red" if ng in [top_ja, top_cn, top_en, top_dnn] else "black"

                def gan_thong_tin(ng=ng):
                    entry_ten.delete(0, tk.END)
                    entry_ten.insert(0, ng["ten"])
                    ten_day_du = {"Ja": "Tiếng Nhật", "Cn": "Tiếng Trung", "En": "Tiếng Anh", "Dnn": "Đa ngôn ngữ"}.get(ng.get("ngonngu", ""), "Đa ngôn ngữ")
                    combo_ngon_ngu.set(ten_day_du)
                    nonlocal entry_level
                    entry_level = hien_level_combobox(ng["level"] + 1)

                row, col = i // so_cot, i % so_cot
                tk.Button(khung_luoi, text=text, font=("Arial", 8), width=30, fg=fg, command=gan_thong_tin).grid(row=row, column=col, padx=2, pady=2)

        except Exception as e:
            print("❌ Lỗi hiển thị người chơi gần đây:", e)

        def start_game(ten=None, level_raw=None, ngon_ngu_chon=None):
            ten = ten or entry_ten.get().strip()
            if not ten:
                messagebox.showerror("Lỗi", "Vui lòng nhập tên.")
                return
            try:
                level = int(level_raw)
            except:
                messagebox.showerror("Lỗi", "Level phải là số.")
                return

            ma_hoa = {
                "Tiếng Nhật": "Ja",
                "Tiếng Trung": "Cn",
                "Tiếng Anh": "En",
                "Đa ngôn ngữ": "Dnn"
            }
            ngonngu_ma = ma_hoa.get(ngon_ngu_chon, ngon_ngu_chon)

            # 🧠 Kiểm tra người chơi cũ
            lich_su = tai_lich_su()
            muc_gan_dung = [u for u in lich_su if u["ten"] == ten and u.get("ngonngu") == ngonngu_ma]
            nguoi_cu = max(muc_gan_dung, key=lambda u: u["level"], default=None)

            if nguoi_cu:
                level_max = nguoi_cu["level"] + 1
                if level > level_max:
                    messagebox.showwarning("Giới hạn", f"Bạn chỉ có thể chọn level ≤ {level_max}.")
                    return
            else:
                if level != 1:
                    messagebox.showwarning("Người mới", "Người mới phải bắt đầu từ level 1.")
                    return

            try:
                df = doc_du_lieu_game(ngonngu_ma)
            except Exception as e:
                messagebox.showerror("Lỗi", f"Lỗi khi đọc dữ liệu:\n{e}")
                return

            if df.empty:
                messagebox.showerror("Lỗi", "File Game_doan_chu.xlsx trống hoặc lỗi")
                return

            popup_game_start.destroy()
            choi_game_ngay_trong_popup(df, ten, level, ngonngu_ma)


#====
    #==Gọi từ ngoài
    if mo_tu_ben_ngoai:
        popup.after(300, bat_dau_game_popup)
#====
    def tao_video_game_popup():
        # DEAD CODE - remove later. Video game đã bị ẩn khỏi UI.
        print("Disabled: game video export removed in audio-tool version.")
        try:
            messagebox.showinfo("Đã tắt", "Tính năng tạo video game đã được tắt.", parent=popup)
        except Exception:
            pass
        return

        popup_video = tk.Toplevel(popup)
        set_popup_icon(popup_video)
        popup_video.title("Tạo Video Game Đoán Chữ")
        popup_video.geometry("850x600")
        popup_video.grab_set()

        tk.Label(popup_video, text="Chọn ngôn ngữ:").pack()
        combo_ngon_ngu = ttk.Combobox(popup_video, font=("Arial", 14), state="readonly")
        combo_ngon_ngu['values'] = ["Tiếng Nhật", "Tiếng Trung", "Tiếng Anh", "Đa ngôn ngữ"]
        combo_ngon_ngu.set("Tiếng Nhật")
        combo_ngon_ngu.pack(pady=5)

        tk.Label(popup_video, text="Chọn hình nền:").pack()
        bg_path_var = tk.StringVar()
        tk.Entry(popup_video, textvariable=bg_path_var, width=50).pack()
        tk.Button(popup_video, text="Chọn...", command=lambda: bg_path_var.set(filedialog.askopenfilename(title="Chọn hình nền", filetypes=[("Image files", "*.png;*.jpg")]))).pack()

        tk.Label(popup_video, text="Chọn từ câu số:").pack()
        entry_start = tk.Entry(popup_video)
        entry_start.insert(0, "1")
        entry_start.pack()

        tk.Label(popup_video, text="Đến câu số:").pack()
        entry_end = tk.Entry(popup_video)
        entry_end.insert(0, "10")
        entry_end.pack()

        def run_video_game():
            try:
                # Lấy ngôn ngữ được chọn
                ngonngu_raw = combo_ngon_ngu.get()
                ma_hoa = {
                    "Tiếng Nhật": "Ja",
                    "Tiếng Trung": "Cn",
                    "Tiếng Anh": "En",
                    "Đa ngôn ngữ": "Dnn"
                }
                ngonngu = ma_hoa.get(ngonngu_raw, ngonngu_raw)

                # Đọc dữ liệu
                df = doc_du_lieu_game(ngonngu)

                start = int(entry_start.get()) - 1
                end = int(entry_end.get())
                bg = bg_path_var.get()
                if not os.path.exists(bg):
                    messagebox.showerror("Lỗi", "Hình nền không tồn tại!")
                    return
                sub_df = df.iloc[start:end]
                if sub_df.empty:
                    messagebox.showerror("Lỗi", "Không có câu nào trong khoảng đã chọn.")
                    return

                popup_video.destroy()
                tao_video_game_doan_chu(sub_df, bg)
            except Exception as e:
                messagebox.showerror("Lỗi", str(e))

        tk.Button(popup_video, text="Tạo video", command=lambda: threading.Thread(target=run_video_game).start(), bg="lightgreen").pack(pady=10)

# ===== END DEAD CODE: GAME UI / GAME AUDIO / GAME IMAGE =====
#=====


    # Buttons (2 cột cho gọn hơn)
    tk.Label(option_frame, text="Thao tác nhanh", font=("Arial", 10, "bold"), bg="#f8fff8").pack(anchor="w", padx=10, pady=(10, 4))

    button_grid = tk.Frame(option_frame, bg="#f8fff8")
    button_grid.pack(fill="x", padx=10, pady=(0, 8))

    action_buttons = [
        # Disabled: game/video/background/subtitle features removed in audio-tool version.
        # ("🎮 Game Đoán Chữ", "#ccffcc", bat_dau_game_popup),
        ("📥 Import Excel + Deploy Supabase", "#e8ffe8", import_excel_va_deploy_supabase),
        ("📦 HSK 2.0 / 3.0 Vocab ZIP Builder", "#d9eaff", hsk30_vocab_zip_builder),
        ("▶️ Đọc nội dung", "lightgreen", doc_popup),
        ("🎯 Đồng bộ ngôn ngữ dòng", "#ffe6cc", lambda: ep_toan_bo_dong_ve_lang()),
        ("⏸ Dừng đọc", "orange", dung_doc),
        ("🔁 Đọc lại", "lightblue", doc_lai_popup),
        ("🎧 Xuất MP3", "lightyellow", xuat_popup_mp3),
        ("🎧 Xuất M4A", "lightyellow", xuat_popup_m4a),
        # ("📄 Tạo phụ đề + MP3", "lightblue", run_xuat_srt_thread),
        # ("🖼️ Chọn background", "lightgray", chon_bg),
        # ("🎥 Tạo Video", "lightpink", run_xuat_video_thread),
        ("🎧 Xuất M4A MultiFiles", "#ffcccc", xuat_popup_m4a_multifiles),
        ("⬅ Trang trước", "#f0f0f0", prev_page),
        ("Trang tiếp ➡", "#f0f0f0", next_page),
    ]

    for idx, (label_text, bg_color, callback) in enumerate(action_buttons):
        row = idx // 2
        col = idx % 2
        tk.Button(
            button_grid,
            text=label_text,
            bg=bg_color,
            command=callback,
        ).grid(row=row, column=col, padx=3, pady=4, sticky="ew", ipady=2)

    button_grid.columnconfigure(0, weight=1)
    button_grid.columnconfigure(1, weight=1)

    tk.Label(
        option_frame,
        text="Lưu ý: nút 'Đồng bộ ngôn ngữ dòng' chỉ chỉnh ngôn ngữ từng dòng, không đổi engine.",
        bg="#f8fff8",
        fg="#666",
        wraplength=260,
        justify="left",
    ).pack(anchor="w", padx=10, pady=(0, 8))


    
    tk.Label(frame, text="", height=2).pack() # Thêm label trống để tránh bị cut cuối
    frame.update_idletasks()
    canvas.config(scrollregion=canvas.bbox("all"))

    def on_frame_configure(event):
        canvas.configure(scrollregion=canvas.bbox("all"))

    frame.bind("<Configure>", on_frame_configure)

    def on_close():
        global popup_lang_open
        popup_lang_open = False
        popup.destroy()

    popup.protocol("WM_DELETE_WINDOW", on_close)

#===============
#=====phát đa ngôn ngữ

def phat_da_ngon_ngu(danh_sach_cau, giong="Nam", toc_do="Bình thường", close_popup=None, engine=None):
    import threading, tempfile, os, time
    import pygame
    from gtts import gTTS
    from pydub import AudioSegment
    global dung_doc_ngay, dang_doc

    if engine is None:
        engine, configured_voice = _get_runtime_tts_selection()
        if giong == "Nam" and configured_voice:
            giong = configured_voice
    engine = engine or "gTTS"

    if 'dang_doc' not in globals():
        dang_doc = False
    if dang_doc:
        print("⛔ Đang đọc, không thể đọc mới.")
        return

    dung_doc_ngay = False
    dang_doc = True

    def run():
        global dung_doc_ngay, dang_doc
        count = 0

        for dong, lang in danh_sach_cau:
            if dung_doc_ngay:
                print("⛔ Dừng đọc ngay.")
                break

            try:
                dong_sach = lam_sach_van_ban(dong)
                if lang not in ["vi", "en", "ja", "zh"]:
                    lang = "vi"

                # Nếu hội thoại thì xen kẽ giọng nam/nữ
                if giong in {"Hội thoại 1 câu nam - 1 câu nữ", "Hội thoại 1 câu nữ - 1 câu nam"}:
                    first_voice, second_voice = _dialogue_voice_pair(giong)
                    voice = first_voice if count % 2 == 0 else second_voice
                    count += 1
                else:
                    voice = giong

                # ⚡ Ghi log giọng, sau này dùng engine khác có thể đổi file hoặc style
                print(f"📢 Đọc ({lang}) [{voice}]: {dong_sach}")

                slow = True if toc_do == "Chậm" else False

                file_mp3 = tempfile.mktemp(suffix=".mp3")
                tao_file_mp3(dong_sach, lang=lang, voice=voice, toc_do=toc_do, engine=engine, file_out=file_mp3)

                pygame.mixer.init()
                sound = pygame.mixer.Sound(file_mp3)
                channel = pygame.mixer.find_channel()
                if channel:
                    channel.play(sound)
                    while channel.get_busy():
                        if dung_doc_ngay:
                            print("⛔ Dừng đọc giữa dòng.")
                            channel.stop()
                            break
                        time.sleep(0.05)
                os.remove(file_mp3)

            except Exception as e:
                print(f"❌ Lỗi đọc: {e}")

        dang_doc = False
        print("✅ Kết thúc đọc nội dung.")
        if close_popup:
            try:
                close_popup()
            except Exception as e:
                print("Không thể đóng popup:", e)

    threading.Thread(target=run).start()
#=====

#======= Xuất mp3/wav giao diện chính


def xuat_file_mp3():
    from tkinter import messagebox

    print("Disabled: main-window audio export was removed. Use the language selection popup instead.")
    try:
        messagebox.showinfo(
            "Đã tắt xuất audio ở màn hình chính",
            "Vui lòng dùng popup chọn ngôn ngữ để xuất MP3/M4A.",
        )
    except Exception as e:
        print("Không thể hiển thị thông báo tắt xuất audio main UI:", e)
    return

    import os
    from tkinter import filedialog, messagebox
    import pygame

    def export_thread():
        print("🔄 Bắt đầu xuất file MP3...")
        btn_xuat_mp3.config(text="⏳ Đang chuẩn bị...", state="disabled")
        progress_var.set(0)
        progress_bar.pack()

        danh_sach = tach_de_thanh_danh_sach_da_ngon_ngu()
        if not danh_sach:
            messagebox.showwarning("Trống", "Không có nội dung để xuất.")
            btn_xuat_mp3.config(text="🎧 Xuất file MP3", state="normal")
            progress_bar.pack_forget()
            pygame.mixer.init()
            pygame.mixer.music.load(WARNING_SOUND)
            pygame.mixer.music.play()
            return

        # === Hỏi chọn định dạng ===
        res = messagebox.askquestion("Chọn định dạng", "Bạn muốn xuất file MP3 (Yes) hay WAV (No)?", icon="question")
        if res == "yes":
            ext = ".mp3"
            filetypes = [("MP3 files", "*.mp3")]
            fmt = "mp3"
        elif res == "no":
            ext = ".wav"
            filetypes = [("WAV files", "*.wav")]
            fmt = "wav"
        else:
            btn_xuat_mp3.config(text="🎧 Xuất file MP3", state="normal")
            progress_bar.pack_forget()
            return

        file_path = filedialog.asksaveasfilename(defaultextension=ext,
                                                 filetypes=filetypes,
                                                 title=f"Lưu file {fmt.upper()}")
        if not file_path:
            btn_xuat_mp3.config(text="🎧 Xuất file MP3", state="normal")
            progress_bar.pack_forget()
            return

        if not FFMPEG_PATH or not os.path.isfile(FFMPEG_PATH):
            messagebox.showerror("Lỗi", "Không tìm thấy ffmpeg trên máy và cũng không có bản đi kèm trong portable.")
            btn_xuat_mp3.config(text="🎧 Xuất file MP3", state="normal")
            progress_bar.pack_forget()
            pygame.mixer.init()
            pygame.mixer.music.load(WARNING_SOUND)
            pygame.mixer.music.play()
            return

        try:
            full_audio = AudioSegment.silent(duration=500)
            temp_dir = tempfile.gettempdir()
            tong_dong = len(danh_sach)

            for i, (dong, lang_raw) in enumerate(danh_sach):
                try:
                    cleaned = lam_sach_van_ban(dong)
                    toc_do = combo_toc_do.get()
                    slow = True if toc_do == "Chậm" else False

                    lang = doan_ngon_ngu_theo_ky_tu(cleaned)
                    if lang in ["zh-cn", "zh-tw", "zh-hk"]:
                        lang = "zh"
                    if lang not in ["vi", "en", "ja", "zh"]:
                        lang = "vi"

                    btn_xuat_mp3.config(text=f"⏳ Dòng {i + 1}/{tong_dong}...")

                    temp_mp3 = os.path.join(temp_dir, f"temp_{uuid.uuid4().hex}.mp3")
                    try:
                        # Use central tao_file_mp3 which prefers Google Cloud for vi and falls back
                        tao_file_mp3(cleaned, lang=lang, voice=None, toc_do=toc_do, engine=combo_engine.get(), file_out=temp_mp3)
                        segment = AudioSegment.from_mp3(temp_mp3)
                        full_audio += segment + AudioSegment.silent(duration=300)
                    except Exception as e:
                        msg = str(e)
                        if "No text to speak" in msg or "No text to send to TTS API" in msg:
                            full_audio += AudioSegment.silent(duration=300)
                        else:
                            raise
                    finally:
                        try:
                            if os.path.exists(temp_mp3):
                                os.remove(temp_mp3)
                        except Exception:
                            pass

                    progress = int((i + 1) / tong_dong * 100)
                    progress_var.set(progress)
                    print(f"✅ {i + 1}/{tong_dong}: ({lang}) {cleaned[:40]}...")

                except Exception as e:
                    print(f"❌ Lỗi dòng {i}: ({lang_raw}) {dong} → {e}")
                    continue

            # === Xuất file với định dạng đã chọn ===
            full_audio.export(file_path, format=fmt, bitrate="192k")

            pygame.mixer.init()
            pygame.mixer.music.load(SUCCESS_SOUND)
            pygame.mixer.music.play()
            gui_discord_thong_bao(f"🎙️ [TextToMp3] Đã xuất xong: {file_path}") #báo tới điện thoại discor

            print("✅ Xuất xong:", file_path)

            btn_xuat_mp3.config(text="✅ Đã xuất xong!")
            btn_xuat_mp3.after(3000, lambda: btn_xuat_mp3.config(text="🎧 Xuất file MP3", state="normal"))
            progress_bar.pack_forget()

            # ==== Thông báo popup mở file ====
            def open_file():
                try: open_path_cross_platform(file_path)
                except Exception as e: messagebox.showerror("Lỗi", f"Không mở được file:\n{e}")

            def open_folder():
                try: open_path_cross_platform(os.path.dirname(file_path))
                except Exception as e: messagebox.showerror("Lỗi", f"Không mở được thư mục:\n{e}")

            def play_now():
                top = tk.Toplevel()
                set_popup_icon(top)
                top.title("🔊 Đang phát: " + os.path.basename(file_path))
                top.geometry("300x120")
                is_playing = [True]
                def stop_play():
                    is_playing[0] = False
                    top.destroy()
                def run():
                    try:
                        pygame.mixer.init()
                        pygame.mixer.music.load(file_path)
                        pygame.mixer.music.play()
                        while pygame.mixer.music.get_busy() and is_playing[0]:
                            time.sleep(0.1)
                    except Exception as e:
                        print("Lỗi phát:", e)
                threading.Thread(target=run, daemon=True).start()
                tk.Button(top, text="Dừng phát", command=stop_play, fg="red", font=("Arial", 11)).pack(pady=16)

            popup = tk.Toplevel(root)
            set_popup_icon(popup)
            popup.title("Hoàn tất xuất MP3")
            popup.geometry("390x180+{}+{}".format(
                root.winfo_x() + root.winfo_width() // 2 - 195,
                root.winfo_y() + root.winfo_height() // 2 - 90
            ))
            popup.grab_set()
            popup.transient(root)
            tk.Label(popup, text=f"🎉 Đã xuất file {fmt.upper()}:\n" + file_path, font=("Arial", 11, "bold"), fg="green").pack(pady=13)
            frm = tk.Frame(popup)
            frm.pack(pady=3)
            tk.Button(frm, text="Mở file", width=10, command=lambda: [popup.destroy(), open_file()]).pack(side="left", padx=8)
            tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup.destroy(), open_folder()]).pack(side="left", padx=8)
            tk.Button(frm, text="Phát ngay", width=10, command=lambda: [popup.destroy(), play_now()]).pack(side="left", padx=8)
            tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=7)

        except Exception as e:
            print("❌ Xuất MP3 lỗi:", e)
            pygame.mixer.init()
            pygame.mixer.music.load(WARNING_SOUND)
            pygame.mixer.music.play()
            messagebox.showerror("Lỗi", f"Xuất {fmt.upper()} bị lỗi:\n{e}")
            btn_xuat_mp3.config(text="🎧 Xuất file MP3", state="normal")
            progress_bar.pack_forget()

    threading.Thread(target=export_thread, daemon=True).start()

#===ĐỌC LẠI
def doc_lai():
    if not noi_dung_cuoi:
        messagebox.showinfo("Chưa có nội dung", "Không có nội dung nào để đọc lại.")
        return
    txt_de.delete("1.0", tk.END)
    txt_de.insert(tk.END, noi_dung_cuoi)
    doc_noi_dung_de()

def xu_ly_doc_noi_dung():
    if che_do_doc.get() == "Tự động":
        doc_noi_dung_de()
    else:
        danh_sach = tach_de_thanh_danh_sach_da_ngon_ngu()
        engine, voice = _get_runtime_tts_selection()
        phat_da_ngon_ngu(danh_sach, giong=voice, engine=engine)
#================


# Hàm import tài liệu
def import_tai_lieu():
    from tkinter import filedialog
    import docx
    import PyPDF2
    import openpyxl

    file_path = filedialog.askopenfilename(
        title="Chọn file tài liệu",
        filetypes=[("All files", "*.*"), ("Word", "*.docx"), ("Excel", "*.xlsx"), ("PDF", "*.pdf"), ("PowerPoint", "*.pptx"), ("Text", "*.txt")]

    )
    if not file_path:
        return

    noi_dung = ""
    try:
        ext = os.path.splitext(file_path)[-1].lower()

        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                noi_dung = f.read()

        elif ext == ".docx":
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                noi_dung += para.text + "\n"

        elif ext == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    noi_dung += page.extract_text() + "\n"

        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(file_path)
            sheet = wb.active
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        noi_dung += str(cell) + " "
                noi_dung += "\n"

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                tk.messagebox.showerror("Lỗi", "Chưa cài đặt python-pptx. Vui lòng cài đặt: pip install python-pptx")
                return
            
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides, 1):
                noi_dung += f"[Slide {slide_idx}]\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        for line in shape.text.split('\n'):
                            if line.strip():
                                noi_dung += line.strip() + "\n"
                noi_dung += "\n"

        else:
            tk.messagebox.showwarning("Không hỗ trợ", f"Định dạng file {ext} chưa được hỗ trợ.")

    except Exception as e:
        tk.messagebox.showerror("Lỗi", f"Không đọc được file:\n{e}")
        return

    if noi_dung.strip():
        txt_de.delete("1.0", tk.END)
        txt_de.insert(tk.END, noi_dung.strip())
    else:
        tk.messagebox.showinfo("Trống", "Không có nội dung hợp lệ trong tài liệu.")


#hàm chuyển đổi định dạng âm thanh

def convert_mp3_wav():
    import os
    from tkinter import filedialog, messagebox
    import pygame

    file_path = filedialog.askopenfilename(title="Chọn file MP3 hoặc WAV",
                                           filetypes=[("Audio files", "*.mp3 *.wav")])
    if not file_path:
        return

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [".mp3", ".wav"]:
        messagebox.showerror("Lỗi", "Chỉ hỗ trợ MP3 và WAV!")
        if os.path.exists(WARNING_SOUND):
            pygame.mixer.init()
            pygame.mixer.music.load(WARNING_SOUND)
            pygame.mixer.music.play()
        return

    # Xác định định dạng đích
    if ext == ".mp3":
        fmt_out = "wav"
        new_ext = ".wav"
    else:
        fmt_out = "mp3"
        new_ext = ".mp3"

    save_path = filedialog.asksaveasfilename(defaultextension=new_ext,
                                             filetypes=[(f"{fmt_out.upper()} files", f"*{new_ext}")],
                                             title=f"Lưu file {fmt_out.upper()}")
    if not save_path:
        return

    def thread_convert():
        try:
            progress_var.set(0)
            progress_bar.pack()
            progress_bar.update()

            audio = AudioSegment.from_file(file_path)
            # Mô phỏng tiến trình (giả lập % để nhìn thấy)
            for i in range(1, 101):
                progress_var.set(i)
                progress_bar.update()
                time.sleep(0.01)

            audio.export(save_path, format=fmt_out, bitrate="192k")

            pygame.mixer.init()
            pygame.mixer.music.load(SUCCESS_SOUND)
            pygame.mixer.music.play()

            progress_bar.pack_forget()

            # === Popup hoàn tất ===
            popup = tk.Toplevel(root)
            set_popup_icon(popup)
            popup.title("Hoàn tất chuyển đổi")
            popup.geometry("390x180+{}+{}".format(
                root.winfo_x() + root.winfo_width() // 2 - 195,
                root.winfo_y() + root.winfo_height() // 2 - 90
            ))
            popup.grab_set()
            popup.transient(root)
            tk.Label(popup, text=f"🎉 Đã chuyển xong:\n{save_path}", font=("Arial", 11, "bold"), fg="green").pack(pady=13)
            frm = tk.Frame(popup)
            frm.pack(pady=3)
            tk.Button(frm, text="Mở file", width=10, command=lambda: [popup.destroy(), open_path_cross_platform(save_path)]).pack(side="left", padx=8)
            tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup.destroy(), open_path_cross_platform(os.path.dirname(save_path))]).pack(side="left", padx=8)
            tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=7)

        except Exception as e:
            print("❌ Lỗi convert:", e)
            progress_bar.pack_forget()
            pygame.mixer.init()
            pygame.mixer.music.load(WARNING_SOUND)
            pygame.mixer.music.play()
            messagebox.showerror("Lỗi", f"Chuyển đổi bị lỗi:\n{e}")
            progress_var.set(0)

    threading.Thread(target=thread_convert, daemon=True).start()

#=================
# cắt âm thanh

def cutter_sound():
    import os
    from tkinter import filedialog, messagebox
    import pygame
    import numpy as np

    win = tk.Toplevel()
    win.title("Waveform Cutter")
    win.geometry("650x400")
    set_popup_icon(win)
    win.lift()
    win.grab_set()

    cutter_audio_filetypes = [("Audio files", "*.mp3 *.wav *.m4a *.aac"), ("All files", "*.*")]
    cutter_save_filetypes = [
        ("WAV files", "*.wav"),
        ("MP3 files", "*.mp3"),
        ("M4A files", "*.m4a"),
        ("AAC files", "*.aac"),
    ]
    supported_cutter_exts = {".mp3", ".wav", ".m4a", ".aac"}

    def copy_m4a_segment_without_reencode(source_path, save_path, start_ms, end_ms):
        if not FFMPEG_PATH or not os.path.isfile(FFMPEG_PATH):
            raise ValueError("Không tìm thấy ffmpeg để cắt M4A giữ nguyên chất lượng.")

        start_sec = start_ms / 1000
        duration_sec = (end_ms - start_ms) / 1000
        startupinfo = None
        if os.name == "nt":
            startupinfo = subprocess.STARTUPINFO()
            startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW

        cmd = [
            FFMPEG_PATH, "-y",
            "-ss", f"{start_sec:.3f}",
            "-i", source_path,
            "-t", f"{duration_sec:.3f}",
            "-vn",
            "-c", "copy",
            "-movflags", "+faststart",
            save_path,
        ]
        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            startupinfo=startupinfo,
        )
        if result.returncode != 0:
            raise Exception(result.stderr or "FFmpeg không cắt được M4A giữ nguyên chất lượng.")

    def export_cut_segment(segment, save_path, source_path=None, start_ms=None, end_ms=None):
        ext = os.path.splitext(save_path)[1].lower()
        source_ext = os.path.splitext(source_path or "")[1].lower()
        if ext == ".m4a" and source_ext == ".m4a" and start_ms is not None and end_ms is not None:
            copy_m4a_segment_without_reencode(source_path, save_path, start_ms, end_ms)
        elif ext == ".m4a":
            segment.export(save_path, format="ipod", codec="aac", bitrate="192k")
        elif ext == ".aac":
            segment.export(save_path, format="adts", codec="aac", bitrate="192k")
        elif ext == ".mp3":
            segment.export(save_path, format="mp3", bitrate="192k")
        elif ext == ".wav":
            segment.export(save_path, format="wav")
        else:
            raise ValueError("Chỉ hỗ trợ lưu WAV, MP3, M4A hoặc AAC.")

    file_path = filedialog.askopenfilename(title="Chọn file MP3, WAV, M4A hoặc AAC",
                                           filetypes=cutter_audio_filetypes,
                                           parent=win)  # ✅ Thêm parent
    if not file_path:
        win.destroy()
        return

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in supported_cutter_exts:
        messagebox.showerror("Lỗi", "Chỉ hỗ trợ MP3, WAV, M4A và AAC!", parent=win)
        win.destroy()
        return

    audio = AudioSegment.from_file(file_path)
    data = np.array(audio.get_array_of_samples())
    if audio.channels == 2:
        data = data.reshape((-1, 2))
        data = data.sum(axis=1) / 2

    block_size = max(int(len(data) / 500), 1)
    data_r = [np.mean(data[i:i+block_size]) for i in range(0, len(data), block_size)]

    canvas = tk.Canvas(win, bg="white", width=600, height=150)
    canvas.pack(pady=10)

    mid = 75
    scale = max(max(data_r), abs(min(data_r))) or 1

    prev_x = 0
    prev_y = mid
    for i, val in enumerate(data_r):
        x = int(i * (600 / len(data_r)))
        y = int(mid - (val / scale * 70))
        canvas.create_line(prev_x, prev_y, x, y, fill="blue")
        prev_x = x
        prev_y = y

    frm_entry = tk.Frame(win)
    frm_entry.pack(pady=5)

    tk.Label(frm_entry, text="Start (s):").grid(row=0, column=0, padx=5)
    entry_start = tk.Entry(frm_entry, width=10)
    entry_start.insert(0, "0")
    entry_start.grid(row=0, column=1)

    tk.Label(frm_entry, text="End (s):").grid(row=0, column=2, padx=5)
    entry_end = tk.Entry(frm_entry, width=10)
    entry_end.insert(0, str(len(audio) // 1000))
    entry_end.grid(row=0, column=3)

    progress_play = ttk.Progressbar(win, orient="horizontal", length=600, mode="determinate")
    progress_play.pack(pady=5)

    label_time = tk.Label(win, text="0:00 / 0:00")
    label_time.pack()

    is_playing = [False]
    is_paused = [False]
    start_time_ref = [0]
    duration_ref = [0]

    def update_progress():
        if is_playing[0] and not is_paused[0]:
            elapsed = time.time() - start_time_ref[0]
            progress = int(min((elapsed / duration_ref[0]) * 100, 100))
            progress_play["value"] = progress
            m = int(elapsed // 60)
            s = int(elapsed % 60)
            label_time.config(text=f"{m}:{s:02} / {int(duration_ref[0] // 60)}:{int(duration_ref[0] % 60):02}")
            if elapsed < duration_ref[0]:
                win.after(100, update_progress)
            else:
                progress_play["value"] = 0
                label_time.config(text=f"0:00 / {int(duration_ref[0] // 60)}:{int(duration_ref[0] % 60):02}")
                is_playing[0] = False
        else:
            progress_play["value"] = 0

    def play_audio():
        try:
            stop_audio()
            start_sec = int(entry_start.get())
            end_sec = int(entry_end.get())
            start_ms = start_sec * 1000
            end_ms = end_sec * 1000
            seg = audio[start_ms:end_ms]
            temp_file = os.path.join(tempfile.gettempdir(), f"temp_play_{uuid.uuid4().hex}.wav")
            seg.export(temp_file, format="wav")

            pygame.mixer.quit()
            pygame.mixer.init()
            pygame.mixer.music.load(temp_file)
            pygame.mixer.music.set_volume(1.0)
            pygame.mixer.music.play()

            is_playing[0] = True
            is_paused[0] = False
            duration_ref[0] = (end_ms - start_ms) / 1000
            start_time_ref[0] = time.time()

            label_time.config(text=f"0:00 / {int(duration_ref[0] // 60)}:{int(duration_ref[0] % 60):02}")
            win.after(100, update_progress)
        except Exception as e:
            print("Lỗi play:", e)
            messagebox.showerror("Lỗi", f"Lỗi play:\n{e}")

    def pause_audio():
        if is_playing[0] and not is_paused[0]:
            pygame.mixer.music.pause()
            is_paused[0] = True
        elif is_playing[0] and is_paused[0]:
            pygame.mixer.music.unpause()
            is_paused[0] = False
            start_time_ref[0] = time.time() - (progress_play["value"] / 100) * duration_ref[0]
            win.after(100, update_progress)

    def stop_audio():
        is_playing[0] = False
        is_paused[0] = False
        pygame.mixer.music.stop()
        progress_play["value"] = 0
        label_time.config(text="0:00 / 0:00")

    def cut_and_save():
        try:
            start_sec = int(entry_start.get())
            end_sec = int(entry_end.get())
            start_ms = start_sec * 1000
            end_ms = end_sec * 1000
            if start_ms >= end_ms:
                messagebox.showerror("Lỗi", "Start phải nhỏ hơn End.")
                return

            save_path = filedialog.asksaveasfilename(defaultextension=".wav",
                                                     filetypes=cutter_save_filetypes,
                                                     title="Lưu file cắt",
                                                     parent=win)  # ✅ Thêm parent

            if not save_path:
                return

            def thread_cut():
                try:
                    progress_var.set(0)
                    progress_bar.pack()
                    progress_bar.update()

                    segment = audio[start_ms:end_ms]
                    for i in range(1, 101):
                        progress_var.set(i)
                        progress_bar.update()
                        time.sleep(0.01)

                    export_cut_segment(segment, save_path, file_path, start_ms, end_ms)

                    pygame.mixer.quit()
                    pygame.mixer.init()
                    pygame.mixer.music.load(SUCCESS_SOUND)
                    pygame.mixer.music.play()

                    progress_bar.pack_forget()

                    popup = tk.Toplevel(root)
                    set_popup_icon(popup)
                    popup.title("Hoàn tất cắt")
                    popup.geometry("390x180+{}+{}".format(
                        root.winfo_x() + root.winfo_width() // 2 - 195,
                        root.winfo_y() + root.winfo_height() // 2 - 90
                    ))
                    popup.grab_set()
                    popup.transient(root)
                    tk.Label(popup, text=f"🎉 Đã cắt xong:\n{save_path}", font=("Arial", 11, "bold"), fg="green").pack(pady=13)
                    frm = tk.Frame(popup)
                    frm.pack(pady=3)
                    tk.Button(frm, text="Mở file", width=10, command=lambda: [popup.destroy(), open_path_cross_platform(save_path)]).pack(side="left", padx=8)
                    tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup.destroy(), open_path_cross_platform(os.path.dirname(save_path))]).pack(side="left", padx=8)
                    tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=7)

                except Exception as e:
                    print("❌ Lỗi cắt:", e)
                    progress_bar.pack_forget()
                    pygame.mixer.quit()
                    pygame.mixer.init()
                    pygame.mixer.music.load(WARNING_SOUND)
                    pygame.mixer.music.play()
                    messagebox.showerror("Lỗi", f"Cắt bị lỗi:\n{e}")
                    progress_var.set(0)

            threading.Thread(target=thread_cut, daemon=True).start()

        except Exception as e:
            print("Lỗi cắt & lưu:", e)
            messagebox.showerror("Lỗi", f"Lỗi cắt & lưu:\n{e}")

    frm_btn = tk.Frame(win)
    frm_btn.pack(pady=5)

    tk.Button(frm_btn, text="▶ Play", width=8, command=play_audio).pack(side="left", padx=5)
    tk.Button(frm_btn, text="⏸ Pause", width=8, command=pause_audio).pack(side="left", padx=5)
    tk.Button(frm_btn, text="⏹ Stop", width=8, command=stop_audio).pack(side="left", padx=5)
    tk.Button(frm_btn, text="✂ Cut & Save", width=12, command=cut_and_save).pack(side="left", padx=5)

    tk.Button(win, text="Đóng", command=lambda: [stop_audio(), win.destroy()]).pack(pady=5)



#===
#====Dán âm thanh
def joiner_sound():
    import os
    from tkinter import filedialog, messagebox
    import pygame

    files = filedialog.askopenfilenames(title="Chọn nhiều file âm thanh để nối",
                                        filetypes=[("Audio files", "*.mp3 *.wav")])
    if not files or len(files) < 2:
        messagebox.showwarning("Ít file", "Cần chọn ít nhất 2 file để nối!")
        return

    first_ext = os.path.splitext(files[0])[1].lower()
    fmt_out = first_ext.replace(".", "")

    save_path = filedialog.asksaveasfilename(defaultextension=first_ext,
                                             filetypes=[("Audio files", "*.mp3 *.wav")],
                                             title="Lưu file nối")
    if not save_path:
        return

    def thread_join():
        try:
            progress_var.set(0)
            progress_bar.pack()
            progress_bar.update()

            combined = AudioSegment.empty()
            total = len(files)

            for i, f in enumerate(files):
                audio = AudioSegment.from_file(f)
                combined += audio

                progress = int(((i + 1) / total) * 100)
                progress_var.set(progress)
                progress_bar.update()
                time.sleep(0.05)

            combined.export(save_path, format=fmt_out, bitrate="192k")

            pygame.mixer.init()
            pygame.mixer.music.load(SUCCESS_SOUND)
            pygame.mixer.music.play()

            progress_bar.pack_forget()

            popup = tk.Toplevel(root)
            set_popup_icon(popup)
            popup.title("Hoàn tất nối")
            popup.geometry("390x180+{}+{}".format(
                root.winfo_x() + root.winfo_width() // 2 - 195,
                root.winfo_y() + root.winfo_height() // 2 - 90
            ))
            popup.grab_set()
            popup.transient(root)
            tk.Label(popup, text=f"🎉 Đã nối xong:\n{save_path}", font=("Arial", 11, "bold"), fg="green").pack(pady=13)
            frm = tk.Frame(popup)
            frm.pack(pady=3)
            tk.Button(frm, text="Mở file", width=10, command=lambda: [popup.destroy(), open_path_cross_platform(save_path)]).pack(side="left", padx=8)
            tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup.destroy(), open_path_cross_platform(os.path.dirname(save_path))]).pack(side="left", padx=8)
            tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=7)

        except Exception as e:
            print("❌ Lỗi nối:", e)
            progress_bar.pack_forget()
            pygame.mixer.init()
            pygame.mixer.music.load(WARNING_SOUND)
            pygame.mixer.music.play()
            messagebox.showerror("Lỗi", f"Nối bị lỗi:\n{e}")
            progress_var.set(0)

    threading.Thread(target=thread_join, daemon=True).start()

#====================
# ===== DEAD CODE: SUBTITLE VIDEO EXPORT =====
# DEAD CODE - remove later. Feature disabled in audio-tool version.

def mo_popup_tao_video_phu_de():
    print("Disabled: subtitle video export removed in audio-tool version.")
    try:
        messagebox.showinfo("Đã tắt", "Tính năng tạo video phụ đề đã được tắt.")
    except Exception:
        pass
    return

    popup = tk.Toplevel(root)
    set_popup_icon(popup)
    popup.title("Tạo video có hình nền (chuẩn YouTube)")
    popup.geometry("520x420")
    popup.grab_set()
    popup.transient(root)

    tk.Label(popup, text="Chọn file MP3:", font=("Arial", 11)).pack(pady=(10, 0))
    entry_mp3 = tk.Entry(popup, width=60)
    entry_mp3.pack(pady=2)

    def browse_mp3():
        path = filedialog.askopenfilename(filetypes=[("MP3 files", "*.mp3")])
        if path:
            entry_mp3.delete(0, tk.END)
            entry_mp3.insert(0, path)

    tk.Button(popup, text="Browse", command=browse_mp3).pack()

    tk.Label(popup, text="Chọn hình nền (JPG/PNG):", font=("Arial", 11)).pack(pady=(10, 0))
    entry_bg = tk.Entry(popup, width=60)
    entry_bg.pack(pady=2)

    def browse_bg():
        path = filedialog.askopenfilename(filetypes=[("Image files", "*.png;*.jpg;*.jpeg")])
        if path:
            entry_bg.delete(0, tk.END)
            entry_bg.insert(0, path)

    tk.Button(popup, text="Browse", command=browse_bg).pack()

    tk.Label(popup, text="Độ phân giải (ví dụ: 1280x720):", font=("Arial", 11)).pack(pady=(10, 0))
    entry_resolution = tk.Entry(popup, width=20)
    entry_resolution.insert(0, "1280x720")
    entry_resolution.pack(pady=2)

    def xuat_video():
        mp3_path = entry_mp3.get().strip()
        bg_path = entry_bg.get().strip()
        resolution = entry_resolution.get().strip()

        if not mp3_path or not bg_path or not resolution:
            messagebox.showwarning("Thiếu thông tin", "Vui lòng chọn MP3, hình nền và độ phân giải.", parent=popup)
            return

        mp3_name = os.path.splitext(os.path.basename(mp3_path))[0]
        suggest_name = mp3_name + ".mp4"

        output_path = filedialog.asksaveasfilename(defaultextension=".mp4", initialfile=suggest_name,
                                                   filetypes=[("MP4 files", "*.mp4")], title="Lưu video")
        if not output_path:
            return

        progress_popup = tk.Toplevel(popup)
        set_popup_icon(progress_popup)
        progress_popup.title("Đang xuất video")
        progress_popup.geometry("420x160")
        progress_popup.grab_set()
        progress_popup.transient(popup)

        label_status = tk.Label(progress_popup, text="Đang xử lý, vui lòng chờ...", font=("Arial", 11))
        label_status.pack(pady=5)
        progress = ttk.Progressbar(progress_popup, orient="horizontal", length=300, mode="determinate")
        progress.pack(pady=5)

        btn_cancel = tk.Button(progress_popup, text="❌ Huỷ xuất video", fg="red")
        btn_cancel.pack(pady=5)

        def run_export():
            nonlocal ffmpeg_process
            try:
                mp3_fixed = mp3_path.replace("\\", "/")
                bg_fixed = bg_path.replace("\\", "/")
                output_fixed = output_path.replace("\\", "/")

                cmd = [
                    FFMPEG_PATH, "-y",
                    "-loop", "1",
                    "-i", bg_fixed,
                    "-i", mp3_fixed,
                    "-vf", f"scale={resolution},format=yuv420p",
                    "-c:v", "libx264",
                    "-c:a", "aac",
                    "-shortest",
                    output_fixed
                ]

                startupinfo = subprocess.STARTUPINFO()
                startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW

                ffmpeg_process = subprocess.Popen(
                    cmd,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    text=True,
                    encoding="utf-8",         # ✅ Fix chính ở đây
                    errors="ignore",          # ✅ Bỏ qua ký tự không hợp lệ
                    startupinfo=startupinfo
                )

                #ffmpeg_process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, startupinfo=startupinfo)
                stdout, stderr = ffmpeg_process.communicate()

                if ffmpeg_process.returncode != 0:
                    print("⚠️ FFmpeg stderr:\n", stderr)
                    raise Exception(f"❌ Lỗi ffmpeg: {stderr}")

                progress["value"] = 100
                label_status.config(text="✅ Hoàn tất!")

                if not os.path.exists(output_path):
                    progress_popup.destroy()
                    messagebox.showerror("Lỗi", "Không tìm thấy file xuất ra!")
                    return

                def open_file_fix():
                    if os.path.exists(output_path):
                        open_path_cross_platform(os.path.normpath(output_path))
                    else:
                        messagebox.showerror("Lỗi", "Không tìm thấy file!")

                def open_folder_fix():
                    folder = os.path.dirname(output_path)
                    if os.path.exists(folder):
                        open_path_cross_platform(folder)
                    else:
                        messagebox.showerror("Lỗi", "Không tìm thấy thư mục!")

                def upload_yt():
                    try:
                        popup_google_login(output_path)
                    except Exception as e:
                        messagebox.showerror("Lỗi upload", f"{e}")

                popup_done = tk.Toplevel(popup)
                set_popup_icon(popup_done)
                pygame.mixer.init()
                pygame.mixer.music.load(SUCCESS_SOUND)
                pygame.mixer.music.play()                
                popup_done.title("Hoàn tất")
                popup_done.geometry("450x160")
                popup_done.grab_set()
                popup_done.transient(popup)
                gui_discord_thong_bao(f"🎙️ [TextToMp3] 🎉 Đã xuất video: {output_path}")
                tk.Label(popup_done, text="🎉 Đã xuất video:\n" + output_path, font=("Arial", 11, "bold"), fg="green").pack(pady=13)
                frm = tk.Frame(popup_done)
                frm.pack(pady=3)
                tk.Button(frm, text="Mở file", width=10, command=lambda: [popup_done.destroy(), open_file_fix()]).pack(side="left", padx=6)
                tk.Button(frm, text="Mở thư mục", width=12, command=lambda: [popup_done.destroy(), open_folder_fix()]).pack(side="left", padx=6)
                tk.Button(frm, text="Đăng YouTube", width=12, command=lambda: [popup_done.destroy(), upload_yt()]).pack(side="left", padx=6)
                tk.Button(popup_done, text="Đóng", command=popup_done.destroy).pack(pady=8)

                progress_popup.destroy()

            except Exception as e:
                progress_popup.destroy()

                pygame.mixer.init()
                pygame.mixer.music.load(WARNING_SOUND)
                pygame.mixer.music.play()                
                messagebox.showerror("Lỗi", f"Xuất video bị lỗi:\n{e}")
                gui_discord_thong_bao(f"🎙️ [TextToMp3] 🎉 Xuất video bị lỗi:\n{e}")

        def cancel_export():
            try:
                if ffmpeg_process and ffmpeg_process.poll() is None:
                    ffmpeg_process.terminate()
                    label_status.config(text="⚠️ Đã huỷ!")
                    progress_popup.destroy()
                    messagebox.showinfo("Huỷ", "Đã huỷ xuất video.")
            except Exception as ex:
                messagebox.showerror("Lỗi", f"Không thể huỷ: {ex}")

        btn_cancel.config(command=cancel_export)

        ffmpeg_process = None
        threading.Thread(target=run_export, daemon=True).start()

    tk.Button(popup, text="🎬 Xuất Video", bg="green", fg="white", command=xuat_video).pack(pady=15)

# ===== END DEAD CODE: SUBTITLE VIDEO EXPORT =====


#=======================
#gọi game từ ngoài
def goi_popup_game_tu_ben_ngoai():
    # DEAD CODE - remove later. Entry point game đã bị ẩn khỏi UI.
    print("Disabled: Game Đoán Chữ removed in audio-tool version.")
    try:
        messagebox.showinfo("Đã tắt", "Game Đoán Chữ đã được tắt.")
    except Exception:
        pass
    return

    # Gọi popup chọn ngôn ngữ
    try:
        mo_popup_chon_lang(mo_tu_ben_ngoai=True)
    except Exception as e:
        print("❌ Lỗi khi gọi popup chơi game:", e)
    
    # Sau khi mở popup xong, tự động gọi nút Bắt đầu Game sau 500ms
    def mo_game_tiep():
        try:
            for w in tk._default_root.winfo_children():
                if isinstance(w, tk.Toplevel) and "Chọn ngôn ngữ" in str(w.title()):
                    for child in w.winfo_children():
                        if isinstance(child, tk.Button) and "Chơi Game Đoán Chữ" in child.cget("text"):
                            child.invoke()  # ấn nút "Chơi Game Đoán Chữ" trong popup chọn lang
                            return
        except Exception as e:
            print("⚠ Lỗi gọi game từ ngoài:", e)

    # Trì hoãn chút để popup ngôn ngữ kịp mở xong
    tk._default_root.after(500, mo_game_tiep)
        
#tạo nút - KHUNG HIỂN THỊ NỘI DUNG
def create_frame_noi_dung(parent):
    frame_noi_dung = tk.LabelFrame(parent, text="📋 Phần mềm xuất mp3 hội thoại đa ngôn ngữ - Máy Học Tập", font=("Arial", 11, "bold"), bg="#f8fff8", fg="green")
    frame_noi_dung.place(x=10, y=42, width=1360, height=700)

    # Vùng soạn đề/to ra đề + thanh cuộn
    global txt_de
    frame_text = tk.Frame(frame_noi_dung)
    frame_text.place(x=10, y=8, width=1000, height=540)

    scrollbar = tk.Scrollbar(frame_text)
    scrollbar.pack(side="right", fill="y")

    txt_de = tk.Text(frame_text, font=("Arial", 16), wrap="word", yscrollcommand=scrollbar.set)
    txt_de.pack(side="left", fill="both", expand=True)
    attach_mouse_text_menu(txt_de)

    scrollbar.config(command=txt_de.yview)

    # Khung hỏi GitHub Models/Gemini sát mép dưới trái (dưới txt_de)
    frame_hoi_gpt = tk.LabelFrame(frame_noi_dung, text="🧠 Hỏi GitHub Models hoặc Gemini", font=("Arial", 10, "bold"), bg="#f0fff0", fg="darkgreen")
    frame_hoi_gpt.place(x=10, y=560, width=620, height=105)

    entry_cau_hoi = tk.Entry(frame_hoi_gpt, font=("Arial", 9))
    entry_cau_hoi.place(x=8, y=8, width=395, height=28)
    attach_mouse_text_menu(entry_cau_hoi)
    
    tk.Button(frame_hoi_gpt, text="Hỏi GitHub", bg="#ffe6e6", fg="red", font=("Arial", 10, "bold"),
              command=gui_hoi_github_models).place(x=410, y=8, width=65, height=28)
    tk.Button(frame_hoi_gpt, text="Hỏi Mẹ", bg="#e6ffe6", fg="green", font=("Arial", 10, "bold"),
              command=gui_hoi_gemini).place(x=480, y=8, width=65, height=28)
    tk.Button(frame_hoi_gpt, text="🎙Nói", font=("Arial", 8),
              command=lambda: nhap_giong_noi_advanced(entry_cau_hoi, root=root)).place(x=550, y=8, width=50, height=28)

    lbl_mic_effect = tk.Label(frame_hoi_gpt, textvariable=mic_effect_var, font=("Arial", 14), fg="green")
    lbl_mic_effect.place(x=575, y=42)


    # Khung đọc đề sát mép phải (xuất MP3/M4A nằm trong popup chọn ngôn ngữ)
    frame_doc = tk.LabelFrame(frame_noi_dung, text="🎧 Đọc Nội Dung & Chọn Ngôn Ngữ", font=("Arial", 8, "bold"), bg="#f8fff8")
    frame_doc.place(x=1030, y=8, width=320, height=660)

    btn_doc = tk.Button(frame_doc, text="Đọc / Chọn ngôn ngữ", font=("Arial", 14, "bold"),
                        command=lambda: doc_noi_dung_de() if che_do_doc.get() == "Tự động" else mo_popup_chon_lang(), bg="lightyellow")
    btn_doc.place(x=35, y=10, width=250, height=34)
    btn_tam_dung = tk.Button(frame_doc, text="⏸ Dừng", font=("Arial", 8), command=toggle_tam_dung, bg="lightyellow")
    btn_tam_dung.place(x=55, y=58, width=210, height=30)
    btn_doc_lai = tk.Button(frame_doc, text="🔁 Đọc lại", font=("Arial", 8), command=doc_lai, bg="lightyellow")
    btn_doc_lai.place(x=55, y=96, width=210, height=30)
    # Disabled: xuất MP3/WAV trực tiếp từ main UI.
    # Giữ object để các biến global cũ không bị vỡ, nhưng không place lên UI.
    # Xuất MP3/M4A chính thức nằm trong popup chọn ngôn ngữ.
    btn_xuat_mp3 = tk.Button(frame_doc, text="🎧 Xuất file MP3", font=("Arial", 8), command=xuat_file_mp3, bg="lightblue", state="disabled")
    # btn_xuat_mp3.place(x=55, y=134, width=210, height=30)
    
    


    progress_var = tk.DoubleVar()
    progress_bar = ttk.Progressbar(frame_doc, variable=progress_var, maximum=100, length=280, mode='determinate')
    progress_bar.place(x=20, y=174, width=280)
    progress_bar.place_forget()

   # Chế độ đọc (dòng trên)
    frame_che_do = tk.Frame(frame_doc, bg="#f8fff8")
    frame_che_do.place(x=35, y=180, width=260, height=30)

    che_do_doc = tk.StringVar(value="Chọn tay")
    tk.Label(frame_che_do, text="Chế độ:").pack(side=tk.LEFT, padx=(0,2))
    tk.Radiobutton(frame_che_do, text="Tự động", variable=che_do_doc, value="Tự động", font=("Arial", 9)).pack(side=tk.LEFT)
    tk.Radiobutton(frame_che_do, text="Chọn tay", variable=che_do_doc, value="Chọn tay", font=("Arial", 9)).pack(side=tk.LEFT)

    # Hai combobox nằm ngang nhau (dòng dưới)
    combo_ngon_ngu = ttk.Combobox(frame_doc, values=["Việt", "Anh", "Nhật", "Trung"], font=("Arial", 10))
    combo_ngon_ngu.set("")
    combo_ngon_ngu.place(x=55, y=216, width=95)

    combo_toc_do = ttk.Combobox(frame_doc, values=["Chậm", "Bình thường"], font=("Arial", 10))
    combo_toc_do.set("Bình thường")
    combo_toc_do.place(x=160, y=216, width=105)
    
    # Disabled: Game Đoán Chữ removed in audio-tool version.
    # btn_game = tk.Button(frame_doc, text="🎮 Game Đoán Chữ", font=("Arial", 11, "bold"),
    #                  bg="yellow", fg="red", command=goi_popup_game_tu_ben_ngoai)
    # btn_game.place(x=55, y=258, width=210, height=35)

    # Nút import tài liệu
    tk.Button(frame_doc, text="📥 Import tài liệu", font=("Arial", 9, "bold"),
              command=import_tai_lieu, bg="lightgreen").place(x=55, y=448, width=210, height=30)

    # Nút Convert MP3/WAV
    tk.Button(frame_doc, text="Convert Mp3 🎹 Wav", font=("Arial", 9, "bold"),
              command=convert_mp3_wav, bg="#f0f0f0").place(x=55, y=306, width=210, height=30)
    
    tk.Button(frame_doc, text="✂️ Cutter Sound", font=("Arial", 9, "bold"),
          command=cutter_sound, bg="#f0f0f0").place(x=55, y=342, width=210, height=30)

    tk.Button(frame_doc, text="➕ Joiner Sound", font=("Arial", 9, "bold"),
          command=joiner_sound, bg="#f0f0f0").place(x=55, y=378, width=210, height=30)
    # Disabled: video/subtitle export removed in audio-tool version.
    # tk.Button(frame_doc, text="🎞 Tạo Video Phụ Đề", font=("Arial", 9, "bold"),
    #       command=mo_popup_tao_video_phu_de, bg="#f0f0f0").place(x=55, y=414, width=210, height=30)

    
    



        # Trả về các biến cần dùng ở ngoài hàm
    return (frame_noi_dung, frame_text, frame_hoi_gpt, frame_doc, txt_de, entry_cau_hoi,
            btn_doc, btn_tam_dung, btn_doc_lai, btn_xuat_mp3, progress_var, progress_bar,
            che_do_doc, combo_ngon_ngu, combo_toc_do)

#====================================================================
#=====Tải ảnh tự động cho game theo cột Nghĩa TV
# ===== DEAD CODE: GAME IMAGE DOWNLOADER =====
# DEAD CODE - remove later. Game/image feature disabled in audio-tool version.
def tai_anh_con_thieu_game():
    print("Disabled: game image downloader removed in audio-tool version.")
    try:
        messagebox.showinfo("Đã tắt", "Tính năng tải ảnh game đã được tắt.")
    except Exception:
        pass
    return

    from urllib.parse import quote
    import requests
    from bs4 import BeautifulSoup
    from openpyxl import load_workbook
    from tkinter import messagebox
    import os, time

    try:
        wb = load_workbook(EXCEL_GAME_PATH)
        so_tai_duoc = 0

        SHEETS = ["Ja", "Cn", "En"]

        def get_first_image(keyword):
            q = quote(keyword)
            url = f"https://www.google.com/search?tbm=isch&q={q}"
            try:
                res = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=10)
                soup = BeautifulSoup(res.text, "html.parser")
                img_tags = soup.find_all("img")
                for img in img_tags[1:]:
                    src = img.get("src", "")
                    if src.startswith("http"):
                        return src
            except Exception:
                return None
            return None

        for sheet_name in SHEETS:
            if sheet_name not in wb.sheetnames:
                continue

            ws = wb[sheet_name]
            folder_lang = os.path.join(IMAGE_FOLDER, sheet_name)
            os.makedirs(folder_lang, exist_ok=True)

            header_row = [cell.value for cell in ws[1]]
            if "Nghĩa TV" not in header_row:
                continue

            col_idx_keyword = header_row.index("Nghĩa TV") + 1
            if "images" not in header_row:
                ws.cell(row=1, column=len(header_row) + 1).value = "images"
                col_idx_image = len(header_row) + 1
            else:
                col_idx_image = header_row.index("images") + 1

            # Dọn lại dữ liệu cột ảnh
            total_rows = ws.max_row - 1
            current_row = 0

            for row in range(2, ws.max_row + 1):
                keyword = str(ws.cell(row=row, column=col_idx_keyword).value or "").strip()
                ten_anh = str(ws.cell(row=row, column=col_idx_image).value or "").strip()

                current_row += 1
                print(f"📥 [{sheet_name}] Đang xử lý {current_row}/{total_rows}: {keyword}")

                if not keyword or ten_anh.lower() == "none.jpg" or ten_anh:
                    continue

                print(f"🔍 Tìm ảnh: {keyword} – sheet: {sheet_name}")
                img_url = get_first_image(keyword)
                if not img_url:
                    print("❌ Không tìm thấy ảnh.")
                    continue


                ext = ".jpg"
                file_name = f"{quote(keyword)}{ext}"
                img_path = os.path.join(folder_lang, file_name)

                try:
                    r = requests.get(img_url, timeout=10)
                    with open(img_path, "wb") as f:
                        f.write(r.content)

                    ws.cell(row=row, column=col_idx_image).value = file_name
                    print(f"✅ Đã lưu ảnh: {img_path}")
                    so_tai_duoc += 1
                except Exception as e:
                    print(f"⚠️ Lỗi tải ảnh {keyword}: {e}")

                time.sleep(2)
        unprotect_file(EXCEL_GAME_PATH)
        wb.save(EXCEL_GAME_PATH)
        messagebox.showinfo("Hoàn tất", f"🎉 Đã cập nhật ảnh cho game.\nTải được {so_tai_duoc} ảnh.")
    except Exception as e:
        messagebox.showerror("Lỗi", f"❌ Lỗi tải ảnh: {e}")


def tai_anh_con_thieu_game_thread():
    import threading
    print("Disabled: game image downloader removed in audio-tool version.")
    try:
        messagebox.showinfo("Đã tắt", "Tính năng tải ảnh game đã được tắt.")
    except Exception:
        pass
    return

    from tkinter import Toplevel, Label
    import tkinter as tk
    from urllib.parse import quote
    import requests
    from bs4 import BeautifulSoup
    from openpyxl import load_workbook
    import os, time
    from tkinter import messagebox

    popup = Toplevel()
    set_popup_icon(popup)
    popup.title("Đang tải ảnh thiếu...")
    popup.geometry("520x320")
    popup.resizable(False, False)
    popup.configure(bg="white")
    Label(popup, text="🔄 Đang tải ảnh thiếu từ Google...", font=("Arial", 22), bg="white").pack(pady=20)
    progress_label = Label(popup, text="Vui lòng đợi...", font=("Arial", 16), bg="white", fg="gray")
    progress_label.pack()

    def run_download():
        SHEETS = ["Ja", "Cn", "En"]
        COLUMN_KEYWORD = "Nghĩa TV"
        COLUMN_IMAGE = "images"
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
        so_tai_duoc = 0
        so_dong_can_tai = 0

        try:
            wb = load_workbook(EXCEL_GAME_PATH)

            def get_first_image(keyword):
                q = quote(keyword)
                url = f"https://www.google.com/search?tbm=isch&q={q}"
                try:
                    res = requests.get(url, headers=headers, timeout=10)
                    soup = BeautifulSoup(res.text, "html.parser")
                    img_tags = soup.find_all("img")
                    for img in img_tags[1:]:
                        src = img.get("src", "")
                        if src.startswith("http"):
                            return src
                except Exception:
                    return None
                return None

            # Đếm tổng số dòng cần tải
            for sheet_name in SHEETS:
                if sheet_name not in wb.sheetnames:
                    continue

                ws = wb[sheet_name]
                header_row = [cell.value for cell in ws[1]]
                if COLUMN_KEYWORD not in header_row:
                    continue

                col_idx_keyword = header_row.index(COLUMN_KEYWORD) + 1
                col_idx_image = header_row.index(COLUMN_IMAGE) + 1 if COLUMN_IMAGE in header_row else len(header_row) + 1

                for row in range(2, ws.max_row + 1):
                    keyword = str(ws.cell(row=row, column=col_idx_keyword).value or "").strip()
                    ten_anh = str(ws.cell(row=row, column=col_idx_image).value or "").strip()
                    if keyword and not ten_anh and ten_anh.lower() != "none.jpg":
                        so_dong_can_tai += 1

            # Tiến hành tải
            for sheet_name in SHEETS:
                if sheet_name not in wb.sheetnames:
                    continue

                ws = wb[sheet_name]
                header_row = [cell.value for cell in ws[1]]
                if COLUMN_KEYWORD not in header_row:
                    continue

                col_idx_keyword = header_row.index(COLUMN_KEYWORD) + 1
                if COLUMN_IMAGE not in header_row:
                    ws.cell(row=1, column=len(header_row) + 1).value = COLUMN_IMAGE
                    col_idx_image = len(header_row) + 1
                else:
                    col_idx_image = header_row.index(COLUMN_IMAGE) + 1

                folder_sub = os.path.join(APPDATA_IMAGE_FOLDER, sheet_name)
                os.makedirs(folder_sub, exist_ok=True)

                for row in range(2, ws.max_row + 1):
                    keyword = str(ws.cell(row=row, column=col_idx_keyword).value or "").strip()
                    ten_anh = str(ws.cell(row=row, column=col_idx_image).value or "").strip()
                    if not keyword or ten_anh or ten_anh.lower() == "none.jpg":
                        continue

                    progress_label.config(text=f"[{sheet_name}] 🔍 Đang tìm: {keyword}\nĐã tải {so_tai_duoc}/{so_dong_can_tai} ảnh")
                    popup.update()

                    img_url = get_first_image(keyword)
                    if not img_url:
                        continue

                    ext = ".jpg"
                    file_name = f"{quote(keyword)}{ext}"
                    img_path = os.path.join(folder_sub, file_name)

                    try:
                        r = requests.get(img_url, timeout=10)
                        with open(img_path, "wb") as f:
                            f.write(r.content)
                        ws.cell(row=row, column=col_idx_image).value = file_name
                        so_tai_duoc += 1
                        progress_label.config(text=f"[{sheet_name}] ✅ {keyword} – {so_tai_duoc}/{so_dong_can_tai} ảnh")
                        popup.update()
                    except Exception as e:
                        print(f"⚠️ Lỗi tải ảnh {keyword}: {e}")

                    time.sleep(2)

            unprotect_file(EXCEL_GAME_PATH)
            wb.save(EXCEL_GAME_PATH)
            popup.destroy()
            messagebox.showinfo("Hoàn tất", f"🎉 Đã tải {so_tai_duoc}/{so_dong_can_tai} ảnh cho 3 sheet.")
        except Exception as e:
            popup.destroy()
            messagebox.showerror("Lỗi", f"❌ Lỗi tải ảnh: {e}")

    threading.Thread(target=run_download).start()

#=== đếm ảnh
def dem_so_anh_thieu():
    print("Disabled: game image missing check removed in audio-tool version.")
    return 0

    from openpyxl import load_workbook
    import os
    from urllib.parse import quote

    try:
        wb = load_workbook(EXCEL_GAME_PATH)
        sheets = ["Ja", "Cn", "En"]
        dem = 0

        for sheet in sheets:
            if sheet not in wb.sheetnames:
                continue

            ws = wb[sheet]
            header = [cell.value for cell in ws[1]]
            if "Nghĩa TV" not in header or "images" not in header:
                continue

            col_keyword = header.index("Nghĩa TV") + 1
            col_image = header.index("images") + 1
            folder_lang = os.path.join(IMAGE_FOLDER, sheet)
            os.makedirs(folder_lang, exist_ok=True)

            for row in range(2, ws.max_row + 1):
                keyword = str(ws.cell(row=row, column=col_keyword).value or "").strip()
                image_name = str(ws.cell(row=row, column=col_image).value or "").strip()

                # Nếu chưa có ảnh hoặc ảnh không tồn tại thật sự
                if keyword:
                    if not image_name or image_name.lower() == "none.jpg":
                        dem += 1
                    else:
                        image_path = os.path.join(folder_lang, image_name)
                        if not os.path.exists(image_path):
                            dem += 1

        return dem
    except Exception as e:
        print("❌ Lỗi kiểm tra ảnh thiếu:", e)
        return 0


# ===== END DEAD CODE: GAME IMAGE DOWNLOADER =====

  
#======tẠO CỬA SỔ gui

root = tk.Tk()
#root.state("zoomed")  # Tự full màn hình khi mở

#LOGO APP
logo_path = os.path.join(IMAGES_DIR, "logo.png")

try:
    root.iconphoto(False, tk.PhotoImage(file=logo_path))
except Exception as e:
    print(f"Không tìm thấy logo: {e}")
root.title(f"Text To MP3/M4A ĐA NGÔN NGỮ - Audio Tool - VCJ International School | {APP_BUILD_TAG}")
root.configure(bg="#eef3ee")


def bung_man_hinh():
    root.deiconify()
    root.lift()
    root.focus_force()
    try:
        dang_toan_man_hinh = str(root.attributes("-fullscreen")).lower() in ("1", "true", "yes")
        root.attributes("-fullscreen", not dang_toan_man_hinh)
    except Exception:
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        root.geometry(f"{screen_width}x{screen_height}+0+0")


btn_bung_man_hinh = tk.Button(
    root,
    text="🖥 Toàn màn hình",
    font=("Arial", 10, "bold"),
    bg="#fff2cc",
    fg="#8a4b00",
    command=bung_man_hinh,
)
btn_bung_man_hinh.place(x=10, y=8, width=120, height=28)

def mo_cai_dat_popup(event=None):
    try:
        x = root.winfo_rootx() + 10
        y = root.winfo_rooty() + 40
        menu_cai_dat.tk_popup(x, y)
    finally:
        try:
            menu_cai_dat.grab_release()
        except Exception:
            pass

btn_cai_dat = tk.Button(
    root,
    text="Cai dat",
    font=("Arial", 10, "bold"),
    bg="#e8f0fe",
    fg="#0b57d0",
    command=mo_cai_dat_popup,
)
btn_cai_dat.place(x=138, y=8, width=90, height=28)

# Cố định kích thước và canh giữa màn hình
root.update_idletasks()
width, height = 1380, 760
screen_width = root.winfo_screenwidth()
screen_height = root.winfo_screenheight()
x = (screen_width // 2) - (width // 2)
y = (screen_height // 2) - (height // 2)
root.geometry(f"{width}x{height}+{x}+{y}")
root.resizable(True, True)
root.minsize(1050, 620)
#=======================
#GỌI FRAME
mic_effect_var = tk.StringVar(value="🎙️")
(frame_noi_dung, frame_text, frame_hoi_gpt, frame_doc, txt_de, entry_cau_hoi, btn_doc, btn_tam_dung, btn_doc_lai,
 btn_xuat_mp3, progress_var, progress_bar, che_do_doc, combo_ngon_ngu, combo_toc_do
) = create_frame_noi_dung(root)
#==================

# === Logo và bố cục co giãn của cửa sổ chính ===
frame_logo = tk.Frame(root, bg="#f5f7e8", bd=0, highlightthickness=0)
lbl_logo = tk.Label(frame_logo, text="Audio Tool", bg="#f5f7e8", fg="green", font=("Arial", 14, "bold"))
lbl_logo.pack(fill="both", expand=True, padx=6, pady=6)

try:
    from PIL import Image, ImageTk
    logo_source_image = Image.open(LOGO_PATH).convert("RGBA")
except Exception as e:
    logo_source_image = None
    print(f"Không tải được logo để co giãn: {e}")


def _resize_logo(width, height):
    """Giữ nguyên tỉ lệ logo, luôn nằm gọn trong khung và không bị cắt."""
    if logo_source_image is None or width < 10 or height < 10:
        return
    image = logo_source_image.copy()
    image.thumbnail((max(10, width - 12), max(10, height - 12)), Image.LANCZOS)
    logo_img = ImageTk.PhotoImage(image)
    lbl_logo.configure(image=logo_img, text="")
    lbl_logo.image = logo_img


def _place_doc_controls(panel_width, panel_height):
    """Dàn các nút bên phải theo kích thước panel; không thay đổi command."""
    children = frame_doc.winfo_children()
    buttons = [w for w in children if isinstance(w, tk.Button)]
    by_text = {str(w.cget("text")): w for w in buttons}
    inner_w = max(180, panel_width - 40)
    button_w = min(280, inner_w)
    button_x = max(10, (panel_width - button_w) // 2)
    button_h = 34
    placements = [
        ("Đọc / Chọn ngôn ngữ", 10, button_h),
        ("⏸ Dừng", 58, 30),
        ("🔁 Đọc lại", 96, 30),
        ("Convert Mp3 🎹 Wav", max(250, panel_height - 320), 30),
        ("✂️ Cutter Sound", max(286, panel_height - 284), 30),
        ("➕ Joiner Sound", max(322, panel_height - 248), 30),
        ("📥 Import tài liệu", max(358, panel_height - 172), 30),
    ]
    for text, y_pos, h in placements:
        widget = by_text.get(text)
        if widget is not None:
            widget.place(x=button_x, y=y_pos, width=button_w, height=h)

    frame_che_do = next((w for w in children if isinstance(w, tk.Frame)), None)
    if frame_che_do is not None:
        frame_che_do.place(x=button_x, y=150, width=button_w, height=34)
    combo_ngon_ngu.place(x=button_x, y=194, width=max(80, (button_w - 10) // 2), height=30)
    combo_toc_do.place(x=button_x + (button_w + 10) // 2, y=194,
                       width=max(80, (button_w - 10) // 2), height=30)


_layout_job = None
def apply_main_layout(event=None):
    """Responsive layout cho phần giao diện chính khi resize/maximize."""
    global _layout_job
    if _layout_job is not None:
        try:
            root.after_cancel(_layout_job)
        except Exception:
            pass
    _layout_job = root.after(30, _apply_main_layout_now)


def _apply_main_layout_now():
    global _layout_job
    _layout_job = None
    root.update_idletasks()
    win_w = max(1050, root.winfo_width())
    win_h = max(620, root.winfo_height())
    margin = 10
    top = 42
    main_w = win_w - margin * 2
    main_h = win_h - top - margin
    frame_noi_dung.place(x=margin, y=top, width=main_w, height=main_h)

    right_w = min(340, max(285, int(main_w * 0.235)))
    gap = 14
    left_w = max(520, main_w - right_w - gap - 20)
    question_h = 105
    text_h = max(260, main_h - question_h - 32)
    frame_text.place(x=10, y=8, width=left_w, height=text_h)
    # Khung chat rộng bằng khung nhập nội dung để hai vùng thẳng hàng.
    frame_hoi_gpt.place(x=10, y=text_h + 18, width=left_w, height=question_h)
    frame_doc.place(x=left_w + gap, y=8, width=right_w, height=main_h - 16)

    question_w = left_w
    entry_cau_hoi.place(x=8, y=8, width=max(300, question_w - 255), height=28)
    question_buttons = [w for w in frame_hoi_gpt.winfo_children() if isinstance(w, tk.Button)]
    x = max(210, question_w - 230)
    for index, widget in enumerate(question_buttons):
        widget.place(x=x + index * 72, y=8, width=68 if index < 2 else 54, height=28)
    for widget in frame_hoi_gpt.winfo_children():
        if isinstance(widget, tk.Label):
            widget.place(x=max(10, question_w - 45), y=42)

    panel_h = max(300, main_h - 16)
    _place_doc_controls(right_w, panel_h)
    # Logo nằm trong khoảng trống riêng, phía trên nhóm nút công cụ;
    # không đặt ở đáy panel để tránh che nút Import/Cutter/Joiner.
    action_y = max(250, panel_h - 320)
    logo_x = left_w + gap + 15
    logo_w = min(185, right_w - 30)
    logo_y = top + 8 + 250
    logo_h = min(150, action_y - 270)
    if logo_h >= 70:
        frame_logo.place(x=logo_x + (right_w - logo_w - 30) // 2,
                         y=logo_y, width=logo_w, height=logo_h)
        _resize_logo(logo_w, logo_h)
    else:
        frame_logo.place_forget()


root.bind("<Configure>", apply_main_layout)
root.after(100, _apply_main_layout_now)

def _bring_root_to_front():
    try:
        root.deiconify()
        try:
            root.attributes("-fullscreen", False)
        except Exception:
            pass
        root.lift()
        root.focus_force()
        root.attributes("-topmost", True)
        root.after(1200, lambda: root.attributes("-topmost", False))
    except Exception:
        pass

root.after(300, _bring_root_to_front)
root.after(500, lambda: threading.Thread(target=_load_secret_bundle_after_ui, daemon=True).start())

#============
def cau_hinh_khoi_dong_cung_win():
    import winreg
    import sys, os
    from tkinter import messagebox

    key = r"Software\Microsoft\Windows\CurrentVersion\Run"
    app_name = "AppSmartLearning"

    # ✅ Chỉ dùng được khi chạy file EXE đã build
    if not getattr(sys, 'frozen', False):
        messagebox.showwarning("Không hỗ trợ", "Tính năng này chỉ hoạt động với file .EXE đã đóng gói.")
        return

    exe_path = os.path.abspath(sys.executable)

    def luu():
        try:
            if var_bat.get():
                with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key, 0, winreg.KEY_SET_VALUE) as reg_key:
                    winreg.SetValueEx(reg_key, app_name, 0, winreg.REG_SZ, f'"{exe_path}"')
            else:
                with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key, 0, winreg.KEY_SET_VALUE) as reg_key:
                    try:
                        winreg.DeleteValue(reg_key, app_name)
                    except FileNotFoundError:
                        pass
            messagebox.showinfo("✅ Thành công", "Đã cập nhật khởi động cùng Windows.")
            top.destroy()
        except Exception as e:
            messagebox.showerror("Lỗi", str(e))

    # Kiểm tra trạng thái hiện tại
    try:
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key) as reg_key:
            current_val, _ = winreg.QueryValueEx(reg_key, app_name)
            bat_dau = True
    except:
        bat_dau = False

    # Giao diện
    top = tk.Toplevel(root)
    set_popup_icon(top)
    top.title("🖥️ Khởi động cùng Windows")
    top.geometry("350x160")
    top.resizable(False, False)
    top.grab_set()

    var_bat = tk.BooleanVar(value=bat_dau)
    tk.Checkbutton(top, text="Tự động khởi động ứng dụng khi bật máy tính",
                   variable=var_bat, font=("Arial", 10)).pack(pady=20)
    tk.Button(top, text="💾 Áp dụng", font=("Arial", 11, "bold"),
              bg="green", fg="white", command=luu).pack(pady=10)


#===========================
# Disabled: không kiểm tra/tải ảnh game khi khởi động audio-tool version.
# try:
#     so_thieu = dem_so_anh_thieu()
#     if so_thieu > 0:
#         if messagebox.askyesno("Thiếu ảnh", f"Phát hiện {so_thieu} ảnh còn thiếu.\nBạn có muốn tải từ Google không?"):
#             tai_anh_con_thieu_game_thread()
# except Exception as e:
#     print("❌ Lỗi kiểm tra ảnh thiếu:", e)

#=========Sửa từng phần key=================

def sua_key_don(loai, key_field, label_hientai, label_moi, show_pw=False):
    def thuc_hien(ok):
        if not ok:
            return
        global config
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                config = json.load(f)
            old_value = config.get(key_field, "")

            top = tk.Toplevel(root)
            top.protocol("WM_DELETE_WINDOW", top.destroy)
            set_popup_icon(top)
            top.title(f"Sửa {loai}")
            top.geometry("560x340" if key_field == "GITHUB_MODELS_TOKEN" else "560x260")
            top.grab_set()
            top.resizable(False, False)

            tk.Label(top, text=label_hientai, font=("Arial", 10, "bold")).pack(pady=(10, 0))
            ent_old = tk.Entry(top, font=("Arial", 10), width=60)
            ent_old.insert(0, old_value)
            ent_old.configure(state='readonly')
            ent_old.pack(pady=2)

            tk.Label(top, text=label_moi, font=("Arial", 10)).pack()
            ent_new = tk.Entry(top, font=("Arial", 11), width=60, show="*" if show_pw else None)
            ent_new.pack(pady=2)

            if key_field == "GITHUB_MODELS_TOKEN":
                tk.Label(
                    top,
                    text=(
                        "Dùng GitHub Personal Access Token (PAT).\n"
                        "Fine-grained PAT cần quyền Models: Read; classic PAT cần scope models.\n"
                        "Không chia sẻ token cho người khác."
                    ),
                    justify="left",
                    wraplength=510,
                    fg="#444444",
                ).pack(padx=18, pady=(8, 2), anchor="w")
                tk.Button(
                    top,
                    text="Mở GitHub để tạo token",
                    command=lambda: open_path_cross_platform("https://github.com/settings/personal-access-tokens/new"),
                    fg="blue",
                ).pack(pady=(0, 2))

            def luu():
                new_value = _normalize_secret_text(ent_new.get())
                if not new_value:
                    messagebox.showwarning("Thiếu", f"Chưa nhập {loai}", parent=top)
                    return

                config[key_field] = new_value

                try:
                    unprotect_file(CONFIG_FILE)  # Gỡ bảo vệ
                    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                        json.dump(config, f, indent=2, ensure_ascii=False)
                    global GITHUB_MODELS_TOKEN, GEMINI_API_KEY, GOOGLE_TTS_API_KEY, DISCORD_WEBHOOK_URL, client
                    if key_field == "GITHUB_MODELS_TOKEN":
                        GITHUB_MODELS_TOKEN = new_value
                        client = _create_github_models_client(GITHUB_MODELS_TOKEN)
                    elif key_field == "GEMINI_API_KEY":
                        GEMINI_API_KEY = new_value
                    elif key_field == "GOOGLE_TTS_API_KEY":
                        GOOGLE_TTS_API_KEY = new_value
                    elif key_field == "DISCORD_WEBHOOK_URL":
                        DISCORD_WEBHOOK_URL = new_value
                    messagebox.showinfo("OK", f"Đã cập nhật {loai}!", parent=top)
                    top.destroy()

                except Exception as e:
                    try:
                        if top.winfo_exists():
                            messagebox.showerror("Lỗi", f"Không ghi được file:\n{e}", parent=top)
                        else:
                            messagebox.showerror("Lỗi", f"Không ghi được file:\n{e}")
                    except:
                        messagebox.showerror("Lỗi", f"Không ghi được file:\n{e}")

            def kiem_tra():
                import threading
                def run():
                    new_value = _normalize_secret_text(ent_new.get())
                    if not new_value:
                        messagebox.showwarning("Thiếu", f"Chưa nhập {loai}", parent=top)
                        return
                    try:
                        import requests
                        if key_field == "GITHUB_MODELS_TOKEN":
                            if not is_github_models_token(new_value):
                                raise Exception("❌ Token GitHub không đúng định dạng. Token PAT phải bắt đầu bằng 'ghp_' hoặc 'github_pat_'.")
                            try:
                                response = requests.get(
                                    GITHUB_MODELS_CATALOG_URL,
                                    headers={
                                        "Accept": "application/vnd.github+json",
                                        "Authorization": f"Bearer {new_value}",
                                        "X-GitHub-Api-Version": GITHUB_MODELS_API_VERSION,
                                    },
                                    timeout=15,
                                )
                                if response.status_code == 401:
                                    raise Exception("❌ Token GitHub không hợp lệ hoặc đã bị thu hồi.")
                                if response.status_code == 403:
                                    raise Exception("❌ Token chưa có quyền GitHub Models. Fine-grained PAT cần Models: Read; classic PAT cần scope 'models'.")
                                if response.status_code == 429:
                                    raise Exception("❌ GitHub Models đang giới hạn lượt dùng. Hãy chờ quota được làm mới rồi thử lại.")
                                response.raise_for_status()
                            except Exception as e:
                                msg = str(e)
                                if "quota" in msg.lower() or "rate limit" in msg.lower():
                                    raise Exception("❌ Token hợp lệ nhưng GitHub Models đã hết lượt miễn phí / bị giới hạn. Hãy chờ rồi thử lại.")
                                elif "Token GitHub" in msg or "quyền GitHub Models" in msg:
                                    raise
                                else:
                                    raise Exception(f"❌ Không kiểm tra được GitHub Models: {e}")
                            messagebox.showinfo("OK", "✅ GitHub Models token hoạt động!", parent=top)

                        elif key_field == "GEMINI_API_KEY":
                            temp_client = genai.Client(api_key=new_value)
                            temp_client.models.generate_content(
                                model="gemini-1.5-flash",
                                contents="ping",
                            )
                            messagebox.showinfo("OK", "✅ Gemini API key hoạt động!", parent=top)

                        elif key_field == "DISCORD_WEBHOOK_URL":
                            if new_value.startswith("ps://"):
                                new_value = new_value.replace("ps://", "https://", 1)
                            if "api/webhooks/" not in new_value:
                                raise Exception("⚠ Webhook không đúng định dạng.\nURL phải chứa /api/webhooks/")
                            test_message = (
                                "✅ Test webhook từ Text To MP3/M4A.\n"
                                f"Thời gian: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n"
                                "Nếu bạn thấy tin nhắn này thì webhook đang hoạt động."
                            )
                            r = requests.post(new_value, json={"content": test_message}, timeout=10)
                            if r.status_code in [200, 204] or r.ok:
                                messagebox.showinfo("OK", "✅ Webhook Discord hoạt động và đã gửi tin nhắn test!", parent=top)
                            elif r.status_code == 401:
                                raise Exception("❌ Webhook sai hoặc đã bị xoá (401 Unauthorized).")
                            elif r.status_code == 404:
                                raise Exception("❌ Webhook không tồn tại (404 Not Found).")
                            else:
                                raise Exception(f"Lỗi HTTP {r.status_code}\nPhản hồi: {r.text}")

                        elif key_field == "GOOGLE_TTS_API_KEY":
                            from google.api_core.client_options import ClientOptions
                            from google.cloud import texttospeech
                            test_client = texttospeech.TextToSpeechClient(
                                client_options=ClientOptions(api_key=new_value)
                            )
                            test_client.list_voices(language_code="vi-VN")
                            messagebox.showinfo("OK", "✅ Google TTS API key hoạt động!", parent=top)

                        else:
                            messagebox.showinfo("Thông báo", f"⚠ Chưa hỗ trợ kiểm tra loại '{loai}'.", parent=top)

                    except Exception as e:
                        if not thong_bao_loi_api(e, loai):
                            messagebox.showerror("Lỗi kiểm tra", f"❌ Không kiểm tra được:\n{e}")


                threading.Thread(target=run, daemon=True).start()

            frame = tk.Frame(top)
            frame.pack(pady=10)
            tk.Button(frame, text="💾 Lưu", command=luu, bg="green", fg="white", width=10).pack(side="left", padx=10)
            tk.Button(frame, text="✅ Kiểm tra", command=kiem_tra, bg="#0066cc", fg="white", width=12).pack(side="left", padx=10)

        except Exception as e:
            print("Lỗi tạo popup:", e)

    ask_password_with_keyboard(thuc_hien)


API_KEY_SETTINGS = {
    "GitHub Models": ("GitHub Models Token", "GITHUB_MODELS_TOKEN", "GitHub PAT hiện tại:", "Nhập GitHub PAT mới:"),
    "Gemini": ("Gemini API KEY cho khung chat", "GEMINI_API_KEY", "Gemini Key hiện tại:", "Nhập Gemini Key mới:"),
    "Google TTS": ("API key cho Google TTS", "GOOGLE_TTS_API_KEY", "Google TTS Key hiện tại:", "Nhập Google TTS Key mới:"),
    "Discord": ("Discord Webhook", "DISCORD_WEBHOOK_URL", "Webhook Discord hiện tại:", "Nhập Webhook Discord mới:"),
    "Email": ("Email & App Password", "EMAIL", "Email/App Password hiện tại:", "Nhập lại trong mục Email & App Password:"),
    "AWS": ("AWS Keys", "AWS", "AWS Keys hiện tại:", "Nhập lại trong mục AWS Keys:"),
}


def _api_key_service_from_error(error_text, context=""):
    """Return the settings item that should be opened for an auth/quota error."""
    text = f"{context} {error_text}".lower()
    auth_words = (
        "api key", "apikey", "invalid_api_key", "unauthorized", "forbidden", "401", "403", "429",
        "quota", "billing", "resource_exhausted", "rate limit", "credential", "accessdenied", "authentication", "permission",
        "service_disabled", "expired", "hết hạn", "không hợp lệ",
    )
    if not any(word in text for word in auth_words):
        return None
    if any(word in text for word in ("gemini", "genai", "generate_content")):
        return "Gemini"
    if any(word in text for word in ("text-to-speech", "texttospeech", "google cloud tts", "google tts")):
        return "Google TTS"
    if any(word in text for word in ("github models", "models.github.ai", "github pat", "github token")):
        return "GitHub Models"
    if any(word in text for word in ("discord", "webhook")):
        return "Discord"
    if any(word in text for word in ("smtp", "email", "app password", "535")):
        return "Email"
    if any(word in text for word in ("aws", "amazon", "polly", "botocore", "unrecognizedclient")):
        return "AWS"
    return None


def _api_error_hint(error_text, context=""):
    """Classify the error so the popup can explain what to check next."""
    text = f"{context} {error_text}".lower()
    if any(word in text for word in (
        "api key not valid",
        "invalid api key",
        "invalid_api_key",
        "incorrect api key",
        "api key không đúng",
        "không đúng định dạng",
    )):
        return "invalid"
    if any(word in text for word in ("quota", "billing", "resource_exhausted", "insufficient_quota", "out of quota", "exceeded your current quota")):
        return "quota"
    if any(word in text for word in ("permission", "forbidden", "service disabled", "disabled", "access denied", "unauthorized", "401", "403", "restricted")):
        return "restricted"
    if any(word in text for word in ("expired", "hết hạn", "revoked", "deleted", "not found", "404")):
        return "expired"
    return "general"


def hien_popup_loi_api_key(service, error_text, context=""):
    """Show an actionable popup with a button to the matching key setting."""
    setting = API_KEY_SETTINGS.get(service)
    if not setting:
        return False

    hint = _api_error_hint(error_text, context)
    if hint == "invalid":
        headline = f"Key của {service} đang không hợp lệ hoặc đã bị dán sai."
        guidance = "Hãy kiểm tra lại key bạn copy có bị dính khoảng trắng, xuống dòng, dấu nháy, hoặc bị dán nhầm sang key khác."
    elif hint == "quota":
        headline = f"Key của {service} hợp lệ nhưng đã hết quota / billing."
        guidance = "Hãy kiểm tra hạn mức, thanh toán và quyền dùng API trong tài khoản dịch vụ."
    elif hint == "restricted":
        headline = f"Key của {service} hợp lệ nhưng đang bị chặn quyền truy cập."
        if service == "GitHub Models":
            guidance = "Hãy tạo lại GitHub PAT: fine-grained cần Models: Read, còn classic cần scope 'models'; rồi dán vào Cài đặt."
        else:
            guidance = "Hãy kiểm tra API đã bật đúng chưa, key có bị giới hạn sai dịch vụ, IP, referrer hoặc project không."
    elif hint == "expired":
        headline = f"Key của {service} có thể đã bị xoá, thu hồi hoặc hết hạn."
        guidance = "Hãy tạo key mới hoặc kiểm tra lại key đang dùng có còn tồn tại không."
    else:
        headline = f"Dịch vụ {service} đang báo lỗi key, quyền truy cập hoặc quota."
        guidance = "Vui lòng kiểm tra key còn hiệu lực, billing/quota và thay key trong Cài đặt."

    top = tk.Toplevel(root)
    set_popup_icon(top)
    top.title("⚠️ Kiểm tra API key")
    top.geometry("640x280")
    top.grab_set()
    tk.Label(
        top,
        text=headline,
        font=("Arial", 11, "bold"),
        fg="#b00020",
        wraplength=590,
    ).pack(padx=18, pady=(18, 8))
    tk.Label(
        top,
        text=guidance,
        justify="left",
        wraplength=590,
    ).pack(padx=18, pady=4)
    tk.Label(
        top,
        text=str(error_text)[:900],
        justify="left",
        wraplength=590,
        fg="#555555",
    ).pack(padx=18, pady=4)

    def open_setting():
        top.destroy()
        if service == "Email":
            sua_key_nhom(
                "Email & App Password",
                ["SENDER_EMAIL", "SENDER_NAME", "APP_PASSWORD"],
                ["Email gửi", "Tên gửi", "App Password"],
                show_pws=[False, False, True],
            )
            return
        if service == "AWS":
            sua_key_nhom(
                "AWS Keys",
                ["AWS_ACCESS_KEY_ID", "AWS_SECRET_ACCESS_KEY", "AWS_REGION"],
                ["Access Key ID", "Secret Access Key", "Region"],
                show_pws=[False, True, False],
            )
            return
        sua_key_don(*setting, show_pw=True)

    frame = tk.Frame(top)
    frame.pack(pady=14)
    tk.Button(frame, text=f"🔧 Mở thay key {service}", command=open_setting, bg="#0066cc", fg="white", width=22).pack(side="left", padx=8)
    tk.Button(frame, text="Đóng", command=top.destroy, width=12).pack(side="left", padx=8)
    return True


def thong_bao_loi_api(error, context=""):
    service = _api_key_service_from_error(str(error), context)
    if service:
        hien_popup_loi_api_key(service, str(error), context)
    return service


def _normalize_secret_text(value):
    """Normalize pasted secrets by removing invisible whitespace and outer quotes."""
    text = (value or "").strip()
    text = text.strip('"').strip("'")
    text = re.sub(r"\s+", "", text)
    return text

 
#==============
def gioi_thieu_ung_dung():
    top = tk.Toplevel(root)
    set_popup_icon(top)
    top.title("ℹ️ Giới thiệu Ứng dụng")
    top.geometry("1000x500")
    top.attributes('-topmost', True)
    thong_tin = """
    🌟 Ứng dụng Text to mp3 - Đa NGÔN NGỮ🌟 thuộc dự án  Smart Learning SmL - Học Tập Đa NGÔN NGỮ

    🎯 Chức năng chính Ứng dụng Smart Learning SmL - Học Tập Đa NGÔN NGỮ:
     Link tải exe : https://tuadenu.github.io/smartlearning/latest.html
    - Hẹn giờ tự động thông minh file nghe, bài đọc cho từng học viên các khung giờ khác nhau, có chuông thông báo và gọi đúng tên học viên
    - Tạo đề thông minh luyện tập: Toán, Ngoại ngữ, Tự luận, Trắc nghiệm...Đa ngôn ngữ các trình độ , phạm vi tuỳ chọn
    - Tích hợp GitHub Models & Gemini tương tác thông minh với học viên, có thể in ra trực tiếp 1 click
    - Hẹn giờ phát nhạc - kết hợp camera gửi ảnh
    - Điều khiển thiết bị Broadlink ,nhà thông minh, học lệnh điều khiển với 1 nút nhấn
    - Gửi ảnh ở bàn học 3 phút/lần (hoặc tuỳ chọn với 3 camera) qua Telegram, Discord cho giáo viên, phụ huynh học viên hoặc chính học viên thông minh
    - Giám sát chặt chẽ quá trình học hoặc làm bài thi
    - Đọc bất kỳ nội dung gì trong khung hiển thị đề kể cả tài liệu đa ngôn ngữ, như Anh, Trung, Nhật, việt... thông minh tự động hoặc thủ công chọn dòng
    - Xuất nội dung bất  ra MP3 bằng đa ngôn ngữ, như Anh, Trung, Nhật, việt...
    - Kết nối trực tiếp máy in và lưu lịch sử in, thư mục chứ ảnh người in và tài liệu PDF tương ứng, gửi 1 bản tới nhiều email tuỳ chọn
    - Hỗ trợ bàn phím mini, giao diện fullscreen thân thiện dễ dùng
    - Tích hợp máy tính mini ngay trên giao diện
    - Trình chiếu ảnh hoặc file ảnh học thuật tự động
    - Tích hợp trình phát nhạc mini để luyện nghe
    - Bảng thông tin thời tiết địa phương
    - Tích hợp 4camera rtps/onvif xem trực tiếp trên App
    - Tự động tải cài đặt update khi có bản mới, có thể call sdt để lấy link khi bị lỗi update nhé!
    💡 Phát triển bởi: Linh Dương - Cty TNHH DU LỊCH VÀ THƯƠNG MẠI QUỐC TẾ VIỆT TRUNG NHẬT 
    📧 Liên hệ: halamchuc@gmail.com 0986183806
"""
    tk.Label(top, text=thong_tin, font=("Arial", 10),fg="green", justify="left").pack(padx=15, pady=10)

# sủa key email, âmazzon
def sua_key_nhom(loai, key_fields, labels, show_pws=None):
    def thuc_hien(ok):
        if not ok:
            return
        global config
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                config = json.load(f)

            top = tk.Toplevel(root)
            top.protocol("WM_DELETE_WINDOW", top.destroy)
            set_popup_icon(top)
            top.title(f"Sửa {loai}")
            top.geometry("640x{}".format(220 + 60 * len(key_fields)))
            top.grab_set()
            top.resizable(False, False)

            entries = {}

            for idx, key in enumerate(key_fields):
                old_value = config.get(key) or config_default.get(key, "")

                tk.Label(top, text=labels[idx] + " (hiện tại):", font=("Arial", 10, "bold")).pack(pady=(5, 0))
                ent_old = tk.Entry(top, font=("Arial", 10), width=60)
                ent_old.insert(0, old_value)
                ent_old.pack(pady=2)

                tk.Label(top, text=labels[idx] + " (mới):", font=("Arial", 10)).pack()
                show_pw = show_pws[idx] if show_pws else False
                ent_new = tk.Entry(top, font=("Arial", 11), width=60, show="*" if show_pw else None)
                ent_new.pack(pady=2)

                entries[key] = ent_new

            # Progress & status
            status_var = tk.StringVar()
            tk.Label(top, textvariable=status_var, font=("Arial", 9), fg="green").pack(pady=(4, 2))
            progress = ttk.Progressbar(top, orient="horizontal", length=280, mode="determinate")
            progress.pack(pady=(2, 4))

            def luu():
                updated = False
                for key, ent in entries.items():
                    new_value = _normalize_secret_text(ent.get())
                    if new_value:
                        config[key] = new_value
                        updated = True
                if not updated:
                    messagebox.showwarning("Thiếu", "Chưa nhập giá trị mới.", parent=top)
                    return

                try:
                    unprotect_file(CONFIG_FILE)
                    with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                        json.dump(config, f, indent=2, ensure_ascii=False)
                    global SENDER_EMAIL, SENDER_NAME, APP_PASSWORD, AWS_ACCESS_KEY_ID, AWS_SECRET_ACCESS_KEY, AWS_REGION
                    SENDER_EMAIL = config.get("SENDER_EMAIL", SENDER_EMAIL)
                    SENDER_NAME = config.get("SENDER_NAME", SENDER_NAME)
                    APP_PASSWORD = config.get("APP_PASSWORD", APP_PASSWORD)
                    AWS_ACCESS_KEY_ID = config.get("AWS_ACCESS_KEY_ID", AWS_ACCESS_KEY_ID)
                    AWS_SECRET_ACCESS_KEY = config.get("AWS_SECRET_ACCESS_KEY", AWS_SECRET_ACCESS_KEY)
                    AWS_REGION = config.get("AWS_REGION", AWS_REGION)
                    messagebox.showinfo("OK", f"Đã cập nhật {loai}!", parent=top)
                    top.destroy()
                except Exception as e:
                    messagebox.showerror("Lỗi", f"Không ghi được file:\n{e}", parent=top)

            def co_mang_internet():
                try:
                    import socket
                    socket.create_connection(("8.8.8.8", 53), timeout=3)
                    return True
                except:
                    return False

            def kiem_tra():
                if not co_mang_internet():
                    status_var.set("")
                    progress["value"] = 0
                    messagebox.showerror("Mất kết nối", "Không có kết nối Internet.")
                    return

                # ==== Kiểm tra Email + App Password ====
                if "SENDER_EMAIL" in entries and "APP_PASSWORD" in entries:
                    email = entries["SENDER_EMAIL"].get().strip()
                    pw = entries["APP_PASSWORD"].get().strip()
                    status_var.set("🔄 Đang kiểm tra email, vui lòng đợi...")
                    progress["value"] = 0
                    top.update()

                    try:
                        import smtplib
                        from email.message import EmailMessage

                        msg = EmailMessage()
                        msg['Subject'] = "✅ Kiểm tra cấu hình Gmail gửi thành công"
                        msg['From'] = email
                        msg['To'] = email
                        msg.set_content("Bạn đã cấu hình đúng email gửi + mật khẩu ứng dụng.")

                        progress["value"] = 20
                        top.update()

                        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
                            progress["value"] = 50
                            top.update()
                            smtp.login(email, pw)
                            progress["value"] = 80
                            top.update()
                            smtp.send_message(msg)
                            progress["value"] = 100
                            top.update()

                        status_var.set("✅ Kiểm tra thành công!")
                        messagebox.showinfo("✅ OK", "Đã gửi email test thành công đến chính bạn!")

                    except Exception as e:
                        status_var.set("")
                        progress["value"] = 0
                        if not thong_bao_loi_api(e, "Email"):
                            messagebox.showerror("❌ Lỗi", f"Không gửi được email test:\n{e}")

                # ==== Kiểm tra AWS ====
                elif "AWS_ACCESS_KEY_ID" in entries and "AWS_SECRET_ACCESS_KEY" in entries and "AWS_REGION" in entries:
                    key_id = entries["AWS_ACCESS_KEY_ID"].get().strip()
                    secret = entries["AWS_SECRET_ACCESS_KEY"].get().strip()
                    region = entries["AWS_REGION"].get().strip()
                    status_var.set("🔄 Đang kiểm tra AWS, vui lòng đợi...")
                    progress["value"] = 0
                    top.update()

                    try:
                        import boto3
                        progress["value"] = 20
                        top.update()

                        polly_client = boto3.Session(
                            aws_access_key_id=key_id,
                            aws_secret_access_key=secret,
                            region_name=region
                        ).client("polly")

                        progress["value"] = 60
                        top.update()

                        polly_client.describe_voices()

                        progress["value"] = 100
                        top.update()

                        status_var.set("✅ AWS Polly hoạt động tốt!")
                        messagebox.showinfo("✅ OK", "AWS Polly hoạt động tốt, key & region hợp lệ.")

                    except Exception as e:
                        status_var.set("")
                        progress["value"] = 0
                        if not thong_bao_loi_api(e, "AWS"):
                            messagebox.showerror("❌ Lỗi", f"Không kiểm tra được AWS Polly:\n{e}")

                else:
                    messagebox.showinfo("Thông báo", f"⚡ Chưa hỗ trợ kiểm tra online cho nhóm {loai}.", parent=top)

            frame = tk.Frame(top)
            frame.pack(pady=10)
            tk.Button(frame, text="💾 Lưu", command=luu, bg="green", fg="white", width=10).pack(side="left", padx=10)
            tk.Button(frame, text="✅ Kiểm tra", command=kiem_tra, bg="#0066cc", fg="white", width=12).pack(side="left", padx=10)

        except Exception as e:
            print("Lỗi tạo popup:", e)

    ask_password_with_keyboard(thuc_hien)
#==============
# các ứng dụng khác
def mo_popup_ung_dung_khac():
    danh_sach_app = [
        ("Máy học tập các phiên bản ", "https://1drv.ms/f/c/86031e0f977fa7e6/Elmiqga_boFIjz7MQ_6GmvcBt96I2nle_UhpLYdn7GAZmg?e=4yFG3R"),
        ("Smart Learning  bản mới nhất", "https://tuadenu.github.io/smartlearning/latest.html"),
        ("Youtube", "https://www.youtube.com/@tiengtrunglinhduong"),
        ("Gọi điện thoại", "https://sites.google.com/view/tiengtrunglinhduong?fbclid=IwY2xjawLduOFleHRuA2FlbQIxMABicmlkETE4WE9ZM1Njdm1zckZobnBuAR7U7401p0yATuofu5gjqB5M1QR37Ait5AHquytcY9uq5nx1MjF7K4djnRmK-Q_aem_E05MPn-wzny_zd1XiUIKvw"),
        ("Facebook", "https://www.facebook.com/tiengtrunglinhduong")
    ]

    popup = tk.Toplevel(root)
    set_popup_icon(popup)
    popup.title("📦 Tải các ứng dụng khác")
    popup.geometry("400x300")
    popup.grab_set()

    tk.Label(popup, text="Danh sách ứng dụng:", font=("Arial", 12, "bold")).pack(pady=10)

    for ten_app, link in danh_sach_app:
        frame = tk.Frame(popup)
        frame.pack(fill="x", padx=10, pady=5)
        tk.Label(frame, text=ten_app, anchor="w").pack(side="left", expand=True)
        tk.Button(frame, text="Truy cập link", fg="blue", command=lambda l=link: open_path_cross_platform(l)).pack(side="right")

    tk.Button(popup, text="Đóng", command=popup.destroy).pack(pady=10)

#===============


# MENU CẤU HÌNH ===

menu_cai_dat = tk.Menu(root, tearoff=0)
menu_cai_dat.add_command(label="Khoi dong cung Windows", command=cau_hinh_khoi_dong_cung_win)
menu_cai_dat.add_separator()
menu_cai_dat.add_command(label="Gioi thieu ung dung", command=gioi_thieu_ung_dung)
menu_cai_dat.add_separator()
menu_cai_dat.add_separator()
menu_cai_dat.add_command(label="Doi mat khau toan ung dung", command=doi_mat_khau)
menu_cai_dat.add_command(label="Sua GitHub Models Token", command=lambda: sua_key_don("GitHub Models Token", "GITHUB_MODELS_TOKEN", "GitHub PAT hiện tại:", "Nhập GitHub PAT mới:", show_pw=True))
menu_cai_dat.add_command(label="Sua Gemini API KEY cho khung chat", command=lambda: sua_key_don("Gemini API KEY cho khung chat", "GEMINI_API_KEY", "Gemini Key hiện tại:", "Nhập Gemini Key mới:", show_pw=True))
menu_cai_dat.add_command(label="API key cho Google TTS", command=lambda: sua_key_don("API key cho Google TTS", "GOOGLE_TTS_API_KEY", "Google TTS Key hiện tại:", "Nhập Google TTS Key mới:", show_pw=True))
menu_cai_dat.add_command(label="Sua Discord Webhook", command=lambda: sua_key_don("Discord Webhook", "DISCORD_WEBHOOK_URL", "Webhook Discord hiện tại:", "Nhập Webhook Discord mới:"))
menu_cai_dat.add_separator()
menu_cai_dat.add_command(label="Tai cac ung dung khac", command=mo_popup_ung_dung_khac)
# Disabled: game image downloader menu removed in audio-tool version.
# menu_cai_dat.add_command(label="📥 Tải ảnh còn thiếu cho Game Đoán Chữ", command=tai_anh_con_thieu_game_thread)



menu_cai_dat.add_command(
    label="Sua Email va App Password",
    command=lambda: sua_key_nhom(
        "Email & App Password",
        ["SENDER_EMAIL", "SENDER_NAME", "APP_PASSWORD"],
        ["Email gửi", "Tên gửi", "App Password"],
        show_pws=[False, False, True]
    )
)

menu_cai_dat.add_command(
    label="Sua AWS Keys",
    command=lambda: sua_key_nhom(
        "AWS Keys",
        ["AWS_ACCESS_KEY_ID", "AWS_SECRET_ACCESS_KEY", "AWS_REGION"],
        ["Access Key ID", "Secret Access Key", "Region"],
        show_pws=[False, True, False]
    )
)


# Chỉ cần tạo 1 menubar, gắn menu_cai_dat:
menubar = tk.Menu(root)
menubar.add_cascade(label="Cai dat", menu=menu_cai_dat)
if sys.platform != "darwin":
    root.config(menu=menubar)

#+++++++=====
# PHÍM TẮT MÀN HÌNH TẠM THỜI
#Tắt màn hàm

def tat_man_hinh():
    try:
        if os.name == "nt":
            import ctypes
            HWND_BROADCAST = 0xFFFF
            WM_SYSCOMMAND = 0x0112
            SC_MONITORPOWER = 0xF170
            ctypes.windll.user32.PostMessageW(HWND_BROADCAST, WM_SYSCOMMAND, SC_MONITORPOWER, 2)
            return

        if sys.platform == "darwin":
            subprocess.run(["pmset", "displaysleepnow"], check=True)
            return

        raise RuntimeError("Chức năng tắt màn hình chưa hỗ trợ trên hệ điều hành này.")
    except Exception as e:
        messagebox.showerror("Lỗi", f"Không tắt được màn hình:\n{e}")


# Tạo cửa sổ nổi nhỏ chỉ chứa nút tắt màn hình trên Windows
if os.name == "nt":
    top_btn = tk.Toplevel(root)
    top_btn.overrideredirect(True)  # Ẩn khung cửa sổ
    top_btn.attributes("-topmost", True)  # Luôn trên cùng
    top_btn.geometry("+{}+{}".format(root.winfo_screenwidth()-250, root.winfo_screenheight()-80))  # Góc dưới phải

    btn_tatman = tk.Button(
        top_btn,
        text="🖥Ấn Để Tắt màn🖥",
        width=15, height=1,
        command=tat_man_hinh,
        font=("Arial", 11),
        fg="red", bg="#f0f8ff"
    )
    btn_tatman.pack()

    # Đảm bảo nút di chuyển theo cửa sổ cha (tùy chọn)
    def keep_on_top():
        # Lấy tọa độ root
        x = root.winfo_x() + root.winfo_width() - 150
        y = root.winfo_y() + root.winfo_height() - 30
        top_btn.geometry(f"+{x}+{y}")
        root.after(300, keep_on_top)

    keep_on_top()

#=====================


root.mainloop()
